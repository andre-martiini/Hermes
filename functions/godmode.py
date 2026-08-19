"""
Hermes Godmode — modo estratégico do Copiloto, rodando sobre Claude
(Anthropic) em vez de Gemini. Módulo aditivo: não altera nada do fluxo
existente de askCopilotoHermes; expõe seu próprio callable e seu próprio
registro de ferramentas.

Escopo: leitura ampla sobre todas as superfícies do Hermes onde o usuário
deixa rastro (tarefas, metas estratégicas, finanças, saúde, diário pessoal,
agenda, pessoas, WhatsApp, conhecimento) — o mesmo panorama que já alimenta
o diário pessoal diário (ver `personal_diary.py:_collect_diary_material`) —
e escrita restrita ao módulo de Estratégia (`estrategia_pessoal`), via os
mesmos módulos compartilhados usados pelo Copiloto padrão
(`strategy_tools.py`, `health_tools.py`, `tools/telegram_extended.py`,
`whatsapp_ingest.py`).
"""

import json
import os

from firebase_admin import firestore
from firebase_functions import https_fn, options

import strategy_tools
from health_tools import build_health_summary
from llm_providers import claude_provider

GODMODE_MODEL = os.environ.get("GODMODE_MODEL", "claude-fable-5")
GODMODE_FALLBACK_MODEL = os.environ.get("GODMODE_FALLBACK_MODEL", "claude-opus-4-8")
GODMODE_MAX_TOKENS = int(os.environ.get("GODMODE_MAX_TOKENS", "4096"))
GODMODE_FUNCTION_TIMEOUT_SEC = 300

GODMODE_PERSONA = (
    "## MODO GODMODE — Diretor de Estratégia e Operações Virtual\n"
    "Você é o Hermes Godmode: um conselheiro adversário, não um assistente condescendente. "
    "Sua função é auditar viabilidade e confrontar premissas do usuário com dados reais do "
    "sistema — nunca validar por educação.\n\n"
    "Escopo: você enxerga o panorama completo do usuário no Hermes — tarefas/ações, metas "
    "estratégicas, finanças, saúde (telemetria diária e relatório semanal), diário pessoal, "
    "agenda, pessoas/contatos, conversas de WhatsApp indexadas e base de conhecimento. Cruze "
    "essas fontes: uma meta financeira parada, uma dor lombar recorrente atrapalhando a agenda, "
    "um objetivo estratégico sem indicador com progresso real — esse tipo de conexão entre "
    "domínios é exatamente o valor que você entrega e que uma consulta isolada não mostra.\n\n"
    "Módulo Estratégia: além de consultar, você pode criar, editar e excluir objetivos "
    "estratégicos e seus indicadores/marcos. Sempre apresente o que vai fazer (ou o que vai "
    "excluir) e peça confirmação explícita do usuário antes de chamar uma ferramenta de escrita "
    "— em especial excluir_objetivo_estrategico, que é irreversível. Todas as demais ferramentas "
    "são somente-leitura.\n\n"
    "Regras de conduta:\n"
    "- Nunca bajule. Discorde com fundamento quando os fatos não sustentarem o pedido.\n"
    "- Toda afirmação de risco, custo ou prazo deve estar ancorada em dado concreto obtido via "
    "ferramenta. Se não houver dado suficiente, diga isso explicitamente em vez de estimar sem base.\n"
    "- Aponte o custo de oportunidade: onde o tempo analisado traria mais alavancagem.\n\n"
    "Regras de formatação (a leitura é majoritariamente em celular):\n"
    "- Nunca use tabelas Markdown — são ilegíveis em tela estreita. Prefira listas de tópicos, "
    "com o item em negrito seguido de uma explicação curta.\n"
    "- Frases curtas e diretas. Evite parágrafos longos."
)


def _get_persona(db) -> str:
    parts = [GODMODE_PERSONA]
    try:
        core_doc = db.collection("system").document("copilot_core").get()
        if core_doc.exists:
            content = (core_doc.to_dict() or {}).get("content")
            if content:
                parts.append(f"## CORE ESTÁTICO DO COPILOTO\n{content}")
    except Exception:
        pass
    return "\n\n".join(parts)


def _get_user_profile_text(db, uid: str | None) -> str:
    """Mesmo perfil completo (dados básicos + preferências + personalidade
    destilada semanalmente do diário pessoal) lido pelo copiloto web
    (`main.py:_format_ai_profile_for_prompt`) e pela ponte de voz
    (`context.py:_format_user_profile`) — ver `docs/okf/arquitetura/schema-firestore.md`."""
    if not uid:
        return "(usuário não autenticado)"
    from main import _format_ai_profile_for_prompt  # import tardio: evita import circular com main.py

    try:
        snap = db.collection("usuarios").document(uid).get()
        profile = (snap.to_dict() or {}).get("ai_profile") if snap.exists else {}
    except Exception:
        return "(perfil indisponível)"
    return _format_ai_profile_for_prompt(profile)


_GAPPS_EXPORT_MIME = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}

_PLAIN_TEXT_EXTENSIONS = (".txt", ".csv", ".md", ".markdown", ".json", ".xml", ".html", ".htm", ".eml")


def _extract_attached_file_text(drive_file_id: str, gemini_key: str | None) -> tuple[str, str]:
    """Baixa do Drive o arquivo anexado à mensagem e extrai seu texto. Retorna (nome_real, texto)."""
    from main import get_drive_service  # import tardio: evita import circular com main.py

    import io as _io

    from googleapiclient.http import MediaIoBaseDownload

    drive_service = get_drive_service()
    meta = drive_service.files().get(fileId=drive_file_id, fields="name,mimeType").execute()
    real_name = meta.get("name", "arquivo anexado")
    mime = meta.get("mimeType", "application/octet-stream")

    if mime in _GAPPS_EXPORT_MIME:
        request = drive_service.files().export_media(fileId=drive_file_id, mimeType=_GAPPS_EXPORT_MIME[mime])
    else:
        request = drive_service.files().get_media(fileId=drive_file_id)

    buffer = _io.BytesIO()
    downloader = MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    file_bytes = buffer.getvalue()

    lower_name = real_name.lower()
    if mime.startswith("image/"):
        text = (
            "(Este anexo é uma imagem — o Godmode ainda não realiza leitura visual/OCR. "
            "Peça ao usuário para descrever em texto o conteúdo relevante da imagem.)"
        )
    elif mime in _GAPPS_EXPORT_MIME or lower_name.endswith(_PLAIN_TEXT_EXTENSIONS):
        text = file_bytes.decode("utf-8", errors="replace")
    elif lower_name.endswith(".docx") or mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import mammoth

        text = mammoth.extract_raw_text(_io.BytesIO(file_bytes)).value or ""
    elif lower_name.endswith(".pptx") or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation

        prs = Presentation(_io.BytesIO(file_bytes))
        parts = [
            shape.text.strip()
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        text = "\n".join(parts)
    elif lower_name.endswith(".xlsx") or mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        import pandas as pd

        sheets = pd.read_excel(_io.BytesIO(file_bytes), sheet_name=None)
        text = "\n\n".join(f"ABA: {name}\n{sheet.to_csv(index=False)}" for name, sheet in sheets.items())
    elif lower_name.endswith(".pdf") or mime == "application/pdf":
        from pdf_precision import extract_pdf_text_with_fallback

        result = extract_pdf_text_with_fallback(
            file_bytes, real_name, api_key=gemini_key, allow_gemini_fallback=bool(gemini_key)
        )
        text = result.get("text") or ""
    else:
        text = ""

    return real_name, text.strip()


def _build_tools(
    db,
    user_uid: str | None,
    gemini_key: str | None,
    attached_drive_file_id: str | None = None,
    attached_file_name: str | None = None,
):
    """Registro de ferramentas do Godmode: leitura ampla sobre as superfícies do
    Hermes, e escrita restrita ao módulo de Estratégia (via strategy_tools)."""

    def consultar_tarefas(status: str = "", area_tematica: str = "", incluir_concluidas: bool = False, limite: int = 50) -> dict:
        limite = max(1, min(int(limite or 50), 150))
        query = db.collection("tarefas")
        # Só um filtro no Firestore por vez (evita depender de índice composto);
        # o restante do filtro é aplicado em memória.
        if area_tematica:
            query = query.where("area_tematica", "==", area_tematica)
        elif status:
            query = query.where("status", "==", status)

        tarefas = []
        for d in query.limit(limite * 3).stream():
            data = d.to_dict() or {}
            if status and data.get("status") != status:
                continue
            if not status and not incluir_concluidas and data.get("status") == "concluído":
                continue
            if area_tematica and data.get("area_tematica") != area_tematica:
                continue
            tarefas.append({
                "id": d.id,
                "titulo": data.get("titulo"),
                "status": data.get("status"),
                "area_tematica": data.get("area_tematica"),
                "projeto": data.get("projeto"),
                "data_limite": data.get("data_limite"),
                "prazo_final": data.get("prazo_final"),
                "contabilizar_meta": data.get("contabilizar_meta"),
                "tags": data.get("tags"),
            })
            if len(tarefas) >= limite:
                break
        return {"tarefas": tarefas}

    def consultar_metas_estrategicas(apenas_ativas: bool = True) -> dict:
        if not user_uid:
            return {"error": "Sem usuário autenticado — não é possível ler metas pessoais."}
        query = db.collection("estrategia_pessoal").where("userId", "==", user_uid)
        metas = [dict(d.to_dict() or {}, id=d.id) for d in query.limit(80).stream()]
        if apenas_ativas:
            metas = [
                m for m in metas
                if str(m.get("status", "")).lower() not in ("concluido", "concluído", "cancelado", "arquivado")
            ]
        return {"metas": metas}

    def criar_objetivo_estrategico(
        objetivoMacro: str,
        pilar: str = "carreira",
        tipoMeta: str = "relativa_qualitativa",
        status: str = "ativo",
        diretrizes: list | None = None,
        indicadores: list | None = None,
        marcos: list | None = None,
        metrica_valor_inicial: float | None = None,
        metrica_valor_atual: float | None = None,
        metrica_valor_objetivo: float | None = None,
        metrica_unidade: str = "",
    ) -> dict:
        return strategy_tools.criar_objetivo_estrategico(
            db, user_uid, objetivoMacro, pilar, tipoMeta, status,
            diretrizes, indicadores, marcos,
            metrica_valor_inicial, metrica_valor_atual, metrica_valor_objetivo, metrica_unidade,
        )

    def editar_objetivo_estrategico(
        objetivo_id: str,
        objetivoMacro: str | None = None,
        pilar: str | None = None,
        tipoMeta: str | None = None,
        status: str | None = None,
        diretrizes: list | None = None,
        metrica_valor_inicial: float | None = None,
        metrica_valor_atual: float | None = None,
        metrica_valor_objetivo: float | None = None,
        metrica_unidade: str | None = None,
    ) -> dict:
        return strategy_tools.editar_objetivo_estrategico(
            db, user_uid, objetivo_id, objetivoMacro, pilar, tipoMeta, status,
            diretrizes, metrica_valor_inicial, metrica_valor_atual, metrica_valor_objetivo, metrica_unidade,
        )

    def gerenciar_item_estrategico(
        objetivo_id: str,
        tipo: str,
        acao: str,
        descricao: str | None = None,
        item_id: str | None = None,
    ) -> dict:
        return strategy_tools.gerenciar_item_estrategico(db, user_uid, objetivo_id, tipo, acao, descricao, item_id)

    def excluir_objetivo_estrategico(objetivo_id: str) -> dict:
        return strategy_tools.excluir_objetivo_estrategico(db, user_uid, objetivo_id)

    def consultar_financas(mes: int | None = None, ano: int | None = None) -> dict:
        try:
            from tools.telegram_extended import execute as _execute_telegram_tool
            return json.loads(_execute_telegram_tool("consultar_financas_v2", {"mes": mes, "ano": ano}, db))
        except Exception as exc:
            return {"error": str(exc)}

    def consultar_saude(ultimos_dias: int = 7, data_especifica: str = "") -> dict:
        try:
            return build_health_summary(db, ultimos_dias, data_especifica or None)
        except Exception as exc:
            return {"error": str(exc)}

    def consultar_relatorio_semanal_saude(semana: str = "") -> dict:
        try:
            week_id = (semana or "").strip()
            if week_id:
                doc = db.collection("health_weekly_reports").document(week_id).get()
                if not doc.exists:
                    return {"error": f"Relatório semanal '{week_id}' não encontrado."}
                return dict(doc.to_dict() or {}, semana=week_id)

            # Sem semana especificada: pega o mais recente já gerado. IDs no formato
            # ISO 'YYYY-Www' ordenam corretamente por comparação lexicográfica (ano
            # primeiro) — ordenar pelo nome do documento resolve direto, sem depender
            # de caminhar semana a semana (o que devolveria "não encontrado" à toa se
            # o scheduler tivesse ficado parado por mais de uma semana).
            docs = list(
                db.collection("health_weekly_reports")
                .order_by("__name__", direction=firestore.Query.DESCENDING)
                .limit(1)
                .stream()
            )
            if not docs:
                return {"error": "Nenhum relatório semanal de saúde foi gerado ainda."}
            doc = docs[0]
            return dict(doc.to_dict() or {}, semana=doc.id)
        except Exception as exc:
            return {"error": str(exc)}

    def consultar_diario_pessoal(dias: int = 7) -> dict:
        try:
            from datetime import datetime as _dt, timedelta as _timedelta
            from zoneinfo import ZoneInfo as _ZoneInfo

            n = max(1, min(int(dias or 7), 30))
            today_local = _dt.now(_ZoneInfo("America/Sao_Paulo")).date()
            diarios = []
            for i in range(n):
                data_str = (today_local - _timedelta(days=i)).strftime("%Y-%m-%d")
                doc = db.collection("diario_pessoal").document(data_str).get()
                if not doc.exists:
                    continue
                d = doc.to_dict() or {}
                if d.get("sem_material") or not d.get("texto"):
                    continue
                diarios.append({
                    "data": d.get("data", data_str),
                    "texto": d.get("texto"),
                    "fontes": d.get("fontes"),
                    "editado": bool(d.get("editado")),
                })
            return {"diarios": diarios}
        except Exception as exc:
            return {"error": str(exc)}

    def consultar_agenda(data_inicio: str, data_fim: str) -> dict:
        # `google_calendar_events.data_inicio` é gravado cru a partir da API do Calendar
        # (main.py, `event['start'].get('dateTime', ...)`) — offset LOCAL do evento (ex.
        # "-03:00"), nunca normalizado para UTC, ou uma data pura em eventos de dia inteiro.
        # Comparar essas strings direto contra data_inicio/data_fim por ordem lexicográfica
        # é incorreto perto das bordas do intervalo (mesmo problema já corrigido para
        # data_fim em email_action_linker.py:CALENDAR_QUERY_SLACK_MINUTES) — a query abaixo
        # é só um pré-filtro barato com folga de 1 dia; quem decide de fato é a data local
        # recalculada a partir do datetime normalizado.
        try:
            from datetime import datetime as _dt, timedelta as _timedelta
            from zoneinfo import ZoneInfo as _ZoneInfo
            from main import parse_iso_datetime

            tz = _ZoneInfo("America/Sao_Paulo")
            prefilter_start = (_dt.strptime(data_inicio, "%Y-%m-%d") - _timedelta(days=1)).strftime("%Y-%m-%d")
            prefilter_end = (_dt.strptime(data_fim, "%Y-%m-%d") + _timedelta(days=1)).strftime("%Y-%m-%d") + "T23:59:59"

            eventos = []
            query = (
                db.collection("google_calendar_events")
                .where("data_inicio", ">=", prefilter_start)
                .where("data_inicio", "<=", prefilter_end)
            )
            for d in query.stream():
                ev = d.to_dict() or {}
                raw_inicio = ev.get("data_inicio")
                start_dt = parse_iso_datetime(raw_inicio)
                if start_dt is None:
                    continue
                # Com horário: normaliza para o fuso local antes de extrair a data.
                # Sem horário (dia inteiro): a própria string já é a data local.
                data_local = start_dt.astimezone(tz).strftime("%Y-%m-%d") if start_dt.tzinfo else str(raw_inicio)[:10]
                if not (data_inicio <= data_local <= data_fim):
                    continue
                eventos.append({
                    "titulo": ev.get("titulo"),
                    "inicio": ev.get("data_inicio"),
                    "fim": ev.get("data_fim"),
                    "criado_pelo_hermes": ev.get("criado_pelo_hermes"),
                })
            eventos.sort(key=lambda x: x.get("inicio") or "")
            return {"eventos": eventos}
        except Exception as exc:
            return {"error": str(exc)}

    def buscar_contato(termo: str, limite: int = 5) -> dict:
        try:
            termo_lower = (termo or "").strip().lower()
            if not termo_lower:
                return {"error": "Termo de busca vazio."}
            candidatos = []
            for d in db.collection("perfil_pessoas").limit(500).stream():
                pdata = d.to_dict() or {}
                nome = (pdata.get("nome") or "").lower()
                email = (pdata.get("email") or "").lower()
                tags = [str(t).lower() for t in (pdata.get("tags") or [])]
                score = 0.0
                if nome == termo_lower or email == termo_lower:
                    score = 1.0
                elif termo_lower in nome:
                    score = 0.8
                elif termo_lower in email:
                    score = 0.7
                elif any(termo_lower in t for t in tags):
                    score = 0.5
                if score > 0:
                    candidatos.append({
                        "pessoa_id": d.id,
                        "nome": pdata.get("nome", ""),
                        "email": pdata.get("email", ""),
                        "tags": pdata.get("tags", []),
                        "score": score,
                    })
            candidatos.sort(key=lambda x: -x["score"])
            return {"candidatos": candidatos[: max(1, int(limite or 5))]}
        except Exception as exc:
            return {"error": str(exc)}

    def consultar_interacoes_pessoa(pessoa_id: str, limite: int = 20) -> dict:
        try:
            if not pessoa_id:
                return {"error": "pessoa_id é obrigatório."}
            interacoes = [
                dict(d.to_dict() or {}, id=d.id)
                for d in db.collection("interacoes_pessoas").where("pessoa_id", "==", pessoa_id).limit(max(1, min(int(limite or 20), 60))).stream()
            ]
            interacoes.sort(key=lambda x: str(x.get("data") or ""), reverse=True)
            return {"interacoes": interacoes}
        except Exception as exc:
            return {"error": str(exc)}

    def buscar_conversas_whatsapp(query: str, limite: int = 5) -> dict:
        try:
            from whatsapp_ingest import buscar_conversas_whatsapp as _buscar_whatsapp
            return _buscar_whatsapp(db, query, limite)
        except Exception as exc:
            return {"error": str(exc)}

    def consultar_dados_cadastrais() -> dict:
        try:
            from dados_cadastrais import get_dados_cadastrais
            return get_dados_cadastrais(db, user_uid)
        except Exception as exc:
            return {"error": str(exc)}

    def buscar_conhecimento(consulta: str, area_tematica: str = "", tags: list | None = None) -> dict:
        if not gemini_key:
            return {"error": "RAG indisponível: chave Gemini não configurada (necessária para embeddings)."}
        try:
            from knowledge_graph import extract_kg_rag_context
            nodes, context_str = extract_kg_rag_context(
                db=db,
                api_key=gemini_key,
                area_tematica=area_tematica or consulta,
                tags=tags or [],
            )
            return {"nodes": nodes[:20], "contexto": context_str[:6000]}
        except Exception as exc:
            return {"error": str(exc)}

    def ler_arquivo_anexado(query_especifica: str = "") -> dict:
        if not attached_drive_file_id:
            return {"erro": "Nenhum arquivo foi anexado a esta mensagem."}
        try:
            real_name, text = _extract_attached_file_text(attached_drive_file_id, gemini_key)
        except Exception as exc:
            return {"erro": f"Falha ao ler o arquivo anexado ({attached_file_name or attached_drive_file_id}): {exc}"}
        if not text:
            return {"arquivo": real_name, "aviso": "Não foi possível extrair texto deste arquivo."}
        return {"arquivo": real_name, "conteudo": text[:7000]}

    function_map = {
        "consultar_tarefas": consultar_tarefas,
        "consultar_metas_estrategicas": consultar_metas_estrategicas,
        "criar_objetivo_estrategico": criar_objetivo_estrategico,
        "editar_objetivo_estrategico": editar_objetivo_estrategico,
        "gerenciar_item_estrategico": gerenciar_item_estrategico,
        "excluir_objetivo_estrategico": excluir_objetivo_estrategico,
        "consultar_financas": consultar_financas,
        "consultar_saude": consultar_saude,
        "consultar_relatorio_semanal_saude": consultar_relatorio_semanal_saude,
        "consultar_diario_pessoal": consultar_diario_pessoal,
        "consultar_agenda": consultar_agenda,
        "buscar_contato": buscar_contato,
        "consultar_interacoes_pessoa": consultar_interacoes_pessoa,
        "buscar_conversas_whatsapp": buscar_conversas_whatsapp,
        "consultar_dados_cadastrais": consultar_dados_cadastrais,
        "buscar_conhecimento": buscar_conhecimento,
    }
    if attached_drive_file_id:
        function_map["ler_arquivo_anexado"] = ler_arquivo_anexado

    tools_schema = [
        {
            "name": "consultar_tarefas",
            "description": (
                "Lista tarefas/ações do usuário (título, status, área temática, prazos). "
                "Por padrão oculta tarefas concluídas — use incluir_concluidas para trazê-las. "
                "Atenção: tarefas com campo 'projeto' preenchido tratam de projetos institucionais, "
                "não de assuntos pessoais do usuário."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "status": {
                        "type": "string",
                        "enum": ["em andamento", "stand-by", "concluído"],
                        "description": "Filtra por status exato. Deixe vazio para todos.",
                    },
                    "area_tematica": {
                        "type": "string",
                        "description": "Filtra por área temática exata, se conhecida.",
                    },
                    "incluir_concluidas": {
                        "type": "boolean",
                        "description": "Se true, inclui tarefas com status 'concluído' quando nenhum status específico for informado.",
                    },
                    "limite": {
                        "type": "integer",
                        "description": "Número máximo de tarefas retornadas (padrão 50, máximo 150).",
                    },
                },
            },
        },
        {
            "name": "consultar_metas_estrategicas",
            "description": (
                "Lista os objetivos/metas estratégicas pessoais do usuário autenticado "
                "(equivalente a metas de PGD/plano de gestão)."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "apenas_ativas": {
                        "type": "boolean",
                        "description": "Se true (padrão), oculta metas concluídas/canceladas/arquivadas.",
                    },
                },
            },
        },
        {
            "name": "criar_objetivo_estrategico",
            "description": (
                "[ESCRITA] Cria um novo objetivo estratégico pessoal em estrategia_pessoal. "
                "Use APENAS quando o usuário pedir explicitamente para criar/cadastrar um objetivo, "
                "meta ou pilar estratégico. Apresente um rascunho ao usuário e só chame após "
                "confirmação explícita."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "objetivoMacro": {"type": "string", "description": "Enunciado do objetivo macro (obrigatório)."},
                    "pilar": {
                        "type": "string",
                        "enum": ["carreira", "financas", "saude", "intelectual", "estilo_vida"],
                        "description": "Pilar estratégico. Padrão: carreira.",
                    },
                    "tipoMeta": {
                        "type": "string",
                        "enum": ["absoluta", "relativa_qualitativa"],
                        "description": "'absoluta' (com métrica numérica) ou 'relativa_qualitativa'.",
                    },
                    "status": {
                        "type": "string",
                        "enum": ["ativo", "revisar", "concluido"],
                        "description": "Padrão: ativo.",
                    },
                    "diretrizes": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Diretrizes derivadas (frases que orientam a IA) — obrigatório ao menos uma.",
                    },
                    "indicadores": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Descrições de indicadores contínuos de sucesso.",
                    },
                    "marcos": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Descrições de marcos pontuais.",
                    },
                    "metrica_valor_inicial": {"type": "number", "description": "Só para tipoMeta 'absoluta'."},
                    "metrica_valor_atual": {"type": "number", "description": "Só para tipoMeta 'absoluta'."},
                    "metrica_valor_objetivo": {"type": "number", "description": "Só para tipoMeta 'absoluta'."},
                    "metrica_unidade": {"type": "string", "description": "Só para tipoMeta 'absoluta'."},
                },
                "required": ["objetivoMacro", "diretrizes"],
            },
        },
        {
            "name": "editar_objetivo_estrategico",
            "description": (
                "[ESCRITA] Edita um objetivo estratégico existente (objetivo_id visível em "
                "consultar_metas_estrategicas). Só passe os campos que devem mudar. Para indicadores "
                "ou marcos individuais use gerenciar_item_estrategico. Use APENAS após confirmação do usuário."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "objetivo_id": {"type": "string", "description": "ID do objetivo (obrigatório)."},
                    "objetivoMacro": {"type": "string"},
                    "pilar": {"type": "string", "enum": ["carreira", "financas", "saude", "intelectual", "estilo_vida"]},
                    "tipoMeta": {"type": "string", "enum": ["absoluta", "relativa_qualitativa"]},
                    "status": {"type": "string", "enum": ["ativo", "revisar", "concluido"]},
                    "diretrizes": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Se fornecida, SUBSTITUI a lista completa de diretrizes.",
                    },
                    "metrica_valor_inicial": {"type": "number"},
                    "metrica_valor_atual": {"type": "number"},
                    "metrica_valor_objetivo": {"type": "number"},
                    "metrica_unidade": {"type": "string"},
                },
                "required": ["objetivo_id"],
            },
        },
        {
            "name": "gerenciar_item_estrategico",
            "description": (
                "[ESCRITA] Gerencia um indicador ou marco dentro de um objetivo estratégico, preservando "
                "IDs. Use APENAS após confirmação do usuário."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "objetivo_id": {"type": "string"},
                    "tipo": {"type": "string", "enum": ["indicador", "marco"]},
                    "acao": {
                        "type": "string",
                        "enum": ["adicionar", "editar", "remover", "concluir"],
                        "description": "'adicionar' precisa de descricao; 'editar' precisa item_id+descricao; 'remover'/'concluir' precisam item_id.",
                    },
                    "descricao": {"type": "string"},
                    "item_id": {"type": "string", "description": "ID do item existente (visível no snapshot da meta)."},
                },
                "required": ["objetivo_id", "tipo", "acao"],
            },
        },
        {
            "name": "excluir_objetivo_estrategico",
            "description": (
                "[ESCRITA — IRREVERSÍVEL] Exclui definitivamente um objetivo estratégico e seus "
                "indicadores/marcos/diretrizes. Use APENAS após confirmação explícita e inequívoca do usuário."
            ),
            "input_schema": {
                "type": "object",
                "properties": {"objetivo_id": {"type": "string"}},
                "required": ["objetivo_id"],
            },
        },
        {
            "name": "consultar_financas",
            "description": (
                "Resumo financeiro unificado do mês (rendas, contas fixas, metas financeiras, reserva de "
                "emergência e balancete previsto/atual). Sem argumentos, usa o mês corrente."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "mes": {"type": "integer", "description": "Mês 0-11 (0=janeiro). Padrão: mês atual."},
                    "ano": {"type": "integer", "description": "Ano (YYYY). Padrão: ano atual."},
                },
            },
        },
        {
            "name": "consultar_saude",
            "description": (
                "Telemetria diária de saúde: peso, caminhada/passos (com nível frente à meta mínima/ideal), "
                "calorias, sono e dor. Cobre um intervalo de dias ou um dia específico."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "ultimos_dias": {"type": "integer", "description": "Número de dias a consultar (padrão 7, máximo 30)."},
                    "data_especifica": {"type": "string", "description": "YYYY-MM-DD — sobrepõe ultimos_dias."},
                },
            },
        },
        {
            "name": "consultar_relatorio_semanal_saude",
            "description": (
                "Relatório semanal de saúde já computado em código (placa de resultado, regra de ajuste "
                "decidida e auditoria do ajuste da semana anterior — o texto só redige por cima de números "
                "já calculados, não é opinião do modelo). Sem argumento, traz a semana mais recente disponível."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "semana": {"type": "string", "description": "Formato ISO 'YYYY-Www' (ex.: '2026-W33'). Deixe vazio para a mais recente."},
                },
            },
        },
        {
            "name": "consultar_diario_pessoal",
            "description": (
                "Lê o diário pessoal do usuário (texto em primeira pessoa gerado diariamente a partir de "
                "todas as superfícies do Hermes — ações, saúde, finanças, agenda, conversas, pessoas). É a "
                "fonte mais rica para entender o panorama e o estado recente do usuário."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "dias": {"type": "integer", "description": "Quantos dias recentes trazer (padrão 7, máximo 30)."},
                },
            },
        },
        {
            "name": "consultar_agenda",
            "description": "Lista eventos do Google Calendar sincronizados num intervalo de datas.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "data_inicio": {"type": "string", "description": "YYYY-MM-DD (obrigatório)."},
                    "data_fim": {"type": "string", "description": "YYYY-MM-DD (obrigatório)."},
                },
                "required": ["data_inicio", "data_fim"],
            },
        },
        {
            "name": "buscar_contato",
            "description": "Busca contatos por nome, e-mail ou tag em perfil_pessoas.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "termo": {"type": "string", "description": "Nome, e-mail ou tag a buscar."},
                    "limite": {"type": "integer", "description": "Padrão 5."},
                },
                "required": ["termo"],
            },
        },
        {
            "name": "consultar_interacoes_pessoa",
            "description": "Histórico de interações registradas com uma pessoa específica (reuniões, menções em tarefas/diário/copiloto, WhatsApp).",
            "input_schema": {
                "type": "object",
                "properties": {
                    "pessoa_id": {"type": "string", "description": "ID da pessoa (obtido via buscar_contato)."},
                    "limite": {"type": "integer", "description": "Padrão 20, máximo 60."},
                },
                "required": ["pessoa_id"],
            },
        },
        {
            "name": "buscar_conversas_whatsapp",
            "description": "Busca semântica nas conversas de WhatsApp indexadas (digests). Use quando a pergunta envolver algo discutido no WhatsApp.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Texto da busca."},
                    "limite": {"type": "integer", "description": "Padrão 5."},
                },
                "required": ["query"],
            },
        },
        {
            "name": "consultar_dados_cadastrais",
            "description": (
                "Consulta os dados cadastrais pessoais completos do usuário (documentos — CPF, RG, "
                "título de eleitor, PIS/PASEP, CTPS, CNH —, contato, família, formação acadêmica, "
                "carreira, dados bancários, plano de saúde etc.). Use quando o usuário pedir ajuda "
                "para preencher um formulário/documento oficial, ou perguntar um dado cadastral "
                "específico que esqueceu."
            ),
            "input_schema": {"type": "object", "properties": {}},
        },
        {
            "name": "buscar_conhecimento",
            "description": (
                "Busca semântica (RAG) na base de conhecimento/procedimentos já registrada no Hermes "
                "(knowledge_graph) para embasar a análise com precedentes reais."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "consulta": {"type": "string", "description": "Texto da busca."},
                    "area_tematica": {
                        "type": "string",
                        "description": "Área temática para restringir a busca, se conhecida.",
                    },
                    "tags": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Tags adicionais para refinar a busca.",
                    },
                },
                "required": ["consulta"],
            },
        },
    ]

    if attached_drive_file_id:
        tools_schema.append({
            "name": "ler_arquivo_anexado",
            "description": (
                f"Lê o conteúdo do arquivo '{attached_file_name or 'anexo'}' que o usuário anexou "
                "a esta mensagem e retorna seu texto extraído. Use sempre que a pergunta se referir "
                "ao conteúdo desse anexo antes de responder."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "query_especifica": {
                        "type": "string",
                        "description": "O que exatamente você precisa encontrar no arquivo (opcional, ajuda a focar a leitura).",
                    },
                },
            },
        })

    return tools_schema, function_map


def _load_history(db, session_id: str) -> list[dict]:
    if not session_id:
        return []
    docs = (
        db.collection("sessoes_godmode")
        .document(session_id)
        .collection("mensagens")
        .order_by("timestamp")
        .limit(40)
        .stream()
    )
    history = []
    for d in docs:
        data = d.to_dict() or {}
        role = "assistant" if data.get("role") == "assistant" else "user"
        content = data.get("content") or ""
        if content:
            history.append({"role": role, "content": content})
    return history


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.GB_2,
    timeout_sec=GODMODE_FUNCTION_TIMEOUT_SEC,
)
def askHermesGodmode(req: https_fn.CallableRequest):
    data = req.data or {}
    prompt = (data.get("prompt") or "").strip()
    session_id = (data.get("sessionId") or "").strip()
    drive_file_id = (data.get("driveFileId") or "").strip() or None
    drive_file_name = (data.get("driveFileName") or "").strip() or None
    drive_files_raw = data.get("driveFiles") or []
    drive_files = []
    if isinstance(drive_files_raw, list) and drive_files_raw:
        for df in drive_files_raw[:10]:
            if isinstance(df, dict) and df.get("driveFileId"):
                drive_files.append({
                    "driveFileId": str(df["driveFileId"]).strip(),
                    "driveFileName": str(df.get("driveFileName") or "documento").strip()
                })
    elif drive_file_id:
        drive_files.append({"driveFileId": drive_file_id, "driveFileName": drive_file_name or "documento"})

    first_file_id = drive_files[0]["driveFileId"] if drive_files else None
    first_file_name = drive_files[0]["driveFileName"] if drive_files else None

    from main import _require_internal_user  # import tardio: evita import circular com main.py

    # Godmode roda com o Admin SDK, que ignora firestore.rules — sem esta checagem,
    # bastaria qualquer conta Firebase autenticada (não o dono verificado que as
    # regras exigem para estas coleções) para ler tarefas, finanças, saúde, diário,
    # agenda, pessoas e WhatsApp através dele.
    _require_internal_user(req)
    user_uid = req.auth.uid

    if not prompt and not drive_files:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Prompt é obrigatório.",
        )
    if not prompt:
        names = ", ".join(f"'{df['driveFileName']}'" for df in drive_files)
        prompt = f"Arquivos anexados: {names}"

    try:
        import anthropic
    except ImportError:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            message="Dependência 'anthropic' não instalada no ambiente de functions.",
        )

    db = firestore.client()
    keys_doc = db.collection("system").document("api_keys").get()
    keys_data = keys_doc.to_dict() if keys_doc.exists else {}
    claude_key = (keys_data or {}).get("claude_api_key")
    gemini_key = (keys_data or {}).get("gemini_api_key")

    if not claude_key:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            message="Chave Claude não configurada (system/api_keys.claude_api_key).",
        )

    if not session_id:
        session_ref = db.collection("sessoes_godmode").document()
        session_id = session_ref.id
        session_ref.set(
            {
                "userId": user_uid,
                "createdAt": firestore.SERVER_TIMESTAMP,
                "lastMessageAt": firestore.SERVER_TIMESTAMP,
                "titulo": prompt[:80],
            }
        )
    else:
        session_ref = db.collection("sessoes_godmode").document(session_id)

    mensagens_ref = session_ref.collection("mensagens")
    # Carrega o histórico ANTES de gravar a mensagem atual do usuário — o loop
    # de tool-calling já injeta `prompt` como o turno mais recente, então
    # gravar antes duplicaria esse turno (e quebraria a alternância obrigatória
    # de roles da Claude Messages API).
    history = _load_history(db, session_id)
    mensagens_ref.add(
        {"role": "user", "content": prompt, "timestamp": firestore.SERVER_TIMESTAMP}
    )

    client = anthropic.Anthropic(api_key=claude_key)
    system_instruction = _get_persona(db) + "\n\n## PERFIL DO USUÁRIO\n" + _get_user_profile_text(db, user_uid)
    tools, function_map = _build_tools(db, user_uid, gemini_key, first_file_id, first_file_name)

    llm_user_message = prompt
    if drive_files:
        file_notices = []
        for df in drive_files:
            file_notices.append(f"'{df['driveFileName']}' (ID: {df['driveFileId']})")
        llm_user_message = (
            f"[O usuário anexou os seguintes arquivos a esta mensagem: {', '.join(file_notices)}. "
            "Use a ferramenta ler_arquivo_anexado para consultar o conteúdo antes de responder.]\n\n"
            f"{prompt}"
        )

    try:
        result = claude_provider.run_tool_loop(
            client=client,
            model=GODMODE_MODEL,
            system_instruction=system_instruction,
            tools=tools,
            function_map=function_map,
            history=history,
            user_message=llm_user_message,
            max_tokens=GODMODE_MAX_TOKENS,
            fallback_model=GODMODE_FALLBACK_MODEL,
        )
    except anthropic.APIStatusError as exc:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.UNAVAILABLE,
            message=f"Erro na Claude API (modelo principal e fallback indisponíveis): {exc}",
        )

    mensagens_ref.add(
        {
            "role": "assistant",
            "content": result["text"],
            "toolsUsed": result["tools_used"],
            "modelUsed": result["model_used"],
            "timestamp": firestore.SERVER_TIMESTAMP,
        }
    )
    session_ref.set({"lastMessageAt": firestore.SERVER_TIMESTAMP}, merge=True)

    return {
        "result": result["text"],
        "sessionId": session_id,
        "toolsUsed": result["tools_used"],
        "usage": result["usage"],
        "modelUsed": result["model_used"],
        "fallbackUsed": result["fallback_used"],
    }
