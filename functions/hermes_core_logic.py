"""
Hermes Core Logic — Telegram Integration
Webhook receiver + Firestore-triggered async processor.
"""
import json
import html
import os
import re
import base64
import tempfile
import threading
import time
import unicodedata
import wave
from contextlib import contextmanager
from datetime import datetime, timedelta, timezone
from io import BytesIO
from typing import Optional

import requests as _requests
from firebase_admin import firestore, get_app, initialize_app, storage
from firebase_functions import firestore_fn, https_fn, options
from firebase_functions.firestore_fn import Event, Change, DocumentSnapshot
from google.cloud.firestore_v1 import DocumentReference
from gemini_cost_controls import (
    GEMINI_BALANCED_MODEL,
    GEMINI_TRANSCRIPTION_MODEL,
    GEMINI_TTS_MODEL,
    generate_content_logged,
    send_message_logged,
)

try:
    get_app()
except ValueError:
    initialize_app()

_MAX_HISTORY_TURNS = 20
_MAX_FILE_BYTES = 20 * 1024 * 1024  # 20 MB
_MAX_INLINE_MEDIA_BYTES = 700 * 1024  # base64 stays safely below Firestore field limit
_TTS_MODEL_ID = os.environ.get("TELEGRAM_TTS_MODEL", GEMINI_TTS_MODEL)
_TEXT_MODEL_ID = os.environ.get("TELEGRAM_TEXT_MODEL", GEMINI_BALANCED_MODEL)
_DEFAULT_MALE_VOICE = "Charon"
_MAX_TTS_TRANSCRIPT_CHARS = 1500

# ---------------------------------------------------------------------------
# In-process Firestore document cache (TTL = 60 s per Cloud Function instance)
# ---------------------------------------------------------------------------
_DOC_CACHE: dict = {}
_DOC_CACHE_TTL = 60
_ACTION_SNAPSHOT_CACHE: dict = {}
_ACTION_SNAPSHOT_CACHE_TTL = 45
_TELEGRAM_USER_CACHE: dict = {}
_TELEGRAM_USER_CACHE_TTL = 300
_COPILOT_SESSION_HISTORY_LIMIT = 12

def _cached_doc_get(db, collection: str, document: str):
    key = f"{collection}/{document}"
    now = time.monotonic()
    cached = _DOC_CACHE.get(key)
    if cached and (now - cached[0]) < _DOC_CACHE_TTL:
        return cached[1]
    doc = db.collection(collection).document(document).get()
    _DOC_CACHE[key] = (now, doc)
    return doc


def _safe_telegram_session_id(chat_id: str) -> str:
    safe_chat_id = re.sub(r"[^a-zA-Z0-9_-]+", "_", str(chat_id or "").strip())
    return f"telegram_{safe_chat_id or 'unknown'}"


def _resolve_user_id_for_telegram_chat(db, chat_id: str) -> str | None:
    now = time.monotonic()
    cached = _TELEGRAM_USER_CACHE.get(chat_id)
    if cached and (now - cached[0]) < _TELEGRAM_USER_CACHE_TTL:
        return cached[1]

    user_id = None
    try:
        chat_candidates = [str(chat_id)]
        try:
            chat_candidates.append(int(chat_id))
        except Exception:
            pass
        for field in ("telegram_chat_id", "telegramChatId", "chat_id", "telegram_id"):
            for candidate in chat_candidates:
                docs = (
                    db.collection("usuarios")
                    .where(field, "==", candidate)
                    .limit(1)
                    .get()
                )
                if docs:
                    user_id = docs[0].id
                    break
            if user_id:
                break
    except Exception as exc:
        print(f"[SessionBridge] Falha ao resolver usuario do Telegram chat_id={chat_id}: {exc}")

    _TELEGRAM_USER_CACHE[chat_id] = (now, user_id)
    return user_id


def _ensure_copilot_session(db, chat_id: str, session: dict, first_text: str = "", from_user: dict | None = None) -> str:
    session_id = session.get("copilot_session_id") or _safe_telegram_session_id(chat_id)
    session["copilot_session_id"] = session_id

    session_ref = db.collection("sessoes_copiloto").document(session_id)
    user_id = session.get("userId") or _resolve_user_id_for_telegram_chat(db, chat_id)
    if user_id:
        session["userId"] = user_id

    from_user = from_user or {}
    try:
        snap = session_ref.get()
        base = {
            "channel": "telegram",
            "telegramChatId": str(chat_id),
            "lastMessageAt": firestore.SERVER_TIMESTAMP,
        }
        if user_id:
            base["userId"] = user_id
        if from_user:
            base["telegramUser"] = {
                "id": from_user.get("id"),
                "username": from_user.get("username"),
                "first_name": from_user.get("first_name"),
                "last_name": from_user.get("last_name"),
            }
        if not snap.exists:
            title_seed = (first_text or "Conversa via Telegram").strip().replace("\n", " ")
            base.update({
                "title": (title_seed[:40] + "...") if len(title_seed) > 40 else title_seed,
                "createdAt": firestore.SERVER_TIMESTAMP,
                "taskId": session.get("acao_id") or None,
                "systemId": None,
            })
        session_ref.set(base, merge=True)
    except Exception as exc:
        print(f"[SessionBridge] Falha ao garantir sessao do Copiloto {session_id}: {exc}")

    return session_id


def _persist_copilot_message(
    db,
    session_id: str | None,
    role: str,
    content: str,
    *,
    source: str = "telegram",
    tools_used: list | None = None,
    extra: dict | None = None,
):
    if not session_id or not content:
        return
    try:
        payload = {
            "role": role,
            "content": content,
            "source": source,
            "timestamp": firestore.SERVER_TIMESTAMP,
        }
        if tools_used:
            payload["toolsUsed"] = tools_used
        if extra:
            payload.update(extra)
        db.collection("sessoes_copiloto").document(session_id).collection("mensagens").add(payload)
        db.collection("sessoes_copiloto").document(session_id).set({
            "lastMessageAt": firestore.SERVER_TIMESTAMP,
            "channel": "telegram",
        }, merge=True)
    except Exception as exc:
        print(f"[SessionBridge] Falha ao persistir mensagem na sessao {session_id}: {exc}")


def _load_copilot_session_history(db, session_id: str | None, types, limit: int = _COPILOT_SESSION_HISTORY_LIMIT) -> list:
    if not session_id:
        return []
    try:
        msg_docs = (
            db.collection("sessoes_copiloto").document(session_id)
            .collection("mensagens")
            .order_by("timestamp", direction=firestore.Query.DESCENDING)
            .limit(limit)
            .get()
        )
        history = []
        for doc in reversed(list(msg_docs)):
            data = doc.to_dict() or {}
            content = (data.get("content") or "").strip()
            if not content:
                continue
            role = data.get("role")
            if role == "assistant":
                role = "model"
            if role not in ("user", "model"):
                continue
            history.append(types.Content(role=role, parts=[types.Part(text=content)]))
        return history
    except Exception as exc:
        print(f"[SessionBridge] Falha ao carregar historico da sessao {session_id}: {exc}")
        return []


def _get_hermes_storage_bucket():
    project_id = (
        os.environ.get("GCLOUD_PROJECT")
        or os.environ.get("GCP_PROJECT")
        or os.environ.get("GOOGLE_CLOUD_PROJECT")
        or "gestao-hermes"
    )
    candidates = [
        os.environ.get("FIREBASE_STORAGE_BUCKET"),
        f"{project_id}.firebasestorage.app",
        f"{project_id}.appspot.com",
        os.environ.get("TELEGRAM_MEDIA_BUCKET"),
        os.environ.get("HERMES_STORAGE_BUCKET"),
        f"{project_id}-telegram-media-us-central1",
        f"{project_id}-slides-us-central1",
    ]

    checked = []
    for bucket_name in [name for name in candidates if name]:
        bucket = storage.bucket(bucket_name)
        checked.append(bucket_name)
        try:
            if bucket.exists():
                return bucket
        except Exception as exc:
            print(f"[Storage] Falha ao verificar bucket {bucket_name}: {exc}")

    raise RuntimeError(
        "Nenhum bucket de Storage disponivel para midias do Telegram. "
        f"Buckets testados: {', '.join(checked)}"
    )


def _format_web_copilot_text_for_telegram(text: str) -> str:
    safe = html.escape(text or "")
    safe = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", safe, flags=re.DOTALL)
    safe = re.sub(r"__(.+?)__", r"<b>\1</b>", safe, flags=re.DOTALL)
    safe = re.sub(r"\*(.+?)\*", r"<i>\1</i>", safe, flags=re.DOTALL)
    safe = re.sub(r"\[(.*?)\]\((task:[^)]+)\)", r"\1 (\2)", safe)
    safe = re.sub(r"\[(.*?)\]\((tool:diagnosis:[^)]+)\)", r"\1 (\2)", safe)
    return safe.strip()


def _call_web_callable(
    *,
    function_name: str,
    data: dict,
    user_uid: str | None = None,
    timeout: int = 480,
) -> dict:
    project_id = (
        os.environ.get("GCLOUD_PROJECT")
        or os.environ.get("GCP_PROJECT")
        or os.environ.get("GOOGLE_CLOUD_PROJECT")
        or "gestao-hermes"
    )
    region = os.environ.get("FUNCTION_REGION") or os.environ.get("GOOGLE_CLOUD_REGION") or "us-central1"
    env_key = re.sub(r"[^A-Z0-9]+", "_", function_name.upper()) + "_URL"
    legacy_env_key = "ASK_COPILOTO_HERMES_URL" if function_name == "askCopilotoHermes" else None
    url = (
        (os.environ.get(legacy_env_key) if legacy_env_key else None)
        or os.environ.get(env_key)
        or f"https://{region}-{project_id}.cloudfunctions.net/{function_name}"
    )

    payload = {"data": data}
    headers = {"Content-Type": "application/json"}
    if user_uid:
        headers["X-Hermes-User-Uid"] = user_uid

    resp = _requests.post(url, json=payload, headers=headers, timeout=timeout)
    if not resp.ok:
        raise RuntimeError(f"{function_name} HTTP {resp.status_code}: {resp.text[:300]}")

    response_data = resp.json()
    if response_data.get("error"):
        raise RuntimeError(json.dumps(response_data["error"], ensure_ascii=False))
    result = response_data.get("result", response_data)
    return result if isinstance(result, dict) else {"result": str(result or "")}


def _run_web_copilot_engine(
    *,
    prompt: str,
    session_id: str,
    task_id: str | None = None,
    system_id: str | None = None,
    user_uid: str | None = None,
) -> dict:
    return _call_web_callable(
        function_name="askCopilotoHermes",
        data={
            "prompt": prompt,
            "sessionId": session_id,
            "taskId": task_id,
            "systemId": system_id,
            "routingIndex": [],
        },
        user_uid=user_uid,
        timeout=480,
    )


def _extract_task_ids_from_text(text: str) -> list[str]:
    seen = set()
    task_ids = []
    for match in re.finditer(r"\btask:([a-zA-Z0-9_-]{8,})", text or ""):
        task_id = match.group(1).strip()
        if task_id not in seen:
            seen.add(task_id)
            task_ids.append(task_id)
    return task_ids


def _find_latest_copilot_card_message_id(db, session_id: str | None, field_name: str) -> str | None:
    if not session_id:
        return None
    try:
        docs = (
            db.collection("sessoes_copiloto").document(session_id)
            .collection("mensagens")
            .order_by("timestamp", direction=firestore.Query.DESCENDING)
            .limit(12)
            .get()
        )
        for doc in docs:
            data = doc.to_dict() or {}
            if data.get(field_name):
                return doc.id
    except Exception as exc:
        print(f"[TelegramCards] Falha ao localizar card {field_name}: {exc}")
    return None


def _set_latest_copilot_card_status(db, session_id: str | None, field_name: str, status: str):
    message_id = _find_latest_copilot_card_message_id(db, session_id, field_name)
    if not message_id:
        return
    try:
        db.collection("sessoes_copiloto").document(session_id).collection("mensagens").document(message_id).update({
            f"{field_name}.status": status,
            f"{field_name}.status_ui": status,
        })
    except Exception as exc:
        print(f"[TelegramCards] Falha ao atualizar status {field_name}: {exc}")


def _summarize_pending_edit(pending_edit: dict) -> str:
    titulo = html.escape(str(pending_edit.get("titulo") or pending_edit.get("task_id") or "acao"))
    lines = ["", "<b>Card de edicao pendente</b>", f"Acao: {titulo}"]
    alteracoes = pending_edit.get("alteracoes") or {}
    for campo, change in list(alteracoes.items())[:5]:
        change = change if isinstance(change, dict) else {}
        original = html.escape(str(change.get("original") or "-"))
        novo = html.escape(str(change.get("novo") or change.get("novo_raw") or "-"))
        lines.append(f"- {html.escape(str(campo))}: {original} -> {novo}")
    if len(alteracoes) > 5:
        lines.append(f"- mais {len(alteracoes) - 5} alteracao(oes)")
    lines.append("No Telegram, use os botoes abaixo para confirmar ou cancelar.")
    return "\n".join(lines)


def _summarize_pending_batch_reschedule(batch_reschedule: dict) -> str:
    items = batch_reschedule.get("items") or []
    lines = ["", "<b>Card de reagendamento em lote</b>", f"Acoes afetadas: {len(items)}"]
    justificativa = batch_reschedule.get("justificativa")
    if justificativa:
        lines.append(f"Motivo: {html.escape(str(justificativa)[:300])}")
    for item in items[:6]:
        titulo = html.escape(str(item.get("titulo") or item.get("task_id") or "acao"))
        original = html.escape(str(item.get("data_limite_original") or "-"))
        nova = html.escape(str(item.get("nova_data_limite") or "-"))
        lines.append(f"- {titulo}: {original} -> {nova}")
    if len(items) > 6:
        lines.append(f"- mais {len(items) - 6} acao(oes)")
    lines.append("No Telegram, use os botoes abaixo para confirmar ou cancelar.")
    return "\n".join(lines)


def _summarize_memory_conflict(conflict: dict) -> str:
    lines = [
        "",
        "<b>Card de conflito de memoria</b>",
        "O Hermes encontrou duas memorias parecidas e precisa de uma decisao.",
    ]
    if conflict.get("existing_text"):
        lines.append(f"Atual: {html.escape(str(conflict.get('existing_text'))[:300])}")
    if conflict.get("proposed_text"):
        lines.append(f"Nova: {html.escape(str(conflict.get('proposed_text'))[:300])}")
    return "\n".join(lines)


def _build_web_copilot_telegram_adaptation(session: dict, delegated: dict, delegated_text: str) -> tuple[str, list | None]:
    rows = []
    extra_lines = []
    pending_cards = session.setdefault("pending_web_cards", {})

    pending_edit = delegated.get("pendingEdit")
    if isinstance(pending_edit, dict) and pending_edit.get("status", "pending") == "pending":
        pending_cards["edit"] = pending_edit
        extra_lines.append(_summarize_pending_edit(pending_edit))
        rows.append([
            {"text": "Confirmar edicao", "callback_data": "webedit_confirm"},
            {"text": "Cancelar", "callback_data": "webedit_cancel"},
        ])

    pending_batch = delegated.get("pendingBatchReschedule")
    if isinstance(pending_batch, dict) and pending_batch.get("status", "pending") == "pending":
        pending_cards["batch_reschedule"] = pending_batch
        extra_lines.append(_summarize_pending_batch_reschedule(pending_batch))
        rows.append([
            {"text": "Confirmar tudo", "callback_data": "webbatch_confirm"},
            {"text": "Cancelar", "callback_data": "webbatch_cancel"},
        ])

    memory_conflict = delegated.get("pendingMemoryConflict")
    if isinstance(memory_conflict, dict):
        pending_cards["memory_conflict"] = memory_conflict
        extra_lines.append(_summarize_memory_conflict(memory_conflict))
        rows.append([
            {"text": "Manter antiga", "callback_data": "webmem_keep_old"},
            {"text": "Manter nova", "callback_data": "webmem_keep_new"},
        ])

    tool_invocation = delegated.get("toolInvocation")
    if isinstance(tool_invocation, dict):
        tool_id = html.escape(str(tool_invocation.get("tool_id") or "ferramenta"))
        extra_lines.append(
            "\n<b>Ferramenta visual</b>\n"
            f"O Copiloto acionou a ferramenta <code>{tool_id}</code>. "
            "No Telegram ainda nao ha uma versao visual completa para esse componente; mantive o resultado textual acima."
        )

    linked_task_ids = _extract_task_ids_from_text(delegated_text)
    for task_id in linked_task_ids[:3]:
        callback_data = f"lock:{task_id}"
        if len(callback_data) <= 64:
            rows.append([{"text": f"Entrar no contexto {task_id[:8]}", "callback_data": callback_data}])

    if pending_cards:
        session["_pending_web_card_updated_at"] = datetime.now(timezone.utc).isoformat()

    return "\n".join(extra_lines), rows or None

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_db():
    return firestore.client()


def _perf_now_ms() -> int:
    return int(time.perf_counter() * 1000)


def _perf_mark(perf_state: dict, name: str):
    started_at = perf_state.get("_last_ms", perf_state["start_ms"])
    now_ms = _perf_now_ms()
    perf_state.setdefault("steps", []).append({
        "name": name,
        "duration_ms": max(0, now_ms - started_at),
    })
    perf_state["_last_ms"] = now_ms


def _perf_log(prefix: str, perf_state: dict, extra: dict | None = None):
    payload = {
        "prefix": prefix,
        "total_ms": max(0, _perf_now_ms() - perf_state["start_ms"]),
        "steps": perf_state.get("steps", []),
    }
    if perf_state.get("tool_calls"):
        payload["tool_calls"] = perf_state["tool_calls"]
    if perf_state.get("tool_rounds"):
        payload["tool_rounds"] = perf_state["tool_rounds"]
    if extra:
        payload.update(extra)
    try:
        print(f"[Perf] {json.dumps(payload, ensure_ascii=False)}")
    except Exception:
        print(f"[Perf] {prefix} total_ms={payload['total_ms']}")


def _get_api_keys(db=None):
    db = db or _get_db()
    doc = _cached_doc_get(db, "system", "api_keys")
    return doc.to_dict() or {} if doc.exists else {}


def _get_telegram_token(db=None) -> str:
    keys = _get_api_keys(db)
    token = keys.get("telegram_bot_token") or os.environ.get("TELEGRAM_BOT_TOKEN", "")
    if not token:
        raise RuntimeError("telegram_bot_token não configurado em system/api_keys.")
    return token


def _get_allowed_chat_id() -> Optional[str]:
    return os.environ.get("ALLOWED_TELEGRAM_CHAT_ID")


_HTML_ANCHOR_RE = re.compile(r"<a\s+[^>]*href=[\"']([^\"']+)[\"'][^>]*>(.*?)</a>", re.IGNORECASE | re.DOTALL)


def _neutralize_hidden_links(text: str) -> str:
    """Keep URLs visible in Telegram instead of allowing hidden anchor text."""
    def repl(match: re.Match) -> str:
        href = html.unescape((match.group(1) or "").strip())
        label = re.sub(r"<[^>]+>", "", match.group(2) or "").strip() or href
        safe_label = html.escape(label)
        if href.startswith(("http://", "https://", "task:")):
            safe_href = html.escape(href, quote=False)
            return safe_label if href in label else f"{safe_label} ({safe_href})"
        return safe_label

    return _HTML_ANCHOR_RE.sub(repl, text or "")


def _send_telegram_message(token: str, chat_id: str | int, text: str, parse_mode: str = "HTML"):
    """POST direto à Telegram Bot API."""
    text = _neutralize_hidden_links(text)
    # Telegram HTML: truncate at 4096 chars
    if len(text) > 4096:
        text = text[:4090] + "\n..."
    resp = _requests.post(
        f"https://api.telegram.org/bot{token}/sendMessage",
        json={"chat_id": chat_id, "text": text, "parse_mode": parse_mode},
        timeout=30,
    )
    if not resp.ok:
        print(f"[Telegram] sendMessage failed: {resp.status_code} {resp.text[:300]}")
        if parse_mode:
            plain_text = html.unescape(re.sub(r"<[^>]+>", "", text))
            retry_resp = _requests.post(
                f"https://api.telegram.org/bot{token}/sendMessage",
                json={"chat_id": chat_id, "text": plain_text[:4096]},
                timeout=30,
            )
            if not retry_resp.ok:
                print(f"[Telegram] sendMessage plain fallback failed: {retry_resp.status_code} {retry_resp.text[:300]}")
                return None
            try:
                return retry_resp.json().get("result", {}).get("message_id")
            except Exception:
                return None
        return None
    try:
        return resp.json().get("result", {}).get("message_id")
    except Exception:
        return None


def _send_telegram_message_with_keyboard(
    token: str, chat_id: str | int, text: str, inline_keyboard: list, parse_mode: str = "HTML"
):
    """POST sendMessage com InlineKeyboardMarkup."""
    text = _neutralize_hidden_links(text)
    if len(text) > 4096:
        text = text[:4090] + "\n..."
    resp = _requests.post(
        f"https://api.telegram.org/bot{token}/sendMessage",
        json={
            "chat_id": chat_id,
            "text": text,
            "parse_mode": parse_mode,
            "reply_markup": {"inline_keyboard": inline_keyboard},
        },
        timeout=30,
    )
    if not resp.ok:
        print(f"[Telegram] sendMessage+keyboard failed: {resp.status_code} {resp.text[:300]}")
        if parse_mode:
            plain_text = html.unescape(re.sub(r"<[^>]+>", "", text))
            retry_resp = _requests.post(
                f"https://api.telegram.org/bot{token}/sendMessage",
                json={
                    "chat_id": chat_id,
                    "text": plain_text[:4096],
                    "reply_markup": {"inline_keyboard": inline_keyboard},
                },
                timeout=30,
            )
            if not retry_resp.ok:
                print(f"[Telegram] sendMessage+keyboard plain fallback failed: {retry_resp.status_code} {retry_resp.text[:300]}")
                return None
            try:
                return retry_resp.json().get("result", {}).get("message_id")
            except Exception:
                return None
        return None
    try:
        return resp.json().get("result", {}).get("message_id")
    except Exception:
        return None


def _is_action_context_locked(session: dict | None) -> bool:
    session = session or {}
    return session.get("contexto_ativo") == "acao" and bool(session.get("acao_id") or session.get("acao_titulo"))


def _merge_inline_keyboards(inline_keyboard: list | None, trailing_keyboard: list | None) -> list | None:
    keyboards = []
    seen_callbacks = set()

    for keyboard in (inline_keyboard or [], trailing_keyboard or []):
        if not isinstance(keyboard, list):
            continue
        if keyboard and isinstance(keyboard[0], dict):
            keyboard = [keyboard]
        for row in keyboard:
            if not isinstance(row, list):
                continue
            row_copy = []
            for button in row:
                if not isinstance(button, dict):
                    continue
                callback_data = button.get("callback_data")
                if callback_data and callback_data in seen_callbacks:
                    continue
                if callback_data:
                    seen_callbacks.add(callback_data)
                row_copy.append(button)
            if row_copy:
                keyboards.append(row_copy)

    return keyboards or None


def _send_telegram_session_message(
    db,
    token: str,
    chat_id: str | int,
    text: str,
    session: dict | None = None,
    inline_keyboard: list | None = None,
    parse_mode: str = "HTML",
):
    session = session or _get_session(db, str(chat_id))
    final_text = text
    final_keyboard = inline_keyboard

    if _is_action_context_locked(session):
        acao_titulo = session.get("acao_titulo") or session.get("acao_id") or "acao"
        if not final_text.lstrip().startswith("<b>[Contexto:"):
            final_text = f"<b>[Contexto: {acao_titulo}]</b>\n\n{final_text}"
        final_keyboard = _merge_inline_keyboards(final_keyboard, _EXIT_KEYBOARD)

    # --- Medidor de Contexto ---
    ctx = session.get("contexto_ativo", "geral")
    hist_key = "history_acao" if ctx == "acao" else "history"
    raw_history = session.get(hist_key, [])
    # _MAX_HISTORY_TURNS * 2 é o limite de itens (user+model)
    turns = len(raw_history) // 2
    
    if turns >= (_MAX_HISTORY_TURNS * 0.75):  # Começa a avisar a partir de 75% (15 turnos)
        meter_icon = "🔴" if turns >= _MAX_HISTORY_TURNS else "🟡"
        warning = f"\n\n{meter_icon} <b>Contexto: {turns}/{_MAX_HISTORY_TURNS} turnos</b>"
        if turns >= _MAX_HISTORY_TURNS:
            warning += " (Limite atingido, mensagens antigas serão esquecidas. Use /limpar se necessário)."
        # Evita duplicar se já houver o aviso (importante para redifusões ou mensagens de sistema)
        if warning.strip() not in final_text:
            final_text += warning

    if final_keyboard:
        return _send_telegram_message_with_keyboard(token, chat_id, final_text, final_keyboard, parse_mode=parse_mode)
    return _send_telegram_message(token, chat_id, final_text, parse_mode=parse_mode)


def _answer_callback_query(token: str, callback_query_id: str, text: str = None):
    """Confirma o recebimento de um callback_query (remove indicador de loading no Telegram)."""
    payload: dict = {"callback_query_id": callback_query_id}
    if text:
        payload["text"] = text[:200]
    try:
        _requests.post(
            f"https://api.telegram.org/bot{token}/answerCallbackQuery",
            json=payload,
            timeout=5,
        )
    except Exception:
        pass


def _send_telegram_chat_action(token: str, chat_id: str | int, action: str):
    try:
        _requests.post(
            f"https://api.telegram.org/bot{token}/sendChatAction",
            json={"chat_id": chat_id, "action": action},
            timeout=5,
        )
    except Exception:
        pass


def _delete_telegram_message(token: str, chat_id: str | int, message_id: int):
    if not message_id: return
    try:
        _requests.post(
            f"https://api.telegram.org/bot{token}/deleteMessage",
            json={"chat_id": chat_id, "message_id": message_id},
            timeout=5,
        )
    except Exception:
        pass


def _send_telegram_typing(token: str, chat_id: str | int):
    _send_telegram_chat_action(token, chat_id, "typing")


@contextmanager
def _telegram_action_heartbeat(token: str, chat_id: str | int, action: str, interval_seconds: int = 4):
    stop_event = threading.Event()

    def _loop():
        while not stop_event.is_set():
            _send_telegram_chat_action(token, chat_id, action)
            stop_event.wait(interval_seconds)

    thread = threading.Thread(target=_loop, daemon=True)
    thread.start()
    try:
        yield
    finally:
        stop_event.set()
        thread.join(timeout=1)


def _send_telegram_voice(token: str, chat_id: str | int, audio_bytes: bytes, filename: str, mime_type: str, caption: str = "") -> bool:
    files = {
        "voice": (filename, audio_bytes, mime_type),
    }
    data = {
        "chat_id": str(chat_id),
    }
    if caption:
        data["caption"] = caption[:1024]
    resp = _requests.post(
        f"https://api.telegram.org/bot{token}/sendVoice",
        data=data,
        files=files,
        timeout=120,
    )
    if not resp.ok:
        print(f"[Telegram] sendVoice failed: {resp.status_code} {resp.text[:500]}")
    return bool(resp.ok)


def _send_telegram_document(token: str, chat_id: str | int, file_bytes: bytes, filename: str, mime_type: str, caption: str = "") -> bool:
    files = {
        "document": (filename, file_bytes, mime_type),
    }
    data = {
        "chat_id": str(chat_id),
    }
    if caption:
        data["caption"] = caption[:1024]
    resp = _requests.post(
        f"https://api.telegram.org/bot{token}/sendDocument",
        data=data,
        files=files,
        timeout=120,
    )
    if not resp.ok:
        print(f"[Telegram] sendDocument failed: {resp.status_code} {resp.text[:500]}")
    return bool(resp.ok)


def _send_telegram_photo(token: str, chat_id: str | int, photo_url: str, caption: str = "") -> bool:
    data = {
        "chat_id": str(chat_id),
        "photo": photo_url,
    }
    if caption:
        data["caption"] = caption[:1024]
    resp = _requests.post(
        f"https://api.telegram.org/bot{token}/sendPhoto",
        json=data,
        timeout=120,
    )
    if not resp.ok:
        print(f"[Telegram] sendPhoto failed: {resp.status_code} {resp.text[:500]}")
    return bool(resp.ok)



def _get_telegram_file(token: str, file_id: str) -> dict:
    """Calls getFile and returns the file metadata dict."""
    resp = _requests.get(
        f"https://api.telegram.org/bot{token}/getFile",
        params={"file_id": file_id},
        timeout=15,
    )
    resp.raise_for_status()
    data = resp.json()
    return data.get("result", {})


def _download_telegram_file(token: str, file_path: str) -> bytes:
    url = f"https://api.telegram.org/file/bot{token}/{file_path}"
    resp = _requests.get(url, timeout=60)
    resp.raise_for_status()
    return resp.content


# ---------------------------------------------------------------------------
# Session state machine
# ---------------------------------------------------------------------------

def _get_session(db, chat_id: str) -> dict:
    doc = db.collection("telegram_sessions").document(chat_id).get()
    return doc.to_dict() or {"chat_id": chat_id, "contexto_ativo": "geral", "history": []}


def _save_session(db, chat_id: str, session: dict):
    session["updated_at"] = firestore.SERVER_TIMESTAMP
    db.collection("telegram_sessions").document(chat_id).set(session)
    copilot_session_id = session.get("copilot_session_id")
    if copilot_session_id:
        try:
            db.collection("sessoes_copiloto").document(copilot_session_id).set({
                "channel": "telegram",
                "telegramChatId": str(chat_id),
                "taskId": session.get("acao_id") or None,
                "lastMessageAt": firestore.SERVER_TIMESTAMP,
            }, merge=True)
        except Exception as exc:
            print(f"[SessionBridge] Falha ao sincronizar metadados da sessao {copilot_session_id}: {exc}")


def _session_matches_processing_state(session: dict, contexto_ativo: str, acao_id: str | None) -> bool:
    current_contexto = session.get("contexto_ativo", "geral")
    if current_contexto != (contexto_ativo or "geral"):
        return False
    if (contexto_ativo or "geral") == "acao":
        return session.get("acao_id") == acao_id
    return True


def _handle_command(text: str, session: dict) -> Optional[str]:
    """
    Returns a reply string if the message is a state-machine command, else None.
    Side-effects: mutates session.
    """
    text = (text or "").strip()
    m = re.match(r"^/contexto\s+(.+)$", text, re.IGNORECASE)
    if m:
        ctx = m.group(1).strip()
        session["contexto_ativo"] = ctx
        session["history"] = []  # fresh history for new context
        return f"Contexto ativo definido como <b>{ctx}</b>. Histórico reiniciado."

    if re.match(r"^/sair$", text, re.IGNORECASE):
        previous = session.get("acao_titulo") or session.get("contexto_ativo", "geral")
        session["contexto_ativo"] = "geral"
        session["history"] = []
        session["acao_id"] = None
        session["acao_titulo"] = None
        session["acao_context_snapshot"] = None
        session["history_acao"] = []
        return f"Saindo do contexto <b>{previous}</b>. Voltando ao modo geral."

    if re.match(r"^/start$", text, re.IGNORECASE):
        return (
            "👋 Olá! Sou o <b>Hermes Copiloto</b>.\n\n"
            "Comandos disponíveis:\n"
            "• <code>/entrar [termo]</code> — busca ações e trava o contexto nelas\n"
            "• <code>/sair</code> — sai do contexto trancado, retorna ao modo geral\n"
            "• <code>/status</code> — mostra o contexto e histórico atuais\n\n"
            "Envie texto, áudio ou arquivos. Tamanho máximo: 20 MB."
        )

    if re.match(r"^/status$", text, re.IGNORECASE):
        ctx = session.get("contexto_ativo", "geral")
        acao_titulo = session.get("acao_titulo", "")
        hist_key = "history_acao" if ctx == "acao" else "history"
        turns = len(session.get(hist_key, [])) // 2
        ctx_display = f"ação trancada — <b>{acao_titulo}</b>" if ctx == "acao" else f"<b>{ctx}</b>"
        return f"Contexto ativo: {ctx_display}\nTurnos no histórico: {turns}"

    if re.match(r"^/limpar$", text, re.IGNORECASE):
        ctx = session.get("contexto_ativo", "geral")
        hist_key = "history_acao" if ctx == "acao" else "history"
        session[hist_key] = []
        return f"✅ Histórico do contexto <b>{ctx}</b> foi limpo. Podemos recomeçar!"

    return None


# ---------------------------------------------------------------------------
# Action context helpers (session locking)
# ---------------------------------------------------------------------------

_DRIVE_ID_PATTERNS = [
    re.compile(r"/d/([a-zA-Z0-9_-]{10,})"),
    re.compile(r"[?&]id=([a-zA-Z0-9_-]{10,})"),
    re.compile(r"/folders/([a-zA-Z0-9_-]{10,})"),
]


def _extract_drive_file_id(value: str | None) -> str | None:
    value = (value or "").strip()
    if re.fullmatch(r"[a-zA-Z0-9_-]{10,}", value):
        return value
    for pattern in _DRIVE_ID_PATTERNS:
        match = pattern.search(value)
        if match:
            return match.group(1)
    return None


def _parse_diary_file_note(note: str | None) -> dict | None:
    """Parse FILE:: diary rich notes created by the web diary UI."""
    note = (note or "").strip()
    if not note.startswith("FILE::"):
        return None

    payload = note[len("FILE::"):]
    name = ""
    value = ""
    if payload.startswith("JSON::"):
        try:
            parsed = json.loads(payload[len("JSON::"):])
            if isinstance(parsed, dict):
                name = str(parsed.get("n") or "")
                value = str(parsed.get("v") or "")
        except Exception:
            return None
    else:
        separator_index = payload.find("::")
        if separator_index == -1:
            value = payload
        else:
            name = payload[:separator_index]
            value = payload[separator_index + 2:]

    if not value:
        return None
    return {"nome": name or "Arquivo", "url": value}


def _drive_url_from_file(file_info: dict) -> str:
    url = (file_info.get("url") or "").strip()
    if url.startswith("http://") or url.startswith("https://"):
        return url
    fid = file_info.get("drive_file_id")
    return f"https://drive.google.com/file/d/{fid}/view" if fid else ""


_ACTION_SEARCH_STOPWORDS = {
    "de", "a", "o", "que", "e", "do", "da", "em", "um", "uma",
    "os", "as", "no", "na", "com", "por", "para", "dos", "das",
    "nos", "nas", "ao", "se", "ou", "acao", "acoes", "tarefa", "tarefas",
    "pesquisa", "pesquisar", "pesquise", "busca", "buscar", "busque",
    "procura", "procurar", "procure", "localiza", "localizar", "localize",
    "ative", "ativar", "ativa", "contexto", "hermes", "por", "favor",
}


def _action_query_terms(text: str) -> list[str]:
    normalized = _normalize_for_matching(text)
    return [
        term
        for term in re.findall(r"[\w./-]+", normalized, flags=re.UNICODE)
        if term not in _ACTION_SEARCH_STOPWORDS and len(term) > 2
    ]


def _clean_action_search_query(text: str) -> str:
    query = _normalize_for_matching(text)
    query = re.sub(r"\b(?:hermes|por favor)\b", " ", query, flags=re.IGNORECASE)
    query = re.sub(r"\b(?:pesquise|pesquisa|pesquisar|busque|busca|buscar|procure|procura|procurar|localize|localiza|localizar)\b", " ", query, flags=re.IGNORECASE)
    query = re.sub(r"\b(?:a|uma|o|um)?\s*(?:acao|acoes|tarefa|tarefas)\b", " ", query, flags=re.IGNORECASE)
    query = re.sub(r"\b(?:e\s+)?(?:ative|ativa|ativar|entre|entra|entrar)\b.*\bcontexto\b", " ", query, flags=re.IGNORECASE)
    query = re.sub(r"\b(?:no|na|em|o)?\s*contexto\b", " ", query, flags=re.IGNORECASE)
    query = re.sub(r"\s+", " ", query)
    return query.strip(" .,:;-")


def _extract_action_search_context_query(text: str) -> str | None:
    lowered = _normalize_for_matching(text).strip()
    if not lowered:
        return None
    # Long messages (forwarded texts, context info) are never intent commands
    if len(text.strip()) > 500:
        return None
    wants_search = any(re.search(r'\b' + marker + r'\b', lowered) for marker in (
        "pesquisa", "pesquise", "buscar", "busca", "busque",
        "procura", "procure", "localiza", "localize",
    ))
    mentions_action = any(marker in lowered for marker in ("acao", "acoes", "tarefa", "tarefas"))
    wants_context = "contexto" in lowered and any(re.search(r'\b' + marker + r'\b', lowered) for marker in (
        "ative", "ativa", "ativar", "entre", "entra", "entrar",
    ))
    if not (wants_search and mentions_action and wants_context):
        return None
    cleaned = _clean_action_search_query(text)
    return cleaned or None


def _extract_action_lookup_query(text: str) -> str | None:
    """Detects requests that only ask to find/list an action, without locking context."""
    lowered = _normalize_for_matching(text).strip()
    if not lowered:
        return None
    if len(text.strip()) > 500:
        return None
    wants_search = any(marker in lowered for marker in (
        "pesquisa", "pesquise", "pesquisar", "buscar", "busca", "busque",
        "procura", "procurar", "procure", "localiza", "localizar", "localize",
        "encontra", "encontrar", "encontre", "ache", "achar",
    ))
    mentions_action = any(marker in lowered for marker in ("acao", "acoes", "tarefa", "tarefas"))
    wants_context = "contexto" in lowered and any(marker in lowered for marker in (
        "ative", "ativa", "ativar", "entre", "entra", "entrar",
    ))
    if not (wants_search and mentions_action) or wants_context:
        return None
    cleaned = _clean_action_search_query(text)
    cleaned = re.sub(r"\b(?:relacionad[ao]s?|referente|sobre|chamad[ao]s?|nomead[ao]s?)\b", " ", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" .,:;-")
    return cleaned or None


def _extract_natural_context_query(text: str) -> str | None:
    """Detects natural-language requests such as 'entre no contexto da acao X'."""
    lowered = _normalize_for_matching(text).strip()
    if not lowered:
        return None

    patterns = [
        r"^(?:entre|entra|entrar)\s+(?:no|na|em)?\s*contexto\s*(?:da|do|de)?\s*(?:acao|tarefa)?\s*(?:da|do|de)?\s*(.*)$",
        r"^(?:ative|ativa|ativar)\s+(?:o\s+)?contexto\s*(?:da|do|de)?\s*(?:acao|tarefa)?\s*(?:da|do|de)?\s*(.*)$",
    ]
    for pattern in patterns:
        match = re.match(pattern, lowered, flags=re.IGNORECASE)
        if match:
            return (match.group(1) or "").strip(" .,:;-")
    return None


def _is_list_actions_request(text: str) -> tuple[bool, str | None]:
    """Detects 'liste as ações de hoje' style requests.
    Returns (is_match, date_iso_or_None) — date_iso is None when no date was specified.
    """
    lowered = _normalize_for_matching(text).strip()
    if not lowered or len(text.strip()) > 250:
        return False, None

    list_markers = (
        "liste", "listar", "me liste",
        "mostre", "mostrar", "me mostre", "mostra", "me mostra",
        "me da", "me de",
        "quais sao",
    )
    action_markers = ("acoes", "tarefas", "atividades")

    has_list = any(marker in lowered for marker in list_markers)
    has_action = any(marker in lowered for marker in action_markers)

    wants_context = "contexto" in lowered and any(m in lowered for m in ("ative", "ativa", "entre", "entra"))
    has_file = any(m in lowered for m in ("arquivo", "arquivos", "documento", "documentos", "link", "links", "anexo", "anexos"))

    if not (has_list and has_action) or wants_context or has_file:
        return False, None

    date_iso = _requested_date_from_text(text)
    return True, date_iso


def _fetch_actions_by_date(date_iso: str | None) -> list:
    """Fetches actions filtered by data_limite == date_iso, or all if date_iso is None."""
    from tools.busca_grafo import buscar_tarefas
    if date_iso:
        res = buscar_tarefas(
            "",
            match_mode="any",
            limite=200,
            data_limite_inicio=date_iso,
            data_limite_fim=date_iso,
        )
    else:
        res = buscar_tarefas("", match_mode="any", limite=200)
    return res.get("resultados", [])


def _format_actions_simple_list(results: list, date_iso: str | None) -> str:
    """Formats a plain bullet-point list of action titles."""
    if date_iso:
        try:
            d = datetime.fromisoformat(date_iso)
            date_label = d.strftime("%d/%m/%Y")
        except Exception:
            date_label = date_iso
        header = f"Ações para <b>{date_label}</b>:"
    else:
        header = "Ações cadastradas:"

    if not results:
        return header + "\n\n<i>Nenhuma ação encontrada.</i>"

    lines = [header, ""]
    for r in results:
        titulo = html.escape(r.get("titulo") or "sem título")
        lines.append(f"• {titulo}")

    return "\n".join(lines)


def _is_reset_request(text: str) -> bool:
    """Detecta se o usuário quer resetar/limpar a sessão via linguagem natural."""
    lowered = _normalize_for_matching(text).strip()
    if not lowered:
        return False
    
    # Marcadores de ação
    reset_markers = ["limpar", "limpa", "resetar", "reset", "recomecar", "recomeca", "reiniciar", "reinicia"]
    # Marcadores de objeto
    context_markers = ["chat", "conversa", "historico", "sessao", "memoria"]
    
    has_reset = any(m in lowered for m in reset_markers)
    has_context = any(m in lowered for m in context_markers)
    
    # Se disser apenas "reset" ou "limpa tudo", ou combinar reset + contexto
    return (has_reset and has_context) or (has_reset and len(lowered.split()) <= 2)


def _select_context_result(query: str, results: list) -> dict | None:
    if not results:
        return None
    if len(results) == 1:
        return results[0]

    norm_query = _normalize_for_matching(query).strip()
    if not norm_query:
        return None

    exact = [r for r in results if _normalize_for_matching(r.get("titulo", "")) == norm_query]
    if len(exact) == 1:
        return exact[0]

    contained = [r for r in results if norm_query in _normalize_for_matching(r.get("titulo", ""))]
    if len(contained) == 1:
        return contained[0]

    query_terms = set(_action_query_terms(query))
    if query_terms:
        scored = []
        for r in results:
            title_terms = set(_action_query_terms(r.get("titulo", "")))
            overlap = len(query_terms & title_terms)
            scored.append((overlap, r))
        scored.sort(key=lambda item: item[0], reverse=True)
        best_score = scored[0][0]
        second_score = scored[1][0] if len(scored) > 1 else 0
        if best_score >= 2 and best_score > second_score:
            return scored[0][1]
    return None


def _lock_action_session(session: dict, task_id: str, snapshot: dict) -> None:
    session["contexto_ativo"] = "acao"
    session["acao_id"] = task_id
    session["acao_titulo"] = snapshot.get("titulo") or task_id
    session["acao_context_snapshot"] = snapshot
    session["history_acao"] = []


def _format_context_locked_message(snapshot: dict) -> str:
    titulo = html.escape(snapshot.get("titulo") or "acao")
    arquivos_count = len(snapshot.get("arquivos_disponiveis") or [])
    return (
        f"🔒 <b>[Contexto: {titulo}]</b>\n\n"
        "Contexto trancado. Estou focado exclusivamente nesta ação.\n"
        f"Arquivos detectados no contexto: <b>{arquivos_count}</b>.\n\n"
        "Use <i>Sair do Contexto</i> para retornar ao modo geral."
    )


def _is_context_file_request(text: str) -> bool:
    lowered = _normalize_for_matching(text)
    file_markers = (
        "link",
        "links",
        "documento",
        "documentos",
        "arquivo",
        "arquivos",
        "anexo",
        "anexos",
    )
    request_markers = (
        "qual",
        "quais",
        "mande",
        "manda",
        "envie",
        "enviar",
        "me manda",
        "pode me mandar",
        "listar",
        "liste",
        "mostre",
        "mostrar",
    )
    return any(marker in lowered for marker in file_markers) and any(marker in lowered for marker in request_markers)


def _requested_date_from_text(text: str) -> str | None:
    lowered = _normalize_for_matching(text)
    try:
        from zoneinfo import ZoneInfo
        today = datetime.now(ZoneInfo("America/Sao_Paulo")).date()
    except Exception:
        today = datetime.now(timezone.utc).date()

    if "ontem" in lowered:
        return (today - timedelta(days=1)).isoformat()
    if "hoje" in lowered:
        return today.isoformat()

    match = re.search(r"\b(\d{1,2})/(\d{1,2})(?:/(\d{2,4}))?\b", lowered)
    if not match:
        return None
    day = int(match.group(1))
    month = int(match.group(2))
    year_raw = match.group(3)
    year = int(year_raw) if year_raw else today.year
    if year < 100:
        year += 2000
    try:
        return datetime(year, month, day).date().isoformat()
    except ValueError:
        return None


def _format_context_files_response(acao_snapshot: dict, request_text: str) -> str:
    arquivos = list(acao_snapshot.get("arquivos_disponiveis") or [])
    date_filter = _requested_date_from_text(request_text)
    if date_filter:
        arquivos = [
            a for a in arquivos
            if str(a.get("data_diario") or a.get("data_criacao") or "").startswith(date_filter)
        ]

    lowered = _normalize_for_matching(request_text)
    requested_terms = []
    if "deferid" in lowered:
        requested_terms.append("deferid")
    if "indeferid" in lowered:
        requested_terms.append("indeferid")
    if requested_terms:
        term_filtered = [
            a for a in arquivos
            if any(term in _normalize_for_matching(a.get("nome", "")) for term in requested_terms)
        ]
        if term_filtered:
            arquivos = term_filtered

    titulo = html.escape(acao_snapshot.get("titulo") or "acao")
    if not arquivos:
        when = f" em {date_filter}" if date_filter else ""
        return (
            f"<b>[Contexto: {titulo}]</b>\n\n"
            f"Nao encontrei anexos ou links registrados{when} nesse contexto.\n"
            "Nao vou criar links de exemplo. Se o arquivo estiver visivel no diario mas nao aparecer aqui, "
            "provavelmente ele nao foi gravado no pool/snapshot da tarefa."
        )

    lines = [
        f"<b>[Contexto: {titulo}]</b>",
        "",
        "Encontrei estes arquivos registrados no contexto. Estou mostrando a URL completa para evitar link oculto incorreto:",
        "",
    ]
    for file_info in arquivos[:12]:
        name = html.escape(file_info.get("nome") or "Arquivo sem nome")
        url = html.escape(_drive_url_from_file(file_info), quote=False)
        date_label = (file_info.get("data_diario") or file_info.get("data_criacao") or "")[:10]
        date_suffix = f" ({html.escape(date_label)})" if date_label else ""
        if url:
            lines.append(f"- <b>{name}</b>{date_suffix}\n  {url}")
        else:
            fid = html.escape(file_info.get("drive_file_id") or "sem drive_file_id")
            lines.append(f"- <b>{name}</b>{date_suffix}\n  DRIVE_FILE_ID: <code>{fid}</code>")

    if len(arquivos) > 12:
        lines.append(f"\nMais {len(arquivos) - 12} arquivo(s) omitido(s) para caber no Telegram.")
    return "\n".join(lines)


def _fetch_acao_snapshot(db, task_id: str) -> dict | None:
    """Busca e formata dados de uma tarefa para persistir no snapshot da sessão."""
    try:
        from main import get_drive_service

        doc = db.collection("tarefas").document(task_id).get()
        if not doc.exists:
            return None
        data = doc.to_dict() or {}

        acomp_raw = data.get("acompanhamento") or []
        acomp_entries = []
        diario_full = []

        def _obter_data_ordenacao(entry):
            d = entry.get('data')
            if not d:
                return ""
            if isinstance(d, str):
                return d
            import datetime as _datetime
            if isinstance(d, _datetime.datetime):
                if d.tzinfo is None:
                    d = d.replace(tzinfo=_datetime.timezone.utc)
                return d.isoformat()
            if isinstance(d, _datetime.date):
                return d.isoformat()
            if hasattr(d, 'isoformat'):
                try:
                    return d.isoformat()
                except Exception:
                    pass
            return str(d)

        def _formatar_data_curta(entry):
            d_val = entry.get('data')
            if not d_val:
                return ""
            if isinstance(d_val, str):
                return d_val[:10]
            elif hasattr(d_val, "strftime"):
                try:
                    return d_val.strftime("%Y-%m-%d")
                except Exception:
                    return str(d_val)[:10]
            else:
                return str(d_val)[:10]

        for entry in sorted(acomp_raw, key=_obter_data_ordenacao):
            nota = (entry.get("nota") or "").strip()
            data_entry = _formatar_data_curta(entry)
            if nota:
                diario_full.append(f"[{data_entry}] {nota[:500]}")
        for entry in sorted(acomp_raw, key=_obter_data_ordenacao, reverse=True)[:3]:
            nota = (entry.get("nota") or "").strip()
            data_entry = _formatar_data_curta(entry)
            if nota:
                acomp_entries.append(f"[{data_entry}] {nota[:300]}")

        plano_raw = data.get("plano_acao") or []
        plano_entries = []
        for passo in plano_raw:
            if isinstance(passo, dict):
                texto = (passo.get("text") or passo.get("titulo") or "").strip()
                concluido = passo.get("completed", False)
                if texto:
                    marcador = "✓" if concluido else "○"
                    plano_entries.append(f"{marcador} {texto[:200]}")
            elif isinstance(passo, str) and passo.strip():
                plano_entries.append(f"○ {passo.strip()[:200]}")

        tags = data.get("kg_tags") or data.get("tags") or []
        arquivos_disponiveis = []
        arquivos_seen = set()
        drive_service = None

        def add_file_candidate(nome: str, url: str, data_criacao: str = "", data_diario: str = ""):
            nonlocal drive_service
            fid = _extract_drive_file_id(url)
            if not fid:
                return
            key = fid or f"{nome}|{url}"
            if key in arquivos_seen:
                for existing in arquivos_disponiveis:
                    if existing.get("drive_file_id") == fid:
                        if data_diario and not existing.get("data_diario"):
                            existing["data_diario"] = data_diario
                        if data_criacao and not existing.get("data_criacao"):
                            existing["data_criacao"] = data_criacao
                        if url and not existing.get("url"):
                            existing["url"] = url
                        break
                return
            if drive_service is None:
                try:
                    drive_service = get_drive_service()
                except Exception:
                    drive_service = False
            if drive_service:
                try:
                    meta = drive_service.files().get(fileId=fid, fields="id,trashed").execute()
                    if meta.get("trashed"):
                        return
                except Exception:
                    return
            arquivos_seen.add(key)
            arquivos_disponiveis.append({
                "nome": nome or "Arquivo sem nome",
                "drive_file_id": fid,
                "url": url,
                "data_criacao": data_criacao or "",
                "data_diario": data_diario or "",
            })

        for item in data.get("pool_dados", []):
            if item.get("tipo") != "arquivo":
                continue
            fid = item.get("drive_file_id")
            raw_url = item.get("valor", "") or ""
            url = raw_url if _extract_drive_file_id(raw_url) else (f"https://drive.google.com/file/d/{fid}/view" if fid else raw_url)
            add_file_candidate(
                item.get("nome", "Arquivo sem nome"),
                url,
                data_criacao=str(item.get("data_criacao") or ""),
            )

        for entry in acomp_raw:
            parsed_file = _parse_diary_file_note(entry.get("nota"))
            if not parsed_file:
                continue
            add_file_candidate(
                parsed_file.get("nome") or "Arquivo",
                parsed_file.get("url") or "",
                data_diario=str(entry.get("data") or ""),
            )

        return {
            "id": task_id,
            "titulo": data.get("titulo", "sem titulo"),
            "status": data.get("status", ""),
            "area": data.get("area_tematica", ""),
            "area_tematica": data.get("area_tematica", ""),
            "descricao": (data.get("descricao") or "")[:500],
            "notas": (data.get("notas") or "")[:400],
            "sintese": (data.get("sintese_demanda") or data.get("demanda") or "")[:400],
            "plano_atual": data.get("plano_acao", []),
            "plano_acao": plano_entries,
            "diario_integral": "\n".join(diario_full)[-4000:],
            "acompanhamento_recente": acomp_entries,
            "tags": tags,
            "arquivos_disponiveis": arquivos_disponiveis[:20],
        }
    except Exception as e:
        print(f"[Session] Falha ao buscar tarefa {task_id}: {e}")
        return None


def _cached_acao_snapshot(db, task_id: str) -> dict | None:
    now = time.monotonic()
    cached = _ACTION_SNAPSHOT_CACHE.get(task_id)
    if cached and (now - cached[0]) < _ACTION_SNAPSHOT_CACHE_TTL:
        return cached[1]

    snapshot = _fetch_acao_snapshot(db, task_id)
    if snapshot:
        _ACTION_SNAPSHOT_CACHE[task_id] = (now, snapshot)
    return snapshot


def _search_actions_for_context(db, query: str) -> list:
    """Busca tarefas usando buscar_tarefas e retorna até 5 resultados para seleção."""
    from tools.busca_grafo import buscar_tarefas

    if query:
        terms = _action_query_terms(query)
        mode = "all" if len(terms) >= 2 else "any"
        res = buscar_tarefas(query, match_mode=mode, limite=5)
        if mode == "all" and not res.get("resultados"):
            res = buscar_tarefas(query, match_mode="any", limite=5)
    else:
        res = buscar_tarefas("", match_mode="any", limite=5)

    return res.get("resultados", [])


def _build_action_keyboard(results: list) -> list:
    """Constrói InlineKeyboard com uma linha por resultado de ação."""
    keyboard = []
    for r in results:
        task_id = r["id"]
        titulo = (r.get("titulo") or "sem titulo")[:35]
        area = r.get("area", "")
        label = titulo + (f" [{area}]" if area else "")
        keyboard.append([{"text": label[:60], "callback_data": f"lock:{task_id}"}])
    return keyboard


def _format_action_lookup_results(query: str, results: list) -> str:
    query_label = html.escape(query or "sua busca")
    if not results:
        return (
            f"Nenhuma acao encontrada para <i>{query_label}</i> no historico do Hermes.\n"
            "Nao pesquisei e-mails, acervo ou internet porque voce pediu uma acao do sistema."
        )

    if len(results) == 1:
        header = f"Encontrei esta acao para <i>{query_label}</i>:"
    else:
        header = f"Encontrei {len(results)} acoes mais provaveis para <i>{query_label}</i>:"

    lines = [header]
    for idx, r in enumerate(results[:5], 1):
        titulo = html.escape(r.get("titulo") or "sem titulo")
        task_id = html.escape(str(r.get("id") or ""))
        status = html.escape(r.get("status") or "nao informado")
        area = html.escape(r.get("area") or "nao informada")
        prazo = html.escape(str(r.get("data_limite") or "nao informado"))
        lines.append(f"\n<b>{idx}. {titulo}</b>")
        lines.append(f"ID: <code>{task_id}</code>")
        lines.append(f"Status: {status} | Area: {area} | Prazo: {prazo}")
        if r.get("processo_sei"):
            lines.append(f"Processo SEI: {html.escape(str(r.get('processo_sei')))}")
        lines.append(f"Abrir: task:{task_id}")
    return "\n".join(lines)


_EXIT_KEYBOARD = [[{"text": "🔓 Sair do Contexto", "callback_data": "exit_context"}]]
_RESET_CONFIRM_KEYBOARD = [
    [
        {"text": "✅ Sim, limpar", "callback_data": "reset_session_confirm"},
        {"text": "❌ Não, manter", "callback_data": "reset_session_cancel"}
    ]
]

# Novos teclados para fluxo de "Proposta -> Confirmação"
_CONFIRM_ACAO_KEYBOARD = [
    [
        {"text": "✅ Confirmar Registro", "callback_data": "confirm_acao"},
        {"text": "❌ Cancelar", "callback_data": "cancel_acao"}
    ]
]

_CONFIRM_FINANCEIRO_KEYBOARD = [
    [
        {"text": "✅ Confirmar Lançamento", "callback_data": "confirm_financeiro"},
        {"text": "❌ Cancelar", "callback_data": "cancel_financeiro"}
    ]
]

_CONFIRM_WHATSAPP_KEYBOARD = [
    [
        {"text": "✅ Confirmar Envio", "callback_data": "confirm_whatsapp"},
        {"text": "❌ Cancelar", "callback_data": "cancel_whatsapp"}
    ]
]



def _handle_telegram_callback(db, token: str, callback_query: dict) -> "https_fn.Response":
    """Processa um callback_query de botão inline de forma síncrona (sem LLM)."""
    query_id = callback_query.get("id", "")
    data = (callback_query.get("data") or "").strip()
    message = callback_query.get("message") or {}
    chat_id = str((message.get("chat") or {}).get("id", ""))

    if not chat_id or not data:
        _answer_callback_query(token, query_id)
        return https_fn.Response("OK", status=200)

    allowed = _get_allowed_chat_id()
    if allowed and chat_id != allowed:
        _answer_callback_query(token, query_id)
        return https_fn.Response("OK", status=200)

    session = _get_session(db, chat_id)
    copilot_session_id = _ensure_copilot_session(
        db,
        chat_id,
        session,
        first_text=f"[Botão Telegram] {data}",
        from_user=callback_query.get("from") or {},
    )
    _save_session(db, chat_id, session)

    def _persist_callback_turn(user_label: str, assistant_text: str):
        _persist_copilot_message(db, copilot_session_id, "user", user_label, source="telegram_callback")
        _persist_copilot_message(db, copilot_session_id, "assistant", assistant_text, source="telegram_callback")

    def _pending_web_card(card_type: str) -> dict | None:
        card = (session.get("pending_web_cards") or {}).get(card_type)
        return card if isinstance(card, dict) else None

    def _clear_pending_web_card(card_type: str):
        cards = session.get("pending_web_cards") or {}
        cards.pop(card_type, None)
        if cards:
            session["pending_web_cards"] = cards
        else:
            session.pop("pending_web_cards", None)



    if data.startswith("health_pain:"):
        parts = data.split(":")
        pain_date = parts[1] if len(parts) > 1 else ""
        pain_raw = parts[2] if len(parts) > 2 else ""
        try:
            pain_score = max(0, min(10, int(pain_raw)))
            if not re.match(r"^\d{4}-\d{2}-\d{2}$", pain_date):
                raise ValueError("data inválida")
            doc_ref = db.collection("health_exercise_logs").document(pain_date)
            current = doc_ref.get()
            current_pain = ((current.to_dict() or {}).get("pain") or {}) if current.exists else {}
            current_pain["evening"] = pain_score
            current_pain["telegram_checked_at"] = datetime.now(timezone.utc).isoformat()
            doc_ref.set({"pain": current_pain}, merge=True)
            _answer_callback_query(token, query_id, f"Dor registrada: {pain_score}/10")
            response_text = f"✅ Check-in lombar registrado: <b>{pain_score}/10</b> em {pain_date}."
            _persist_callback_turn(f"Check-in lombar: {pain_score}/10", response_text)
            _send_telegram_message(token, chat_id, response_text)
        except Exception as exc:
            print(f"[HealthCheckin] Falha ao registrar dor: {exc}")
            _answer_callback_query(token, query_id, "Não consegui registrar.")
            _send_telegram_message(token, chat_id, "⚠️ Não consegui registrar esse check-in lombar. Tente responder novamente.")

    elif data == "exit_context":
        acao_titulo = session.get("acao_titulo") or "anterior"
        response_text = f"✅ Saindo do contexto <b>{acao_titulo}</b>. Voltando ao modo geral."
        _answer_callback_query(token, query_id, "Contexto liberado.")
        session["contexto_ativo"] = "geral"
        session["acao_id"] = None
        session["acao_titulo"] = None
        session["acao_context_snapshot"] = None
        session["history_acao"] = []
        _save_session(db, chat_id, session)
        _persist_callback_turn("Botão: Sair do Contexto", response_text)
        _send_telegram_message(
            token, chat_id,
            response_text
        )

    elif data.startswith("lock:"):
        task_id = data[len("lock:"):]
        _answer_callback_query(token, query_id, "Carregando contexto...")
        snapshot = _cached_acao_snapshot(db, task_id)
        if not snapshot:
            response_text = f"⚠️ Ação <code>{task_id}</code> não encontrada."
            _persist_callback_turn(f"Botão: entrar no contexto task:{task_id}", response_text)
            _send_telegram_session_message(
                db,
                token, chat_id,
                response_text
            )
            return https_fn.Response("OK", status=200)

        titulo = snapshot.get("titulo", task_id)
        response_text = (
            f"🔒 <b>[Contexto: {titulo}]</b>\n\n"
            f"Contexto trancado. Estou focado exclusivamente nesta ação.\n"
            f"Histórico anterior isolado — nenhum ruído de conversas passadas.\n\n"
            f"Use <i>Sair do Contexto</i> para retornar ao modo geral."
        )
        session["contexto_ativo"] = "acao"
        session["acao_id"] = task_id
        session["acao_titulo"] = titulo
        session["acao_context_snapshot"] = snapshot
        session["history_acao"] = []
        _save_session(db, chat_id, session)
        _persist_callback_turn(f"Botão: entrar no contexto task:{task_id}", response_text)

        _send_telegram_message_with_keyboard(
            token, chat_id,
            response_text,
            _EXIT_KEYBOARD,
        )

    elif data == "reset_session_confirm":
        ctx = session.get("contexto_ativo", "geral")
        hist_key = "history_acao" if ctx == "acao" else "history"
        session[hist_key] = []
        _save_session(db, chat_id, session)
        _answer_callback_query(token, query_id, "Histórico limpo!")
        response_text = f"✅ Histórico do contexto <b>{ctx}</b> foi limpo. Podemos recomeçar!"
        _persist_callback_turn("Botão: confirmar limpeza de histórico", response_text)
        _send_telegram_message(token, chat_id, response_text)

    elif data == "reset_session_cancel":
        _answer_callback_query(token, query_id, "Ação cancelada.")
        response_text = "Ok, mantive o histórico atual."
        _persist_callback_turn("Botão: cancelar limpeza de histórico", response_text)
        _send_telegram_message(token, chat_id, response_text)

    elif data == "webedit_confirm":
        _answer_callback_query(token, query_id, "Aplicando edição...")
        pending = _pending_web_card("edit")
        if not pending:
            response_text = "Não encontrei uma edição pendente para confirmar."
            _persist_callback_turn("Botão: confirmar edição", response_text)
            _send_telegram_message(token, chat_id, response_text)
            return https_fn.Response("OK", status=200)
        try:
            alteracoes = {}
            for campo, change in (pending.get("alteracoes") or {}).items():
                change = change if isinstance(change, dict) else {}
                alteracoes[campo] = change.get("novo_raw") if "novo_raw" in change else change.get("novo")
            message_id = _find_latest_copilot_card_message_id(db, copilot_session_id, "pendingEdit")
            result = _call_web_callable(
                function_name="confirmarEdicaoAcao",
                data={
                    "sessionId": copilot_session_id,
                    "messageId": message_id,
                    "taskId": pending.get("task_id"),
                    "alteracoes": alteracoes,
                    "snapshotTs": pending.get("snapshot_ts"),
                },
                user_uid=session.get("userId"),
                timeout=60,
            )
            _clear_pending_web_card("edit")
            _save_session(db, chat_id, session)
            status = result.get("status") or "completed"
            response_text = "Edição confirmada e aplicada no Hermes." if status == "completed" else (result.get("message") or f"Edição retornou status: {status}")
            _persist_callback_turn("Botão: confirmar edição", response_text)
            _send_telegram_message(token, chat_id, response_text)
        except Exception as exc:
            response_text = f"Erro ao confirmar edição: {exc}"
            _persist_callback_turn("Botão: confirmar edição", response_text)
            _send_telegram_message(token, chat_id, response_text)

    elif data == "webedit_cancel":
        _answer_callback_query(token, query_id, "Edição cancelada.")
        _clear_pending_web_card("edit")
        _set_latest_copilot_card_status(db, copilot_session_id, "pendingEdit", "cancelled")
        _save_session(db, chat_id, session)
        response_text = "Edição cancelada. Nada foi alterado."
        _persist_callback_turn("Botão: cancelar edição", response_text)
        _send_telegram_message(token, chat_id, response_text)

    elif data == "webbatch_confirm":
        _answer_callback_query(token, query_id, "Reagendando...")
        pending = _pending_web_card("batch_reschedule")
        if not pending:
            response_text = "Não encontrei um reagendamento em lote pendente para confirmar."
            _persist_callback_turn("Botão: confirmar reagendamento em lote", response_text)
            _send_telegram_message(token, chat_id, response_text)
            return https_fn.Response("OK", status=200)
        try:
            message_id = _find_latest_copilot_card_message_id(db, copilot_session_id, "pendingBatchReschedule")
            result = _call_web_callable(
                function_name="confirmarReagendamentoEmLote",
                data={
                    "sessionId": copilot_session_id,
                    "messageId": message_id,
                    "items": pending.get("items") or [],
                    "justificativa": pending.get("justificativa") or "Reagendamento em lote via Telegram.",
                },
                user_uid=session.get("userId"),
                timeout=60,
            )
            _clear_pending_web_card("batch_reschedule")
            _save_session(db, chat_id, session)
            response_text = f"Reagendamento confirmado. Ações atualizadas: {result.get('count', 0)}."
            _persist_callback_turn("Botão: confirmar reagendamento em lote", response_text)
            _send_telegram_message(token, chat_id, response_text)
        except Exception as exc:
            response_text = f"Erro ao confirmar reagendamento: {exc}"
            _persist_callback_turn("Botão: confirmar reagendamento em lote", response_text)
            _send_telegram_message(token, chat_id, response_text)

    elif data == "webbatch_cancel":
        _answer_callback_query(token, query_id, "Reagendamento cancelado.")
        _clear_pending_web_card("batch_reschedule")
        _set_latest_copilot_card_status(db, copilot_session_id, "pendingBatchReschedule", "cancelled")
        _save_session(db, chat_id, session)
        response_text = "Reagendamento em lote cancelado. Nada foi alterado."
        _persist_callback_turn("Botão: cancelar reagendamento em lote", response_text)
        _send_telegram_message(token, chat_id, response_text)

    elif data in ("webmem_keep_old", "webmem_keep_new"):
        decision = "manter_existente" if data == "webmem_keep_old" else "substituir_pelo_novo"
        _answer_callback_query(token, query_id, "Resolvendo memória...")
        conflict = _pending_web_card("memory_conflict")
        if not conflict:
            response_text = "Não encontrei um conflito de memória pendente."
            _persist_callback_turn("Botão: resolver conflito de memória", response_text)
            _send_telegram_message(token, chat_id, response_text)
            return https_fn.Response("OK", status=200)
        try:
            message_id = _find_latest_copilot_card_message_id(db, copilot_session_id, "pendingMemoryConflict")
            result = _call_web_callable(
                function_name="confirmarConflitoMemoria",
                data={
                    "sessionId": copilot_session_id,
                    "messageId": message_id,
                    "memoriaId": conflict.get("memoria_id") or conflict.get("memory_id") or conflict.get("id"),
                    "decisao": decision,
                    "fatoAtualizado": conflict.get("proposed_text") or "",
                    "categoria": conflict.get("categoria") or "fato_isolado",
                },
                user_uid=session.get("userId"),
                timeout=60,
            )
            _clear_pending_web_card("memory_conflict")
            _save_session(db, chat_id, session)
            response_text = f"Conflito de memória resolvido: {result.get('decision', decision)}."
            _persist_callback_turn("Botão: resolver conflito de memória", response_text)
            _send_telegram_message(token, chat_id, response_text)
        except Exception as exc:
            response_text = f"Erro ao resolver conflito de memória: {exc}"
            _persist_callback_turn("Botão: resolver conflito de memória", response_text)
            _send_telegram_message(token, chat_id, response_text)

    elif data == "confirm_acao":
        _answer_callback_query(token, query_id, "Processando registro...")
        pending = session.get("pending_confirmations", {}).get("acao")
        if not pending:
            response_text = "⚠️ Nenhuma ação pendente de confirmação encontrada ou o prazo expirou."
            _persist_callback_turn("Botão: confirmar registro de ação", response_text)
            _send_telegram_message(token, chat_id, response_text)
            return https_fn.Response("OK", status=200)
            
        # Executa a criação (lógica de criar_acao_no_sistema)
        try:
            import uuid as _uuid
            now_iso = datetime.now(timezone.utc).isoformat()
            task_id = str(_uuid.uuid4())[:20]
            
            # Reuso da lógica de reagendamento se houver horários
            try:
                from main import get_calendar_service, get_target_calendar_id
                import hermes_calendar_tools as hc_tools
                c_service = get_calendar_service()
                c_id = get_target_calendar_id(db)
                if c_service and c_id and pending.get("horario_inicio") and pending.get("horario_fim"):
                    hc_tools.reagendar_acoes_hermes(db, c_service, c_id, pending.get("data_limite"), pending.get("horario_inicio"), pending.get("horario_fim"))
            except Exception: pass

            plano_convertido = [
                {"id": str(_uuid.uuid4())[:8], "text": str(p), "completed": False}
                for p in (pending.get("plano_acao") or []) if str(p).strip()
            ]
            
            doc = {
                "id": task_id,
                "titulo": pending["titulo"].strip(),
                "descricao": pending.get("descricao") or "",
                "area_tematica": pending.get("area_tematica") or "GERAL",
                "data_limite": pending.get("data_limite"),
                "horario_inicio": pending.get("horario_inicio"),
                "horario_fim": pending.get("horario_fim"),
                "tipo_acao": pending.get("tipo_acao") or "fast",
                "tags": pending.get("tags") or [],
                "notas": pending.get("notas") or "",
                "plano_acao": plano_convertido,
                "status": "em andamento",
                "criado_em": now_iso,
                "data_criacao": now_iso,
                "origem_ingestao": "telegram",
                "sync_status": "new",
            }
            db.collection("tarefas").document(task_id).set(doc)
            
            # Limpa pendência
            session.get("pending_confirmations", {}).pop("acao", None)
            if session.get("_pending_confirm_type") == "acao":
                session.pop("_pending_confirm_type", None)
            _save_session(db, chat_id, session)

            response_text = f"✅ <b>Ação registrada com sucesso!</b>\nID: <code>{task_id}</code>\nTítulo: {pending['titulo']}"
            _persist_callback_turn("Botão: confirmar registro de ação", response_text)
            _send_telegram_message(token, chat_id, response_text)
        except Exception as e:
            response_text = f"❌ Erro ao registrar ação: {e}"
            _persist_callback_turn("Botão: confirmar registro de ação", response_text)
            _send_telegram_message(token, chat_id, response_text)

    elif data == "cancel_acao":
        _answer_callback_query(token, query_id, "Cancelado.")
        session.get("pending_confirmations", {}).pop("acao", None)
        if session.get("_pending_confirm_type") == "acao":
            session.pop("_pending_confirm_type", None)
        _save_session(db, chat_id, session)
        response_text = "❌ Registro de ação cancelado pelo usuário."
        _persist_callback_turn("Botão: cancelar registro de ação", response_text)
        _send_telegram_message(token, chat_id, response_text)

    elif data == "confirm_financeiro":
        _answer_callback_query(token, query_id, "Processando lançamento...")
        pending = session.get("pending_confirmations", {}).get("financeiro")
        if not pending:
            response_text = "⚠️ Nenhum lançamento financeiro pendente."
            _persist_callback_turn("Botão: confirmar lançamento financeiro", response_text)
            _send_telegram_message(token, chat_id, response_text)
            return https_fn.Response("OK", status=200)
            
        try:
            from tools.telegram_extended import execute
            res = execute("registrar_item_financeiro_v2", pending, db)
            
            session.get("pending_confirmations", {}).pop("financeiro", None)
            if session.get("_pending_confirm_type") == "financeiro":
                session.pop("_pending_confirm_type", None)
            _save_session(db, chat_id, session)

            response_text = f"✅ <b>Lançamento financeiro realizado!</b>\n{res}"
            _persist_callback_turn("Botão: confirmar lançamento financeiro", response_text)
            _send_telegram_message(token, chat_id, response_text)
        except Exception as e:
            response_text = f"❌ Erro no financeiro: {e}"
            _persist_callback_turn("Botão: confirmar lançamento financeiro", response_text)
            _send_telegram_message(token, chat_id, response_text)

    elif data == "cancel_financeiro":
        _answer_callback_query(token, query_id, "Cancelado.")
        session.get("pending_confirmations", {}).pop("financeiro", None)
        if session.get("_pending_confirm_type") == "financeiro":
            session.pop("_pending_confirm_type", None)
        _save_session(db, chat_id, session)
        response_text = "❌ Lançamento financeiro descartado."
        _persist_callback_turn("Botão: cancelar lançamento financeiro", response_text)
        _send_telegram_message(token, chat_id, response_text)

    elif data == "confirm_whatsapp":
        _answer_callback_query(token, query_id, "Enfileirando mensagem...")
        pending = session.get("pending_confirmations", {}).get("whatsapp")
        if not pending:
            response_text = "⚠️ Nenhuma mensagem de WhatsApp pendente."
            _persist_callback_turn("Botão: confirmar WhatsApp", response_text)
            _send_telegram_message(token, chat_id, response_text)
            return https_fn.Response("OK", status=200)

        try:
            from tools.schedule_whatsapp_message import schedule_whatsapp_message as _schedule_whatsapp
            res = _schedule_whatsapp(
                db,
                pending.get("contact_number") or "",
                pending.get("message") or "",
                pending.get("scheduled_time") or "",
            )

            session.get("pending_confirmations", {}).pop("whatsapp", None)
            if session.get("_pending_confirm_type") == "whatsapp":
                session.pop("_pending_confirm_type", None)
            _save_session(db, chat_id, session)

            response_text = f"✅ <b>WhatsApp confirmado.</b>\n{html.escape(str(res))}"
            _persist_callback_turn("Botão: confirmar WhatsApp", response_text)
            _send_telegram_message(token, chat_id, response_text)
        except Exception as e:
            response_text = f"❌ Erro ao enfileirar WhatsApp: {html.escape(str(e))}"
            _persist_callback_turn("Botão: confirmar WhatsApp", response_text)
            _send_telegram_message(token, chat_id, response_text)

    elif data == "cancel_whatsapp":
        _answer_callback_query(token, query_id, "Cancelado.")
        session.get("pending_confirmations", {}).pop("whatsapp", None)
        if session.get("_pending_confirm_type") == "whatsapp":
            session.pop("_pending_confirm_type", None)
        _save_session(db, chat_id, session)
        response_text = "❌ Envio de WhatsApp descartado."
        _persist_callback_turn("Botão: cancelar WhatsApp", response_text)
        _send_telegram_message(token, chat_id, response_text)

    else:

        _answer_callback_query(token, query_id)

    return https_fn.Response("OK", status=200)


# ---------------------------------------------------------------------------
# Audio transcription (reusing Groq / Gemini pattern from main.py)
# ---------------------------------------------------------------------------

def _transcribe_audio_bytes(audio_bytes: bytes, extension: str, db) -> str:
    """Transcreve áudio usando Groq (Whisper) com fallback para Gemini."""
    import base64 as _b64
    import tempfile
    import os
    keys = _get_api_keys(db)
    
    # 1. Tentativa com Groq (Whisper-Large-V3-Turbo) - Alta qualidade para PT-BR
    groq_key = keys.get("groq_api_key")
    if groq_key:
        try:
            from groq import Groq
            client = Groq(api_key=groq_key)
            suffix = extension if extension.startswith(".") else f".{extension}"
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(audio_bytes)
                tmp_path = tmp.name
            try:
                with open(tmp_path, "rb") as f:
                    transcription = client.audio.transcriptions.create(
                        file=(os.path.basename(tmp_path), f),
                        model="whisper-large-v3-turbo",
                        response_format="json",
                        language="pt",
                        temperature=0.0,
                    )
                if transcription and hasattr(transcription, "text") and transcription.text:
                    return transcription.text.strip()
            finally:
                try:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)
                except OSError:
                    pass
        except Exception as groq_err:
            print(f"[Transcription] Groq failed: {groq_err}")

    # 2. Fallback: Gemini (Multimodal)
    gemini_key = keys.get("gemini_api_key")
    if gemini_key:
        try:
            from google import genai
            from google.genai import types
            client = genai.Client(api_key=gemini_key)
            
            # Mapeamento básico de extensões para mime-types
            mime_map = {
                "ogg": "audio/ogg", "oga": "audio/ogg", "wav": "audio/wav",
                "mp3": "audio/mpeg", "m4a": "audio/mp4", "flac": "audio/flac"
            }
            clean_ext = extension.lower().strip(".")
            mime_type = mime_map.get(clean_ext, "audio/ogg")
            
            response = generate_content_logged(
                client,
                model=GEMINI_TRANSCRIPTION_MODEL,
                contents=[
                    types.Part.from_bytes(data=audio_bytes, mime_type=mime_type),
                    "Transcreva este áudio literalmente. Responda apenas com o texto transcrito, sem introduções ou explicações."
                ],
                feature="telegram.audio_transcription_fallback",
            )
            if response and response.text:
                return response.text.strip()
        except Exception as gemini_err:
            print(f"[Transcription] Gemini fallback failed: {gemini_err}")

    return "[Transcrição indisponível: erro nos motores Groq/Gemini]"


# ---------------------------------------------------------------------------
# Media upload to Firebase Storage
# ---------------------------------------------------------------------------


# ---------------------------------------------------------------------------
# Storage public URL helper (works with Uniform Bucket-Level Access)
# ---------------------------------------------------------------------------

def _blob_public_url(blob) -> str:
    """
    Returns a public URL for a Firebase Storage blob.
    Generates and patches a Firebase Storage download token metadata to ensure the link is publicly accessible.
    """
    import urllib.parse
    import uuid
    
    token = None
    if blob.metadata:
        token = blob.metadata.get("firebaseStorageDownloadTokens")
        
    if not token:
        token = str(uuid.uuid4())
        metadata = dict(blob.metadata or {})
        metadata["firebaseStorageDownloadTokens"] = token
        blob.metadata = metadata
        try:
            blob.patch()
        except Exception as e:
            print(f"[Storage] Warning: Failed to patch blob metadata for token: {e}")
            
    bucket_name = blob.bucket.name
    encoded_name = urllib.parse.quote(blob.name, safe='')
    if token:
        return f"https://firebasestorage.googleapis.com/v0/b/{bucket_name}/o/{encoded_name}?alt=media&token={token}"
    return f"https://firebasestorage.googleapis.com/v0/b/{bucket_name}/o/{encoded_name}?alt=media"


def _upload_to_storage(file_bytes: bytes, file_name: str, mime_type: str, chat_id: str) -> str:
    """Uploads to Firebase Storage and returns the public URL."""
    bucket = _get_hermes_storage_bucket()
    path = f"telegram_uploads/{chat_id}/{file_name}"
    blob = bucket.blob(path)
    blob.upload_from_string(file_bytes, content_type=mime_type)
    return _blob_public_url(blob)


def _store_telegram_media(file_bytes: bytes, file_name: str, mime_type: str, chat_id: str) -> tuple[str | None, str | None]:
    """
    Stores media inline when small enough; otherwise uploads to Storage and returns the blob path.
    Returns: (media_bytes_b64, storage_path)
    """
    import base64

    if len(file_bytes) <= _MAX_INLINE_MEDIA_BYTES:
        return base64.b64encode(file_bytes).decode(), None

    safe_name = os.path.basename(file_name or f"upload_{int(time.time())}")
    bucket = _get_hermes_storage_bucket()
    path = f"telegram_uploads/{chat_id}/{int(time.time())}_{safe_name}"
    blob = bucket.blob(path)
    blob.upload_from_string(file_bytes, content_type=mime_type)
    return None, path


def _load_telegram_media_bytes(media_bytes_b64: str | None, storage_path: str | None) -> bytes | None:
    if media_bytes_b64:
        import base64

        return base64.b64decode(media_bytes_b64)
    if storage_path:
        bucket = _get_hermes_storage_bucket()
        blob = bucket.blob(storage_path)
        return blob.download_as_bytes()
    return None


# ---------------------------------------------------------------------------
# Gemini orchestration
# ---------------------------------------------------------------------------

_AREAS_VALIDAS_CACHE: tuple | None = None  # (monotonic_ts, list[str])
_AREAS_VALIDAS_TTL = 300  # unidades mudam raramente; evita stream da coleção a cada mensagem


def carregar_areas_tematicas_validas(db) -> list[str]:
    """Lista canonica de areas tematicas = nomes das Unidades + GERAL/NAO CLASSIFICADA."""
    global _AREAS_VALIDAS_CACHE
    now = time.monotonic()
    if _AREAS_VALIDAS_CACHE and (now - _AREAS_VALIDAS_CACHE[0]) < _AREAS_VALIDAS_TTL:
        return _AREAS_VALIDAS_CACHE[1]
    nomes = []
    try:
        for d in db.collection('unidades').stream():
            nome = ((d.to_dict() or {}).get('nome') or '').strip().upper()
            if nome:
                nomes.append(nome)
    except Exception as e:
        print(f"[areas_validas] Falha ao carregar unidades: {e}")
    base = ['GERAL', 'NÃO CLASSIFICADA']
    # preserva ordem e remove duplicados
    resultado = base + [n for n in nomes if n not in base]
    if nomes:  # não cacheia resultado de falha
        _AREAS_VALIDAS_CACHE = (now, resultado)
    return resultado


def normalizar_area_tematica(valor: str, areas_validas: list[str]) -> str:
    """Casa case-insensitive contra a lista valida; fallback 'GERAL'."""
    alvo = (valor or '').strip().upper()
    for a in areas_validas:
        if a.upper() == alvo:
            return a
    return 'GERAL'


def _build_system_instruction(copilot_core: str, copilot_soul: str, contexto_ativo: str) -> str:
    from datetime import datetime as _dt
    from zoneinfo import ZoneInfo
    tz = ZoneInfo("America/Sao_Paulo")
    now_local = _dt.now(tz)
    today = now_local.strftime("%Y-%m-%d")
    time_str = now_local.strftime("%H:%M")
    ctx_hint = (
        f"\n\n<b>Contexto ativo:</b> {contexto_ativo}"
        if contexto_ativo != "geral"
        else ""
    )
    return (
        f"Você é o Copiloto Hermes, estrategista sênior de processos. Hoje é {today} e o horário local atual é {time_str}."
        f"{ctx_hint}\n\n"
        "## CORE ESTÁTICO DO COPILOTO\n"
        f"{copilot_core}\n\n"
        "## PERSONALIDADE DINÂMICA ATUAL\n"
        f"{copilot_soul}\n\n"
        "## CANAL DE COMUNICAÇÃO: TELEGRAM\n"
        "REGRA ABSOLUTA DE FORMATAÇÃO: Você está respondendo via Telegram. "
        "Use EXCLUSIVAMENTE as seguintes tags HTML suportadas pelo Telegram: "
        "<b>negrito</b>, <i>itálico</i>, <code>código inline</code>, <pre>bloco de código</pre>. "
        "PROIBIDO usar Markdown (asteriscos, underlines, backticks, #, ##). "
        "Listas: use • ou - como prefixo de linha, sem Markdown. "
        "Mantenha respostas concisas — Telegram tem limite de 4096 caracteres por mensagem.\n\n"
        "## REGRAS ABSOLUTAS\n"
        "1. JAMAIS expanda siglas arbitrariamente.\n"
        "2. Se qualquer ferramenta retornar campo 'erro', reproduza o erro literal.\n"
        "3. Agendamento/Agenda: Por padrão, ao propor ou criar uma ação, NÃO proponha nem defina horários de início/fim (campos horario_inicio e horario_fim devem ser nulos/vazios), definindo apenas o dia (data_limite). SÓ preencha horario_inicio e horario_fim se o usuário pedir explicitamente para agendar um horário específico. Nesse caso, você DEVE usar consultar_agenda ou encontrar_slot_livre ANTES de agendar. Horário de funcionamento: 08:00 às 19:00, janela D+7. Se o agendamento for para hoje, o horário inicial DEVE ser sempre posterior ao horário local atual. Se houver conflito em horário específico, pergunte se força inserção ou busca outro slot.\n"
        "4. Para criar uma nova ação, você DEVE usar obrigatoriamente a ferramenta propor_acao_para_confirmacao. Ela gerará botões de ✅/❌ para o usuário confirmar o registro.\n"
        "5. Links de tarefas: use o formato task:{ID} no texto (ex: 'Ação task:abc123').\n"
        "6. Acione salvar_memoria_global apenas para fatos duráveis e preferências estáveis.\n"
        "7. ACESSO FINANCEIRO: para consultas sobre balanço, rendas, obrigações ou metas, use consultar_financas_v2. Para novos registros, use obrigatoriamente propor_lancamento_financeiro para que o usuário receba botões de confirmação. NUNCA invente números.\n"
    )

def _build_system_instruction_guarded(
    copilot_core: str,
    copilot_soul: str,
    contexto_ativo: str,
    acao_snapshot: dict | None = None,
) -> str:
    from datetime import datetime as _dt
    from zoneinfo import ZoneInfo
    tz = ZoneInfo("America/Sao_Paulo")
    now_local = _dt.now(tz)
    today = now_local.strftime("%Y-%m-%d")
    time_str = now_local.strftime("%H:%M")
    ctx_hint = (
        f"\n\n<b>Contexto ativo:</b> {contexto_ativo}"
        if contexto_ativo != "geral"
        else ""
    )
    base = (
        f"Voce e o Copiloto Hermes, estrategista senior de processos. Hoje e {today} e o horario local atual e {time_str}."
        f"{ctx_hint}\n\n"
        "## CORE ESTATICO DO COPILOTO\n"
        f"{copilot_core}\n\n"
        "## PERSONALIDADE DINAMICA ATUAL\n"
        f"{copilot_soul}\n\n"
        "## CANAL DE COMUNICACAO: TELEGRAM\n"
        "REGRA ABSOLUTA DE FORMATACAO: Voce esta respondendo via Telegram. "
        "Use EXCLUSIVAMENTE as seguintes tags HTML suportadas pelo Telegram: "
        "<b>negrito</b>, <i>italico</i>, <code>codigo inline</code>, <pre>bloco de codigo</pre>. "
        "PROIBIDO usar Markdown. "
        "Listas: use • ou - como prefixo de linha, sem Markdown. "
        "Mantenha respostas concisas; Telegram tem limite de 4096 caracteres por mensagem.\n\n"
        "## REGRAS ABSOLUTAS\n"
        "1. JAMAIS expanda siglas arbitrariamente.\n"
        "2. Se qualquer ferramenta retornar campo 'erro', reproduza o erro literal.\n"
        "4. Se o pedido for sobre dados internos do Hermes, tarefas, acoes, historico, agenda ou documentos do sistema, NUNCA use internet como fallback. Nesses casos, use apenas ferramentas internas. Se nao encontrar nada apos usar as ferramentas, admita explicitamente que nao encontrou nos registros do sistema.\n"
        "5. Agendamento/Agenda: Por padrao, ao propor ou criar uma acao, NAO proponha nem defina horarios de inicio/fim (campos horario_inicio e horario_fim devem ser nulos/vazios), definindo apenas o dia (data_limite). SO preencha horario_inicio e horario_fim se o usuario pedir explicitamente para agendar um horario especifico. Nesse caso, voce DEVE usar consultar_agenda ou encontrar_slot_livre ANTES de agendar. Horario de funcionamento: 08:00 as 19:00, janela D+7. Se o agendamento for para hoje, o horario inicial DEVE ser sempre posterior ao horario local atual. Se houver conflito em horario especifico, pergunte se forca insercao ou busca outro slot.\n"
        "6. Para criar uma nova ação, você DEVE usar obrigatoriamente a ferramenta propor_acao_para_confirmacao. Isso apresentará os botões de ✅/❌ ao usuário.\n"
        "7. Links de tarefas: use o formato task:{ID} no texto (ex: 'Acao task:abc123').\n"
        "8. Acione salvar_memoria_global apenas para fatos duraveis e preferencias estaveis.\n"
        "9. GOVERNANCA DE FONTES: ao descrever uma tarefa encontrada por consultar_historico_acoes, use SOMENTE os campos retornados por essa ferramenta. Se um campo nao constar no retorno da ferramenta, diga 'nao informado' em vez de inventar.\n"
        "10. ACESSO FINANCEIRO: use exclusivamente consultar_financas_v2 e para novos registros use propor_lancamento_financeiro para lidar com dados financeiros internos (rendas, contas, metas). Proibido inventar valores.\n"
    )

    if not acao_snapshot:
        return base

    titulo = acao_snapshot.get("titulo", "")
    plano_lines = "\n".join(f"  {p}" for p in acao_snapshot.get("plano_acao", [])) or "  (sem passos registrados)"
    diario_lines = "\n".join(f"  {d}" for d in acao_snapshot.get("acompanhamento_recente", [])) or "  (sem entradas recentes)"
    tags_str = ", ".join(str(t) for t in acao_snapshot.get("tags", [])) or "nenhuma"

    acao_section = (
        "\n\n## ACAO VINCULADA — MODO CONTEXTO TRANCADO\n"
        f"ID: {acao_snapshot.get('id', '')}\n"
        f"Titulo: {titulo}\n"
        f"Area: {acao_snapshot.get('area', '')}\n"
        f"Status: {acao_snapshot.get('status', '')}\n"
        f"Descricao: {acao_snapshot.get('descricao', '') or 'nao informada'}\n"
        f"Sintese: {acao_snapshot.get('sintese', '') or 'nao informada'}\n"
        f"Plano de Acao:\n{plano_lines}\n"
        f"Diario Recente:\n{diario_lines}\n"
        f"Notas: {acao_snapshot.get('notas', '') or 'nenhuma'}\n"
        f"Tags: {tags_str}\n\n"
        "MODO CONTEXTO TRANCADO ATIVADO — REGRAS ABSOLUTAS:\n"
        "1. O historico de conversas anteriores foi isolado. Concentre-se exclusivamente nesta acao.\n"
        "2. Use SOMENTE os dados desta secao ao descrever a acao. PROIBIDO usar RAG ou inferencia.\n"
        "3. Ao receber perguntas sobre a acao, responda com base nos dados acima.\n"
        "4. Se o usuario pedir para atualizar, criar passos ou registrar progresso, use as ferramentas disponiveis.\n"
    )
    return base + acao_section


def _build_system_instruction_guarded_v2(
    copilot_core: str,
    copilot_soul: str,
    contexto_ativo: str,
    acao_snapshot: dict | None = None,
) -> str:
    from datetime import datetime as _dt
    from zoneinfo import ZoneInfo
    tz = ZoneInfo("America/Sao_Paulo")
    now_local = _dt.now(tz)
    today = now_local.strftime("%Y-%m-%d")
    time_str = now_local.strftime("%H:%M")
    ctx_hint = (
        f"\n\n<b>Contexto ativo:</b> {contexto_ativo}"
        if contexto_ativo != "geral"
        else ""
    )
    base = (
        f"Voce e o Copiloto Hermes, estrategista senior de processos. Hoje e {today} e o horario local atual e {time_str}."
        f"{ctx_hint}\n\n"
        "## CORE ESTATICO DO COPILOTO\n"
        f"{copilot_core}\n\n"
        "## PERSONALIDADE DINAMICA ATUAL\n"
        f"{copilot_soul}\n\n"
        "## CANAL DE COMUNICACAO: TELEGRAM\n"
        "REGRA ABSOLUTA DE FORMATACAO: Voce esta respondendo via Telegram. "
        "Use EXCLUSIVAMENTE as seguintes tags HTML suportadas pelo Telegram: "
        "<b>negrito</b>, <i>italico</i>, <code>codigo inline</code>, <pre>bloco de codigo</pre>. "
        "PROIBIDO usar Markdown. "
        "Listas: use - ou bullet simples como prefixo de linha, sem Markdown. "
        "Mantenha respostas concisas; Telegram tem limite de 4096 caracteres por mensagem.\n\n"
        "## REGRAS ABSOLUTAS\n"
        "1. JAMAIS expanda siglas arbitrariamente.\n"
        "2. Se qualquer ferramenta retornar campo 'erro', reproduza o erro literal.\n"
        "4. Se o pedido for sobre dados internos do Hermes, tarefas, acoes, historico, agenda ou documentos do sistema, NUNCA use internet como fallback. Nesses casos, use apenas ferramentas internas. Se nao encontrar nada apos usar as ferramentas, admita explicitamente que nao encontrou nos registros do sistema.\n"
        "5. Agendamento/Agenda: Por padrao, ao propor ou criar uma acao, NAO proponha nem defina horarios de inicio/fim (campos horario_inicio e horario_fim devem ser nulos/vazios), definindo apenas o dia (data_limite). SO preencha horario_inicio e horario_fim se o usuario pedir explicitamente para agendar um horario especifico. Nesse caso, voce DEVE usar consultar_agenda ou encontrar_slot_livre ANTES de agendar. Horario de funcionamento: 08:00 as 19:00, janela D+7. Se o agendamento for para hoje, o horario inicial DEVE ser sempre posterior ao horario local atual. Se houver conflito em horario especifico, pergunte se forca insercao ou busca outro slot.\n"
        "6. Para criar uma nova ação, você DEVE usar obrigatoriamente a ferramenta propor_acao_para_confirmacao. Isso apresentará os botões de ✅/❌ ao usuário.\n"
        "7. Links de tarefas: use o formato task:{ID} no texto (ex: 'Acao task:abc123').\n"
        "8. Acione salvar_memoria_global apenas para fatos duraveis e preferencias estaveis.\n"
        "9. GOVERNANCA DE FONTES: ao descrever uma tarefa encontrada por consultar_historico_acoes, use SOMENTE os campos retornados por essa ferramenta. Se um campo nao constar no retorno da ferramenta, diga 'nao informado' em vez de inventar.\n"
        "10. LINKS E ARQUIVOS: nunca crie hiperlinks, texto-ancora ou URLs que nao aparecam literalmente no contexto, em uma ferramenta ou em um DRIVE_FILE_ID real. Se o usuario pedir links e eles nao estiverem disponiveis, diga que nao encontrou.\n"
        "11. PESQUISA DE ACOES: se o usuario pedir para pesquisar/localizar uma acao ou tarefa, use consultar_historico_acoes primeiro. Nao substitua resultado ausente por acervo, email ou internet, salvo se o usuario pedir explicitamente essa ampliacao.\n"
        "12. ACESSO FINANCEIRO: para qualquer dado sobre rendas, contas, metas ou balanco interno, use consultar_financas_v2. Para novos registros, use obrigatoriamente propor_lancamento_financeiro para que o usuário receba os botões de confirmação. Detalhe os valores com precisao absoluta conforme retornado pelo sistema.\n"
        "13. EFICIENCIA: quando precisar de varias consultas independentes, solicite todas na mesma rodada de ferramentas. Evite rodadas sequenciais se uma unica rodada paralela resolver. Nunca chame mais de uma ferramenta de escrita/registro no mesmo turno; proponha uma confirmacao por vez.\n"
        "14. WHATSAPP: para enviar ou agendar mensagem de WhatsApp, use schedule_whatsapp_message. A ferramenta deve apenas preparar a proposta; o envio real depende de confirmacao por botao.\n"
    )

    if not acao_snapshot:
        return base

    titulo = acao_snapshot.get("titulo", "")
    plano_lines = "\n".join(f"  {p}" for p in acao_snapshot.get("plano_acao", [])) or "  (sem passos registrados)"
    diario_lines = "\n".join(f"  {d}" for d in acao_snapshot.get("acompanhamento_recente", [])) or "  (sem entradas recentes)"
    diario_integral = acao_snapshot.get("diario_integral", "") or "(sem diario integral)"
    arquivos = acao_snapshot.get("arquivos_disponiveis", []) or []
    arquivos_lines = (
        "\n".join(f"  - {a.get('nome', 'Arquivo sem nome')} | DRIVE_FILE_ID: {a.get('drive_file_id', '')}" for a in arquivos)
        or "  (sem arquivos disponiveis)"
    )
    tags_str = ", ".join(str(t) for t in acao_snapshot.get("tags", [])) or "nenhuma"

    acao_section = (
        "\n\n## ACAO VINCULADA - MODO CONTEXTO TRANCADO\n"
        f"ID: {acao_snapshot.get('id', '')}\n"
        f"Titulo: {titulo}\n"
        f"Area: {acao_snapshot.get('area_tematica', acao_snapshot.get('area', ''))}\n"
        f"Status: {acao_snapshot.get('status', '')}\n"
        f"Descricao: {acao_snapshot.get('descricao', '') or 'nao informada'}\n"
        f"Sintese: {acao_snapshot.get('sintese', '') or 'nao informada'}\n"
        f"Plano de Acao:\n{plano_lines}\n"
        f"Plano Atual Estruturado: {json.dumps(acao_snapshot.get('plano_atual', []), ensure_ascii=False)}\n"
        f"Diario Recente:\n{diario_lines}\n"
        f"Diario Integral:\n{diario_integral}\n"
        f"Notas: {acao_snapshot.get('notas', '') or 'nenhuma'}\n"
        f"Tags: {tags_str}\n\n"
        f"Arquivos Disponiveis:\n{arquivos_lines}\n\n"
        "MODO CONTEXTO TRANCADO ATIVADO - REGRAS ABSOLUTAS:\n"
        "1. O historico de conversas anteriores foi isolado. Concentre-se exclusivamente nesta acao.\n"
        "2. Use SOMENTE os dados desta secao ao descrever a acao. PROIBIDO usar RAG ou inferencia.\n"
        "3. Este payload replica a estrutura de contexto da interface web: plano atual, diario integral, tags e arquivos disponiveis.\n"
        "4. Ao receber perguntas sobre a acao, responda com base nos dados acima.\n"
        "5. Se o usuario pedir para atualizar, criar passos ou registrar progresso, use as ferramentas disponiveis.\n"
        "6. Se o usuario pedir links ou anexos, cite apenas as URLs/DRIVE_FILE_ID listados em Arquivos Disponiveis. Nunca use rotulos clicaveis como 'Acessar Pasta' sem mostrar a URL real.\n"
        "7. PROIBIDO usar consultar_historico_acoes neste modo — a acao ja esta carregada acima. Use apenas ferramentas de atualizacao (plano, diario, notas).\n"
    )
    return base + acao_section


def _normalize_for_matching(text: str) -> str:
    normalized = unicodedata.normalize("NFKD", text or "")
    normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    return normalized.lower()


def _is_explicit_web_request(text: str) -> bool:
    lowered = _normalize_for_matching(text)
    web_markers = [
        "na internet",
        "na web",
        "pesquise na internet",
        "busque na internet",
        "pesquise no google",
        "busque no google",
        "noticia",
        "noticias",
        "site oficial",
        "link externo",
        "fonte externa",
        "fonte oficial",
    ]
    return any(marker in lowered for marker in web_markers)


def _extract_document_text(file_bytes: bytes, filename: str, mime_type: str) -> str | None:
    import io
    name = filename.lower()
    mime = mime_type.lower()
    
    # Word DOCX
    if mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or name.endswith('.docx'):
        try:
            import mammoth
            result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
            return (result.value or "").strip()
        except Exception as e:
            print(f"[Parser] Erro DOCX: {e}")
            
    # Excel XLSX / XLS
    elif mime in ('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel') or name.endswith(('.xlsx', '.xls')):
        try:
            import pandas as pd
            df_list = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            sheets_text = []
            for sheet_name, df in df_list.items():
                sheets_text.append(f"ABA: {sheet_name}\n{df.to_csv(index=False)}")
            return "\n\n".join(sheets_text).strip()
        except Exception as e:
            print(f"[Parser] Erro Excel: {e}")
            
    # PowerPoint PPTX
    elif mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or name.endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slides_text.append(shape.text.strip())
            return '\n'.join(slides_text).strip()
        except Exception as e:
            print(f"[Parser] Erro PPTX: {e}")
            
    # Plain text / CSV
    elif mime.startswith('text/') or name.endswith(('.txt', '.csv', '.json', '.xml', '.yaml', '.yml', '.md')):
        try:
            return file_bytes.decode('utf-8', errors='replace')
        except Exception as e:
            print(f"[Parser] Erro Texto: {e}")
            
    return None


def _is_internal_hermes_request(text: str) -> bool:
    lowered = _normalize_for_matching(text)
    internal_markers = [
        "hermes",
        "no sistema",
        "no copilot",
        "copiloto",
        "acao",
        "acoes",
        "tarefa",
        "tarefas",
        "agendamento",
        "documentacao",
        "deferido",
        "deferidos",
        "assistencia estudantil",
        "historico",
        "cadastro",
    ]
    return any(marker in lowered for marker in internal_markers)


def _extract_response_mode(text: str, session: dict) -> tuple[str, str]:
    raw_text = (text or "").strip()
    if not raw_text:
        return "texto", raw_text

    lowered = _normalize_for_matching(raw_text)
    explicit_audio_markers = (
        "/audio",
        "#audio",
        "[audio]",
        "mensagem de voz",
        "em voz",
        "por voz",
    )
    audio_keywords = ("audio", "voz")
    request_verbs = ("responda", "resposta", "mande", "manda", "envie", "quero")
    wants_audio = any(marker in lowered for marker in explicit_audio_markers)
    if not wants_audio:
        wants_audio = any(v in lowered for v in request_verbs) and any(k in lowered for k in audio_keywords)

    if wants_audio:
        cleaned = raw_text
        cleaned = re.sub(r"^/audio\b", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"#audio|\[audio\]", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"(me\s+)?responda\s+(em|por|para|pro?)\s+(o\s+)?[aá]udio", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"resposta\s+(em|por|para|pro?)\s+(o\s+)?[aá]udio", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"envie\s+(em|por|para|pro?)\s+(o\s+)?[aá]udio", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"(me\s+)?mande\s+(em|por|para|pro?)\s+(o\s+)?[aá]udio", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"quero\s+(em|por|para|pro?)\s+(o\s+)?[aá]udio", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"mensagem de voz", "", cleaned, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"\b(em|por)\s+voz\b", "", cleaned, flags=re.IGNORECASE).strip()
        return "audio", cleaned.strip(" ,.-")

    explicit_text_markers = ("/texto", "#texto", "[texto]")
    if any(marker in lowered for marker in explicit_text_markers):
        cleaned = re.sub(r"^/texto\b|#texto|\[texto\]", "", raw_text, flags=re.IGNORECASE).strip()
        return "texto", cleaned

    return "texto", raw_text


def _extract_voice_profile(text: str, default_voice: str = "masculina") -> tuple[str, str]:
    raw_text = (text or "").strip()
    if not raw_text:
        return default_voice, raw_text

    lowered = _normalize_for_matching(raw_text)
    masculine_markers = (
        "#voz:masculina",
        "[voz:masculina]",
        "/voz masculina",
        "voz masculina",
        "na voz masculina",
        "com voz masculina",
    )

    if any(marker in lowered for marker in masculine_markers):
        cleaned = re.sub(r"#voz:masculina|\[voz:masculina\]|/voz masculina", "", raw_text, flags=re.IGNORECASE).strip()
        cleaned = re.sub(r"\b(na|com)?\s*voz masculina\b", "", cleaned, flags=re.IGNORECASE).strip()
        return "masculina", cleaned

    return default_voice, raw_text


def _run_gemini_text(gemini_key: str, system_instruction: str, user_prompt: str, model_id: str = _TEXT_MODEL_ID) -> str:
    from google import genai
    from google.genai import types

    client = genai.Client(api_key=gemini_key)
    response = generate_content_logged(
        client,
        model=model_id,
        contents=user_prompt,
        feature="telegram.tts_script",
        config=types.GenerateContentConfig(system_instruction=system_instruction),
    )
    return (response.text or "").strip()


def _build_tts_director_instruction(voice_profile: str) -> str:
    return (
        "Voce e o diretor de TTS do Hermes. "
        "Recebera uma resposta factual pronta e deve apenas converte-la em um roteiro curto para fala natural. "
        "Nao adicione fatos novos. "
        "Nao inclua links, nomes tecnicos longos nem listas densas. "
        "Nao use tags, colchetes nem anotacoes de palco. "
        "Entregue apenas o texto a ser falado. "
        "Mantenha tom de assistente operacional, conversa um a um, sem soar como narrador. "
        "Priorize objetividade, proximidade e sobriedade. "
        f"Perfil de voz desejado: {voice_profile}. "
        "Responda apenas com o roteiro final."
    )


def _run_gemini_tts(gemini_key: str, script_text: str, voice_profile: str) -> tuple[bytes, str]:
    from google import genai
    from google.genai import types

    voice_name = _DEFAULT_MALE_VOICE
    client = genai.Client(api_key=gemini_key)
    style_prompt = (
        "### DIRECTOR'S NOTES\n"
        "Style: concise operational assistant, one-to-one, practical and calm.\n"
        "Pacing: natural conversational pace, with short pauses only when meaning changes.\n"
        "Tone: professional, direct, grounded, helpful and discreet.\n"
        "Delivery: do not sound like a narrator, announcer, storyteller, presenter or commercial voice-over.\n"
        "Delivery: sound like a senior assistant speaking directly to one person.\n\n"
        "### SCRIPT\n"
        f"\"{script_text[:_MAX_TTS_TRANSCRIPT_CHARS]}\""
    )
    response = generate_content_logged(
        client,
        model=_TTS_MODEL_ID,
        contents=style_prompt,
        feature="telegram.tts_audio",
        config=types.GenerateContentConfig(
            response_modalities=["AUDIO"],
            speech_config=types.SpeechConfig(
                voice_config=types.VoiceConfig(
                    prebuilt_voice_config=types.PrebuiltVoiceConfig(voice_name=voice_name)
                )
            ),
        ),
    )

    for candidate in (response.candidates or []):
        for part in (candidate.content.parts or []):
            inline_data = getattr(part, "inline_data", None)
            if inline_data and getattr(inline_data, "data", None):
                return inline_data.data, getattr(inline_data, "mime_type", "audio/L16;rate=24000")
    raise RuntimeError("Gemini TTS nÃ£o retornou Ã¡udio.")


def _transcode_audio_for_telegram_voice(audio_bytes: bytes, mime_type: str) -> tuple[bytes, str, str] | tuple[None, None, None]:
    import shutil
    import subprocess

    sample_rate = 24000
    rate_match = re.search(r"rate=(\d+)", (mime_type or "").lower())
    if rate_match:
        sample_rate = int(rate_match.group(1))

    ffmpeg_bin = shutil.which("ffmpeg")
    if not ffmpeg_bin:
        return None, None, None

    with tempfile.TemporaryDirectory() as tmpdir:
        wav_path = os.path.join(tmpdir, "input.wav")
        ogg_path = os.path.join(tmpdir, "output.ogg")
        with wave.open(wav_path, "wb") as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(audio_bytes)
        cmd = [
            ffmpeg_bin,
            "-y",
            "-i",
            wav_path,
            "-c:a",
            "libopus",
            "-b:a",
            "32k",
            ogg_path,
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=120, check=False)
        if proc.returncode != 0 or not os.path.exists(ogg_path):
            print(f"[TTS] ffmpeg transcode failed: {proc.stderr[:400]}")
            return None, None, None
        with open(ogg_path, "rb") as f:
            return f.read(), "audio/ogg", "hermes_voice.ogg"


def _wrap_pcm_audio_as_wav(audio_bytes: bytes, mime_type: str) -> tuple[bytes, str, str]:
    sample_rate = 24000
    rate_match = re.search(r"rate=(\d+)", (mime_type or "").lower())
    if rate_match:
        sample_rate = int(rate_match.group(1))

    wav_buffer = BytesIO()
    with wave.open(wav_buffer, "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(sample_rate)
        wav_file.writeframes(audio_bytes)
    return wav_buffer.getvalue(), "audio/wav", "hermes_audio.wav"


def _send_tts_failure_notice(token: str, chat_id: str | int):
    _send_telegram_message(
        token,
        chat_id,
        "⚠️ Houve uma falha tÃ©cnica ao gerar o Ã¡udio. Repita a solicitaÃ§Ã£o ou peÃ§a a resposta em texto.",
    )


def _send_tts_failure_notice_contextual(db, token: str, chat_id: str | int, session: dict | None = None):
    _send_telegram_session_message(
        db,
        token,
        chat_id,
        "⚠️ Houve uma falha técnica ao gerar o áudio. Repita a solicitação ou peça a resposta em texto.",
        session=session,
    )


def _send_contextual_response(
    db,
    token: str,
    chat_id: str | int,
    text: str,
    session: dict | None = None,
    response_mode: str = "texto",
    gemini_key: str | None = None,
    voice_profile: str = "masculina",
    inline_keyboard: list | None = None,
    perf_state: dict | None = None,
):
    """Envia resposta factual pronta e, se pedido, converte o mesmo texto em audio."""
    message_id = _send_telegram_session_message(
        db,
        token,
        chat_id,
        text,
        session=session,
        inline_keyboard=inline_keyboard,
    )
    if response_mode != "audio":
        return message_id
    if not gemini_key:
        _send_tts_failure_notice_contextual(db, token, chat_id, session=session)
        return message_id

    try:
        with _telegram_action_heartbeat(token, chat_id, "record_voice"):
            tts_script = _run_gemini_text(
                gemini_key=gemini_key,
                system_instruction=_build_tts_director_instruction(voice_profile),
                user_prompt=text,
            )
            audio_bytes, audio_mime = _run_gemini_tts(
                gemini_key=gemini_key,
                script_text=tts_script or text,
                voice_profile=voice_profile,
            )
            tg_audio_bytes, tg_audio_mime, tg_audio_name = _transcode_audio_for_telegram_voice(audio_bytes, audio_mime)
            if tg_audio_bytes:
                if not _send_telegram_voice(token, chat_id, tg_audio_bytes, tg_audio_name, tg_audio_mime):
                    raise RuntimeError("Falha ao enviar voice note pelo Telegram.")
            else:
                wav_bytes, wav_mime, wav_name = _wrap_pcm_audio_as_wav(audio_bytes, audio_mime)
                if not _send_telegram_document(token, chat_id, wav_bytes, wav_name, wav_mime, caption="Resposta em audio"):
                    raise RuntimeError("Falha ao enviar arquivo de audio.")
        if perf_state is not None:
            _perf_mark(perf_state, "telegram.audio_response")
    except Exception as tts_err:
        print(f"[Core] Deterministic TTS error: {tts_err}")
        if perf_state is not None:
            _perf_log("telegram.request.deterministic_tts_error", perf_state, {"chat_id": str(chat_id), "error": str(tts_err)})
        _send_tts_failure_notice_contextual(db, token, chat_id, session=session)
    return message_id


def _run_gemini_turn(
    db,
    gemini_key: str,
    system_instruction: str,
    history: list,
    user_message_parts: list,
    tools_list: list,
    function_map: dict,
    perf_state: dict | None = None,
    sistema_id_contexto: str = None,
) -> str:
    """Runs a full Gemini multi-turn exchange and returns the final text response."""
    from concurrent.futures import ThreadPoolExecutor
    from google import genai
    from google.genai import types

    client = genai.Client(api_key=gemini_key)
    model_id = _TEXT_MODEL_ID
    read_only_parallel_tools = {
        "consultar_historico_acoes",
        "buscar_arquivos_acervo",
        "pesquisar_internet",
        "ler_pagina_web",
        "consultar_agenda",
        "encontrar_slot_livre",
        "consultar_financas_v2",
        "buscar_e_analisar_email",
    }

    def _is_parallel_safe_tool(tool_name: str) -> bool:
        return tool_name in read_only_parallel_tools

    def _execute_tool_call(fc, *, parallel: bool = False):
        fn = function_map.get(fc.name)
        tool_start_ms = _perf_now_ms()
        if fn is None:
            result_text = f"Ferramenta '{fc.name}' não encontrada."
        else:
            try:
                kwargs = dict(fc.args or {})
                result = fn(**kwargs)
                result_text = result if isinstance(result, str) else json.dumps(result, ensure_ascii=False)
            except Exception as tool_err:
                result_text = f"Erro ao executar {fc.name}: {tool_err}"
        if perf_state is not None:
            perf_state.setdefault("tool_calls", []).append({
                "name": fc.name,
                "duration_ms": max(0, _perf_now_ms() - tool_start_ms),
                "parallel": parallel,
            })
        return types.Part.from_function_response(
            name=fc.name,
            response={"result": result_text},
        )

    def _tool_call_groups(func_calls: list) -> list[list]:
        groups = []
        current_read_group = []
        for fc in func_calls:
            if _is_parallel_safe_tool(fc.name):
                current_read_group.append(fc)
                continue
            if current_read_group:
                groups.append(current_read_group)
                current_read_group = []
            groups.append([fc])
        if current_read_group:
            groups.append(current_read_group)
        return groups

    chat = client.chats.create(
        model=model_id,
        config=types.GenerateContentConfig(
            system_instruction=system_instruction,
            tools=tools_list,
            automatic_function_calling=types.AutomaticFunctionCallingConfig(disable=True),
            # Desabilita thinking: evita parts com thought=True que confundem a extração de texto
            thinking_config=types.ThinkingConfig(thinking_budget=0),
        ),
        history=history,
    )
    if perf_state is not None:
        _perf_mark(perf_state, "telegram.chat_create")

    response = send_message_logged(
        chat,
        user_message_parts,
        model=model_id,
        feature="telegram.first_turn",
        db=db,
    )
    if perf_state is not None:
        _perf_mark(perf_state, "telegram.first_model_response")

    # Agentic loop — processa tool calls até obter resposta de texto
    for _ in range(6):
        if not response.candidates:
            break
        candidate = response.candidates[0]
        if candidate.finish_reason and candidate.finish_reason.name not in ("STOP", "MAX_TOKENS", ""):
            break

        # Coleta function calls neste turno
        func_calls = []
        for part in (candidate.content.parts or []):
            if hasattr(part, "function_call") and part.function_call and part.function_call.name:
                func_calls.append(part.function_call)

        if not func_calls:
            break

        # Executa todas as tool calls e coleta resultados
        tool_results = []
        for group in _tool_call_groups(func_calls):
            if len(group) == 1:
                tool_results.append(_execute_tool_call(group[0], parallel=False))
                continue

            max_workers = min(len(group), 6)
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = [executor.submit(_execute_tool_call, fc, parallel=True) for fc in group]
                for future in futures:
                    tool_results.append(future.result())

        if perf_state is not None:
            parallel_count = sum(1 for fc in func_calls if _is_parallel_safe_tool(fc.name))
            serial_count = len(func_calls) - parallel_count
            perf_state.setdefault("tool_rounds", []).append({
                "total_calls": len(func_calls),
                "parallel_safe_calls": parallel_count,
                "serialized_calls": serial_count,
            })
        response = send_message_logged(
            chat,
            tool_results,
            model=model_id,
            feature="telegram.tool_roundtrip",
            db=db,
        )
        if perf_state is not None:
            _perf_mark(perf_state, "telegram.tool_roundtrip")
    else:
        # Se saiu do loop por limite de iterações, força um turno final de texto
        last_chance_msg = [types.Part(text="Limite de pesquisas atingido. Por favor, apresente uma resposta final ao usuário com base no que você encontrou (ou informe que não encontrou).")]
        response = send_message_logged(
            chat,
            last_chance_msg,
            model=model_id,
            feature="telegram.safety_final_turn",
            db=db,
        )
        if perf_state is not None:
            _perf_mark(perf_state, "telegram.safety_final_turn")

    # Extrai texto final — ignora parts de thinking (thought=True) que não são resposta
    text_parts = []
    if response.candidates:
        for part in (response.candidates[0].content.parts or []):
            # Filtra parts de raciocínio interno (thinking) — não são texto de resposta
            if getattr(part, "thought", False):
                continue
            if hasattr(part, "text") and part.text:
                text_parts.append(part.text)
    return "\n".join(text_parts).strip() or "Peço desculpas, mas não consegui formular uma resposta. Tente reformular sua pergunta."


# ---------------------------------------------------------------------------
# Main processing logic
# ---------------------------------------------------------------------------

def _process_telegram_message(db, data: dict):
    """Full processing pipeline for one incoming Telegram message."""
    from google.genai import types
    from tools.busca_grafo import buscar_tarefas
    from tools.busca_acervo import buscar_acervo
    perf_state = {"start_ms": _perf_now_ms(), "_last_ms": _perf_now_ms(), "steps": [], "tool_calls": []}

    chat_id = str(data.get("chat_id", ""))
    text = (data.get("text") or "").strip()
    file_info = data.get("file_info")   # {file_id, file_unique_id, file_size, mime_type, file_name}
    audio_info = data.get("audio_info") # {file_id, file_size, mime_type, duration}
    media_bytes_b64 = data.get("media_bytes_b64")  # base64 of already-downloaded file
    media_storage_path = data.get("media_storage_path")

    if not chat_id:
        return

    keys = _get_api_keys(db)
    gemini_key = keys.get("gemini_api_key")
    token = keys.get("telegram_bot_token") or os.environ.get("TELEGRAM_BOT_TOKEN", "")
    session = _get_session(db, chat_id)
    if not gemini_key or not token:
        _send_telegram_session_message(db, token, chat_id, "⚠️ Configuração incompleta. Contate o administrador.", session=session)
        return
    copilot_session_id = _ensure_copilot_session(
        db,
        chat_id,
        session,
        first_text=text,
        from_user=data.get("from_user") or {},
    )
    _save_session(db, chat_id, session)

    copilot_user_message_persisted = False

    def _persist_turn_to_copilot(
        user_content: str,
        assistant_content: str,
        *,
        tools_used: list | None = None,
        persist_user: bool = True,
    ):
        if persist_user:
            _persist_copilot_message(db, copilot_session_id, "user", user_content, source="telegram")
        _persist_copilot_message(
            db,
            copilot_session_id,
            "assistant",
            assistant_content,
            source="telegram",
            tools_used=tools_used,
        )

    _perf_mark(perf_state, "telegram.bootstrap")
    response_mode, text = _extract_response_mode(text, session)
    voice_profile, text = _extract_voice_profile(text, "masculina")
    print(f"[Core] initial response_mode={response_mode} voice_profile={voice_profile}")
    _perf_mark(perf_state, "telegram.session_load")

    # --- /entrar command — busca semântica de ações para travamento de contexto ---
    if re.match(r"^/entrar(\s|$)", text, re.IGNORECASE):
        query = text[len("/entrar"):].strip()
        results = _search_actions_for_context(db, query)
        if not results:
            reply_text = (
                f"Nenhuma ação encontrada para <i>{query}</i>." if query
                else "Nenhuma ação cadastrada no sistema."
            )
            _persist_turn_to_copilot(text, reply_text, tools_used=["buscar_tarefas"])
            _send_telegram_session_message(
                db,
                token,
                chat_id,
                reply_text,
                session=session,
            )
        else:
            header = (
                f"Resultados para <i>{query}</i>. Toque para entrar no contexto:" if query
                else "Ações recentes. Toque para entrar no contexto:"
            )
            _persist_turn_to_copilot(text, header, tools_used=["buscar_tarefas"])
            _send_telegram_session_message(
                db,
                token,
                chat_id,
                header,
                session=session,
                inline_keyboard=_build_action_keyboard(results),
            )
        return

    # --- Natural-language action search + context lock
    # Ex: "pesquisa a acao de X e ative o contexto"
    action_search_context_query = _extract_action_search_context_query(text)
    if action_search_context_query:
        results = _search_actions_for_context(db, action_search_context_query)
        selected = _select_context_result(action_search_context_query, results)
        if selected:
            snapshot = _cached_acao_snapshot(db, selected["id"])
            if not snapshot:
                response_text = "Encontrei a acao, mas nao consegui carregar o snapshot real do contexto. Tente novamente em instantes."
                _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
                _send_contextual_response(
                    db,
                    token,
                    chat_id,
                    response_text,
                    session=session,
                    response_mode=response_mode,
                    gemini_key=gemini_key,
                    voice_profile=voice_profile,
                    perf_state=perf_state,
                )
                return
            _lock_action_session(session, selected["id"], snapshot)
            _save_session(db, chat_id, session)
            response_text = _format_context_locked_message(snapshot)
            _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
            _send_contextual_response(
                db,
                token,
                chat_id,
                response_text,
                session=session,
                response_mode=response_mode,
                gemini_key=gemini_key,
                voice_profile=voice_profile,
                perf_state=perf_state,
            )
            return

        if results:
            header = (
                f"Encontrei algumas acoes para <i>{html.escape(action_search_context_query)}</i>. "
                "Toque na correta para eu travar o contexto real:"
            )
            _persist_turn_to_copilot(text, header, tools_used=["buscar_tarefas"])
            _send_contextual_response(
                db,
                token,
                chat_id,
                header,
                session=session,
                inline_keyboard=_build_action_keyboard(results),
                response_mode=response_mode,
                gemini_key=gemini_key,
                voice_profile=voice_profile,
                perf_state=perf_state,
            )
            return

        response_text = f"Nenhuma acao encontrada para <i>{html.escape(action_search_context_query)}</i>. Nao ativei contexto sem uma acao real."
        _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
        _send_contextual_response(
            db,
            token,
            chat_id,
            response_text,
            session=session,
            response_mode=response_mode,
            gemini_key=gemini_key,
            voice_profile=voice_profile,
            perf_state=perf_state,
        )
        return

    # --- Natural-language context lock ("entre no contexto da acao X") ---
    natural_context_query = _extract_natural_context_query(text)
    if natural_context_query is not None:
        results = _search_actions_for_context(db, natural_context_query)
        selected = _select_context_result(natural_context_query, results)
        if selected:
            snapshot = _cached_acao_snapshot(db, selected["id"])
            if not snapshot:
                response_text = "Encontrei a acao, mas nao consegui carregar o snapshot real do contexto. Tente novamente em instantes."
                _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
                _send_contextual_response(
                    db,
                    token,
                    chat_id,
                    response_text,
                    session=session,
                    response_mode=response_mode,
                    gemini_key=gemini_key,
                    voice_profile=voice_profile,
                    perf_state=perf_state,
                )
                return
            _lock_action_session(session, selected["id"], snapshot)
            _save_session(db, chat_id, session)
            response_text = _format_context_locked_message(snapshot)
            _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
            _send_contextual_response(
                db,
                token,
                chat_id,
                response_text,
                session=session,
                response_mode=response_mode,
                gemini_key=gemini_key,
                voice_profile=voice_profile,
                perf_state=perf_state,
            )
            return

        if results:
            header = (
                f"Encontrei mais de uma acao para <i>{html.escape(natural_context_query or 'sua busca')}</i>. "
                "Toque na correta para eu travar o contexto real:"
            )
            _persist_turn_to_copilot(text, header, tools_used=["buscar_tarefas"])
            _send_contextual_response(
                db,
                token,
                chat_id,
                header,
                session=session,
                inline_keyboard=_build_action_keyboard(results),
                response_mode=response_mode,
                gemini_key=gemini_key,
                voice_profile=voice_profile,
                perf_state=perf_state,
            )
            return

        response_text = f"Nenhuma acao encontrada para <i>{html.escape(natural_context_query or text)}</i>. Nao ativei contexto sem uma acao real."
        _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
        _send_contextual_response(
            db,
            token,
            chat_id,
            response_text,
            session=session,
            response_mode=response_mode,
            gemini_key=gemini_key,
            voice_profile=voice_profile,
            perf_state=perf_state,
        )
        return

    # --- Deterministic action list ("liste as ações de hoje") ---
    is_list_req, list_date = _is_list_actions_request(text)
    if is_list_req:
        results = _fetch_actions_by_date(list_date)
        response_text = _format_actions_simple_list(results, list_date)
        _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
        _send_contextual_response(
            db,
            token,
            chat_id,
            response_text,
            session=session,
            response_mode=response_mode,
            gemini_key=gemini_key,
            voice_profile=voice_profile,
            perf_state=perf_state,
        )
        return

    # --- Deterministic action lookup ("busque a acao relacionada a X") ---
    action_lookup_query = _extract_action_lookup_query(text)
    if action_lookup_query:
        results = _search_actions_for_context(db, action_lookup_query)
        response_text = _format_action_lookup_results(action_lookup_query, results)
        _persist_turn_to_copilot(text, response_text, tools_used=["buscar_tarefas"])
        _send_contextual_response(
            db,
            token,
            chat_id,
            response_text,
            session=session,
            response_mode=response_mode,
            gemini_key=gemini_key,
            voice_profile=voice_profile,
            perf_state=perf_state,
        )
        return

    # --- Command handler ---
    if text.startswith("/"):
        reply = _handle_command(text, session)
        if reply:
            _save_session(db, chat_id, session)
            _persist_turn_to_copilot(text, reply)
            _send_telegram_session_message(db, token, chat_id, reply, session=session)
            return
        # Unknown command — fall through to Gemini

    # --- Natural-language reset request ---
    if _is_reset_request(text):
        ctx = session.get("contexto_ativo", "geral")
        response_text = f"Você tem certeza que deseja limpar o histórico desta sessão (<b>{ctx}</b>)? Isto ajudará a evitar alucinações, mas eu esquecerei o que acabamos de conversar."
        _persist_turn_to_copilot(text, response_text)
        _send_telegram_message_with_keyboard(
            token, chat_id,
            response_text,
            _RESET_CONFIRM_KEYBOARD,
        )
        return

    # --- Load copilot personality ---
    try:
        core_doc = _cached_doc_get(db, "system", "copilot_core")
        copilot_core = (core_doc.to_dict() or {}).get("content", "") if core_doc.exists else ""
    except Exception:
        copilot_core = ""
    try:
        soul_doc = _cached_doc_get(db, "system", "copilot_soul")
        copilot_soul = (soul_doc.to_dict() or {}).get("content", "") if soul_doc.exists else ""
    except Exception:
        copilot_soul = ""

    contexto_ativo = session.get("contexto_ativo", "geral")
    request_acao_id = session.get("acao_id")
    # When locked to an action, pull its data snapshot for LLM injection
    acao_snapshot = session.get("acao_context_snapshot") if contexto_ativo == "acao" else None
    if contexto_ativo == "acao" and request_acao_id:
        fresh_snapshot = _cached_acao_snapshot(db, request_acao_id)
        if fresh_snapshot:
            acao_snapshot = fresh_snapshot
            session["acao_context_snapshot"] = fresh_snapshot
            session["acao_titulo"] = fresh_snapshot.get("titulo") or session.get("acao_titulo")
        _perf_mark(perf_state, "telegram.action_snapshot")
    system_instruction = _build_system_instruction_guarded_v2(copilot_core, copilot_soul, contexto_ativo, acao_snapshot)

    # Áreas temáticas válidas: o copiloto só pode SELECIONAR uma existente, nunca inventar.
    _areas_validas = carregar_areas_tematicas_validas(db)
    system_instruction = (
        system_instruction
        + "\n\nÁREAS TEMÁTICAS VÁLIDAS (ao criar ações, escolha EXATAMENTE UMA desta lista; "
        + "NUNCA invente outra): "
        + ", ".join(_areas_validas)
        + ". Se nenhuma se encaixar, use 'GERAL'."
    )

    # --- Restore history (trimmed) — isolated buffer when action-locked ---
    hist_key = "history_acao" if contexto_ativo == "acao" else "history"
    raw_history = session.get(hist_key, [])
    # Keep last N turns (each turn = user + model = 2 items)
    if len(raw_history) > _MAX_HISTORY_TURNS * 2:
        raw_history = raw_history[-((_MAX_HISTORY_TURNS * 2)):]
    telegram_history = [
        types.Content(role=h["role"], parts=[types.Part(text=p["text"]) for p in h.get("parts", [])])
        for h in raw_history
        if h.get("role") and h.get("parts")
    ]
    copilot_history = _load_copilot_session_history(db, copilot_session_id, types)
    history = copilot_history or telegram_history
    if copilot_history:
        _perf_mark(perf_state, "telegram.copilot_history_load")

    # --- Resolve user message parts ---
    user_parts = []
    file_context_text = ""

    # Audio transcription
    if audio_info and (media_bytes_b64 or media_storage_path):
        audio_bytes = _load_telegram_media_bytes(media_bytes_b64, media_storage_path)
        if not audio_bytes:
            user_parts.append(types.Part(text="[Audio recebido mas o arquivo nao foi encontrado para transcricao]"))
            audio_bytes = None
        if not audio_bytes:
            pass
        else:
            mime = audio_info.get("mime_type", "audio/ogg")
            raw_ext = mime.split("/")[-1]
            ext = "ogg" if raw_ext in ("ogg", "oga") else raw_ext
            _send_telegram_typing(token, chat_id)
            try:
                transcription = _transcribe_audio_bytes(audio_bytes, ext, db)
            except Exception as transcribe_err:
                print(f"[Core] Transcription error: {transcribe_err}")
                import traceback; traceback.print_exc()
                transcription = None
            if transcription:
                if response_mode != "audio":
                    transcribed_mode, cleaned_transcription = _extract_response_mode(transcription, session)
                    if transcribed_mode == "audio":
                        response_mode = "audio"
                        transcription = cleaned_transcription
                voice_profile, transcription = _extract_voice_profile(transcription, voice_profile)
                print(f"[Core] transcription response_mode={response_mode} voice_profile={voice_profile} transcription={transcription[:160]}")
                file_context_text = f"[Transcri??o de ?udio]: {transcription}"
                user_parts.append(types.Part(text=file_context_text))
            else:
                user_parts.append(types.Part(text="[?udio recebido mas transcri??o falhou]"))
            _perf_mark(perf_state, "telegram.audio_transcription")

    # File/document
    elif file_info and (media_bytes_b64 or media_storage_path):
        import uuid as _uuid
        import mimetypes
        file_bytes = _load_telegram_media_bytes(media_bytes_b64, media_storage_path)
        if not file_bytes:
            user_parts.append(types.Part(text="[Arquivo recebido mas o conteudo nao foi encontrado]"))
        else:
            fname = file_info.get("file_name") or f"upload_{_uuid.uuid4().hex[:8]}"
            mime = file_info.get("mime_type", "application/octet-stream")
            
            # Corrigir Mime Type se for genérico ou ausente
            if mime == "application/octet-stream" or not mime:
                guessed, _ = mimetypes.guess_type(fname)
                if guessed:
                    mime = guessed
            
            _GEMINI_NATIVE_MIMES = {
                # Images
                'image/jpeg', 'image/jpg', 'image/png', 'image/webp', 'image/heic', 'image/heif',
                # Audio
                'audio/wav', 'audio/mp3', 'audio/mpeg', 'audio/aac', 'audio/flac', 'audio/ogg', 'audio/m4a',
                # Video
                'video/mp4', 'video/mpeg', 'video/quicktime', 'video/webm',
                # Documents
                'application/pdf', 'text/plain', 'text/csv', 'text/html', 'text/markdown',
            }
            
            if mime in _GEMINI_NATIVE_MIMES:
                user_parts.append(types.Part.from_bytes(data=file_bytes, mime_type=mime))
                user_parts.append(types.Part(text=(
                    f"[Arquivo recebido]: {fname} ({mime})\n"
                    "Analise este arquivo, identifique o que ele mostra e sugira como ele se relaciona ao contexto atual."
                )))
            else:
                # Tenta extrair texto localmente para formatos não suportados nativamente (Excel, Word, PowerPoint, etc.)
                extracted_text = _extract_document_text(file_bytes, fname, mime)
                if extracted_text:
                    user_parts.append(types.Part(text=(
                        f"[Conteúdo extraído do arquivo recebido '{fname}']:\n\n{extracted_text}\n\n"
                        "Analise as informações do documento acima e sugira como elas se relacionam ao contexto atual."
                    )))
                else:
                    # Se não puder extrair e não for nativo, informa o erro de forma limpa
                    user_parts.append(types.Part(text=(
                        f"[Arquivo recebido]: {fname} ({mime})\n"
                        f"⚠️ Nota: O formato deste arquivo não é suportado diretamente e não foi possível extrair seu conteúdo. "
                        "Por favor, envie o conteúdo em formato compatível (PDF, imagens, planilhas CSV ou texto simples)."
                    )))
            _perf_mark(perf_state, "telegram.file_upload")

    # Text
    if text:
        user_parts.append(types.Part(text=text))

    if not user_parts:
        user_parts.append(types.Part(text="[Mensagem sem conteúdo processável]"))

    request_text_for_routing = " ".join(
        p.text for p in user_parts if hasattr(p, "text") and getattr(p, "text", "")
    )
    user_content_for_copilot = "\n\n".join(
        p.text for p in user_parts if hasattr(p, "text") and getattr(p, "text", "")
    ).strip() or text or "[Mensagem sem conteúdo processável]"
    internal_hermes_request = _is_internal_hermes_request(request_text_for_routing)
    explicit_web_request = _is_explicit_web_request(request_text_for_routing)

    # Deterministic path for links/docs in a locked action context. This avoids
    # the model inventing anchor text or Drive URLs when the user asks for files.
    if contexto_ativo == "acao" and acao_snapshot and _is_context_file_request(request_text_for_routing):
        response_text = _format_context_files_response(acao_snapshot, request_text_for_routing)
        _persist_turn_to_copilot(user_content_for_copilot, response_text)
        user_turn = {"role": "user", "parts": [{"text": p.text} for p in user_parts if hasattr(p, "text") and p.text]}
        model_turn = {"role": "model", "parts": [{"text": response_text}]}
        new_history = raw_history + [user_turn, model_turn]
        if len(new_history) > _MAX_HISTORY_TURNS * 2:
            new_history = new_history[-(_MAX_HISTORY_TURNS * 2):]
        session[hist_key] = new_history
        _save_session(db, chat_id, session)
        _send_contextual_response(
            db,
            token,
            chat_id,
            response_text,
            session=session,
            response_mode=response_mode,
            gemini_key=gemini_key,
            voice_profile=voice_profile,
            perf_state=perf_state,
        )
        _perf_log(
            "telegram.request.direct_context_files",
            perf_state,
            {"chat_id": chat_id, "acao_id": request_acao_id},
        )
        return

    if os.environ.get("TELEGRAM_USE_WEB_COPILOT", "1") != "0" and not file_info:
        processing_msg_id = None
        try:
            _send_telegram_typing(token, chat_id)
            processing_msg_id = _send_telegram_session_message(
                db,
                token,
                chat_id,
                "⏳ <i>Estou processando pelo Copiloto Hermes...</i>",
                session=session,
            )
            _persist_copilot_message(db, copilot_session_id, "user", user_content_for_copilot, source="telegram")
            copilot_user_message_persisted = True
            delegated = _run_web_copilot_engine(
                prompt=request_text_for_routing,
                session_id=copilot_session_id,
                task_id=request_acao_id,
                user_uid=session.get("userId"),
            )
            delegated_text = (delegated.get("result") or "").strip()
            if not delegated_text:
                raise RuntimeError("Motor web nao retornou texto.")

            # Intercept markdown images
            images_to_send = re.findall(r'!\[.*?\]\((https?://[^\)]+)\)', delegated_text)
            for img_url in images_to_send:
                _send_telegram_photo(token, chat_id, img_url)
            
            # Clean up the text before sending to telegram
            clean_delegated_text = re.sub(r'!\[.*?\]\((https?://[^\)]+)\)', '', delegated_text).strip()
            
            response_text = _format_web_copilot_text_for_telegram(clean_delegated_text)
            card_text, inline_keyboard = _build_web_copilot_telegram_adaptation(session, delegated, delegated_text)
            if card_text:
                response_text = f"{response_text}\n{card_text}"
            if processing_msg_id:
                _delete_telegram_message(token, chat_id, processing_msg_id)
            user_turn = {"role": "user", "parts": [{"text": p.text} for p in user_parts if hasattr(p, "text") and p.text]}
            model_turn = {"role": "model", "parts": [{"text": delegated_text}]}
            new_history = raw_history + [user_turn, model_turn]
            if len(new_history) > _MAX_HISTORY_TURNS * 2:
                new_history = new_history[-(_MAX_HISTORY_TURNS * 2):]
            session[hist_key] = new_history
            _save_session(db, chat_id, session)
            _send_contextual_response(
                db,
                token,
                chat_id,
                response_text,
                session=session,
                response_mode=response_mode,
                gemini_key=gemini_key,
                voice_profile=voice_profile,
                inline_keyboard=inline_keyboard,
                perf_state=perf_state,
            )
            _perf_mark(perf_state, "telegram.web_copilot_response")
            _perf_log(
                "telegram.request.web_copilot_complete",
                perf_state,
                {
                    "chat_id": chat_id,
                    "copilot_session_id": copilot_session_id,
                    "response_mode": response_mode,
                },
            )
            return
        except Exception as web_copilot_err:
            if processing_msg_id:
                _delete_telegram_message(token, chat_id, processing_msg_id)
            print(f"[SessionBridge] Fallback para motor Telegram: {web_copilot_err}")
            _perf_mark(perf_state, "telegram.web_copilot_fallback")

    # --- Memory context ---
    # Skips Firestore read for short/trivial messages (< 4 meaningful words)
    try:
        query_text = text or file_context_text or ""
        _STOPWORDS_MEM = {"de", "a", "o", "que", "e", "do", "da", "em", "um", "uma",
                          "os", "as", "no", "na", "com", "por", "para", "sim", "nao",
                          "ok", "crie", "faça", "faz", "ola", "oi"}
        _meaningful_words = [w for w in query_text.lower().split()
                             if w not in _STOPWORDS_MEM and len(w) > 2]
        if len(_meaningful_words) >= 4:
            memory_docs = (
                db.collection("knowledge_nodes")
                .order_by("data_criacao", direction=firestore.Query.DESCENDING)
                .limit(4)
                .stream()
            )
            mem_lines = []
            for m_doc in memory_docs:
                m = m_doc.to_dict() or {}
                fato = m.get("fato", "")
                cat = m.get("categoria", "")
                if fato:
                    mem_lines.append(f"- [{cat}] {fato}")
            if mem_lines:
                mem_text = "## MEMÓRIA GLOBAL ATIVA\n" + "\n".join(mem_lines)
                system_instruction = system_instruction + "\n\n" + mem_text
    except Exception:
        pass
    _perf_mark(perf_state, "telegram.memory_context")

    # --- Tool definitions ---
    def consultar_historico_acoes(
        query: str,
        area_tematica: str = None,
        data_limite_inicio: str = None,
        data_limite_fim: str = None,
        status: str = None,
    ):
        """Busca ações e tarefas no Hermes. Use status para filtrar por estado (ex: 'em andamento', 'concluída', 'cancelada'). Use data_limite_inicio/fim (YYYY-MM-DD) para filtrar por prazo."""
        if contexto_ativo == "acao" and acao_snapshot:
            titulo_acao = acao_snapshot.get("titulo") or request_acao_id or "ação atual"
            return (
                f"Contexto trancado na ação '{titulo_acao}'. "
                "Os dados desta ação já estão carregados no contexto — use-os diretamente. "
                "Para pesquisar outras ações, o usuário deve sair primeiro com /sair."
            )
        from tools.busca_grafo import buscar_tarefas

        _STOPWORDS_TG = {"de", "a", "o", "que", "e", "do", "da", "em", "um", "uma",
                         "os", "as", "no", "na", "com", "por", "para", "dos", "das",
                         "nos", "nas", "ao", "se", "ou", "acao", "acoes", "tarefa",
                         "tarefas", "pesquisa", "pesquisar", "pesquise", "busca",
                         "buscar", "busque", "procura", "procurar", "procure",
                         "localiza", "localizar", "localize", "ative", "ativar",
                         "ativa", "contexto"}
        _tg_terms = [w for w in re.findall(r"[\w./-]+", _normalize_for_matching(query), flags=re.UNICODE) if w not in _STOPWORDS_TG and len(w) > 2]
        _tg_mode = "all" if len(_tg_terms) >= 2 else "any"
        res = buscar_tarefas(query, area_tematica=area_tematica, match_mode=_tg_mode,
                             data_limite_inicio=data_limite_inicio, data_limite_fim=data_limite_fim,
                             status=status)
        if _tg_mode == "all" and not res.get("resultados"):
            res = buscar_tarefas(query, area_tematica=area_tematica, match_mode="any",
                                 data_limite_inicio=data_limite_inicio, data_limite_fim=data_limite_fim,
                                 status=status)
        if res.get("erro"):
            return f"⚠️ [ERRO] {res['erro']}"
        resultados = res.get("resultados", [])
        if not resultados:
            filtros_desc = []
            if query and query.strip():
                filtros_desc.append(f"query='{query.strip()}'")
            if status:
                filtros_desc.append(f"status='{status}'")
            if area_tematica:
                filtros_desc.append(f"area='{area_tematica}'")
            if data_limite_inicio or data_limite_fim:
                filtros_desc.append(f"prazo=[{data_limite_inicio or '*'} a {data_limite_fim or '*'}]")
            filtros_str = ", ".join(filtros_desc) if filtros_desc else "(sem filtros)"
            return (
                f"NENHUMA TAREFA ENCONTRADA com os filtros: {filtros_str}.\n"
                "INSTRUCAO OBRIGATORIA: Informe ao usuario que nao encontrou. "
                "NAO invente titulos, status ou dados. NAO use RAG para compensar."
            )
        lines = [
            "=== TAREFAS REAIS ENCONTRADAS NO BANCO DE DADOS ===",
            "REGRA: Use EXCLUSIVAMENTE os campos abaixo. Nao invente, nao complete, nao use RAG.",
            "",
        ]
        if res.get("aviso"):
            lines.append(f"AVISO TECNICO: {res['aviso']}")
        for r in resultados:
            lines.append(f"ID: {r['id']}")
            lines.append(f"Titulo: {r['titulo']}")
            lines.append(f"Status: {r['status']} | Tipo: {r.get('tipo_acao') or 'nao informado'}")
            lines.append(f"Prazo: {r.get('data_limite', 'N/A')} | Area: {r['area']}")
            if r.get('processo_sei'):
                lines.append(f"Processo SEI: {r['processo_sei']}")
            lines.append(f"Responsavel: {r['responsavel'] or 'nao informado'}")
            if r.get('tags'):
                lines.append(f"Tags: {', '.join(r['tags'])}")
            if r.get('sintese_demanda'):
                lines.append(f"Sintese da Demanda: {r['sintese_demanda']}")
            if r.get('descricao'):
                lines.append(f"Descricao: {r['descricao']}")
            if r.get('notas'):
                lines.append(f"Notas: {r['notas']}")
            plano = r.get('plano_acao', [])
            if plano:
                lines.append("Plano de Acao:")
                for passo in plano:
                    lines.append(f"  {passo}")
            acomp = r.get('acompanhamento_recente', [])
            if acomp:
                lines.append("Diario de Bordo (ultimas entradas):")
                for entrada in acomp:
                    lines.append(f"  {entrada}")
            lines.append("---")
        return "\n".join(lines)

    def buscar_arquivos_acervo(query: str):
        """Busca documentos, manuais e arquivos no Acervo Global do Hermes."""
        from tools.busca_acervo import buscar_acervo
        res = buscar_acervo(query)
        if res.get("erro"):
            return f"⚠️ [ERRO] {res['erro']}"
        resultados = res.get("resultados", [])
        if not resultados:
            return "Nenhum documento encontrado."
        return "\n\n".join(
            f"DOC: {r['titulo']} | FONTE: {r['fonte']}\nTRECHO: {r['trecho']}"
            for r in resultados
        )

    def pesquisar_internet(query: str):
        """Busca informações recentes na internet via Tavily."""
        try:
            if internal_hermes_request and not explicit_web_request:
                return (
                    '{"error": "Bloqueado: o pedido atual e interno do Hermes. '
                    'Nao use internet como substituto para consultas do sistema."}'
                )
            tavily_key = _get_api_keys(db).get("tavily_api_key")
            if not tavily_key:
                return '{"error": "Tavily não configurado."}'
            resp = _requests.post(
                "https://api.tavily.com/search",
                json={"api_key": tavily_key, "query": query, "search_depth": "advanced",
                      "include_answer": True, "include_raw_content": False, "max_results": 5},
                timeout=20,
            )
            resp.raise_for_status()
            data = resp.json()
            parts = []
            if data.get("answer"):
                parts.append(f"RESPOSTA: {data['answer']}\n")
            for r in data.get("results", []):
                parts.append(f"FONTE: {r.get('title','')} ({r.get('url','')})\n{r.get('content','')}")
            return "\n\n".join(parts) or "Sem resultados."
        except Exception as e:
            return f'{{"error": "{e}"}}'

    def ler_pagina_web(url: str):
        """Lê o conteúdo de uma URL via Jina Reader."""
        try:
            resp = _requests.get(f"https://r.jina.ai/{url}",
                                 headers={"Accept": "text/markdown"}, timeout=25)
            if resp.status_code in (401, 403, 429):
                return '{"error": "Acesso bloqueado pela página de destino."}'
            resp.raise_for_status()
            content = resp.text.strip()
            return content[:12000] + "\n[truncado]" if len(content) > 12000 else content
        except Exception as e:
            return f'{{"error": "{e}"}}'

    def consultar_agenda(data_inicio: str, data_fim: str):
        """Retorna eventos ocupados no período para verificação de disponibilidade (YYYY-MM-DD)."""
        try:
            from main import get_calendar_service, get_target_calendar_id
            import hermes_calendar_tools as hc_tools
            c_service = get_calendar_service()
            # Na main.py talvez get_db e db global não funcionem direto dentro da tool, mas 'db' é capturado!
            c_id = get_target_calendar_id(db)
            if not c_service or not c_id:
                return "Google Calendar não configurado."
            events = hc_tools.consultar_eventos(c_service, c_id, data_inicio, data_fim)
            return hc_tools.formatar_eventos_para_llm(events)
        except Exception as e:
            return f"Erro ao consultar agenda: {e}"

    def encontrar_slot_livre(a_partir_de: str, duracao_min: int = 30):
        """Encontra o próximo horário livre na agenda. a_partir_de = YYYY-MM-DD. Retorna JSON com data, horario_inicio, horario_fim."""
        try:
            from main import get_calendar_service, get_target_calendar_id
            import hermes_calendar_tools as hc_tools
            c_service = get_calendar_service()
            c_id = get_target_calendar_id(db)
            if not c_service or not c_id:
                return "Erro: Google Calendar não configurado."
            slot = hc_tools.encontrar_proximo_slot(c_service, c_id, a_partir_de, duracao_min)
            if slot:
                import json as _js
                return _js.dumps(slot, ensure_ascii=False)
            return "Nenhum slot livre encontrado."
        except Exception as e:
            return f"Erro ao buscar slot livre: {e}"

    def criar_acao_no_sistema(
        titulo: str,
        descricao: str = "",
        area_tematica: str = "GERAL",
        data_limite: str = None,
        tipo_acao: str = "fast",
        tags: list[str] = None,
        notas: str = "",
        plano_acao: list[str] = None,
        horario_inicio: str = None,
        horario_fim: str = None,
    ):
        """
        Cria uma nova ação no Hermes. Apresente draft ao usuário antes de chamar.
        Retorna 'OK|{ID}' em caso de sucesso ou 'ERRO|{detalhe}'.
        IMPORTANTE: area_tematica deve ser EXATAMENTE UMA das áreas temáticas válidas
        listadas no contexto do sistema. Nunca invente uma nova; se nenhuma se encaixar, use 'GERAL'.
        """
        import uuid as _uuid
        # Garante que o copiloto só use áreas temáticas existentes (fallback 'GERAL').
        area_tematica = normalizar_area_tematica(area_tematica, _areas_validas)
        now_iso = datetime.now(timezone.utc).isoformat()

        # Normalização de horários
        def normalize_hhmm(t_str: str) -> str | None:
            if not t_str:
                return None
            t_str = str(t_str).strip()
            if ":" not in t_str:
                return None
            try:
                h, m = t_str.split(":")
                return f"{int(h):02d}:{int(m):02d}"
            except:
                return t_str

        horario_inicio = normalize_hhmm(horario_inicio)
        horario_fim = normalize_hhmm(horario_fim)

        from zoneinfo import ZoneInfo
        from datetime import datetime as _dt
        tz = ZoneInfo("America/Sao_Paulo")
        now_local = _dt.now(tz)
        today_local = now_local.strftime("%Y-%m-%d")

        if not data_limite or str(data_limite) < today_local:
            data_limite = today_local

        if data_limite == today_local and horario_inicio:
            current_time_str = now_local.strftime("%H:%M")
            if horario_inicio < current_time_str:
                return f"ERRO|Não é possível agendar um horário anterior ao horário atual ({current_time_str}). Por favor, escolha um horário posterior."

        task_id = str(_uuid.uuid4())[:20]
        plano_convertido = [
            {"id": str(_uuid.uuid4())[:8], "text": str(p), "completed": False}
            for p in (plano_acao or []) if str(p).strip()
        ]
        try:
            from main import get_calendar_service, get_target_calendar_id
            import hermes_calendar_tools as hc_tools
            c_service = get_calendar_service()
            c_id = get_target_calendar_id(db)
            if c_service and c_id and horario_inicio and horario_fim:
                hc_tools.reagendar_acoes_hermes(db, c_service, c_id, data_limite, horario_inicio, horario_fim)
        except Exception as e:
            print(f"[Core] Erro ao reagendar iterativo: {e}")

        doc = {
            "id": task_id,
            "titulo": titulo.strip(),
            "descricao": descricao or "",
            "area_tematica": area_tematica or "GERAL",
            "data_limite": data_limite,
            "horario_inicio": horario_inicio,
            "horario_fim": horario_fim,
            "tipo_acao": tipo_acao or "fast",
            "tags": tags or [],
            "notas": notas or "",
            "plano_acao": plano_convertido,
            "status": "em andamento",
            "criado_em": now_iso,
            "data_criacao": now_iso,
            "data_atualizacao": now_iso,
            "origem_ingestao": "telegram",
            "acompanhamento": [],
            "sync_status": "new",
        }
        try:
            db.collection("tarefas").document(task_id).set(doc)
            return f"OK|{task_id}"
        except Exception as e:
            return f"ERRO|{e}"

    def reagendar_acoes_em_lote(
        nova_data_inicio: str,
        max_por_semana: int = 5,
        estrategia: str = "data_criacao",
        filtro_data: str = None,
        task_ids: list[str] = None,
        justificativa: str = "",
    ):
        """
        Reagenda múltiplas ações de uma vez, redistribuindo-as a partir de uma data de início.
        Executa imediatamente após confirmação do usuário.

        Parâmetros:
        - nova_data_inicio: YYYY-MM-DD — primeiro dia útil a partir do qual distribuir as ações
        - max_por_semana: máximo de ações por semana (padrão 5)
        - estrategia: "data_criacao" (padrão) | "tipo_acao" (fast primeiro) | "alfa" (alfabética)
        - filtro_data: YYYY-MM-DD — seleciona ações com data_limite igual a essa data
        - task_ids: lista explícita de IDs de tarefas (alternativa ao filtro_data)
        - justificativa: motivo gravado no diário de cada ação reagendada

        Retorna resumo textual com as ações reagendadas ou mensagem de erro.
        """
        try:
            from datetime import timedelta as _td

            if not filtro_data and not task_ids:
                return "ERRO|Forneça filtro_data (YYYY-MM-DD) ou task_ids."

            tasks = []
            if task_ids:
                for tid in (task_ids or []):
                    tdoc = db.collection('tarefas').document(str(tid)).get()
                    if tdoc.exists:
                        t = tdoc.to_dict()
                        if t.get('status') not in ('concluído', 'cancelado'):
                            tasks.append({'_id': str(tid), **t})
            else:
                q = db.collection('tarefas')\
                    .where('data_limite', '==', filtro_data)\
                    .where('status', 'in', ['em andamento', 'stand-by'])\
                    .get()
                for qdoc in q:
                    tasks.append({'_id': qdoc.id, **qdoc.to_dict()})

            if not tasks:
                return "Nenhuma ação encontrada com os critérios informados."

            if estrategia == 'tipo_acao':
                tasks.sort(key=lambda x: (0 if x.get('tipo_acao') == 'fast' else 1, x.get('data_criacao', '')))
            elif estrategia == 'alfa':
                tasks.sort(key=lambda x: x.get('titulo', '').lower())
            else:
                tasks.sort(key=lambda x: x.get('data_criacao', ''))

            try:
                start_date = datetime.strptime(nova_data_inicio, "%Y-%m-%d").date()
                today_date = datetime.now(timezone.utc).date()
                if start_date < today_date:
                    start_date = today_date
            except ValueError:
                return f"ERRO|Formato de data inválido: '{nova_data_inicio}'. Use YYYY-MM-DD."

            def _next_weekday(d):
                while d.weekday() >= 5:
                    d += _td(days=1)
                return d

            day_cursor = _next_weekday(start_date)

            batch = db.batch()
            now_iso = datetime.now(timezone.utc).isoformat()
            diary_entry = {
                'data': now_iso,
                'nota': f"[Copiloto Hermes] {justificativa or 'Reagendamento em lote.'}"
            }

            count_this_week = 0
            updated_lines = []

            for task in tasks:
                if count_this_week >= max_por_semana:
                    days_to_monday = 7 - day_cursor.weekday()
                    day_cursor += _td(days=days_to_monday)
                    day_cursor = _next_weekday(day_cursor)
                    count_this_week = 0

                nova_data_str = day_cursor.strftime("%Y-%m-%d")
                task_ref = db.collection('tarefas').document(task['_id'])
                batch.update(task_ref, {
                    'data_limite': nova_data_str,
                    'data_inicio': nova_data_str,
                    'data_atualizacao': now_iso,
                    'acompanhamento': firestore.ArrayUnion([diary_entry]),
                })
                updated_lines.append(f"• [{task.get('titulo', task['_id'])}](task:{task['_id']}) → {nova_data_str}")

                count_this_week += 1
                day_cursor += _td(days=1)
                day_cursor = _next_weekday(day_cursor)

            batch.commit()
            return f"✅ {len(updated_lines)} ações reagendadas:\n" + "\n".join(updated_lines)

        except Exception as _e:
            return f"ERRO|{_e}"

    def salvar_memoria_global(fato: str, categoria: str):
        """Persiste fato durável na memória global do Hermes. Apenas para regras estáveis e preferências permanentes."""
        try:
            import uuid as _uuid
            node_id = str(_uuid.uuid4())[:16]
            db.collection("knowledge_nodes").document(node_id).set({
                "id": node_id,
                "fato": fato,
                "categoria": categoria,
                "data_criacao": datetime.now(timezone.utc).isoformat(),
                "origem": "telegram",
            })
            return json.dumps({"status": "saved", "id": node_id}, ensure_ascii=False)
        except Exception as e:
            return json.dumps({"status": "error", "reason": str(e)}, ensure_ascii=False)

    def registrar_correcao_procedimento(
        area_tematica: str,
        titulo_procedimento: str,
        correcao_descrita: str,
        novo_conteudo_proposto: str,
        justificativa: str,
    ):
        """[FERRAMENTA OCULTA] Registra correção de procedimento silenciosamente."""
        try:
            import uuid as _uuid
            cid = str(_uuid.uuid4())[:12]
            db.collection("correcoes_pendentes").document(cid).set({
                "id": cid,
                "area_tematica": area_tematica,
                "titulo_procedimento": titulo_procedimento,
                "correcao_descrita": correcao_descrita,
                "novo_conteudo_proposto": novo_conteudo_proposto,
                "justificativa_usuario": justificativa,
                "status": "pendente",
                "data_criacao": firestore.SERVER_TIMESTAMP,
                "origem": "telegram",
            })
            return f"Correção registrada (ID: {cid})."
        except Exception as e:
            return f"Erro: {e}"

    def agendar_lembrete_acao(data: str, horario: str, task_id: str = None, texto: str = ""):
        """
        Agenda um lembrete para uma acao do Hermes.
        data: Data do lembrete no formato YYYY-MM-DD.
        horario: Horario do lembrete no formato HH:MM.
        task_id: ID da acao. Opcional quando ja existe uma acao em contexto.
        texto: Texto personalizado opcional que aparecera no lembrete.
        """
        from tools.telegram_extended import execute
        actual_task_id = task_id or request_acao_id or session.get("current_task_id")
        if not actual_task_id:
            return "ERRO|Nenhuma ação ativa em contexto e nenhum task_id foi informado para agendar o lembrete."
        slots = {"task_id": actual_task_id, "data": data, "horario": horario, "texto": texto}
        return execute("agendar_lembrete_acao", slots, db)

    def consultar_financas_v2(mes: int = None, ano: int = None):
        """Retorna resumo financeiro (balancete, metas, extrato). mes (0-11), ano (YYYY). Se mes for omitido, assume o mes anterior ao atual."""
        from tools.telegram_extended import execute
        return execute("consultar_financas_v2", {"mes": mes, "ano": ano}, db)

    def registrar_item_financeiro_v2(tipo: str, descricao: str, valor: float, mes: int = None, ano: int = None, data: str = None):
        """
        Registra nova renda, obrigacao_fixa ou transacao_avulsa no sistema financeiro.
        tipo: 'renda' | 'obrigacao_fixa' | 'transacao_avulsa'.
        Obrigatório apresentar rascunho completo ao usuário para confirmação antes de persistir.
        Categoria não deve ser inferida nem solicitada por enquanto; o sistema grava "Geral" internamente.
        """
        from tools.telegram_extended import execute
        slots = {"tipo": tipo, "descricao": descricao, "valor": valor, "categoria": "Geral", "mes": mes, "ano": ano, "data": data}
        return execute("registrar_item_financeiro_v2", slots, db)

    def buscar_e_analisar_email(query: str, max_results: int = 5):
        """Busca e analisa e-mails no Gmail. Use query padrão do Gmail (ex: 'from:x@y.com newer_than:2d')."""
        try:
            from tools.buscar_e_analisar_email import buscar_e_analisar_email as _fn
            return _fn(query=query, max_results=min(int(max_results), 5))
        except Exception as e:
            return f"Erro: {e}"

    def propor_acao_para_confirmacao(
        titulo: str,
        descricao: str = "",
        area_tematica: str = "GERAL",
        data_limite: str = None,
        tipo_acao: str = "fast",
        tags: list[str] = None,
        notas: str = "",
        plano_acao: list[str] = None,
        horario_inicio: str = None,
        horario_fim: str = None,
    ):
        """
        Gera uma proposta de criação de ação para o usuário confirmar via botões.
        Use esta ferramenta SEMPRE antes de criar uma ação.
        """
        # Normalização de horários
        def normalize_hhmm(t_str: str) -> str | None:
            if not t_str:
                return None
            t_str = str(t_str).strip()
            if ":" not in t_str:
                return None
            try:
                h, m = t_str.split(":")
                return f"{int(h):02d}:{int(m):02d}"
            except:
                return t_str

        horario_inicio = normalize_hhmm(horario_inicio)
        horario_fim = normalize_hhmm(horario_fim)

        from zoneinfo import ZoneInfo
        from datetime import datetime as _dt
        tz = ZoneInfo("America/Sao_Paulo")
        now_local = _dt.now(tz)
        today_local = now_local.strftime("%Y-%m-%d")

        eff_limit = data_limite or today_local
        if eff_limit < today_local:
            eff_limit = today_local

        if eff_limit == today_local and horario_inicio:
            current_time_str = now_local.strftime("%H:%M")
            if horario_inicio < current_time_str:
                return f"ERRO|Não é possível propor um horário anterior ao horário atual ({current_time_str}). Por favor, escolha um horário posterior."

        pending_data = {
            "titulo": titulo,
            "descricao": descricao,
            "area_tematica": area_tematica,
            "data_limite": eff_limit,
            "tipo_acao": tipo_acao,
            "tags": tags or [],
            "notas": notas,
            "plano_acao": plano_acao or [],
            "horario_inicio": horario_inicio,
            "horario_fim": horario_fim,
        }
        session.setdefault("pending_confirmations", {})["acao"] = pending_data
        session["_pending_confirm_type"] = "acao"
        
        draft = (
            f"📝 <b>PROPOSTA DE AÇÃO</b>\n"
            f"• Título: {titulo}\n"
            f"• Área: {area_tematica}\n"
            f"• Prazo: {pending_data['data_limite']}\n"
        )
        if tags: draft += f"• Tags: {', '.join(tags)}\n"
        if plano_acao: draft += f"• Passos: {len(plano_acao)}\n"
        
        return f"Proposta gerada com sucesso. Draft: {draft}\n\n[SISTEMA: Os botões de confirmação serão anexados automaticamente a esta resposta.]"

    def propor_lancamento_financeiro(tipo: str, descricao: str, valor: float, mes: int = None, ano: int = None, data: str = None):
        """
        Gera uma proposta de lançamento financeiro para o usuário confirmar via botões.
        tipo: 'renda' | 'obrigacao_fixa' | 'transacao_avulsa'.
        Categoria não deve ser inferida nem exibida por enquanto.
        """
        pending_data = {
            "tipo": tipo,
            "descricao": descricao,
            "valor": valor,
            "categoria": "Geral",
            "mes": mes,
            "ano": ano,
            "data": data
        }
        session.setdefault("pending_confirmations", {})["financeiro"] = pending_data
        session["_pending_confirm_type"] = "financeiro"
        
        draft = (
            f"💰 <b>PROPOSTA DE LANÇAMENTO</b>\n"
            f"• Tipo: {tipo}\n"
            f"• Descrição: {descricao}\n"
            f"• Valor: R$ {valor:.2f}\n"
        )
        return f"Proposta financeira gerada. Draft: {draft}\n\n[SISTEMA: Os botões de confirmação serão anexados automaticamente a esta resposta.]"

    def schedule_whatsapp_message(contact_number: str, message: str, scheduled_time: str) -> str:
        """
        Prepara uma mensagem de WhatsApp para confirmacao por botoes.
        O enfileiramento real acontece apenas no callback confirm_whatsapp.
        """
        pending_data = {
            "contact_number": str(contact_number or "").strip(),
            "message": str(message or "").strip(),
            "scheduled_time": str(scheduled_time or "").strip(),
        }
        session.setdefault("pending_confirmations", {})["whatsapp"] = pending_data
        session["_pending_confirm_type"] = "whatsapp"

        draft = (
            f"📲 <b>PROPOSTA DE WHATSAPP</b>\n"
            f"• Destinatário: {html.escape(pending_data['contact_number'])}\n"
            f"• Agendamento: {html.escape(pending_data['scheduled_time'])}\n"
            f"• Mensagem: {html.escape(pending_data['message'][:700])}"
        )
        return f"Proposta de WhatsApp gerada. Draft: {draft}\n\n[SISTEMA: Os botões de confirmação serão anexados automaticamente a esta resposta.]"

    tools_list = [

        consultar_historico_acoes,
        buscar_arquivos_acervo,
        pesquisar_internet,
        ler_pagina_web,
        consultar_agenda,
        encontrar_slot_livre,
        criar_acao_no_sistema,
        reagendar_acoes_em_lote,
        salvar_memoria_global,
        registrar_correcao_procedimento,
        buscar_e_analisar_email,
        consultar_financas_v2,
        registrar_item_financeiro_v2,
        propor_acao_para_confirmacao,
        propor_lancamento_financeiro,
        schedule_whatsapp_message,
        agendar_lembrete_acao,
    ]


    function_map = {fn.__name__: fn for fn in tools_list}

    # --- Gemini call ---
    _send_telegram_typing(token, chat_id)
    processing_msg_id = _send_telegram_session_message(
        db,
        token,
        chat_id,
        "⏳ <i>Estou processando seu pedido, aguarde um minuto...</i>",
        session=session,
    )

    try:
        response_text = _run_gemini_turn(
            db=db,
            gemini_key=gemini_key,
            system_instruction=system_instruction,
            history=history,
            user_message_parts=user_parts,
            tools_list=tools_list,
            function_map=function_map,
            perf_state=perf_state,
            sistema_id_contexto=None,
        )
    except Exception as gemini_err:
        if processing_msg_id:
            _delete_telegram_message(token, chat_id, processing_msg_id)
        print(f"[Core] Gemini error: {gemini_err}")
        _perf_log("telegram.request.error", perf_state, {"chat_id": chat_id, "error": str(gemini_err)})
        latest_session = _get_session(db, chat_id)
        _send_telegram_session_message(
            db,
            token,
            chat_id,
            f"⚠️ Erro ao processar: {gemini_err}",
            session=latest_session,
        )
        return

    if processing_msg_id:
        _delete_telegram_message(token, chat_id, processing_msg_id)

    # --- Persist history — write to isolated buffer when action-locked ---
    user_turn = {"role": "user", "parts": [{"text": p.text} for p in user_parts if hasattr(p, "text") and p.text]}
    model_turn = {"role": "model", "parts": [{"text": response_text}]}
    new_history = raw_history + [user_turn, model_turn]
    if len(new_history) > _MAX_HISTORY_TURNS * 2:
        new_history = new_history[-(_MAX_HISTORY_TURNS * 2):]
    latest_session = _get_session(db, chat_id)
    
    # Sincroniza mudanças feitas pelas ferramentas (closures) no objeto 'session' original
    if "pending_confirmations" in session:
        latest_session["pending_confirmations"] = session["pending_confirmations"]
    if "_pending_confirm_type" in session:
        latest_session["_pending_confirm_type"] = session["_pending_confirm_type"]

    if not _session_matches_processing_state(latest_session, contexto_ativo, request_acao_id):
        print(
            f"[Session] Discarding stale response for chat_id={chat_id} "
            f"contexto={contexto_ativo} acao_id={request_acao_id!r}"
        )
        _perf_log(
            "telegram.request.discarded_stale_session",
            perf_state,
            {"chat_id": chat_id, "contexto_ativo": contexto_ativo, "acao_id": request_acao_id},
        )
        return
    latest_session[hist_key] = new_history

    # --- Detecção de propostas pendentes para anexar botões ---
    inline_keyboard = None
    if latest_session.get("_pending_confirm_type") == "acao":
        inline_keyboard = _CONFIRM_ACAO_KEYBOARD
    elif latest_session.get("_pending_confirm_type") == "financeiro":
        inline_keyboard = _CONFIRM_FINANCEIRO_KEYBOARD
    elif latest_session.get("_pending_confirm_type") == "whatsapp":
        inline_keyboard = _CONFIRM_WHATSAPP_KEYBOARD

    # --- Send response ---
    # Text goes out first; session is persisted after so it doesn't block the user.
    
    # Intercept markdown images
    images_to_send = re.findall(r'!\[.*?\]\((https?://[^\)]+)\)', response_text)
    for img_url in images_to_send:
        _send_telegram_photo(token, chat_id, img_url)
    
    # Clean up the text before sending to telegram
    clean_response_text = re.sub(r'!\[.*?\]\((https?://[^\)]+)\)', '', response_text).strip()
    
    _send_telegram_session_message(db, token, chat_id, clean_response_text, session=latest_session, inline_keyboard=inline_keyboard)

    _perf_mark(perf_state, "telegram.text_response")

    tools_used = []
    for call in perf_state.get("tool_calls", []):
        name = call.get("name")
        if name and name not in tools_used:
            tools_used.append(name)
    _persist_turn_to_copilot(
        user_content_for_copilot,
        response_text,
        tools_used=tools_used,
        persist_user=not copilot_user_message_persisted,
    )
    _perf_mark(perf_state, "telegram.copilot_history_persist")

    _save_session(db, chat_id, latest_session)
    _perf_mark(perf_state, "telegram.history_persist")

    if response_mode == "audio":
        try:
            with _telegram_action_heartbeat(token, chat_id, "record_voice"):
                print(f"[Core] sending audio with voice_profile={voice_profile}")
                tts_script = _run_gemini_text(
                    gemini_key=gemini_key,
                    system_instruction=_build_tts_director_instruction(voice_profile),
                    user_prompt=response_text,
                )
                audio_bytes, audio_mime = _run_gemini_tts(
                    gemini_key=gemini_key,
                    script_text=tts_script or response_text,
                    voice_profile=voice_profile,
                )
                tg_audio_bytes, tg_audio_mime, tg_audio_name = _transcode_audio_for_telegram_voice(audio_bytes, audio_mime)
                if tg_audio_bytes:
                    if not _send_telegram_voice(token, chat_id, tg_audio_bytes, tg_audio_name, tg_audio_mime):
                        raise RuntimeError("Falha ao enviar voice note pelo Telegram.")
                else:
                    wav_bytes, wav_mime, wav_name = _wrap_pcm_audio_as_wav(audio_bytes, audio_mime)
                    if not _send_telegram_document(token, chat_id, wav_bytes, wav_name, wav_mime, caption="Resposta em audio"):
                        raise RuntimeError("Falha ao enviar arquivo de audio pelo Telegram.")
            _perf_mark(perf_state, "telegram.audio_response")
        except Exception as tts_err:
            print(f"[Core] TTS error: {tts_err}")
            _perf_log("telegram.request.tts_error", perf_state, {"chat_id": chat_id, "error": str(tts_err)})
            _send_tts_failure_notice_contextual(db, token, chat_id, session=latest_session)
    _perf_log(
        "telegram.request.complete",
        perf_state,
        {
            "chat_id": chat_id,
            "response_mode": response_mode,
            "history_turns": len(raw_history) // 2,
        },
    )


# ---------------------------------------------------------------------------
# Cloud Functions
# ---------------------------------------------------------------------------

@https_fn.on_request(
    cors=options.CorsOptions(cors_origins=["https://app.hermes.com", "http://localhost:3001", "http://localhost:5173", "http://localhost:3025", "http://127.0.0.1:3001", "http://127.0.0.1:5173", "http://127.0.0.1:3025"], cors_methods=["GET", "POST"]),
    timeout_sec=10,
    memory=options.MemoryOption.GB_1,
)
def telegramWebhook(req: https_fn.Request) -> https_fn.Response:
    """
    Porteiro: recebe POST do Telegram, valida, enfileira no Firestore e retorna 200.
    Nunca bloqueia aguardando processamento Gemini.
    """
    if req.method != "POST":
        return https_fn.Response("OK", status=200)

    try:
        update = req.get_json(silent=True) or {}
    except Exception:
        return https_fn.Response("OK", status=200)

    db = _get_db()
    try:
        token = _get_telegram_token(db)
    except Exception:
        return https_fn.Response("OK", status=200)

    # --- Handle callback_query (botões inline — travamento/destravamento de contexto) ---
    callback_query = update.get("callback_query")
    if callback_query:
        return _handle_telegram_callback(db, token, callback_query)

    # --- Extrair mensagem ---
    message = update.get("message") or update.get("edited_message") or {}
    if not message:
        return https_fn.Response("OK", status=200)

    chat = message.get("chat", {})
    chat_id = str(chat.get("id", ""))
    if not chat_id:
        return https_fn.Response("OK", status=200)

    # --- Validação single-owner ---
    allowed = _get_allowed_chat_id()
    if allowed and chat_id != allowed:
        return https_fn.Response("OK", status=200)  # silent ignore

    text = message.get("text") or message.get("caption") or ""
    from_user = message.get("from", {})
    message_id = str(message.get("message_id", ""))

    # --- Detectar mídia ---
    audio_info = None
    file_info = None
    media_bytes_b64 = None
    media_storage_path = None

    # Audio / Voice
    audio = message.get("audio") or message.get("voice")
    if audio:
        file_size = audio.get("file_size", 0)
        if file_size > _MAX_FILE_BYTES:
            _send_telegram_session_message(
                db,
                token,
                chat_id,
                "⚠️ Áudio muito grande (máximo 20 MB). Use o portal Web para arquivos maiores.",
            )
            return https_fn.Response("OK", status=200)
        try:
            file_meta = _get_telegram_file(token, audio["file_id"])
            file_bytes = _download_telegram_file(token, file_meta["file_path"])
            media_bytes_b64, media_storage_path = _store_telegram_media(
                file_bytes,
                f"audio_{message_id}.ogg",
                audio.get("mime_type", "audio/ogg"),
                chat_id,
            )
            audio_info = {
                "file_id": audio["file_id"],
                "file_size": file_size,
                "mime_type": audio.get("mime_type", "audio/ogg"),
                "duration": audio.get("duration", 0),
            }
        except Exception as e:
            print(f"[Webhook] Audio media prepare error: {e}")
            _send_telegram_session_message(
                db,
                token,
                chat_id,
                f"⚠️ Não consegui preparar o áudio para processamento: {e}\nTente novamente ou envie o texto diretamente.",
            )
            return https_fn.Response("OK", status=200)

    # Document / Photo
    doc = message.get("document")
    photos = message.get("photo")
    if not audio_info:
        media_obj = doc or (photos[-1] if photos else None)
        if media_obj:
            file_size = media_obj.get("file_size", 0)
            if file_size > _MAX_FILE_BYTES:
                _send_telegram_session_message(
                    db,
                    token,
                    chat_id,
                    "⚠️ Arquivo muito grande (máximo 20 MB). Use o portal Web para arquivos maiores.",
                )
                return https_fn.Response("OK", status=200)
            try:
                file_meta = _get_telegram_file(token, media_obj["file_id"])
                file_bytes = _download_telegram_file(token, file_meta["file_path"])
                file_info = {
                    "file_id": media_obj["file_id"],
                    "file_size": file_size,
                    "file_name": doc.get("file_name", "arquivo") if doc else "foto.jpg",
                    "mime_type": (doc.get("mime_type") if doc else "image/jpeg") or "application/octet-stream",
                }
                media_bytes_b64, media_storage_path = _store_telegram_media(
                    file_bytes,
                    file_info["file_name"],
                    file_info["mime_type"],
                    chat_id,
                )
            except Exception as e:
                print(f"[Webhook] File download error: {e}")

    if not text and not audio_info and not file_info:
        return https_fn.Response("OK", status=200)

    # --- Enfileira no Firestore ---
    payload = {
        "chat_id": chat_id,
        "message_id": message_id,
        "text": text,
        "from_user": from_user,
        "audio_info": audio_info,
        "file_info": file_info,
        "media_bytes_b64": media_bytes_b64,
        "media_storage_path": media_storage_path,
        "received_at": firestore.SERVER_TIMESTAMP,
        "processed": False,
    }
    db.collection("telegram_inbound").document(f"{chat_id}_{message_id}").set(payload)

    return https_fn.Response("OK", status=200)


@firestore_fn.on_document_created(
    document="telegram_inbound/{docId}",
    timeout_sec=540,
    memory=options.MemoryOption.GB_2,
)
def on_telegram_inbound(event: Event[DocumentSnapshot]) -> None:
    """
    Trigger assíncrono: processa a mensagem enfileirada e envia resposta ao Telegram.
    """
    snap = event.data
    if not snap or not snap.exists:
        return

    data = snap.to_dict() or {}
    if data.get("processed"):
        return

    # Marca como em processamento para evitar double-trigger
    snap.reference.update({"processed": True, "processing_started_at": firestore.SERVER_TIMESTAMP})

    try:
        db = _get_db()
        if not data.get("callback_action"):
            _process_telegram_message(db, data)
    except Exception as e:
        print(f"[on_telegram_inbound] Unhandled error: {e}")
        import traceback
        traceback.print_exc()
        snap.reference.update({"error": str(e)})
        try:
            chat_id = data.get("chat_id")
            db = _get_db()
            token = _get_telegram_token(db)
            if chat_id and token:
                _send_telegram_session_message(
                    db,
                    token,
                    chat_id,
                    f"⚠️ Erro interno ao processar mensagem: {e}",
                )
        except Exception:
            pass
