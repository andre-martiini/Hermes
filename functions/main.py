

from firebase_functions import firestore_fn, scheduler_fn, options, https_fn, pubsub_fn, storage_fn

from firebase_admin import initialize_app, firestore, get_app
from google.cloud.firestore_v1.vector import Vector as FsVector
import json
import base64
from datetime import datetime, timedelta, timezone
import math
import time
import random
import threading
import re
import io
import uuid
import secrets
import os
import sys
import unicodedata

CURRENT_DIR = os.path.dirname(os.path.abspath(__file__))
if CURRENT_DIR not in sys.path:
    sys.path.insert(0, CURRENT_DIR)

from security_portals import (
    generatePgdFromDiariesAI,
    generatePgdFromRawTextAI,
    getPublicFinancePortal,
    getPublicScholarshipProject,
    getPublicShoppingPortal,
    matchShoppingItemsAI,
    mutatePublicShoppingPortal,
    mutateShoppingList,
    submitPublicFinanceTransaction,
    submitPublicScholarshipRegistration,
)
from pdf_precision import extract_pdf_text_with_fallback, is_pdf_mime_type
from godmode import (  # noqa: F401 — registra as Cloud Functions
    askHermesGodmode,
)
from mcp_server import mcpServer  # noqa: F401 — registra a Cloud Function

DOCX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def is_docx_mime_type(filename: str | None = None, mime_type: str | None = None) -> bool:
    mime = (mime_type or "").lower().strip()
    name = (filename or "").lower().strip()
    return mime == DOCX_MIME_TYPE or name.endswith(".docx")


def extract_docx_text(file_bytes: bytes) -> tuple[str, dict]:
    import mammoth

    result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
    warnings = [
        getattr(message, "message", str(message))
        for message in (getattr(result, "messages", None) or [])
    ]
    metadata = {"source": "mammoth", "warnings": warnings}
    return (result.value or "").strip(), metadata

from knowledge_graph import (  # noqa: F401 — registra as Cloud Functions
    on_tarefa_created_kg,
    on_tarefa_concluida_kg,
    on_tarefa_written_extract_people,
    buscar_procedimento,
    crystallize_task_manual,
    extract_kg_rag_context,
    processar_artefato_kg,
    monitorar_acervo_global,
    executar_monitoramento_acervo_global,
    sincronizar_acervo_manual,
    smart_search_kg,
    get_artefato_raw_text,
)
from hermes_core_logic import (  # noqa: F401 — registra as Cloud Functions
    telegramWebhook,
    on_telegram_inbound,
    on_health_log_red_flag,
    carregar_areas_tematicas_validas,
    normalizar_area_tematica,
    _get_telegram_token,
    _get_allowed_chat_id,
    _send_telegram_message,
)
from gemini_cost_controls import (
    check_and_increment_limit,
    decrement_limit,
    log_gemini_usage,
    GEMINI_ROUTING_MODEL,
    GEMINI_LIGHT_MODEL,
    GEMINI_BALANCED_MODEL,
    GEMINI_FRONTIER_MODEL,
    generate_content_logged,
)
from llm_providers import openai_provider


# Inicializa o Firebase Admin apenas uma vez no escopo global
try:
    get_app()
except ValueError:
    initialize_app()

DEFAULT_GOOGLE_CALENDAR_ID = 'cf4953b9512ee2e85a7e064f9d5ce4eaf6e3634564c91e5c7ee2bb01fd46782a@group.calendar.google.com'
SYNC_LOCK_DOC_ID = 'sync_lock'
SYNC_LOCK_STALE_SECONDS = 15 * 60
MAX_SYNC_PASSES = 3
# Sincronização Google Tasks <-> Ações (coleção 'tarefas') desativada a pedido do usuário (2026-06-02).
# As duas direções estão desligadas: criar tarefa no Google não cria ação no Hermes e vice-versa,
# e itens já vinculados também deixam de ser sincronizados. A integração com o Google Calendar
# (sync_google_calendar e a parte de Calendar do push) continua ATIVA.
# Para reativar a sincronização de tarefas, basta voltar este flag para True.
SYNC_GOOGLE_TASKS_ENABLED = False
COPILOT_FUNCTION_TIMEOUT_SEC = 540
COPILOT_SOFT_DEADLINE_SEC = 300
COPILOT_MODEL_TIMEOUT_MS = 70000
COPILOT_MODEL_RETRY_TIMEOUT_MS = 30000
COPILOT_TOOL_TIMEOUT_SEC = 45
COPILOT_TOOL_TIMEOUT_ASYNC_SEC = 150
COPILOT_CHAT_MODEL = os.environ.get("COPILOT_CHAT_MODEL", "gemini-3.5-flash-lite")
COPILOT_FALLBACK_MODEL = os.environ.get("COPILOT_FALLBACK_MODEL", "gemini-3.5-flash-lite")
# Temperatura do chat do copiloto. O default do SDK (~1.0) é alto demais para um
# assistente que executa ferramentas e gerencia tarefas — eleva a variância e o
# risco de alucinação. 0.3 mantém alguma fluência nos textos (slides/relatórios)
# sem a aleatoriedade que faz o modelo "lite" escorregar. Ajustável por env.
try:
    COPILOT_TEMPERATURE = float(os.environ.get("COPILOT_TEMPERATURE", "0.3"))
except (TypeError, ValueError):
    COPILOT_TEMPERATURE = 0.3
# Escalonamento por complexidade: comandos que disparam muitas funções de uma vez
# ou descrevem várias tarefas encadeadas saturam o flash-lite. Nesses casos — e só
# neles — subimos para um tier mais forte. Heurística pura, sem chamada extra de LLM.
COPILOT_AUTO_ESCALATE = os.environ.get("COPILOT_AUTO_ESCALATE", "1") != "0"
COPILOT_COMPLEX_MODEL = os.environ.get("COPILOT_COMPLEX_MODEL", GEMINI_FRONTIER_MODEL)
# Limiares do escalonamento por complexidade. Defaults conservadores: o modelo
# frontier custa ~6x o flash-lite por token, então só escalamos com sinais fortes.
# Tudo ajustável por env para calibrar sem deploy.
def _env_int(name: str, default: int) -> int:
    try:
        return int(os.environ.get(name, str(default)))
    except (TypeError, ValueError):
        return default
COPILOT_ESCALATE_RARE_TOOLS = _env_int("COPILOT_ESCALATE_RARE_TOOLS", 3)
COPILOT_ESCALATE_DYNAMIC_TOOLS = _env_int("COPILOT_ESCALATE_DYNAMIC_TOOLS", 6)
# Teto de caracteres por mensagem ao recompor o histórico da sessão. Mensagens de
# assistente (relatórios/markdown longos) inflam o input reenviado a cada turno.
COPILOT_HISTORY_MSG_MAXCHARS = _env_int("COPILOT_HISTORY_MSG_MAXCHARS", 4000)
COPILOT_DEADLINE_FALLBACK_TEXT = (
    "O modelo de IA demorou demais para concluir esta resposta e eu interrompi a chamada "
    "antes de estourar o tempo da conversa. Tente dividir o pedido em partes menores, "
    "pedir primeiro um levantamento objetivo ou anexar apenas o trecho essencial."
)


def _is_copilot_deadline_error(error: Exception) -> bool:
    code = getattr(error, "code", None)
    status = str(getattr(error, "status", "") or "").upper()
    message = str(getattr(error, "message", "") or "")
    details = str(getattr(error, "details", "") or "")
    full_text = " ".join([str(error), status, message, details]).upper()
    return (
        code == 504
        or "DEADLINE_EXCEEDED" in full_text
        or "DEADLINE EXPIRED" in full_text
        or "READTIMEOUT" in full_text
        or "TIMED OUT" in full_text
        or "TIMEOUT" in full_text
    )


def _is_retryable_gemini_server_error(error: Exception) -> bool:
    code = getattr(error, "code", None)
    status = str(getattr(error, "status", "") or "").upper()
    full_text = f"{error} {status}".upper()
    return code in {500, 502, 503} or any(token in full_text for token in ("INTERNAL", "UNAVAILABLE", "BAD_GATEWAY"))


def _is_retryable_google_api_error(error: Exception) -> bool:
    """Erros transitórios da API do Google (ex.: 503 "service is currently
    unavailable") que devem ser reexecutados com backoff em vez de tratados
    como falha definitiva da sincronização."""
    status_code = getattr(getattr(error, "resp", None), "status", None)
    if status_code in {429, 500, 502, 503, 504}:
        return True
    full_text = str(error).upper()
    return any(token in full_text for token in ("UNAVAILABLE", "INTERNAL ERROR", "BAD GATEWAY", "RATE LIMIT", "QUOTA"))

# ---------------------------------------------------------------------------
# In-process Firestore document cache (TTL = 60 s per Cloud Function instance)
# Evita leituras repetidas de system/api_keys, system/copilot_core, etc.
# ---------------------------------------------------------------------------
_DOC_CACHE: dict = {}
_DOC_CACHE_TTL = 60  # seconds

# Cache da coleção pops_diretrizes. Um POP recém-criado passa a valer na
# próxima janela de TTL — sem revalidação em caso de miss (a revalidação
# anulava o cache, já que a maioria das mensagens não aciona POP nenhum).
_POPS_DATA_CACHE: tuple | None = None  # (monotonic_ts, list[dict])
_POPS_DATA_TTL = 60

# Cache genérico de coleções-catálogo que mudam raramente (sistemas, unidades).
_COLLECTION_CACHE: dict = {}  # collection -> (monotonic_ts, list[dict])
_COLLECTION_CACHE_TTL = 300  # seconds

# Cache de _bootstrap_user_ai_profile por UID (TTL 60 s)
_PROFILE_CACHE: dict = {}  # uid -> (monotonic_ts, dict)

def _cached_doc_get(db, collection: str, document: str):
    key = f"{collection}/{document}"
    now = time.monotonic()
    cached = _DOC_CACHE.get(key)
    if cached and (now - cached[0]) < _DOC_CACHE_TTL:
        return cached[1]
    doc = db.collection(collection).document(document).get()
    _DOC_CACHE[key] = (now, doc)
    return doc


def _cached_collection_list(db, collection: str, ttl: int = _COLLECTION_CACHE_TTL) -> list[dict]:
    """Lista documentos de uma coleção-catálogo com cache por instância (TTL padrão 5 min)."""
    now = time.monotonic()
    cached = _COLLECTION_CACHE.get(collection)
    if cached and (now - cached[0]) < ttl:
        return cached[1]
    docs = [{"id": d.id, **(d.to_dict() or {})} for d in db.collection(collection).stream()]
    _COLLECTION_CACHE[collection] = (now, docs)
    return docs


def get_genai_module():
    from google import genai
    return genai


def get_gemini_api_key() -> str | None:
    db = get_db()
    keys_doc = _cached_doc_get(db, 'system', 'api_keys')
    return keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None


def get_embedding(text: str, api_key: str = None, task_type: str = "RETRIEVAL_DOCUMENT") -> list:
    """Get text embedding via Gemini SDK (gemini-embedding-001, 768 dims).
    task_type should be RETRIEVAL_DOCUMENT for indexing and RETRIEVAL_QUERY for searching."""
    from google import genai
    from google.genai import types
    if not api_key:
        api_key = get_gemini_api_key()
    if not api_key:
        raise ValueError("Chave Gemini não configurada.")
    client = genai.Client(api_key=api_key)
    response = client.models.embed_content(
        model="gemini-embedding-001",
        contents=text[:8000],
        config=types.EmbedContentConfig(
            task_type=task_type,
            output_dimensionality=768  # trava dimensional: compatível com vetores existentes no Firestore
        )
    )
    return response.embeddings[0].values


def _cosine_similarity(a: list[float], b: list[float]) -> float:
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = math.sqrt(sum(x * x for x in a))
    norm_b = math.sqrt(sum(x * x for x in b))
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


def get_db():

    """Retorna a instância do Firestore de forma lazy"""

    return firestore.client()


DESCRIPTION_SYNTHESIS_MODEL = os.environ.get("DESCRIPTION_SYNTHESIS_MODEL", "gemini-3.5-flash-lite")
DESCRIPTION_SYNTHESIS_BATCH_LIMIT = 50


def _clean_meaningful_text(value, max_chars: int = 1600) -> str:
    text = str(value or "").strip()
    text = re.sub(r"\s+", " ", text)
    if len(text) > max_chars:
        return text[:max_chars].rstrip() + "..."
    return text


def _is_blank_action_description(value) -> bool:
    if value is None:
        return True
    if not isinstance(value, str):
        return False
    return re.fullmatch(r"[\s\-_–—]*", value or "") is not None


def _extract_action_diary_entries(task_data: dict, limit: int = 40) -> list[dict]:
    entries = []
    raw_entries = task_data.get("acompanhamento") or []
    if not isinstance(raw_entries, list):
        return entries
    for entry in raw_entries:
        if isinstance(entry, dict):
            note = _clean_meaningful_text(entry.get("nota"), 1800)
            date = _clean_meaningful_text(entry.get("data"), 80)
        else:
            note = _clean_meaningful_text(entry, 1800)
            date = ""
        if note:
            entries.append({"data": date, "nota": note})
    return entries[-limit:]


def _extract_action_plan_items(task_data: dict, limit: int = 80) -> list[dict]:
    items = []
    raw_items = task_data.get("plano_acao") or []
    if not isinstance(raw_items, list):
        return items
    for item in raw_items:
        if isinstance(item, dict):
            text = _clean_meaningful_text(item.get("text") or item.get("titulo") or item.get("descricao"), 1200)
            completed = bool(item.get("completed"))
        else:
            text = _clean_meaningful_text(item, 1200)
            completed = False
        if text:
            items.append({"text": text, "completed": completed})
    return items[:limit]


def _is_task_eligible_for_description_synthesis(task_data: dict) -> tuple[bool, str]:
    if not _is_blank_action_description(task_data.get("descricao")):
        return False, "descricao_existente"
    diary_entries = _extract_action_diary_entries(task_data)
    plan_items = _extract_action_plan_items(task_data)
    if not diary_entries and not plan_items:
        return False, "sem_contexto"
    return True, "elegivel"


def _build_description_synthesis_prompt(task_data: dict) -> str:
    title = _clean_meaningful_text(task_data.get("titulo") or "Acao sem titulo", 500)
    status = _clean_meaningful_text(task_data.get("status"), 120) or "Nao informado"
    area = _clean_meaningful_text(task_data.get("area_tematica") or task_data.get("projeto"), 200) or "Nao informada"
    diary_entries = _extract_action_diary_entries(task_data)
    plan_items = _extract_action_plan_items(task_data)

    diary_text = "\n".join(
        f"- [{entry.get('data') or 'sem data'}] {entry.get('nota')}"
        for entry in diary_entries
    ) or "- Sem registros significativos."
    plan_text = "\n".join(
        f"- [{'X' if item.get('completed') else ' '}] {item.get('text')}"
        for item in plan_items
    ) or "- Sem checklist significativo."

    return f"""Voce e um redator executivo senior. Recebera o titulo de uma acao, seus registros cronologicos de diario de bordo e as etapas do seu plano de acao.
Sua missao e elaborar uma sintese descritiva coesa (de 1 a 2 paragrafos concisos) que explique de forma cristalina:
1. O objetivo principal e escopo da demanda.
2. O progresso atual resumido (o que ja foi feito e os proximos passos cruciais).

Regras de Formatacao:
- Seja impessoal, profissional e direto (pt-BR).
- Nao utilize saudacoes, introducoes ou meta-comentarios ("Aqui esta a descricao").
- Retorne APENAS o texto da descricao resultante.

CONTEXTO DA ACAO
Titulo: {title}
Status: {status}
Area/Projeto: {area}

DIARIO DE BORDO
{diary_text}

CHECKLIST / PLANO DE ACAO
{plan_text}"""


def _sanitize_generated_description(text: str) -> str:
    cleaned = str(text or "").strip()
    cleaned = re.sub(r"^```(?:markdown|text)?", "", cleaned, flags=re.IGNORECASE).strip()
    cleaned = re.sub(r"```$", "", cleaned).strip()
    cleaned = re.sub(r"^(aqui esta|segue|descricao gerada|sintese executiva)\s*:?\s*", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned[:2500].strip()


def _generate_action_description(client, task_data: dict) -> str:
    response = client.models.generate_content(
        model=DESCRIPTION_SYNTHESIS_MODEL,
        contents=_build_description_synthesis_prompt(task_data),
        config={"temperature": 0.25, "max_output_tokens": 700}
    )
    description = _sanitize_generated_description(getattr(response, "text", "") or "")
    if len(description) < 40:
        raise ValueError("Descricao sintetizada ficou vazia ou curta demais.")
    return description


def _count_eligible_description_synthesis_tasks(db_ref) -> int:
    count = 0
    for task_doc in db_ref.collection("tarefas").stream():
        task_data = task_doc.to_dict() or {}
        eligible, _reason = _is_task_eligible_for_description_synthesis(task_data)
        if eligible:
            count += 1
    return count

def _perf_now_ms() -> int:
    return int(time.perf_counter() * 1000)


def _perf_mark(perf_state: dict, name: str):
    started_at = perf_state.get("_last_ms", perf_state["start_ms"])
    now_ms = _perf_now_ms()
    perf_state.setdefault("steps", []).append({
        "name": name,
        "duration_ms": max(0, now_ms - started_at),
    })
    perf_state["_last_ms"] = now_ms


def _perf_log(prefix: str, perf_state: dict, extra: dict | None = None):
    payload = {
        "prefix": prefix,
        "total_ms": max(0, _perf_now_ms() - perf_state["start_ms"]),
        "steps": perf_state.get("steps", []),
    }
    if perf_state.get("tool_calls"):
        payload["tool_calls"] = perf_state["tool_calls"]
    if extra:
        payload.update(extra)
    try:
        print(f"[Perf] {json.dumps(payload, ensure_ascii=False)}")
    except Exception:
        print(f"[Perf] {prefix} total_ms={payload['total_ms']}")


GOOGLE_BASE_SCOPES = [
    'https://www.googleapis.com/auth/tasks',
    'https://www.googleapis.com/auth/gmail.modify',
    'https://www.googleapis.com/auth/calendar',
    'https://www.googleapis.com/auth/drive',
    'https://www.googleapis.com/auth/contacts',
]

GOOGLE_FORMS_SCOPES = GOOGLE_BASE_SCOPES + [
    'https://www.googleapis.com/auth/forms.body'
]


class GoogleAuthRevokedError(Exception):
    """Raised when Google OAuth credentials need a new user consent flow."""


GOOGLE_REAUTH_MESSAGE = (
    "Autenticacao Google expirada ou revogada. Execute setup_credentials.bat "
    "na raiz do projeto para refazer o login Google e atualizar system/google_credentials."
)


def is_google_invalid_grant_error(exc):
    text = str(exc).lower()
    if "invalid_grant" in text or "token has been expired or revoked" in text:
        return True
    cause = getattr(exc, "__cause__", None)
    if cause and cause is not exc:
        return is_google_invalid_grant_error(cause)
    context = getattr(exc, "__context__", None)
    if context and context is not exc:
        return is_google_invalid_grant_error(context)
    return False


def get_google_creds(scopes=None):
    """Busca as credenciais OAuth2 do Firestore e renova se necessário"""
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request

    db = get_db()
    creds_doc = db.collection('system').document('google_credentials').get()

    if not creds_doc.exists:
        raise Exception("Credenciais não encontradas no Firestore.")

    creds_data = creds_doc.to_dict()
    stored_scopes = creds_data.get('scopes') or GOOGLE_BASE_SCOPES
    required_scopes = scopes or GOOGLE_BASE_SCOPES
    missing_scopes = sorted(set(required_scopes) - set(stored_scopes))
    if missing_scopes:
        db.collection('system').document('google_credentials').set({
            'auth_status': 'reauth_required',
            'auth_error': 'missing_scopes',
            'auth_error_message': GOOGLE_REAUTH_MESSAGE,
            'missing_scopes': missing_scopes,
            'updated_at': firestore.SERVER_TIMESTAMP
        }, merge=True)
        raise GoogleAuthRevokedError(
            f"{GOOGLE_REAUTH_MESSAGE} Escopos ausentes: {', '.join(missing_scopes)}"
        )

    expiry_val = creds_data.get('expiry_date') or creds_data.get('expiry')
    parsed_expiry = None
    if expiry_val:
        from datetime import datetime, timezone
        if isinstance(expiry_val, datetime):
            parsed_expiry = expiry_val
        elif isinstance(expiry_val, (int, float)):
            if expiry_val > 1e11:  # milliseconds
                expiry_val = expiry_val / 1000.0
            parsed_expiry = datetime.fromtimestamp(expiry_val, timezone.utc)
        elif isinstance(expiry_val, str):
            try:
                parsed_expiry = datetime.fromisoformat(expiry_val.replace('Z', '+00:00'))
            except ValueError:
                pass

        if parsed_expiry and parsed_expiry.tzinfo is not None:
            parsed_expiry = parsed_expiry.astimezone(timezone.utc).replace(tzinfo=None)

    creds = Credentials(
        token=creds_data.get('token'),
        refresh_token=creds_data.get('refresh_token'),
        token_uri=creds_data.get('token_uri'),
        client_id=creds_data.get('client_id'),
        client_secret=creds_data.get('client_secret'),
        scopes=stored_scopes,
        expiry=parsed_expiry
    )

    # Verifica se o token expirou e tenta renovar
    if creds and creds.expired and creds.refresh_token:
        try:
            creds.refresh(Request())
            # Salva o NOVO token de volta no Firestore para evitar falhas futuras
            db.collection('system').document('google_credentials').update({
                'token': creds.token,
                'expiry_date': creds.expiry,
                'updated_at': firestore.SERVER_TIMESTAMP
            })
            print("Token Google renovado e salvo no Firestore com sucesso.")
        except Exception as e:
            print(f"Erro ao renovar token do Google: {e}")
            if is_google_invalid_grant_error(e):
                db.collection('system').document('google_credentials').set({
                    'auth_status': 'reauth_required',
                    'auth_error': 'invalid_grant',
                    'auth_error_message': GOOGLE_REAUTH_MESSAGE,
                    'updated_at': firestore.SERVER_TIMESTAMP
                }, merge=True)
                raise GoogleAuthRevokedError(GOOGLE_REAUTH_MESSAGE) from e
            
    return creds




def get_tasks_service():

    from googleapiclient.discovery import build

    return build('tasks', 'v1', credentials=get_google_creds())



def get_gmail_service():

    from googleapiclient.discovery import build

    return build('gmail', 'v1', credentials=get_google_creds())


def archive_gmail_message(service, msg_id, sync_ref=None, logs=None, reason="financeiro"):
    """Remove o label INBOX depois que o e-mail financeiro foi absorvido pelo Hermes."""
    try:
        service.users().messages().modify(
            userId='me',
            id=msg_id,
            body={'removeLabelIds': ['INBOX']}
        ).execute(num_retries=3)
        if sync_ref is not None and logs is not None:
            log_to_firestore(sync_ref, logs, f"[GMAIL] E-mail financeiro arquivado ({reason}): {msg_id}")
        return True
    except Exception as e:
        if sync_ref is not None and logs is not None:
            log_to_firestore(sync_ref, logs, f"[GMAIL] Aviso: nao foi possivel arquivar {msg_id}: {e}")
        else:
            print(f"[GMAIL] Aviso: nao foi possivel arquivar {msg_id}: {e}")
        return False



def get_calendar_service():

    from googleapiclient.discovery import build

    return build('calendar', 'v3', credentials=get_google_creds())



def get_drive_service():

    from googleapiclient.discovery import build

    return build('drive', 'v3', credentials=get_google_creds())


def get_target_calendar_id(db=None):

    db = db or get_db()

    try:

        cfg_doc = db.collection('system').document('config').get()

        if cfg_doc.exists:

            cfg = cfg_doc.to_dict() or {}

            calendar_id = (
                cfg.get('googleCalendarId')
                or cfg.get('google_calendar_id')
                or cfg.get('calendarId')
            )

            if isinstance(calendar_id, str) and calendar_id.strip():

                return calendar_id.strip()

    except Exception:

        pass

    return DEFAULT_GOOGLE_CALENDAR_ID


def get_sync_calendar_ids(db=None):

    target_calendar_id = get_target_calendar_id(db)
    calendar_ids = ['primary']

    if target_calendar_id and target_calendar_id != 'primary':
        calendar_ids.append(target_calendar_id)

    return calendar_ids


def parse_iso_datetime(value):

    if not value or not isinstance(value, str):

        return None

    try:

        normalized = value.replace('Z', '+00:00')

        return datetime.fromisoformat(normalized)

    except Exception:

        return None


def build_task_calendar_event_id(task_id):

    stable_uuid = uuid.uuid5(uuid.NAMESPACE_URL, f"hermes-task:{task_id}")

    return f"hermes{stable_uuid.hex}"


def queue_sync_request(db, reason=None):

    payload = {
        'pending_request': True,
        'pending_request_at': datetime.now(timezone.utc).isoformat()
    }

    if reason:

        payload['pending_reason'] = reason

    db.collection('system').document('sync').set(payload, merge=True)


def acquire_sync_lock(db, owner_id):

    lock_ref = db.collection('system').document(SYNC_LOCK_DOC_ID)
    now = datetime.now(timezone.utc)
    lock_payload = {
        'owner_id': owner_id,
        'started_at': now.isoformat(),
        'expires_at': (now + timedelta(seconds=SYNC_LOCK_STALE_SECONDS)).isoformat()
    }

    try:

        lock_ref.create(lock_payload)

        return True

    except Exception:

        pass

    try:

        lock_doc = lock_ref.get()

        if lock_doc.exists:

            lock_data = lock_doc.to_dict() or {}
            expires_at = parse_iso_datetime(lock_data.get('expires_at'))

            if expires_at and expires_at > now:

                return False

        lock_ref.delete()
        lock_ref.create(lock_payload)

        return True

    except Exception:

        return False


def release_sync_lock(db, owner_id):

    lock_ref = db.collection('system').document(SYNC_LOCK_DOC_ID)

    try:

        lock_doc = lock_ref.get()

        if not lock_doc.exists:

            return

        lock_data = lock_doc.to_dict() or {}

        if lock_data.get('owner_id') == owner_id:

            lock_ref.delete()

    except Exception:

        pass



def claim_action_dedup_slot(db, titulo, data_limite, horario_inicio, ttl_minutes=15):
    """Reivindica atomicamente a chave (título, data, horário) de uma nova ação para evitar
    duplicatas quando a criação é chamada mais de uma vez para o mesmo pedido (retry do
    modelo, reenvio de webhook do Telegram, duplo toque no botão de confirmação) — sintoma
    relatado como "aparece duplicada no mesmo horário, com evento duplicado na agenda".
    Ao contrário de uma consulta seguida de escrita (sujeita a corrida), usa create() como
    exclusão mútua atômica no nível do documento, no mesmo espírito de acquire_sync_lock.

    Retorna (status, task_id):
      ("proceed", None)     — nenhuma reivindicação concorrente; prossiga com a criação e
                               chame store_action_dedup_result(...) depois de criar a ação.
      ("duplicate", task_id) — outra chamada já concluiu a criação; reaproveite esse ID.
      ("pending", None)      — outra chamada está criando a mesma ação agora; não crie outra.
    """
    import hashlib

    key_raw = f"{(titulo or '').strip().lower()}|{data_limite}|{horario_inicio}"
    key = "aclock_" + hashlib.sha1(key_raw.encode("utf-8")).hexdigest()[:24]
    claim_ref = db.collection("action_creation_claims").document(key)
    now = datetime.now(timezone.utc)

    def _try_create():
        claim_ref.create({"claimed_at": now.isoformat(), "task_id": None})

    try:
        _try_create()
        return "proceed", None
    except Exception:
        pass

    try:
        existing = claim_ref.get()
        if not existing.exists:
            _try_create()
            return "proceed", None

        data = existing.to_dict() or {}
        claimed_at = parse_iso_datetime(data.get("claimed_at"))
        if not claimed_at or (now - claimed_at) >= timedelta(minutes=ttl_minutes):
            claim_ref.delete()
            _try_create()
            return "proceed", None

        task_id = data.get("task_id")
        if task_id:
            return "duplicate", task_id

        for _ in range(3):
            time.sleep(0.4)
            task_id = (claim_ref.get().to_dict() or {}).get("task_id")
            if task_id:
                return "duplicate", task_id

        return "pending", None
    except Exception:
        return "proceed", None


def store_action_dedup_result(db, titulo, data_limite, horario_inicio, task_id):
    import hashlib

    key_raw = f"{(titulo or '').strip().lower()}|{data_limite}|{horario_inicio}"
    key = "aclock_" + hashlib.sha1(key_raw.encode("utf-8")).hexdigest()[:24]
    try:
        db.collection("action_creation_claims").document(key).set({"task_id": task_id}, merge=True)
    except Exception:
        pass


def release_action_dedup_slot(db, titulo, data_limite, horario_inicio):
    """Libera a reivindicação de claim_action_dedup_slot quando a criação falhou, para que
    uma nova tentativa não fique bloqueada como "pending" até o TTL expirar."""
    import hashlib

    key_raw = f"{(titulo or '').strip().lower()}|{data_limite}|{horario_inicio}"
    key = "aclock_" + hashlib.sha1(key_raw.encode("utf-8")).hexdigest()[:24]
    try:
        db.collection("action_creation_claims").document(key).delete()
    except Exception:
        pass


def emit_notification_backend(title, message, n_type='info', link=None):

    from datetime import datetime

    import uuid

    db = get_db()

    notif_id = str(uuid.uuid4())[:9]

    db.collection('notificacoes').document(notif_id).set({

        'id': notif_id,

        'title': title,

        'message': message,

        'type': n_type,

        'timestamp': datetime.now().isoformat(),

        'isRead': False,

        'link': link,

    })



def log_to_firestore(sync_ref, logs, message, force_update=False):

    from datetime import datetime

    timestamp = datetime.now().strftime('%H:%M:%S')

    log_entry = f"[{timestamp}] {message}"

    logs.append(log_entry)

    print(log_entry)

    if force_update:

        sync_ref.update({'logs': logs})



    message_upper = message.upper()

    if "ERRO" in message_upper:

        # Evita falsos positivos com títulos de tarefas ou descrições que contêm a palavra "erro" (ex: [PULL], [-], [X], [+], [^])
        is_false_positive = (
            message_upper.startswith("[PULL]") or
            message_upper.startswith("[-] ") or
            message_upper.startswith("[X] ") or
            message_upper.startswith("[+] ") or
            message_upper.startswith("[^] ") or
            message_upper.startswith("[CAL->HERMES]") or
            message_upper.startswith("[BOLETO]") or
            message_upper.startswith("[PIX] PROCESSADO") or
            message_upper.startswith("[GMAIL] E-MAIL FINANCEIRO")
        )

        if not is_false_positive:

            emit_notification_backend("Erro de Sincronização", message, 'error')

    # Nota: a notificação "Gasto Realizado via Pix" é emitida diretamente em sync_pix_emails
    # (já com valor e saldo disponível), não mais a partir do texto deste log.



def classify_task(title, notes):

    import re

    text = f"{title} {notes}".upper()

    area_tematica, contabilizar_meta = 'NÃO CLASSIFICADA', False

    tags = re.findall(r'\[(.*?)\]|TAG:\s*([\w\-]+)', text)

    tags = [t[0].upper() if t[0] else t[1].upper() for t in tags]

    

    if any(tag in ['CLC', 'LICITACAO'] for tag in tags):

        area_tematica, contabilizar_meta = 'CLC', True

    elif any(tag in ['ASSISTENCIA', 'ESTUDANTIL'] for tag in tags):

        area_tematica, contabilizar_meta = 'ASSISTÊNCIA', True

    elif 'GERAL' in tags:

        area_tematica = 'GERAL'



    # Se não classificou por tag, tenta por palavra-chave no texto

    if area_tematica == 'NÃO CLASSIFICADA':

        clc_keywords = ['LICITAÇÃO', 'LICITACAO', 'PREGÃO', 'PREGAO', 'CONTRATO', 'DISPENSA', 'INEXIGIBILIDADE', 'COMPRA', 'AQUISIÇÃO', 'AQUISICAO', 'PROCESSO']

        assist_keywords = ['ASSISTÊNCIA', 'ASSISTENCIA', 'ESTUDANTIL', 'ALUNO', 'BOLSA', 'AUXÍLIO', 'AUXILIO', 'PERMANÊNCIA', 'PERMANENCIA']



        if any(kw in text for kw in clc_keywords):

            area_tematica, contabilizar_meta = 'CLC', True

        elif any(kw in text for kw in assist_keywords):

            area_tematica, contabilizar_meta = 'ASSISTÊNCIA', True



    return area_tematica, None, contabilizar_meta



def extract_time_from_notes(notes):

    import re

    if not notes: return None, None

    match = re.search(r'\[Horário:\s*(\d{2}:\d{2})\s*-\s*(\d{2}:\d{2})\]', notes)

    if match:

        return match.group(1), match.group(2)

    return None, None



def update_notes_with_time(notes, start, end):

    import re

    if not notes: notes = ""

    pattern = r'\[Horário:\s*\d{2}:\d{2}\s*-\s*\d{2}:\d{2}\]'

    new_block = f"[Horário: {start} - {end}]" if start and end else ""

    

    if re.search(pattern, notes):

        if new_block:

            return re.sub(pattern, new_block, notes)

        else:

            return re.sub(pattern, '', notes).strip()

    else:

        if new_block:

            return f"{notes}\n\n{new_block}".strip()

        else:

            return notes



def normalize_task_title(title):

    import re

    if not isinstance(title, str):
        return title

    compact = re.sub(r'\s+', ' ', title).strip()
    if not compact:
        return compact

    small_words = {
        'de', 'da', 'do', 'das', 'dos',
        'e', 'em', 'na', 'no', 'nas', 'nos',
        'a', 'o', 'as', 'os',
        'para', 'por', 'com'
    }

    def normalize_piece(piece, is_first):
        if not piece:
            return piece
        if re.fullmatch(r'[A-Z0-9]{2,5}', piece):
            return piece
        lower = piece.lower()
        if (not is_first) and lower in small_words:
            return lower
        return lower[:1].upper() + lower[1:]

    words = []
    for word_index, word in enumerate(compact.split(' ')):
        parts = re.split(r'([/-])', word)
        normalized_parts = []
        part_index = 0
        for part in parts:
            if part in ('/', '-'):
                normalized_parts.append(part)
                continue
            normalized_parts.append(normalize_piece(part, word_index == 0 and part_index == 0))
            part_index += 1
        words.append(''.join(normalized_parts))

    return ' '.join(words)


def parse_iso_datetime(value):

    from datetime import datetime

    if not isinstance(value, str):
        return None

    text = value.strip()
    if not text:
        return None

    # Google APIs frequentemente retornam UTC com sufixo "Z"
    if text.endswith('Z'):
        text = f"{text[:-1]}+00:00"

    try:
        return datetime.fromisoformat(text)
    except Exception:
        return None


def is_remote_calendar_newer(remote_updated, local_updated):
    from datetime import timezone
    remote_dt = parse_iso_datetime(remote_updated)
    local_dt = parse_iso_datetime(local_updated)

    if remote_dt and local_dt:
        # Se um for aware e o outro naive, forçamos o naive para UTC para permitir a comparação
        if remote_dt.tzinfo is not None and local_dt.tzinfo is None:
            local_dt = local_dt.replace(tzinfo=timezone.utc)
        elif remote_dt.tzinfo is None and local_dt.tzinfo is not None:
            remote_dt = remote_dt.replace(tzinfo=timezone.utc)
            
        return remote_dt > local_dt
    if remote_dt and not local_dt:
        return True
    return False


def is_iso_after(left_value, right_value):
    from datetime import timezone
    left_dt = parse_iso_datetime(left_value)
    right_dt = parse_iso_datetime(right_value)

    if left_dt and right_dt:
        if left_dt.tzinfo is not None and right_dt.tzinfo is None:
            right_dt = right_dt.replace(tzinfo=timezone.utc)
        elif left_dt.tzinfo is None and right_dt.tzinfo is not None:
            left_dt = left_dt.replace(tzinfo=timezone.utc)
        return left_dt > right_dt
    if left_dt and not right_dt:
        return True
    if right_dt and not left_dt:
        return False
    if left_value and right_value:
        return str(left_value) > str(right_value)
    return False



def extract_schedule_from_calendar_event(event, tz_name='America/Sao_Paulo'):

    from datetime import timedelta
    from zoneinfo import ZoneInfo

    start_info = event.get('start', {}) or {}
    end_info = event.get('end', {}) or {}

    start_dt = parse_iso_datetime(start_info.get('dateTime'))
    end_dt = parse_iso_datetime(end_info.get('dateTime'))

    if start_dt:
        tz = ZoneInfo(tz_name)
        if start_dt.tzinfo is None:
            start_dt = start_dt.replace(tzinfo=tz)
        start_local = start_dt.astimezone(tz)

        if end_dt:
            if end_dt.tzinfo is None:
                end_dt = end_dt.replace(tzinfo=tz)
            end_local = end_dt.astimezone(tz)
        else:
            end_local = start_local + timedelta(hours=1)

        return {
            'data_inicio': start_local.date().isoformat(),
            'data_limite': start_local.date().isoformat(),
            'horario_inicio': start_local.strftime('%H:%M'),
            'horario_fim': end_local.strftime('%H:%M')
        }

    # Eventos "dia inteiro" (sem horário) - mantém data e limpa horários
    start_date = start_info.get('date')
    if isinstance(start_date, str) and start_date.strip():
        single_date = start_date.strip()
        return {
            'data_inicio': single_date,
            'data_limite': single_date,
            'horario_inicio': None,
            'horario_fim': None
        }

    return None


def sync_google_tasks_pull(service, sync_ref, logs):

    from datetime import datetime

    # Sincronização Google Tasks -> Ações desativada a pedido do usuário.
    # Nenhuma tarefa do Google cria ou altera ações do Hermes.
    if not SYNC_GOOGLE_TASKS_ENABLED:
        print("[PULL] Sincronização Google Tasks -> Ações desativada (SYNC_GOOGLE_TASKS_ENABLED=False).")
        return

    db = get_db()

    try:

        results = service.tasklists().list().execute()

        tasklist_id = next((item['id'] for item in results.get('items', []) if 'tarefa' in item['title'].lower()), None)

        if not tasklist_id: return

        

        g_tasks = []

        next_page_token = None

        while True:

            res = service.tasks().list(tasklist=tasklist_id, showCompleted=True, showHidden=True, maxResults=100, pageToken=next_page_token).execute()

            g_tasks.extend(res.get('items', []))

            if not res.get('nextPageToken') or len(g_tasks) >= 200: break

            next_page_token = res.get('nextPageToken')



        local_tasks = {t.to_dict().get('google_id'): (t.id, t.to_dict()) for t in db.collection('tarefas').stream() if t.to_dict().get('google_id')}

        

        for gt in g_tasks:

            g_id = gt['id']
            title_raw = gt.get('title', '(Sem Título)')
            title = normalize_task_title(title_raw)
            title_was_normalized = title != title_raw

            g_updated = gt.get('updated', '')

            status = 'concluído' if gt.get('status') == 'completed' else 'em andamento'

            g_due = gt.get('due', '').split('T')[0] if gt.get('due') else None

            

            # Extração de horários das notas

            g_notes = gt.get('notes', '')

            h_inicio, h_fim = extract_time_from_notes(g_notes)

            

            # Duração padrão de 1h se houver início mas não fim

            if h_inicio and not h_fim:

                try:

                    h, m = map(int, h_inicio.split(':'))

                    h_fim = f"{(h+1)%24:02d}:{m:02d}"

                except: pass



            if g_id in local_tasks:

                doc_id, t_old = local_tasks[g_id]

                if is_iso_after(g_updated, t_old.get('data_atualizacao', '')):

                    applied_updated = datetime.now().isoformat() if title_was_normalized else g_updated

                    update_data = {

                        'titulo': title, 'status': status, 'data_atualizacao': applied_updated,

                        'data_conclusao': gt.get('completed'), 'notas': g_notes,

                        'horario_inicio': h_inicio, 'horario_fim': h_fim

                    }

                    if g_due:
                        update_data['data_limite'] = g_due
                        update_data['data_inicio'] = g_due

                    db.collection('tarefas').document(doc_id).update(update_data)

                    log_to_firestore(sync_ref, logs, f"[-] ATUALIZADA: {title}")

            else:

                # Importação desativada a pedido do usuário (criação de tarefas no Google Tasks não gera mais ações no Hermes)

                print(f"[PULL] Ignorada importação de nova tarefa '{title}' do Google Tasks (funcionalidade desativada).")

    except Exception as e:

        log_to_firestore(sync_ref, logs, f"ERRO PULL: {e}")



from googleapiclient.errors import HttpError



def sync_google_tasks_push(service, calendar_service, sync_ref, logs):

    from datetime import datetime

    db = get_db()
    calendar_id = get_target_calendar_id(db)

    try:

        # Sincronização Ações -> Google Tasks desativada a pedido do usuário.
        # Mantém-se apenas a parte de Google Calendar deste push (ver SYNC_GOOGLE_TASKS_ENABLED).
        tasklist_id = None
        g_tasks_map = {}

        if SYNC_GOOGLE_TASKS_ENABLED:

            results = service.tasklists().list().execute()

            tasklist_id = next((item['id'] for item in results.get('items', []) if 'tarefa' in item['title'].lower()), None)

            if not tasklist_id:
                # Verifica se há list default
                default_list = service.tasklists().get(tasklist='@default').execute()
                tasklist_id = default_list.get('id')

            # Pega todas as tarefas do Google (com paginação) para o mapa

            next_page_token = None

            while True:

                g_results = service.tasks().list(tasklist=tasklist_id, showCompleted=True, showHidden=True, maxResults=100, pageToken=next_page_token).execute()

                for item in g_results.get('items', []):

                    g_tasks_map[item['id']] = item

                next_page_token = g_results.get('nextPageToken')

                if not next_page_token or len(g_tasks_map) >= 500: break



        for doc in db.collection('tarefas').stream():

            t = doc.to_dict()

            cat = t.get('area_tematica', '')

            if cat == 'SISTEMAS': continue



            g_id = t.get('google_id')
            raw_title = t.get('titulo') or '(Sem Título)'
            title = normalize_task_title(raw_title)
            local_updated = t.get('data_atualizacao', '')
            if title != raw_title:
                local_updated = datetime.now().isoformat()
                doc.reference.update({'titulo': title, 'data_atualizacao': local_updated})

            if t.get('status') == 'excluído':

                # 1. Remover do Google Tasks (desativado: ver SYNC_GOOGLE_TASKS_ENABLED)
                if SYNC_GOOGLE_TASKS_ENABLED and g_id:

                    try:

                        service.tasks().delete(tasklist=tasklist_id, task=g_id).execute()

                        log_to_firestore(sync_ref, logs, f"[X] REMOVIDA DO GOOGLE TASKS: {title}")

                    except HttpError as e:

                        if e.resp.status == 404:

                            log_to_firestore(sync_ref, logs, f"[!] Task {g_id} já não existia no Google Tasks.")

                # 2. Remover do Google Calendar
                cal_id = t.get('google_calendar_id')
                if cal_id:
                    try:
                        calendar_service.events().delete(calendarId=calendar_id, eventId=cal_id).execute()
                        log_to_firestore(sync_ref, logs, f"[X] REMOVIDA DO CALENDAR: {title}")
                    except Exception as e:
                        log_to_firestore(sync_ref, logs, f"[!] Erro ao remover do Calendar: {e}")

                doc.reference.delete()

                continue

            

            g_status = 'completed' if t.get('status') == 'concluído' else 'needsAction'
            
            # Decisão: se houver horario_inicio, enviaremos pro CALENDAR como EVENTO também!
            sync_to_calendar = bool(t.get('horario_inicio') and t.get('data_limite') and t.get('data_limite') != '-')

            if t.get('data_limite') and t.get('data_limite') != '-':
                # Pro Tasks, precisa ser 00:00:00.000Z por limitacao da API.
                g_due = f"{t.get('data_limite')}T00:00:00.000Z"
            else:
                g_due = None

            # Atualiza as notas com o horário para garantir a sincronia

            h_inicio, h_fim = t.get('horario_inicio'), t.get('horario_fim')

            # Se não houver fim mas houver início, assume-se 1h de duração

            if h_inicio and not h_fim:

                try:

                    h, m = map(int, h_inicio.split(':'))

                    h_fim = f"{(h+1)%24:02d}:{m:02d}"

                except: pass

            

            updated_notes = update_notes_with_time(t.get('notas', ''), h_inicio, h_fim)

            # --- PARTE 1: Sincronia Padrão do Google Tasks (desativada: ver SYNC_GOOGLE_TASKS_ENABLED) ---
            if SYNC_GOOGLE_TASKS_ENABLED:
                if not g_id:
                    body = {'title': title, 'notes': updated_notes, 'status': g_status}
                    if g_due: body['due'] = g_due
                    new_task = service.tasks().insert(tasklist=tasklist_id, body=body).execute()
                    doc.reference.update({'google_id': new_task['id'], 'data_atualizacao': new_task.get('updated'), 'notas': updated_notes, 'horario_fim': h_fim if not t.get('horario_fim') else t.get('horario_fim')})
                    log_to_firestore(sync_ref, logs, f"[+] ENVIADA TASKS: {title}")
                    g_id = new_task['id'] # Para usar no Calendar se precisar
                elif g_id in g_tasks_map and is_iso_after(local_updated, g_tasks_map[g_id].get('updated', '')):
                    body = {'id': g_id, 'title': title, 'notes': updated_notes, 'status': g_status}
                    if g_due: body['due'] = g_due
                    try:
                        updated_task = service.tasks().update(tasklist=tasklist_id, task=g_id, body=body).execute()
                        log_to_firestore(sync_ref, logs, f"[^] ATUALIZADA NO TASKS: {title}")
                        sync_updates = {}
                        if updated_task.get('updated'):
                            sync_updates['data_atualizacao'] = updated_task.get('updated')
                        if updated_notes != t.get('notas', ''):
                            sync_updates['notas'] = updated_notes
                        if h_fim and not t.get('horario_fim'):
                            sync_updates['horario_fim'] = h_fim
                        if sync_updates:
                            doc.reference.update(sync_updates)
                    except HttpError as e:
                        if e.resp.status == 404:
                            doc.reference.update({'google_id': None})

            # --- PARTE 2: Sincronia Google Calendar (Se possuir horário) ---
            cal_id = t.get('google_calendar_id')
            desired_event_id = build_task_calendar_event_id(doc.id)
            if sync_to_calendar and g_status == 'needsAction': # Só agenda eventos não concluídos
                # Prepara o Evento
                from datetime import datetime, timezone
                try:
                    # Converte pra ISO 8601 string para API (Timezone default da máquina rodando)
                    start_dt = f"{t.get('data_limite')}T{h_inicio}:00-03:00"
                    end_dt = f"{t.get('data_limite')}T{h_fim}:00-03:00"
                    
                    event_body = {
                        'summary': f"Tarefa: {title}",
                        'description': updated_notes,
                        'start': {'dateTime': start_dt, 'timeZone': 'America/Sao_Paulo'},
                        'end': {'dateTime': end_dt, 'timeZone': 'America/Sao_Paulo'},
                        'extendedProperties': {
                            'private': {
                                'hermes_task_id': doc.id
                            }
                        }
                    }
                    insert_event_body = dict(event_body)
                    insert_event_body['id'] = desired_event_id
                    
                    if not cal_id:
                        # Cria novo
                        try:
                            new_event = calendar_service.events().insert(calendarId=calendar_id, body=insert_event_body).execute()
                        except HttpError as cal_err:
                            if cal_err.resp.status != 409:
                                raise
                            new_event = calendar_service.events().get(calendarId=calendar_id, eventId=desired_event_id).execute()
                        doc.reference.update({'google_calendar_id': new_event['id']})
                        log_to_firestore(sync_ref, logs, f"[+] ALOCADA CALENDAR: {title}")
                    else:
                        # Atualiza
                        try:
                            calendar_service.events().update(calendarId=calendar_id, eventId=cal_id, body=event_body).execute()
                        except HttpError as cal_err:
                            if cal_err.resp.status == 404:
                                try:
                                    new_event = calendar_service.events().insert(calendarId=calendar_id, body=insert_event_body).execute()
                                except HttpError as conflict_err:
                                    if conflict_err.resp.status != 409:
                                        raise
                                    new_event = calendar_service.events().get(calendarId=calendar_id, eventId=desired_event_id).execute()
                                doc.reference.update({'google_calendar_id': new_event['id']})
                except Exception as ce:
                    log_to_firestore(sync_ref, logs, f"[CAL][!] Falha ao sincronizar evento da tarefa '{title}': {ce}")
            elif (not sync_to_calendar or g_status == 'completed') and cal_id:
                 # Tem ID no cal, mas perdeu horario ou foi completada - Remove do Calendar
                 try:
                     calendar_service.events().delete(calendarId=calendar_id, eventId=cal_id).execute()
                 except HttpError as ce: pass
                 doc.reference.update({'google_calendar_id': None})
                 


    except Exception as e:

        if is_google_invalid_grant_error(e):
            raise GoogleAuthRevokedError(GOOGLE_REAUTH_MESSAGE) from e

        log_to_firestore(sync_ref, logs, f"ERRO PUSH: {e}")



def sync_google_calendar(service, sync_ref, logs):

    from datetime import datetime, timedelta, timezone

    db = get_db()
    calendar_ids = get_sync_calendar_ids(db)

    try:

        log_to_firestore(sync_ref, logs, "Sincronizando Google Calendar...", True)

        time_min = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat().replace('+00:00', 'Z')

        time_max = (datetime.now(timezone.utc) + timedelta(days=30)).isoformat().replace('+00:00', 'Z')

        count = 0

        seen_ids = set()

        linked_tasks_by_event_id = {}
        for task_doc in db.collection('tarefas').stream():
            task_data = task_doc.to_dict() or {}
            linked_event_id = task_data.get('google_calendar_id')
            deterministic_event_id = build_task_calendar_event_id(task_doc.id)
            linked_tasks_by_event_id.setdefault(deterministic_event_id, []).append((task_doc.reference, task_data))
            if isinstance(linked_event_id, str) and linked_event_id.strip():
                linked_tasks_by_event_id.setdefault(linked_event_id.strip(), []).append((task_doc.reference, task_data))

        for calendar_id in calendar_ids:

            try:

                events_result = service.events().list(

                    calendarId=calendar_id, timeMin=time_min, timeMax=time_max,

                    singleEvents=True, orderBy='startTime'

                ).execute()

                events = events_result.get('items', [])

            except Exception as cal_err:

                if is_google_invalid_grant_error(cal_err):
                    raise GoogleAuthRevokedError(GOOGLE_REAUTH_MESSAGE) from cal_err

                log_to_firestore(sync_ref, logs, f"[CAL][!] Falha ao listar agenda '{calendar_id}': {cal_err}")

                continue

            for event in events:

                event_id = event['id']

                doc_id = f"{calendar_id}__{event_id}"

                seen_ids.add(doc_id)

                summary = event.get('summary', '(Sem titulo)')

                start = event['start'].get('dateTime', event['start'].get('date'))

                end = event['end'].get('dateTime', event['end'].get('date'))

                # Eventos criados pelo próprio Hermes (ver build_task_calendar_event_id /
                # sync_google_tasks_push) levam extendedProperties.private.hermes_task_id.
                # Guardamos essa marca para que produtores de sinal (ex.: link_calendar_events_to_actions)
                # possam ignorá-los — sinalizar de volta um evento que o próprio Hermes criou é redundante.
                hermes_task_id = ((event.get('extendedProperties') or {}).get('private') or {}).get('hermes_task_id')

                db.collection('google_calendar_events').document(doc_id).set({

                    'google_id': event_id,

                    'calendar_id': calendar_id,

                    'titulo': summary,

                    'data_inicio': start,

                    'data_fim': end,

                    'criado_pelo_hermes': bool(hermes_task_id),

                    'last_sync': datetime.now().isoformat()

                }, merge=True)

                count += 1

                # Sincronia inversa: evento do Calendar (criado pelo Hermes) atualiza data/horário da tarefa
                linked_tasks = linked_tasks_by_event_id.get(event_id, [])
                if not linked_tasks:
                    continue

                schedule = extract_schedule_from_calendar_event(event)
                if not schedule:
                    continue

                event_updated = event.get('updated', '')
                for task_ref, task_data in linked_tasks:
                    local_updated = task_data.get('data_atualizacao', '')
                    if not is_iso_after(event_updated, local_updated):
                        continue

                    local_date = task_data.get('data_limite') or task_data.get('data_inicio')
                    local_start = task_data.get('horario_inicio')
                    local_end = task_data.get('horario_fim')
                    has_schedule_change = (
                        local_date != schedule.get('data_limite')
                        or local_start != schedule.get('horario_inicio')
                        or local_end != schedule.get('horario_fim')
                    )

                    if not has_schedule_change:
                        continue

                    updated_notes = update_notes_with_time(
                        task_data.get('notas', ''),
                        schedule.get('horario_inicio'),
                        schedule.get('horario_fim')
                    )

                    task_updates = {
                        **schedule,
                        'notas': updated_notes,
                        # Usa "agora" para garantir que o push subsequente preserve esse ajuste no Tasks/Calendar
                        'data_atualizacao': datetime.now().isoformat()
                    }
                    task_ref.update(task_updates)
                    task_data.update(task_updates)
                    log_to_firestore(sync_ref, logs, f"[CAL->HERMES] Horário atualizado pela agenda: {task_data.get('titulo', '(Sem titulo)')}")

        # Limpeza de eventos deletados no Google Calendar (somente janela sincronizada)

        docs = db.collection('google_calendar_events')\
            .where('data_inicio', '>=', time_min)\
            .where('data_inicio', '<=', time_max)\
            .stream()

        deleted_count = 0

        for doc in docs:

            if doc.id not in seen_ids:

                doc.reference.delete()

                deleted_count += 1

        log_to_firestore(sync_ref, logs, f"[CAL] {count} eventos sincronizados em {len(calendar_ids)} agenda(s). {deleted_count} removidos.")

    except Exception as e:

        if is_google_invalid_grant_error(e):
            raise GoogleAuthRevokedError(GOOGLE_REAUTH_MESSAGE) from e

        log_to_firestore(sync_ref, logs, f"ERRO CAL: {e}")

def sync_pix_emails(service, sync_ref, logs):
    """
    Busca emails de Pix e registra no Financeiro (Versão Cloud Function)
    """
    import re
    import time
    from datetime import datetime, timezone

    db = get_db()

    try:
        log_to_firestore(sync_ref, logs, "Buscando emails de Pix a partir de 01/02/2026...")

        # Query ampliada para capturar e-mails transacionais de qualquer instituição (Pix, Google Pay, PicPay, etc)
        query = 'after:2026/02/01 (subject:(Pix OR "Google Pay" OR "PicPay" OR "Pagamento" OR "Transferência" OR "Comprovante") OR "Pix" OR "Google Pay" OR "PicPay")'

        messages = []
        page_token = None
        while True:
            results = service.users().messages().list(
                userId='me', q=query, maxResults=100, pageToken=page_token
            ).execute()
            batch = results.get('messages', [])
            if batch:
                messages.extend(batch)
            page_token = results.get('nextPageToken')
            if not page_token or len(messages) >= 500:
                break

        if not messages:
            log_to_firestore(sync_ref, logs, "Nenhum Pix/Pagamento encontrado para os critérios de busca.")
            return

        log_to_firestore(sync_ref, logs, f"Encontrados {len(messages)} e-mails potenciais de Pix/Pagamento. Analisando...")

        # Cache de transações existentes para evitar duplicatas (Bloqueio de duplicidade financeira)
        existing_transactions = []
        existing_income = []
        existing_google_ids = set()

        def parse_iso_date(date_str):
            if not date_str: return None
            try:
                dt = datetime.fromisoformat(date_str.replace('Z', '+00:00'))
                if dt.tzinfo is None:
                    dt = dt.replace(tzinfo=timezone.utc)
                return dt
            except:
                return None

        for t in db.collection('finance_transactions').stream():
            data = t.to_dict()
            existing_transactions.append({
                'doc_id': t.id,
                'description': data.get('description', ''),
                'amount': data.get('amount', 0.0),
                'date': parse_iso_date(data.get('date')),
                'pix_id': data.get('pix_id'),
                'google_message_id': data.get('google_message_id'),
                'status': data.get('status')
            })
            if data.get('google_message_id'): existing_google_ids.add(data['google_message_id'])

        for t in db.collection('finance_income').stream():
            data = t.to_dict()
            existing_income.append({
                'doc_id': t.id,
                'description': data.get('description', ''),
                'amount': data.get('amount', 0.0),
                'date': parse_iso_date(data.get('date')),
                'pix_id': data.get('pix_id'),
                'google_message_id': data.get('google_message_id')
            })
            if data.get('google_message_id'): existing_google_ids.add(data['google_message_id'])

        processed_emails_doc = db.collection('system').document('processed_emails').get()
        processed_ids = set(processed_emails_doc.to_dict().get('ids', [])) if processed_emails_doc.exists else set()
        new_processed_ids = []
        
        # Cache de rubricas de renda e rubricas de contas (saídas)
        income_rubrics_cache = []
        for r in db.collection('income_rubrics').stream():
            d = r.to_dict()
            income_rubrics_cache.append({'id': r.id, 'desc': d.get('description', '').lower()})

        bill_rubrics_cache = []
        for r in db.collection('bill_rubrics').stream():
            d = r.to_dict()
            desc = d.get('description', '')
            keywords = [w.strip().lower() for w in re.split(r'[\(\)\s,-]+', desc) if len(w.strip()) > 2]
            bill_rubrics_cache.append({'id': r.id, 'desc': desc.lower(), 'keywords': keywords, 'full_desc': desc})

        # Orçamento do mês corrente e gasto acumulado, para exibir "Saldo disponível" nas notificações de Pix
        now_utc = datetime.now(timezone.utc)
        current_month_key = f"{now_utc.year}-{now_utc.month:02d}"
        finance_settings_doc = db.collection('finance_settings').document('config').get()
        finance_settings_data = finance_settings_doc.to_dict() if finance_settings_doc.exists else {}
        monthly_budget = (finance_settings_data.get('monthlyBudgets') or {}).get(
            current_month_key, finance_settings_data.get('monthlyBudget', 0) or 0
        )
        current_month_spend = sum(
            t['amount'] for t in existing_transactions
            if t.get('status') != 'deleted' and t.get('date')
            and t['date'].year == now_utc.year and t['date'].month == now_utc.month
        )

        for msg in messages:
            msg_id = msg['id']

            if msg_id in processed_ids or msg_id in existing_google_ids:
                archive_gmail_message(service, msg_id, sync_ref, logs, "pix-ja-processado")
                continue

            details = service.users().messages().get(userId='me', id=msg_id).execute()

            internal_date_ms = int(details.get('internalDate', time.time() * 1000))
            dt = datetime.fromtimestamp(internal_date_ms / 1000.0, tz=timezone.utc)

            snippet = details.get('snippet', '')
            sender = ''
            subject = ''
            for header in details.get('payload', {}).get('headers', []):
                h_name = header['name'].lower()
                if h_name == 'from':
                    sender = header['value']
                elif h_name == 'subject':
                    subject = header['value']

            # Descartar e-mails da XP do extrato de conta corrente PicPay
            # (remetente real é "XP Inc." <no_reply@xpi.com.br> — "xp.com" nunca casava com "xpi.com.br")
            if 'xpi.com.br' in sender.lower() or 'xp investimentos' in sender.lower() or 'xpinvestimentos' in sender.lower() or 'xp inc' in sender.lower() or 'xp ' in subject.lower():
                log_to_firestore(sync_ref, logs, f"[GMAIL-PIX] Ignorando e-mail da XP no extrato da conta PicPay (msg_id: {msg_id})")
                archive_gmail_message(service, msg_id, sync_ref, logs, "pix-xp-ignorado")
                continue

            content = f"{subject} {snippet}"
            value_match = re.search(r'R\$\s*([\d\.,]+)', content)
            pix_id_match = re.search(r'\b(E[A-Z0-9]{31})\b', content)
            pix_id = pix_id_match.group(1) if pix_id_match else None

            if value_match:
                val_raw = value_match.group(1).rstrip('.').rstrip(',')
                if '.' in val_raw and ',' in val_raw:
                    # Formato BR completo: "." separador de milhar, "," separador decimal (ex: "2.708,53")
                    val_str = val_raw.replace('.', '').replace(',', '.')
                elif ',' in val_raw:
                    # Só decimal, sem milhar (ex: "51,86")
                    val_str = val_raw.replace(',', '.')
                elif '.' in val_raw:
                    # Só ponto: pode ser separador de milhar sem centavos (ex: "1.343" = R$ 1.343,00)
                    # ou decimal informal (ex: "1.34"). Milhar em pt-BR sempre agrupa em blocos de 3 dígitos;
                    # decimais de moeda têm no máximo 2 casas — usamos isso para desambiguar.
                    integer_part, _, frac_part = val_raw.rpartition('.')
                    if len(frac_part) == 3 and integer_part:
                        val_str = val_raw.replace('.', '')
                    else:
                        val_str = val_raw
                else:
                    val_str = val_raw
                try:
                    amount = float(val_str)
                except ValueError:
                    continue

                if amount <= 0:
                    continue

                is_income = any(word in content.lower() for word in ['recebido', 'recebeu', 'recebida', 'recebimento', 'creditado', 'entrada'])
                description = f"Pix: {subject}"
                iso_date = dt.isoformat()

                # Conciliação com Rubricas de Contas (Contas Fixas / Saídas)
                if not is_income:
                    matched_bill_rubric = None
                    clean_content = content.lower()
                    for rb in bill_rubrics_cache:
                        if rb['desc'] in clean_content or any(kw in clean_content for kw in rb['keywords']):
                            matched_bill_rubric = rb
                            break

                    if matched_bill_rubric:
                        month = dt.month - 1
                        year = dt.year
                        
                        found_bill_doc = None
                        for fb_doc in db.collection('fixed_bills').where('month', '==', month).where('year', '==', year).stream():
                            fb_data = fb_doc.to_dict()
                            if fb_data.get('rubricId') == matched_bill_rubric['id'] or matched_bill_rubric['desc'] in fb_data.get('description', '').lower():
                                found_bill_doc = fb_doc
                                break
                        
                        if found_bill_doc:
                            db.collection('fixed_bills').document(found_bill_doc.id).update({
                                'isPaid': True,
                                'amount': amount,
                                'data_pagamento': iso_date,
                                'google_message_id': msg_id,
                                'pix_id': pix_id
                            })
                            log_to_firestore(sync_ref, logs, f"[CONCILIAÇÃO] Conta Fixa '{matched_bill_rubric['full_desc']}' baixada como PAGA (R$ {amount:.2f}). Ignorada nos lançamentos avulsos.")
                        else:
                            db.collection('fixed_bills').add({
                                'description': matched_bill_rubric['full_desc'],
                                'amount': amount,
                                'dueDay': dt.day,
                                'month': month,
                                'year': year,
                                'isPaid': True,
                                'rubricId': matched_bill_rubric['id'],
                                'google_message_id': msg_id,
                                'pix_id': pix_id,
                                'created_at': iso_date
                            })
                            log_to_firestore(sync_ref, logs, f"[CONCILIAÇÃO] Criada e baixada Conta Fixa '{matched_bill_rubric['full_desc']}' (R$ {amount:.2f}). Ignorada nos lançamentos avulsos.")
                        
                        new_processed_ids.append(msg_id)
                        archive_gmail_message(service, msg_id, sync_ref, logs, "pix-conciliado-conta-fixa")
                        continue

                is_duplicate = False
                target_cache = existing_income if is_income else existing_transactions

                for item in target_cache:
                    if item.get('google_message_id') == msg_id:
                        is_duplicate = True
                        break
                    if pix_id and item.get('pix_id') == pix_id:
                        is_duplicate = True
                        break
                    if item.get('amount') and abs(item['amount'] - amount) < 0.01 and item.get('date'):
                        item_dt = item['date']
                        if item_dt.tzinfo is None:
                            item_dt = item_dt.replace(tzinfo=timezone.utc)
                        diff_seconds = abs((item_dt - dt).total_seconds())
                        
                        if diff_seconds <= 7200:
                            is_duplicate = True
                            existing_desc = item.get('description', '')
                            if ("pagamento realizado via pix" in existing_desc.lower() or "pix:" in existing_desc.lower()) and "google pay" in description.lower():
                                update_fields = {'description': description}
                                if pix_id: update_fields['pix_id'] = pix_id
                                collection_name = 'finance_income' if is_income else 'finance_transactions'
                                db.collection(collection_name).document(item['doc_id']).update(update_fields)
                                item['description'] = description
                                log_to_firestore(sync_ref, logs, f"[PIX] Atualizada descrição do lançamento existente para: '{description}'")
                            break

                if is_duplicate:
                    new_processed_ids.append(msg_id)
                    archive_gmail_message(service, msg_id, sync_ref, logs, "pix-duplicado")
                    continue

                # Salva no banco
                if is_income:
                    matched_rubric_id = None
                    clean_desc = description.replace('Pix: ', '').lower()
                    for rb in income_rubrics_cache:
                        if rb['desc'] in clean_desc or clean_desc in rb['desc']:
                            matched_rubric_id = rb['id']
                            break

                    new_record = {
                        'description': description, 'amount': amount, 'day': dt.day,
                        'month': dt.month - 1, 'year': dt.year,
                        'category': 'Geral', 'isReceived': True, 'date': iso_date,
                        'google_message_id': msg_id, 'pix_id': pix_id,
                        'rubricId': matched_rubric_id
                    }
                    doc_ref = db.collection('finance_income').add(new_record)[1]
                    existing_income.append({'doc_id': doc_ref.id, 'amount': amount, 'date': dt, 'pix_id': pix_id, 'description': description, 'google_message_id': msg_id})
                else:
                    sprint = 1 if dt.day < 8 else 2 if dt.day < 15 else 3 if dt.day < 22 else 4
                    new_record = {
                        'description': description, 'amount': amount, 'date': iso_date,
                        'sprint': sprint, 'category': 'Geral',
                        'google_message_id': msg_id, 'pix_id': pix_id
                    }
                    doc_ref = db.collection('finance_transactions').add(new_record)[1]
                    existing_transactions.append({'doc_id': doc_ref.id, 'amount': amount, 'date': dt, 'pix_id': pix_id, 'description': description, 'google_message_id': msg_id, 'status': None})

                    if dt.year == now_utc.year and dt.month == now_utc.month:
                        current_month_spend += amount
                    saldo_disponivel = monthly_budget - current_month_spend
                    emit_notification_backend(
                        "Gasto Realizado via Pix",
                        f"Valor: R$ {amount:.2f}\nSaldo disponível: R$ {saldo_disponivel:.2f}",
                        'expense', 'financeiro'
                    )

                new_processed_ids.append(msg_id)
                log_to_firestore(sync_ref, logs, f"[PIX] Processado: {description} (R$ {amount:.2f})")
                archive_gmail_message(service, msg_id, sync_ref, logs, "pix-lancado")

        if new_processed_ids:
            updated_ids = list(processed_ids.union(new_processed_ids))[-1000:]
            db.collection('system').document('processed_emails').set({'ids': updated_ids}, merge=True)

        cleanup_retroactive_pix_duplicates(db, sync_ref, logs)

    except Exception as e:
        log_to_firestore(sync_ref, logs, f"ERRO PIX: {e}")


def cleanup_retroactive_pix_duplicates(db, sync_ref=None, logs=None):
    """
    Varre a coleção finance_transactions e finance_income para identificar e consolidar
    duplicatas retroativas (ex: e-mail do Google Pay + e-mail do PicPay/Banco para a mesma compra < 2 horas).
    """
    import re
    from datetime import datetime, timezone

    def parse_dt(d_str):
        if not d_str: return None
        try:
            dt = datetime.fromisoformat(d_str.replace('Z', '+00:00'))
            if dt.tzinfo is None: dt = dt.replace(tzinfo=timezone.utc)
            return dt
        except: return None

    log_to_firestore(sync_ref, logs, "Iniciando verificação retroativa de duplicatas no Firestore...")

    for col in ['finance_transactions', 'finance_income']:
        try:
            docs = list(db.collection(col).stream())
            items = []
            for doc in docs:
                d = doc.to_dict()
                if d.get('status') == 'deleted': continue
                items.append({
                    'id': doc.id,
                    'description': d.get('description', ''),
                    'amount': d.get('amount', 0.0),
                    'date': parse_dt(d.get('date')),
                    'google_message_id': d.get('google_message_id'),
                    'pix_id': d.get('pix_id')
                })

            items.sort(key=lambda x: x['date'] or datetime.min.replace(tzinfo=timezone.utc))

            removed_count = 0
            to_delete = set()
            for i in range(len(items)):
                if items[i]['id'] in to_delete: continue
                for j in range(i + 1, len(items)):
                    if items[j]['id'] in to_delete: continue

                    if items[i]['date'] and items[j]['date']:
                        diff = abs((items[j]['date'] - items[i]['date']).total_seconds())
                        if diff > 900:
                            break

                    # Só considera duplicata quando há um ID de Pix (E2E) confirmado e igual entre os dois,
                    # ou quando um dos lançamentos é claramente a versão "Google Pay" do outro (mesma
                    # transação relatada por dois provedores diferentes). Nunca remove por valor coincidente
                    # isoladamente, pois isso pode apagar dois pagamentos reais e distintos do mesmo valor.
                    if abs(items[i]['amount'] - items[j]['amount']) < 0.01:
                        same_pix_id = bool(items[i]['pix_id']) and items[i]['pix_id'] == items[j]['pix_id']

                        item_i_gpay = 'google pay' in items[i]['description'].lower()
                        item_j_gpay = 'google pay' in items[j]['description'].lower()
                        is_cross_provider_pair = item_i_gpay != item_j_gpay

                        if not (same_pix_id or is_cross_provider_pair):
                            continue

                        if item_j_gpay and not item_i_gpay:
                            to_delete.add(items[i]['id'])
                            removed_count += 1
                            break
                        else:
                            to_delete.add(items[j]['id'])
                            removed_count += 1

            for doc_id in to_delete:
                db.collection(col).document(doc_id).update({'status': 'deleted'})

            if removed_count > 0:
                log_to_firestore(sync_ref, logs, f"[{col}] Limpeza retroativa concluída: {removed_count} duplicata(s) marcada(s) como removida(s).")
        except Exception as e:
            log_to_firestore(sync_ref, logs, f"ERRO LIMPEZA RETROATIVA ({col}): {e}")


def _format_brl(amount) -> str:
    """Formata valores monetários sem depender do locale da Cloud Function."""
    try:
        formatted = f"{float(amount):,.2f}"
    except (TypeError, ValueError):
        return "valor não informado"
    return "R$ " + formatted.replace(",", "_").replace(".", ",").replace("_", ".")


def _gmail_message_headers(message: dict) -> tuple[str, str]:
    """Extrai remetente e assunto do payload retornado pela API do Gmail."""
    headers = (message.get('payload') or {}).get('headers') or []
    values = {
        str(header.get('name') or '').strip().lower(): str(header.get('value') or '').strip()
        for header in headers
    }
    return values.get('from', ''), values.get('subject', '')


def _build_imported_bills_notification(imported_bills: list[dict]) -> str:
    """Monta um resumo legível no Telegram para os boletos da sincronização."""
    count = len(imported_bills)
    header = f"{count} boleto{' foi' if count == 1 else 's foram'} importado{' ' if count == 1 else 's '}do Gmail."
    total = sum(
        float(bill.get('amount') or 0)
        for bill in imported_bills
        if str(bill.get('amount') or '').strip()
    )

    if count == 1:
        bill = imported_bills[0]
        lines = [header, '', f"📄 {bill['description']}", f"💰 Valor: {_format_brl(bill['amount'])}"]
        if bill.get('due_date'):
            lines.append(f"📅 Vencimento: {bill['due_date'].strftime('%d/%m/%Y')}")
        if bill.get('rubric'):
            lines.append(f"🏷️ Rubrica: {bill['rubric']}")
        if bill.get('sender'):
            lines.append(f"✉️ Remetente: {bill['sender']}")
        if bill.get('subject'):
            lines.append(f"📝 Assunto: {bill['subject']}")
        return "\n".join(lines)

    lines = [header, '']
    for index, bill in enumerate(imported_bills, start=1):
        due = bill['due_date'].strftime('%d/%m') if bill.get('due_date') else 'sem vencimento'
        rubric = f" · {bill['rubric']}" if bill.get('rubric') else ''
        lines.append(f"{index}. {bill['description']} — {_format_brl(bill['amount'])} — vence {due}{rubric}")
        if bill.get('sender'):
            lines.append(f"   ✉️ {bill['sender']}")
    lines.extend(['', f"💰 Total: {_format_brl(total)}"])
    return "\n".join(lines)


def sync_boletos_gmail(service, sync_ref, logs):
    """
    Explora o Gmail em busca de boletos, extraia dados via IA e salva no Firestore (fixed_bills).
    """
    db = get_db()
    
    log_to_firestore(sync_ref, logs, "Buscando boletos no Gmail via IA...")
    
    # Query para emails com anexos PDF ou assuntos de fatura/boleto/pagamento
    # Pegamos os mais recentes para a sincronia automática
    query = 'has:attachment filename:pdf (subject:(boleto OR fatura OR bill OR pagamento OR "o seu boleto" OR "sua fatura" OR "vencimento") OR "boleto" OR "fatura")'
    
    try:
        results = service.users().messages().list(userId='me', q=query, maxResults=15).execute(num_retries=3)
        messages = results.get('messages', [])
        
        if not messages:
            log_to_firestore(sync_ref, logs, "Nenhum boleto recente encontrado no Gmail.")
            return

        # Configurar Gemini
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        api_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        if not api_key:
            log_to_firestore(sync_ref, logs, "ERRO: Gemini API Key não encontrada (em system/api_keys).")
            return
        
        from google.genai import types

        genai = get_genai_module()
        client = genai.Client(api_key=api_key)

        processed_count = 0
        imported_bills = []
        
        # Cache de boletos existentes para permitir duplicatas ou vinculação
        existing_bills_cache = []
        existing_bill_google_ids = set()
        for b in db.collection('fixed_bills').stream():
            d = b.to_dict()
            existing_bills_cache.append({
                'id': b.id,
                'desc': d.get('description', '').lower(),
                'amount': d.get('amount'),
                'dueDay': d.get('dueDay'),
                'month': d.get('month'),
                'year': d.get('year'),
                'isPaid': d.get('isPaid', False),
                'rubricId': d.get('rubricId')
            })
            if d.get('google_message_id'):
                existing_bill_google_ids.add(d['google_message_id'])

        # Cache de rubricas para vinculação automática
        rubrics_cache = []
        for r in db.collection('bill_rubrics').stream():
            d = r.to_dict()
            rubrics_cache.append({
                'id': r.id,
                'desc': d.get('description', '').lower(),
                'label': d.get('description', '').strip()
            })

        processed_emails_doc = db.collection('system').document('processed_emails').get()
        processed_ids = processed_emails_doc.to_dict().get('ids', []) if processed_emails_doc.exists else []
        new_processed_ids = []

        for m_info in messages:
            msg_id = m_info['id']
            if msg_id in existing_bill_google_ids:
                archive_gmail_message(service, msg_id, sync_ref, logs, "boleto-ja-lancado")
                continue
            if msg_id in processed_ids: continue
            
            msg = service.users().messages().get(userId='me', id=msg_id).execute(num_retries=3)
            snippet = msg.get('snippet', '')
            sender, subject = _gmail_message_headers(msg)

            # Tentar baixar o primeiro PDF encontrado
            pdf_data = None
            def find_pdf(part):
                nonlocal pdf_data
                if part.get('parts'):
                    for sub in part['parts']: find_pdf(sub)
                if part.get('filename', '').lower().endswith('.pdf') and part.get('body', {}).get('attachmentId'):
                    att = service.users().messages().attachments().get(userId='me', messageId=msg_id, id=part['body']['attachmentId']).execute(num_retries=3)
                    pdf_data = base64.urlsafe_b64decode(att['data'])
                    return
            
            find_pdf(msg['payload'])
            
            # Formata rubricas para o prompt
            rubrics_text = "\n".join([f"- {r['desc']} (ID: {r['id']})" for r in rubrics_cache])
            
            prompt = f"""
            Você é um assistente financeiro de elite. Analise o e-mail/documento anexo e extraia os dados abaixo para um BOLETO ou PAGAMENTO.
            
            Além disso, compare esta conta com a seguinte lista de RUBRICAS RECORRENTES do usuário:
            {rubrics_text}
            
            Se o boleto corresponder a uma dessas rubricas (mesmo que o nome não seja idêntico, ex: "EDP ENERGIA" corresponde a "EDP (energia)"), informe o ID da rubrica.
            
            Campos obrigatórios no JSON:
            - description: Nome curto da conta (ex: VIVO, Sabesp, Condomínio)
            - amount: valor numérico do boleto
            - due_date: data de vencimento (formato YYYY-MM-DD)
            - barcode: linha digitável ou código de barras (apenas números)
            - pix_code: código Pix Copia e Cola (geralmente começa com 000201...)
            - rubric_id: ID da rubrica correspondente (se houver match) ou null

            Responda APENAS em JSON no formato:
            {{
              "description": "...",
              "amount": 123.45,
              "due_date": "YYYY-MM-DD",
              "barcode": "...",
              "pix_code": "...",
              "rubric_id": "..."
            }}
            Se não for um boleto/fatura ou se não encontrar dados, responda {{"error": "not_a_bill"}}.
            """
            
            content_parts = [prompt, f"E-mail Fragment: {snippet}"]
            if pdf_data:
                content_parts.append(types.Part.from_bytes(data=pdf_data, mime_type="application/pdf"))
            
            try:
                response = client.models.generate_content(model="gemini-3.5-flash-lite", contents=content_parts)
                res_text = response.text.strip()
                if "```json" in res_text:
                    res_text = res_text.split("```json")[-1].split("```")[0].strip()
                elif "```" in res_text:
                    res_text = res_text.split("```")[-1].split("```")[0].strip()

                data = json.loads(res_text)
                if data.get('error'): 
                    new_processed_ids.append(msg_id)
                    continue
                
                due_dt = datetime.fromisoformat(data['due_date'])
                month = due_dt.month - 1
                year = due_dt.year
                
                found_existing_id = None
                is_exact_dup = False
                name_extracted = data['description'].lower()
                # Prioriza o match de rubrica feito pela IA
                rubric_id_from_ai = data.get('rubric_id')

                for eb in existing_bills_cache:
                    if eb['month'] == month and eb['year'] == year:
                        # 1. Tenta match por rubricId (IA indicou este card)
                        if rubric_id_from_ai and eb.get('rubricId') == rubric_id_from_ai:
                            found_existing_id = eb['id']
                            found_existing_rubric_id = eb.get('rubricId')
                            # Se valor for igual, é exatamente o mesmo registro
                            if abs(eb['amount'] - data['amount']) < 0.01 and eb.get('dueDay') == due_dt.day:
                                is_exact_dup = True
                            break

                        # 2. Lógica de vinculação via nome (Fallback/Ambiguidade)
                        if name_extracted in eb['desc'] or eb['desc'] in name_extracted:
                            # Se o valor também for igual, é uma duplicata exata
                            if abs(eb['amount'] - data['amount']) < 0.01 and eb.get('dueDay') == due_dt.day:
                                is_exact_dup = True
                                break
                            # Caso contrário, vamos vincular a este card (atualizá-lo)
                            found_existing_id = eb['id']
                            found_existing_rubric_id = eb.get('rubricId')
                            break 

                if is_exact_dup:
                    new_processed_ids.append(msg_id)
                    archive_gmail_message(service, msg_id, sync_ref, logs, "boleto-duplicado")
                    continue
                
                if found_existing_id:
                    # VINCULAÇÃO: Atualiza card existente
                    update_data = {
                        'amount': data['amount'],
                        'dueDay': due_dt.day,
                        'barcode': data.get('barcode', ''),
                        'pixCode': data.get('pix_code', ''),
                        'google_message_id': msg_id,
                        'updated_at': datetime.now().isoformat()
                    }
                    
                    # Se o card existente não tiver rubricId, usa o da IA ou tenta achar
                    if not found_existing_rubric_id:
                        update_data['rubricId'] = rubric_id_from_ai
                        if not update_data['rubricId']:
                            for rb in rubrics_cache:
                                if name_extracted in rb['desc'] or rb['desc'] in name_extracted:
                                    update_data['rubricId'] = rb['id']
                                    break
                                
                    db.collection('fixed_bills').document(found_existing_id).update(update_data)
                    log_to_firestore(sync_ref, logs, f"[BOLETO] Vinculado ao card '{data['description']}': R$ {data['amount']}")
                    processed_count += 1
                    effective_rubric_id = update_data.get('rubricId') or found_existing_rubric_id
                    archive_gmail_message(service, msg_id, sync_ref, logs, "boleto-vinculado")
                else:
                    # BUSCA EM RUBRICAS (Fallback se a IA não retornou rubric_id)
                    matched_rubric_id = rubric_id_from_ai
                    
                    if not matched_rubric_id:
                        for rb in rubrics_cache:
                            if name_extracted in rb['desc'] or rb['desc'] in name_extracted:
                                matched_rubric_id = rb['id']
                                break

                    # CRIAÇÃO: Adiciona novo card
                    db.collection('fixed_bills').add({
                        'description': data['description'],
                        'amount': data['amount'],
                        'dueDay': due_dt.day,
                        'month': month,
                        'year': year,
                        'barcode': data.get('barcode', ''),
                        'pixCode': data.get('pix_code', ''),
                        'isPaid': False,
                        'category': 'Geral',
                        'rubricId': matched_rubric_id,
                        'google_message_id': msg_id,
                        'created_at': datetime.now().isoformat()
                    })
                    log_to_firestore(sync_ref, logs, f"[BOLETO] Importado (Novo Card): {data['description']} (R$ {data['amount']})")
                    processed_count += 1
                    effective_rubric_id = matched_rubric_id
                    archive_gmail_message(service, msg_id, sync_ref, logs, "boleto-lancado")

                rubric_label = next(
                    (rubric['label'] for rubric in rubrics_cache if rubric['id'] == effective_rubric_id),
                    ''
                )
                imported_bills.append({
                    'description': str(data.get('description') or 'Conta sem identificação').strip(),
                    'amount': data.get('amount'),
                    'due_date': due_dt,
                    'rubric': rubric_label,
                    'sender': sender,
                    'subject': subject,
                })
                
                new_processed_ids.append(msg_id)

            except Exception as e:
                error_msg = str(e)
                if "The document has no pages" in error_msg:
                    log_to_firestore(sync_ref, logs, f"Aviso: O PDF da mensagem {msg_id} está vazio. Ignorando.")
                else:
                    log_to_firestore(sync_ref, logs, f"Aviso: Erro ao processar mensagem {msg_id}: {e}")
                new_processed_ids.append(msg_id)

        if new_processed_ids:
            updated_ids = list(set(processed_ids + new_processed_ids))[-500:]
            db.collection('system').document('processed_emails').set({'ids': updated_ids}, merge=True)

        if processed_count > 0:
            log_to_firestore(sync_ref, logs, f"Sincronização de boletos concluída. {processed_count} novos boletos.")
            emit_notification_backend(
                "Novos Boletos",
                _build_imported_bills_notification(imported_bills),
                "success",
                "financeiro"
            )
    
    except Exception as e:
        log_to_firestore(sync_ref, logs, f"ERRO na busca de boletos: {e}")


@https_fn.on_call(memory=options.MemoryOption.GB_1, timeout_sec=540)
def sync_gmail_bills_callable(req: https_fn.CallableRequest):
    """Executa a sincronização de boletos do Gmail manualmente via app"""
    db = get_db()
    sync_ref = db.collection('system').document('sync')
    logs = [f"[{datetime.now().strftime('%H:%M:%S')}] Iniciando sincronização manual via App..."]
    
    try:
        gs = get_gmail_service()
        sync_boletos_gmail(gs, sync_ref, logs)
        
        sync_ref.update({
            'status': 'completed',
            'last_success': datetime.now().isoformat(),
            'logs': logs
        })
        return {"success": True}
    except Exception as e:
        error_msg = f"Erro na sincronização manual: {str(e)}"
        log_to_firestore(sync_ref, logs, error_msg)
        return {"success": False, "error": error_msg}

def run_full_sync(trigger_reason='unspecified'):
    """Executa o processo completo de sincronização"""
    db = get_db()
    sync_ref = db.collection('system').document('sync')
    run_id = uuid.uuid4().hex
    logs = [f"Iniciando sincronização ({datetime.now().strftime('%Y-%m-%d %H:%M:%S')})... Trigger: {trigger_reason}"]

    if not acquire_sync_lock(db, run_id):
        queue_sync_request(db, f"sync-busy:{trigger_reason}")
        print(f"Sincronização já em andamento. Pedido enfileirado: {trigger_reason}")
        return False

    try:
        sync_ref.set({
            'status': 'processing',
            'active_run_id': run_id,
            'pending_request': False,
            'last_trigger': trigger_reason,
            'started_at': datetime.now(timezone.utc).isoformat(),
            'logs': logs
        }, merge=True)

        for current_pass in range(1, MAX_SYNC_PASSES + 1):
            if current_pass > 1:
                log_to_firestore(sync_ref, logs, f"[SYNC] Reexecutando sincronização para consolidar alterações pendentes (passo {current_pass}/{MAX_SYNC_PASSES}).", True)
                sync_ref.set({
                    'status': 'processing',
                    'active_run_id': run_id,
                    'pending_request': False,
                    'logs': logs
                }, merge=True)

            ts, gs, cs = get_tasks_service(), get_gmail_service(), get_calendar_service()
            # Primeiro puxa o Calendar para permitir sincronia inversa (agenda -> Hermes) antes do push
            sync_google_calendar(cs, sync_ref, logs)
            sync_google_tasks_push(ts, cs, sync_ref, logs)
            sync_google_tasks_pull(ts, sync_ref, logs)

            sync_pix_emails(gs, sync_ref, logs)
            
            # Sincronização de Contatos do Google People API
            sync_google_contacts_internal(db, sync_ref, logs)
            
            # Ingestão de Documentos (Acervo Global)
            log_to_firestore(sync_ref, logs, "[SYNC] Verificando novos documentos na Pasta de Deságue (Acervo Global)...", True)
            executar_monitoramento_acervo_global()


            sync_boletos_gmail(gs, sync_ref, logs)

            # Vínculo automático de e-mails a ações em andamento/stand-by (via IA + confirmação Telegram).
            # Protegido por try/except próprio: uma falha aqui nunca deve derrubar o sync financeiro/agenda.
            try:
                from email_action_linker import link_emails_to_actions
                link_emails_to_actions(db, gs, sync_ref, logs)
            except Exception as e_link:
                log_to_firestore(sync_ref, logs, f"[EMAIL-LINK][ERRO] Falha inesperada no vínculo e-mail-ação: {e_link}", True)

            # Vínculo automático de reuniões encerradas a ações (matching determinístico por
            # google_calendar_id, sem IA). Mesma proteção: nunca derruba o restante do sync.
            try:
                from email_action_linker import link_calendar_events_to_actions
                link_calendar_events_to_actions(db, sync_ref, logs)
            except Exception as e_cal_link:
                log_to_firestore(sync_ref, logs, f"[CAL-LINK][ERRO] Falha inesperada no vínculo calendar-ação: {e_cal_link}", True)

            # Triagem de conversas de WhatsApp capturadas por services/whatsapp-capture
            # (propõe vínculo com ações + grava digests vetorizados). Desligada por padrão
            # e sem efeito nenhum enquanto o worker local não estiver rodando/configurado.
            try:
                from whatsapp_ingest import triage_whatsapp_messages
                triage_whatsapp_messages(db, sync_ref, logs)
            except Exception as e_wa_ingest:
                log_to_firestore(sync_ref, logs, f"[WA-INGEST][ERRO] Falha inesperada na triagem de WhatsApp: {e_wa_ingest}", True)

            sync_state = sync_ref.get().to_dict() or {}
            if not sync_state.get('pending_request'):
                break
            if current_pass == MAX_SYNC_PASSES:
                log_to_firestore(sync_ref, logs, "[SYNC][!] Limite de reexecuções atingido; alterações restantes serão processadas na próxima sincronização.", True)

        sync_ref.set({
            'status': 'completed',
            'last_success': datetime.now().isoformat(),
            'pending_request': False,
            'active_run_id': None,
            'logs': logs
        }, merge=True)

        print("Sincronização concluída com sucesso.")
        return True

    except Exception as e:

        if isinstance(e, GoogleAuthRevokedError) or is_google_invalid_grant_error(e):
            error_msg = f"ERRO GOOGLE AUTH: {GOOGLE_REAUTH_MESSAGE}"
            try:
                db.collection('system').document('google_credentials').set({
                    'auth_status': 'reauth_required',
                    'auth_error': 'invalid_grant',
                    'auth_error_message': GOOGLE_REAUTH_MESSAGE,
                    'updated_at': firestore.SERVER_TIMESTAMP
                }, merge=True)
            except Exception as auth_status_err:
                print(f"Falha ao marcar reautenticacao Google: {auth_status_err}")
        else:
            error_msg = f"ERRO na sincronização: {str(e)}"

        print(error_msg)

        sync_ref.set({
            'status': 'error',
            'error_message': error_msg,
            'pending_request': False,
            'active_run_id': None,
            'logs': logs + [error_msg]
        }, merge=True)
        return False
    finally:
        release_sync_lock(db, run_id)



@firestore_fn.on_document_updated(document="system/sync", timeout_sec=540, memory=options.MemoryOption.GB_1)

def on_sync_request(event: firestore_fn.Event[firestore_fn.Change[firestore_fn.DocumentSnapshot]]):

    """Trigger disparado quando system/sync é atualizado manualmente"""

    if not event.data.after.exists: return

    data = event.data.after.to_dict()

    if data.get('status') != 'requested': return

    run_full_sync('firestore-request')



@scheduler_fn.on_schedule(schedule="every 30 minutes", timeout_sec=540, memory=options.MemoryOption.GB_1)

def scheduled_sync(event: scheduler_fn.ScheduledEvent) -> None:

    """Trigger agendado para rodar a cada 30 minutos"""

    run_full_sync('scheduled')


def _page_monitor_html_to_text(html_content: str, seletor_css: str | None = None) -> str:
    """Extrai texto visivel de uma pagina HTML para comparacao por hash."""
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html_content or "", "html.parser")
        for tag in soup(["script", "style", "noscript", "svg"]):
            tag.decompose()
        root = soup
        if seletor_css:
            selected = soup.select_one(seletor_css)
            if selected is not None:
                root = selected
        text = root.get_text(separator="\n")
    except Exception as exc:
        print(f"[PageMonitor] Fallback regex para extrair texto: {exc}")
        cleaned = re.sub(r"(?is)<(script|style|noscript)[^>]*>.*?</\1>", " ", html_content or "")
        cleaned = re.sub(r"(?s)<[^>]+>", " ", cleaned)
        text = cleaned

    lines = [line.strip() for line in str(text or "").splitlines()]
    return "\n".join(line for line in lines if line)


def _page_monitor_fetch_text(url: str, seletor_css: str | None = None, timeout: int = 30) -> str:
    import requests

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/124.0 Safari/537.36"
        ),
        "Accept-Language": "pt-BR,pt;q=0.9,en;q=0.6",
    }
    response = requests.get(url, headers=headers, timeout=timeout)
    response.raise_for_status()
    return _page_monitor_html_to_text(response.text, seletor_css)


def _page_monitor_analyze_change(db, objetivo: str, texto_antigo: str | None, texto_novo: str | None) -> dict | None:
    gemini_key = get_gemini_api_key()
    if not gemini_key:
        return None

    try:
        genai = get_genai_module()
        client = genai.Client(api_key=gemini_key)
        prompt = (
            "Voce compara duas versoes do texto visivel de uma pagina da web e avalia uma mudanca "
            "em relacao a um OBJETIVO.\n"
            f"OBJETIVO DO USUARIO: {objetivo}\n\n"
            "--- VERSAO ANTERIOR ---\n"
            f"{(texto_antigo or '')[:8000]}\n\n"
            "--- VERSAO ATUAL ---\n"
            f"{(texto_novo or '')[:8000]}\n\n"
            "Tarefa: (1) descreva objetivamente o que mudou de relevante, ignorando rodape, banners, "
            "datas de acesso, menus e contadores triviais. (2) decida se a mudanca AVANCA ou ATENDE "
            "o objetivo do usuario.\n"
            "Responda APENAS um JSON: "
            "{\"avanca_objetivo\": true|false, \"resumo\": \"explicacao objetiva em 1-3 frases\"}"
        )
        response = generate_content_logged(
            client,
            model=GEMINI_LIGHT_MODEL,
            contents=prompt,
            feature="page_monitor.change_analysis",
            db=db,
        )
        raw = (getattr(response, "text", None) or "").strip()
        match = re.search(r"\{.*\}", raw, re.S)
        if not match:
            return {"avanca_objetivo": False, "resumo": raw[:500] or "Sem resposta do modelo."}
        parsed = json.loads(match.group(0))
        return {
            "avanca_objetivo": bool(parsed.get("avanca_objetivo")),
            "resumo": str(parsed.get("resumo") or "").strip()[:1500],
        }
    except Exception as exc:
        print(f"[PageMonitor] Falha na analise Gemini: {exc}")
        return None


def _page_monitor_build_message(apelido: str, objetivo: str, url: str, resumo: str | None, fallback: bool = False) -> str:
    if fallback:
        return (
            f"Monitor de Paginas: {apelido}\n\n"
            "A pagina mudou, mas nao foi possivel avaliar o objetivo automaticamente.\n\n"
            f"Objetivo: {objetivo}\n"
            f"{url}"
        )[:4090]

    return (
        f"Monitor de Paginas: {apelido}\n\n"
        f"{resumo or 'Mudanca relevante detectada.'}\n\n"
        f"Objetivo: {objetivo}\n"
        f"{url}"
    )[:4090]


@scheduler_fn.on_schedule(
    schedule="every 30 minutes",
    timeout_sec=540,
    memory=options.MemoryOption.GB_1,
)
def scheduled_page_monitor(event: scheduler_fn.ScheduledEvent) -> None:
    """Verifica paginas monitoradas e envia alerta por Telegram quando o objetivo avancar."""
    import hashlib

    db = get_db()
    try:
        docs = list(db.collection("paginas_monitoradas").where("ativo", "==", True).stream())
    except Exception as exc:
        print(f"[PageMonitor] Falha ao listar paginas monitoradas: {exc}")
        return

    print(f"[PageMonitor] Verificando {len(docs)} pagina(s) ativa(s).")
    for doc_snap in docs:
        data = doc_snap.to_dict() or {}
        url = str(data.get("url") or "").strip()
        apelido = str(data.get("apelido") or url or doc_snap.id).strip()
        objetivo = str(data.get("objetivo") or "").strip()
        if not url:
            continue

        now_iso = datetime.now(timezone.utc).isoformat()
        update_payload = {"ultima_verificacao": now_iso}
        try:
            texto_novo = _page_monitor_fetch_text(url, data.get("seletor_css"))
            novo_hash = hashlib.sha256(texto_novo.encode("utf-8")).hexdigest()
        except Exception as exc:
            print(f"[PageMonitor] Falha ao buscar '{apelido}': {exc}")
            update_payload["ultima_falha"] = now_iso
            update_payload["erro_ultima_verificacao"] = str(exc)[:500]
            try:
                doc_snap.reference.update(update_payload)
            except Exception as update_exc:
                print(f"[PageMonitor] Falha ao registrar erro de '{apelido}': {update_exc}")
            continue

        hash_antigo = data.get("hash_atual")
        if not hash_antigo:
            update_payload.update({
                "hash_atual": novo_hash,
                "texto_atual": texto_novo[:20000],
                "erro_ultima_verificacao": None,
            })
            doc_snap.reference.update(update_payload)
            print(f"[PageMonitor] Baseline capturado para '{apelido}'.")
            continue

        if novo_hash == hash_antigo:
            update_payload["erro_ultima_verificacao"] = None
            doc_snap.reference.update(update_payload)
            continue

        print(f"[PageMonitor] Mudanca detectada em '{apelido}'.")
        update_payload.update({
            "hash_atual": novo_hash,
            "texto_atual": texto_novo[:20000],
            "ultima_mudanca": now_iso,
            "erro_ultima_verificacao": None,
        })

        analise = _page_monitor_analyze_change(db, objetivo, data.get("texto_atual", ""), texto_novo)
        chat_id = _resolve_telegram_chat_id_for_uid(db, data.get("userId")) or _resolve_default_telegram_chat_id(db)

        if analise is None:
            update_payload["ultima_analise"] = (
                "Mudanca detectada - analise indisponivel; alerta enviado como fallback."
            )
            sent = _send_telegram_message_raw(
                db,
                chat_id,
                _page_monitor_build_message(apelido, objetivo, url, None, fallback=True),
            )
            update_payload["ultimo_alerta_telegram"] = now_iso if sent else None
            update_payload["erro_telegram"] = None if sent else "send_failed_or_chat_id_missing"
            print(f"[PageMonitor] Alerta fallback para '{apelido}': sent={sent}.")
        else:
            update_payload["ultima_analise"] = analise["resumo"]
            if analise["avanca_objetivo"]:
                task_id = str(data.get("task_id") or "").strip()
                linked = False
                if task_id:
                    try:
                        from email_action_linker import _load_candidate_task_by_id, queue_and_maybe_send_suggestion
                        task = _load_candidate_task_by_id(db, task_id)
                        if task:
                            result = queue_and_maybe_send_suggestion(
                                db,
                                f"pagina_{doc_snap.id}_{novo_hash[:16]}",
                                canal="pagina",
                                task=task,
                                titulo_sinal=apelido,
                                origem_sinal=url,
                                resumo=analise["resumo"],
                                nota_sugerida=analise["resumo"],
                                reativar_sugerido=True,
                                chat_id=chat_id,
                                extra={"link_externo": url},
                            )
                            linked = bool(result and result.get("telegram_sent"))
                    except Exception as exc_link:
                        print(f"[PageMonitor] Falha ao vincular '{apelido}' à ação {task_id}: {exc_link}")

                if linked:
                    update_payload["ultimo_alerta_telegram"] = now_iso
                    update_payload["erro_telegram"] = None
                    print(f"[PageMonitor] Alerta de '{apelido}' vinculado à ação {task_id}.")
                else:
                    sent = _send_telegram_message_raw(
                        db,
                        chat_id,
                        _page_monitor_build_message(apelido, objetivo, url, analise["resumo"]),
                    )
                    update_payload["ultimo_alerta_telegram"] = now_iso if sent else None
                    update_payload["erro_telegram"] = None if sent else "send_failed_or_chat_id_missing"
                    print(f"[PageMonitor] Alerta para '{apelido}': sent={sent}.")
            else:
                print(f"[PageMonitor] Mudanca em '{apelido}' nao avancou o objetivo; sem alerta.")

        doc_snap.reference.update(update_payload)

    print("[PageMonitor] Verificacao concluida.")


def _normalize_notification_title(title: str | None) -> str:
    normalized = unicodedata.normalize("NFKD", str(title or ""))
    normalized = "".join(ch for ch in normalized if not unicodedata.combining(ch))
    return normalized.strip().lower()


def _should_mirror_notification_to_telegram(notif: dict) -> bool:
    """Direciona todas as notificações do sistema para o Telegram."""
    return True


def _build_telegram_notification_message(notif: dict) -> str:
    title = str(notif.get('title') or 'Hermes').strip()
    message = str(notif.get('message') or '').strip()
    n_type = str(notif.get('type') or 'info').strip()
    link = str(notif.get('link') or '').strip()

    icons = {
        'success': '✅',
        'warning': '⚠️',
        'error': '🚨',
        'info': '🔔',
        'expense': '💸',
    }
    lines = [f"{icons.get(n_type, '🔔')} Hermes - {title}"]
    assunto = str(notif.get('assunto') or '').strip()
    if assunto:
        lines.extend(["", f"📝 <b>Assunto Detalhado:</b> {assunto}"])
    if message:
        lines.extend(["", message])
    # Links internos (ex: "financeiro") só servem para navegação dentro do app e não são
    # clicáveis no Telegram — exibimos apenas links http(s) reais, que fazem sentido ali.
    if link and link.startswith(("http://", "https://")):
        lines.extend(["", f"Link: {link}"])
    return "\n".join(lines)

@firestore_fn.on_document_created(document="notificacoes/{notification_id}", memory=options.MemoryOption.MB_512)

def on_notificacao_created(event: firestore_fn.Event[firestore_fn.DocumentSnapshot | None]):

    """Trigger disparado quando uma nova notificação é criada"""

    if not event.data: return

    notif = event.data.to_dict()

    if not notif: return

    title = notif.get('title', 'Hermes')

    message = notif.get('message', '')

    db = get_db()

    updates = {}

    # Notificações do scraper SIPAC (functions_node/index.js) tentam primeiro o vínculo
    # determinístico com uma ação por número de processo; se conseguir enviar o cartão de
    # confirmação, pula o espelhamento genérico abaixo para não duplicar o aviso.
    if not notif.get('sent_to_telegram') and notif.get('link') == '@SipacTrackingTool':
        try:
            from email_action_linker import try_link_sipac_notification
            if try_link_sipac_notification(db, event.data.reference.id, notif):
                event.data.reference.update({'sent_to_telegram': True, 'linked_to_acao': True})
                return
        except Exception as exc:
            print(f"[SIPAC-LINK] Falha ao tentar vincular notificação SIPAC a ação: {exc}")

    if not notif.get('sent_to_telegram') and _should_mirror_notification_to_telegram(notif):

        telegram_chat_id = _resolve_default_telegram_chat_id(db)

        if telegram_chat_id:

            sent = _send_telegram_message_raw(db, telegram_chat_id, _build_telegram_notification_message(notif))

            updates['sent_to_telegram'] = bool(sent)

            if not sent:

                updates['telegram_error'] = 'send_failed'

        else:

            updates['sent_to_telegram'] = False

            updates['telegram_error'] = 'chat_id_not_configured'

    if updates:

        event.data.reference.update(updates)



@scheduler_fn.on_schedule(
    schedule="every 1 minutes",
    memory=options.MemoryOption.GB_1,
    timeout_sec=120,
)

def check_and_send_reminders(event: scheduler_fn.ScheduledEvent) -> None:

    """Verifica e envia lembretes agendados (hábitos, pesagem, customizados e ações)"""

    from datetime import datetime, timedelta

    import pytz

    

    db = get_db()

    # Define o fuso horário de Brasília para comparar com as strings de horário do usuário (HH:mm)

    tz = pytz.timezone('America/Sao_Paulo')

    now = datetime.now(tz)

    current_time_str = now.strftime('%H:%M')

    today_str = now.strftime('%Y-%m-%d')

    day_of_week = now.weekday() # 0 = Monday, 1 = Tuesday... 6 = Sunday (Note: Python index matches our dayOfWeek if 0=Mon, but let's check)

    # No helper.tsx: dayOfWeek: 1 // Segunda-feira. Python: 0=Mon, 1=Tue... 

    # Precisamos ajustar para 0=Dom? Não, vamos seguir o padrão do AppSettings.

    # AppSettings weighInReminder dayOfWeek: 0-6 (0=Dom no JS Date.getDay())

    # Python now.strftime('%w') retorna 0 para Domingo.

    js_day_of_week = int(now.strftime('%w'))



    # 1. Carrega Configurações

    settings_doc = db.collection('configuracoes').document('geral').get()

    if settings_doc.exists:

        settings = settings_doc.to_dict()

        notifs_config = settings.get('notifications', {})

        

        # --- Lembrete de Pesagem ---

        weigh_in = notifs_config.get('weighInReminder', {})

        if weigh_in.get('enabled') and weigh_in.get('time') == current_time_str:

            freq = weigh_in.get('frequency', 'weekly')

            target_day = weigh_in.get('dayOfWeek', 1)

            

            should_remind = False

            if js_day_of_week == target_day:

                if freq == 'weekly':

                    should_remind = True

                elif freq == 'biweekly':

                    # Lógica simplificada de biweekly baseada no timestamp da semana

                    week_num = int(now.strftime('%V'))

                    if week_num % 2 == 0: should_remind = True

                elif freq == 'monthly' and now.day == 1:

                    should_remind = True

            

            if should_remind:

                remind_id = f"weighin_{today_str}"

                if not db.collection('system_reminders').document(remind_id).get().exists:

                    emit_notification_backend(

                        "Lembrete de Pesagem",

                        "Hora de registrar seu peso para acompanhar sua evolução no módulo Saúde!",

                        'info',

                        'saude'

                    )

                    db.collection('system_reminders').document(remind_id).set({'sent_at': now.isoformat()})



        # --- Notificações Customizadas ---

        custom_notifs = notifs_config.get('custom', [])

        for cn in custom_notifs:

            if cn.get('enabled') and cn.get('time') == current_time_str:

                freq = cn.get('frequency', 'daily')

                should_send = False

                

                if freq == 'daily':

                    should_send = True

                elif freq == 'weekly' and js_day_of_week in cn.get('daysOfWeek', []):

                    should_send = True

                elif freq == 'monthly' and now.day == cn.get('dayOfMonth', 1):

                    should_send = True

                

                if should_send:

                    remind_id = f"custom_{cn.get('id')}_{today_str}"

                    if not db.collection('system_reminders').document(remind_id).get().exists:

                        emit_notification_backend(

                            "Lembrete Personalizado",

                            cn.get('message', 'Notificação Hermes'),

                            'info'

                        )

                        db.collection('system_reminders').document(remind_id).set({'sent_at': now.isoformat()})



    # 2. Lembretes de saude enviados somente pelo Telegram.
    # Mantenha esta lista em sincronia com DEFAULT_HEALTH_REMINDERS em HealthView.tsx
    # (mesmos id/title/message/time/daysOfWeek/category) — nao ha fonte unica entre
    # Python e TypeScript, entao qualquer mudanca aqui precisa ser replicada la.
    default_health_reminders = [
        {
            "id": "lunch_slow",
            "title": "Almoço com calma",
            "message": "André, lembre de comer devagar no almoço. Ritmo baixo também é estratégia.",
            "time": "11:45",
            "enabled": True,
            "daysOfWeek": [0, 1, 2, 3, 4, 5, 6],
            "category": "nutrition",
        },
        {
            "id": "food_window",
            "title": "Janela alimentar",
            "message": "André, última janela alimentar chegando. Se for comer, mantenha leve.",
            "time": "17:30",
            "enabled": True,
            "daysOfWeek": [0, 1, 2, 3, 4, 5, 6],
            "category": "nutrition",
        },
        {
            "id": "pain_checkin",
            "title": "Check-in lombar",
            "message": "André, check-in rápido: como ficou sua lombar hoje?",
            "time": "21:30",
            "enabled": True,
            "daysOfWeek": [0, 1, 2, 3, 4, 5, 6],
            "category": "pain",
        },
        {
            "id": "strength_training",
            "title": "Treino de força",
            "message": "André, hoje é dia de treino de força (bloco A ou B).",
            "time": "07:00",
            "enabled": True,
            "daysOfWeek": [1, 3, 5],
            "category": "spine",
        },
        {
            "id": "daily_weighin",
            "title": "Pesagem diária",
            "message": "André, pese-se ao acordar, antes do café.",
            "time": "06:30",
            "enabled": True,
            "daysOfWeek": [0, 1, 2, 3, 4, 5, 6],
            "category": "custom",
        },
        {
            "id": "waist_saturday",
            "title": "Cintura da semana",
            "message": "André, meça a circunferência de cintura na altura do umbigo.",
            "time": "07:00",
            "enabled": True,
            "daysOfWeek": [6],
            "category": "custom",
        },
        {
            "id": "batch_cooking_sunday",
            "title": "Batch cooking",
            "message": "André, hora de preparar as refeições da semana.",
            "time": "10:00",
            "enabled": True,
            "daysOfWeek": [0],
            "category": "nutrition",
        },
        {
            "id": "fexofenadina_reminder",
            "title": "Fexofenadina",
            "message": "André, tome a fexofenadina com água — longe de suco e de antiácido, que reduzem a absorção em 30–40%.",
            "time": "08:00",
            "enabled": True,
            "daysOfWeek": [0, 1, 2, 3, 4, 5, 6],
            "category": "custom",
        },
    ]

    health_reminders_by_id = {item["id"]: item for item in default_health_reminders}
    try:
        for reminder_doc in db.collection("health_telegram_reminders").stream():
            data = reminder_doc.to_dict() or {}
            data["id"] = data.get("id") or reminder_doc.id
            health_reminders_by_id[reminder_doc.id] = data
    except Exception as exc:
        print(f"[HealthReminders] Falha ao carregar lembretes: {exc}")

    for reminder in health_reminders_by_id.values():
        reminder_id = str(reminder.get("id") or "").strip()
        if not reminder_id or not reminder.get("enabled", True):
            continue
        if reminder.get("time") != current_time_str:
            continue
        days = reminder.get("daysOfWeek")
        if isinstance(days, list) and days and js_day_of_week not in days:
            continue
        sent_id = f"health_telegram_{reminder_id}_{today_str}"
        if db.collection("system_reminders").document(sent_id).get().exists:
            continue

        owner_uid = reminder.get("created_by_uid")
        telegram_chat_id = _resolve_telegram_chat_id_for_uid(db, owner_uid) or _resolve_default_telegram_chat_id(db)
        if telegram_chat_id:
            title = (reminder.get("title") or "Lembrete de saúde").strip()
            message = (reminder.get("message") or "André, lembrete de saúde configurado no Hermes.").strip()
            keyboard = None
            if reminder.get("category") == "pain":
                keyboard = [
                    [{"text": str(n), "callback_data": f"health_pain:{today_str}:{n}"} for n in range(0, 6)],
                    [{"text": str(n), "callback_data": f"health_pain:{today_str}:{n}"} for n in range(6, 11)],
                ]
            sent = _send_telegram_message_raw_with_keyboard(db, telegram_chat_id, f"{title}\n\n{message}", keyboard)
            db.collection("system_reminders").document(sent_id).set({
                "sent_at": now.isoformat(),
                "sent": bool(sent),
                "type": "health_telegram",
                "reminder_id": reminder_id,
            })
        else:
            print(f"[HealthReminders] Nenhum chat_id encontrado para {reminder_id}")


    # 3. Lembretes de Ações (Specific Task Reminders)

    from google.cloud.firestore import Query

    # Busca tarefas com reminder_at definido e que ainda não foram marcadas como lembradas

    now_iso = now.strftime('%Y-%m-%dT%H:%M:%S')

    tasks_with_reminders = db.collection('tarefas')\
            .where('reminder_sent', '==', False)\
            .where('reminder_at', '<=', now_iso)\
            .stream()



    for task_doc in tasks_with_reminders:

        t = task_doc.to_dict()

        title = t.get('titulo', 'Ação Pendente')
        task_reminders = _normalize_task_reminders(t)
        due_reminder = next((reminder for reminder in task_reminders if not reminder.get('reminder_sent')), None)
        reminder_iso = due_reminder.get('reminder_at') if due_reminder else t.get('reminder_at')

        emit_notification_backend(

            f"Lembrete: {title}",

            "Está na hora de realizar esta ação agendada!",

            'warning',

            'acoes'

        )

        

        owner_uid = t.get('created_by_uid')
        telegram_chat_id = _resolve_telegram_chat_id_for_uid(db, owner_uid) or _resolve_default_telegram_chat_id(db)
        telegram_message = _build_task_reminder_telegram_message(t, reminder_iso)
        if telegram_chat_id:
            _send_telegram_message_raw(db, telegram_chat_id, telegram_message)
        else:
            print(f"[Telegram] Nenhum chat_id encontrado para lembrete da tarefa {task_doc.id}")

        # Criar lembrete correspondente no Google Tasks

        try:
            ts = get_tasks_service()
            results = ts.tasklists().list().execute()
            tasklist_id = next((item['id'] for item in results.get('items', []) if 'tarefa' in item['title'].lower()), None)
            if not tasklist_id:
                default_list = ts.tasklists().get(tasklist='@default').execute()
                tasklist_id = default_list.get('id')
            
            if tasklist_id:
                reminder_text = due_reminder.get('message') if (due_reminder and due_reminder.get('message')) else t.get('reminder_message')
                reminder_body = {
                    'title': f"Lembrete: {title}",
                    'notes': reminder_text or "Está na hora de realizar esta ação agendada!",
                }
                if reminder_iso:
                    date_part = reminder_iso.split('T')[0]
                    reminder_body['due'] = f"{date_part}T00:00:00.000Z"
                    if 'T' in reminder_iso:
                        time_part = reminder_iso.split('T')[1][:5]
                        reminder_body['notes'] = f"Horário agendado: {time_part}\n\n" + reminder_body['notes']
                
                ts.tasks().insert(tasklist=tasklist_id, body=reminder_body).execute()
                print(f"[Google Tasks] Lembrete criado para a tarefa {task_doc.id}")
        except Exception as g_err:
            print(f"[Google Tasks] Erro ao criar lembrete no Google Tasks: {g_err}")

        # Marca como enviado para não repetir

        if due_reminder:
            updated_reminders = []
            matched = False
            for reminder in task_reminders:
                if not matched and reminder.get('id') == due_reminder.get('id'):
                    updated_reminders.append({**reminder, 'reminder_sent': True})
                    matched = True
                else:
                    updated_reminders.append(reminder)
            task_doc.reference.update(_build_task_reminder_state_payload(updated_reminders))
        else:
            task_doc.reference.update({'reminder_sent': True})

    # 4. Notificacoes agendadas pelo planejador proativo de IA e Verificacao de Duplicatas de Contatos
    try:
        from ai_notification_planner import dispatch_pending_ai_notifications, dispatch_scheduled_whatsapp_messages
        from contact_merge_utils import find_and_notify_duplicate_contacts
        dispatch_pending_ai_notifications(db, now)
        dispatch_scheduled_whatsapp_messages(db, now)
    except Exception as exc:
        print(f"[AINotifications] Falha ao despachar notificacoes agendadas: {exc}")



@https_fn.on_call()

def upload_to_drive(req: https_fn.CallableRequest):

    """Realiza o upload de um arquivo para o Google Drive"""

    import base64

    from googleapiclient.http import MediaIoBaseUpload

    import io

    data = req.data

    file_name = data.get('fileName')

    file_content_b64 = data.get('fileContent')

    mime_type = data.get('mimeType', 'application/octet-stream')

    folder_id = data.get('folderId')

    if not file_name or not file_content_b64:

        raise https_fn.HttpsError(

            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,

            message="O nome e o conteúdo do arquivo são obrigatórios."

        )

    try:

        service = get_drive_service()

        file_metadata = {'name': file_name}

        if folder_id:

            file_metadata['parents'] = [folder_id]

        file_content = base64.b64decode(file_content_b64)

        fh = io.BytesIO(file_content)

        media = MediaIoBaseUpload(fh, mimetype=mime_type, resumable=True)

        file = service.files().create(body=file_metadata, media_body=media, fields='id, webViewLink').execute()
        
        # Share the file publicly so frontend preview thumbnails work without 403 errors
        try:
            service.permissions().create(fileId=file.get('id'), body={'type': 'anyone', 'role': 'reader'}).execute()
        except Exception as perm_e:
            print(f"Aviso: Não foi possível definir a permissão pública: {perm_e}")

        return {'fileId': file.get('id'), 'webViewLink': file.get('webViewLink')}

    except Exception as e:

        print(f"Erro no upload para o Drive: {str(e)}")

        if isinstance(e, GoogleAuthRevokedError) or is_google_invalid_grant_error(e):
            try:
                get_db().collection('system').document('google_credentials').set({
                    'auth_status': 'reauth_required',
                    'auth_error': 'invalid_grant',
                    'auth_error_message': GOOGLE_REAUTH_MESSAGE,
                    'updated_at': firestore.SERVER_TIMESTAMP
                }, merge=True)
            except Exception as auth_status_err:
                print(f"Falha ao marcar reautenticacao Google: {auth_status_err}")
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message=GOOGLE_REAUTH_MESSAGE
            )

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INTERNAL, message=str(e))



@firestore_fn.on_document_written(document="tarefas/{taskId}")
def on_tarefa_written(event: firestore_fn.Event[firestore_fn.Change[firestore_fn.DocumentSnapshot]]):
    """Trigger disparado quando uma tarefa é atualizada, para monitorar processo_sei e horário"""
    if not event.data or not event.data.after or not event.data.after.exists: return

    after = event.data.after.to_dict() or {}
    before = event.data.before.to_dict() if event.data.before and event.data.before.exists else {}

    taskId = event.params['taskId']
    db = get_db()

    # Checa alteração de horário de início/fim ou prazo para forçar trigger pro Google Tasks/Calendar
    if (
        after.get('horario_inicio') != before.get('horario_inicio')
        or after.get('horario_fim') != before.get('horario_fim')
        or after.get('data_limite') != before.get('data_limite')
    ):
        sync_ref = db.collection('system').document('sync')
        sync_data = sync_ref.get().to_dict() or {}
        current_status = sync_data.get('status')

        if current_status in ('processing', 'requested'):
            queue_sync_request(db, 'task-schedule-change')
        else:
            sync_ref.set({
                'status': 'requested',
                'requested_at': datetime.now(timezone.utc).isoformat(),
                'last_trigger': 'task-schedule-change'
            }, merge=True)
    
@firestore_fn.on_document_updated(document="tarefas/{taskId}")
def on_processo_updated(event: firestore_fn.Event[firestore_fn.Change[firestore_fn.DocumentSnapshot]]):

    """Trigger disparado quando uma tarefa é atualizada, para monitorar processo_sei"""

    if not event.data.after.exists: return



    before = event.data.before.to_dict() or {}

    after = event.data.after.to_dict() or {}



    # Condição: Se area_tematica == 'CLC' e o campo processo_sei for alterado/inserido.

    if after.get('area_tematica') == 'CLC' and after.get('processo_sei'):

        if before.get('processo_sei') != after.get('processo_sei'):

            taskId = event.params['taskId']

            db = get_db()

            db.collection('tarefas').document(taskId).update({'sync_status': 'processando'})



            # Dispara via PubSub para o Node.js

            from google.cloud import pubsub_v1

            import json

            import os



            try:

                publisher = pubsub_v1.PublisherClient()

                topic_path = publisher.topic_path(os.environ.get('GCLOUD_PROJECT'), 'scrape-sipac')



                message_data = {

                    "taskId": taskId,

                    "processoSei": after.get('processo_sei'),

                    "folderId": db.collection('system').document('config').get().to_dict().get('googleDriveFolderId')

                }



                publisher.publish(topic_path, json.dumps(message_data).encode('utf-8'))

                print(f"Mensagem enviada para tópico scrape-sipac: {taskId}")

            except Exception as e:

                print(f"Erro ao publicar no PubSub: {e}")



@pubsub_fn.on_message_published(topic="vectorize-process")

def on_vectorize_requested(event: pubsub_fn.CloudEvent[pubsub_fn.MessagePublishedData]):

    """Trigger disparado via PubSub para vetorizar documentos"""

    import json

    try:

        # No Gen 2, event.data.message.data contém os bytes
        msg_bytes = event.data.message.data
        if msg_bytes:
            import base64
            if isinstance(msg_bytes, str):
                message_text = base64.b64decode(msg_bytes).decode('utf-8')
            else:
                message_text = msg_bytes.decode('utf-8')
        else:
             message_text = getattr(event.data.message, "text", "")



        data = json.loads(message_text)

        task_id = data.get('taskId')

        if task_id:

            process_vectorization(task_id)

    except Exception as e:

        print(f"Erro ao processar mensagem PubSub: {e}")



@https_fn.on_call(memory=options.MemoryOption.GB_1, timeout_sec=540)

def vectorize_process_docs_callable(req: https_fn.CallableRequest):

    """Versão callable para o frontend ou testes manuais"""

    task_id = req.data.get('taskId')

    if not task_id: return {'success': False, 'error': 'taskId faltante'}

    return process_vectorization(task_id)



def process_vectorization(task_id):

    """Lógica central de extração e vetorização"""

    from google import genai

    db = get_db()

    task_doc = db.collection('tarefas').document(task_id).get()

    if not task_doc.exists: return {'success': False, 'error': 'Tarefa não encontrada'}



    task_data = task_doc.to_dict()

    # Guard: tarefas cristalizadas são indexadas pelo novo pipeline de artefatos (KG).
    # O pipeline legado (processos_conhecimento) não deve criar duplicatas para elas.
    if task_data.get('kg_crystallized'):
        return {'success': True, 'vectorized_count': 0, 'skipped': 'kg_pipeline'}

    pool_dados = task_data.get('pool_dados', [])



    # Buscar chave do Gemini

    keys_doc = _cached_doc_get(db, 'system', 'api_keys')

    GEMINI_API_KEY = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None

    if not GEMINI_API_KEY: return {'success': False, 'error': 'Chave Gemini não configurada'}



    client = genai.Client(api_key=GEMINI_API_KEY)



    count = 0

    for item in pool_dados:

        if item.get('tipo') == 'arquivo' and item.get('drive_file_id'):

            file_id = item['drive_file_id']

            # Verifica se já foi vetorizado

            existing = db.collection('processos_conhecimento').where('file_id', '==', file_id).get()

            if not existing:

                try:

                    # Download do Drive

                    service = get_drive_service()

                    request = service.files().get_media(fileId=file_id)

                    file_content = request.execute()



                    # Determinar MIME type

                    mime_type = "application/pdf" if item.get('nome', '').lower().endswith('.pdf') else "text/html"



                    # Extração de texto via Gemini 1.5 Flash

                    response = client.models.generate_content(model="gemini-3.5-flash-lite", contents=[

                        "Extraia todo o texto relevante deste documento para indexação. Se for HTML, ignore tags. Se for PDF, faça OCR se necessário.",

                        {"mime_type": mime_type, "data": file_content}

                    ])

                    text_content = response.text if response.text else f"Conteúdo de {item.get('nome')}"



                    embedding_vec = get_embedding(text_content, api_key=GEMINI_API_KEY)

                    db.collection('processos_conhecimento').add({

                        'task_id': task_id,

                        'file_id': file_id,

                        'nome': item.get('nome'),

                        'texto': text_content,

                        'embedding': embedding_vec,

                        'data_vetorizacao': firestore.SERVER_TIMESTAMP

                    })

                    count += 1

                except Exception as e:

                    print(f"Erro ao vetorizar {file_id}: {e}")



    return {'success': True, 'vectorized_count': count}



@https_fn.on_call(memory=options.MemoryOption.GB_1, timeout_sec=120)
def vectorizeKnowledgeItemCallable(req: https_fn.CallableRequest):
    """Vetoriza um único item da base de conhecimento."""
    db = get_db()
    
    knowledge_id = req.data.get('knowledgeId')
    if not knowledge_id:
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message="knowledgeId é obrigatório.")

    doc_ref = db.collection('conhecimento').document(knowledge_id)
    doc = doc_ref.get()

    if not doc.exists:
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.NOT_FOUND, message="Item de conhecimento não encontrado.")

    item_data = doc.to_dict()
    text_content = item_data.get('texto_bruto')

    if not text_content:
        return {'success': False, 'message': 'Nenhum texto bruto para vetorizar.'}

    try:
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        GEMINI_API_KEY = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        if not GEMINI_API_KEY:
            raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION, message="Chave Gemini não configurada.")

        embedding_vec = get_embedding(text_content, api_key=GEMINI_API_KEY)
        doc_ref.update({'embedding': embedding_vec})
        
        return {'success': True, 'message': f'Item {knowledge_id} vetorizado.'}
    except Exception as e:
        print(f"Erro ao vetorizar {knowledge_id}: {e}")
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INTERNAL, message=str(e))


@https_fn.on_call(memory=options.MemoryOption.GB_1, timeout_sec=120)
def extractAndVectorizeRAGItem(req: https_fn.CallableRequest):
    """
    Extrai texto de um arquivo (PDF/TXT/MD) e vetoriza o item já existente na coleção 'conhecimento'.
    Chamado automaticamente após o upload de um arquivo para uma base RAG.
    """
    data = req.data
    file_base64 = data.get('fileBase64')
    mime_type = data.get('mimeType', 'application/octet-stream')
    filename = data.get('filename') or data.get('fileName') or 'arquivo'
    knowledge_id = data.get('knowledgeId')

    if not file_base64 or not knowledge_id:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="fileBase64 e knowledgeId são obrigatórios."
        )

    db = get_db()
    doc_ref = db.collection('conhecimento').document(knowledge_id)
    if not doc_ref.get().exists:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.NOT_FOUND,
            message="Item de conhecimento não encontrado."
        )

    # Decodifica o arquivo
    try:
        file_bytes = base64.b64decode(file_base64)
    except Exception as e:
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message=f"Erro ao decodificar arquivo: {e}")

    # Extrai texto
    texto_bruto = ""
    extraction_strategy = "text_decode"
    extraction_metadata = None
    try:
        if is_pdf_mime_type(filename, mime_type):
            pdf_result = extract_pdf_text_with_fallback(
                file_bytes,
                filename,
                api_key=get_gemini_api_key(),
                allow_gemini_fallback=True,
            )
            texto_bruto = pdf_result.get('text', '')
            extraction_strategy = pdf_result.get('strategy', 'none')
            extraction_metadata = pdf_result.get('metadata')
        elif is_docx_mime_type(filename, mime_type):
            texto_bruto, extraction_metadata = extract_docx_text(file_bytes)
            extraction_strategy = "docx_mammoth"
        else:
            texto_bruto = file_bytes.decode('utf-8', errors='replace').strip()
    except Exception as e:
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INTERNAL, message=f"Erro ao extrair texto: {e}")

    if not texto_bruto:
        return {'success': False, 'vectorized': False, 'message': 'Nenhum texto extraído do arquivo.'}

    # Trunca para 500.000 chars antes de salvar no Firestore (limite de 1MB por documento)
    texto_bruto = texto_bruto[:500000]

    # Salva texto no documento existente
    doc_updates = {
        'texto_bruto': texto_bruto,
        'texto_extraido_por': extraction_strategy,
    }
    if extraction_metadata:
        doc_updates['extraction_metadata'] = extraction_metadata
    doc_ref.update(doc_updates)

    # Vetoriza
    try:
        api_key = get_gemini_api_key()
        if not api_key:
            return {'success': True, 'vectorized': False, 'message': 'Texto salvo, mas chave Gemini não configurada.'}

        embedding_vec = get_embedding(texto_bruto, api_key=api_key)
        doc_ref.update({'embedding': embedding_vec})
        return {'success': True, 'vectorized': True}
    except Exception as e:
        print(f"Erro ao vetorizar RAG item {knowledge_id}: {e}")
        return {'success': True, 'vectorized': False, 'message': str(e)}


@https_fn.on_call(memory=options.MemoryOption.GB_1, timeout_sec=120)
def generate_task_with_ia(req: https_fn.CallableRequest):
    """
    Gera os campos de uma tarefa baseada em input de texto/áudio usando Gemini.
    Considera o contexto RAG personalizado se fornecido.
    """
    data = req.data
    content = data.get('content')
    origin = data.get('origin', 'manual')
    rag_context_id = data.get('ragContext') or data.get('base_conhecimento')
    extra_context = data.get('extraContext', '')
    extra_context_id = data.get('extraContextId')
    knowledge_item_ids = data.get('knowledgeItemIds', [])
    available_tags = data.get('availableTags', [])

    if not content:
        return {"error": "Conteúdo não fornecido"}

    today = datetime.now().date().isoformat()

    db = firestore.client()
    keys_doc = _cached_doc_get(db, 'system', 'api_keys')
    api_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None

    if not api_key:
        return {"error": "Gemini API Key não encontrada no sistema."}

    genai = get_genai_module()
    client = genai.Client(api_key=api_key)

    # Busca contexto do RAG da base principal
    rag_retrieved_context = ""
    if rag_context_id and rag_context_id != "Nenhum":
        rag_retrieved_context = retrieve_personalized_rag_context(db, genai, content, rag_context_id)

    # Busca contexto dos arquivos extras desta ação (RAG isolado)
    extra_rag_context = ""
    if extra_context_id or knowledge_item_ids:
        extra_rag_context = retrieve_extra_context_rag(db, genai, content, extra_context_id, knowledge_item_ids)

    tags_string = ", ".join(f'"{tag}"' for tag in available_tags) if available_tags else '"GERAL", "NÃO CLASSIFICADA"'

    # Prompt enriquecido com todos os contextos disponíveis
    prompt = f"""
    Você é o HERMES IA, consultor de produtividade avançada do André.
    Seu objetivo é transformar um fragmento de informação (WhatsApp, Áudio ou Texto) em uma estrutura de Deep Work (Tarefa Planejada) de altíssima qualidade.

    --- BASE DE CONHECIMENTO PRINCIPAL (RAG) ---
    Contexto recuperado da base de conhecimento selecionada pelo usuário:
    {rag_retrieved_context if rag_retrieved_context else "Nenhuma base RAG selecionada."}

    --- DOCUMENTOS EXTRAS DESTA AÇÃO ---
    Documentos carregados especificamente para subsidiar esta demanda:
    {extra_rag_context if extra_rag_context else "Nenhum documento extra carregado."}

    --- CONTEXTO TEXTUAL ADICIONAL ---
    {extra_context if extra_context else "Nenhum."}

    --- CONTEÚDO BRUTO PARA PROCESSAR ---
    Origem: {origin}
    Conteúdo: {content}

    SUA MISSÃO:
    1. Analise TODOS os contextos acima. Priorize os documentos extras e o RAG para definir a forma correta de execução.
    2. Crie um TÍTULO impactante, profissional e específico (reflita exatamente a demanda).
    3. Escreva uma DESCRIÇÃO detalhada: contextualize o André sobre o que é a demanda, por que ela existe e o que precisa ser entregue.
    4. Defina a CATEGORIA escolhendo EXATAMENTE UMA das tags válidas fornecidas abaixo. 
       TAGS DISPONÍVEIS: [{tags_string}]
       IMPORTANTE: Escolha a tag que mais se adeque ao contexto da ação. Se nenhuma for perfeitamente adequada, escolha "GERAL" ou "NÃO CLASSIFICADA". NUNCA INVENTE OUTRA TAG.
    5. Crie um PLANO DE AÇÃO (checklist) com no máximo 5 etapas concretas e sequenciais para resolver a demanda.
       REGRAS DO PLANO: cada etapa deve ser específica e acionável para ESTA demanda. Mencione elementos concretos presentes no conteúdo (nomes, sistemas, processos, documentos). Proibido etapas genéricas como "analisar o processo" sem especificar qual. Se a demanda for simples e não justificar 5 etapas, use menos.
    6. Defina a DATA DE EXECUÇÃO (data_limite) seguindo estas regras obrigatórias:
       - Se houver uma data ou prazo mencionado no conteúdo, use-o — MAS a data gerada DEVE ser igual ou posterior a {today}.
       - Se nenhum prazo for mencionado, use a data de hoje ({today}).
       - NUNCA gere uma data anterior a {today}. Isso é proibido.

    SAÍDA ESPERADA (JSON puro, sem markdown):
    {{
      "titulo": "Título claro e profissional da demanda",
      "descricao": "Descrição detalhada contextualizando a demanda e o que deve ser feito",
      "area_tematica": "NOME_EXATO_DE_UMA_DAS_TAGS_DISPONIVEIS",
      "status": "em andamento",
      "data_limite": "YYYY-MM-DD",
      "plano_acao": ["Passo 1 detalhado", "Passo 2 detalhado", "..."]
    }}
    """

    try:
        response = client.models.generate_content(model="gemini-3.5-flash-lite", contents=prompt)
        text = response.text
        # Limpeza para garantir JSON puro
        if "```json" in text:
            text = text.split("```json")[1].split("```")[0].strip()
        elif "```" in text:
            text = text.split("```")[1].split("```")[0].strip()
        elif "{" in text:
            start = text.find("{")
            end = text.rfind("}") + 1
            text = text[start:end]
            
        result = json.loads(text)
        # Garante que data_limite nunca seja anterior a hoje
        generated_date = result.get('data_limite', '')
        if not generated_date or generated_date < today:
            result['data_limite'] = today
        return result
    except Exception as e:
        print(f"Erro no processamento Gemini RAG: {e}")
        return {"error": str(e), "raw_response": text if 'text' in locals() else None}

def retrieve_personalized_rag_context(db, genai, query_text, base_id):

    """

    Recupera informações de múltiplas fontes para formar o contexto RAG de uma base.

    Inclui busca vetorial e busca estruturada.

    """

    context_parts = []

    

    # --- 1. Vector Search (Semantic) ---

    if query_text:

        try:

            # Re-using the logic from findSimilarKnowledge

            query_embedding = get_embedding(query_text)



            docs_query = db.collection('conhecimento').where('embedding', '!=', None)

            if base_id and base_id != "Nenhum":

                docs_query = docs_query.where('base_id', '==', base_id)

            

            docs = list(docs_query.limit(50).stream()) # Limit to 50 docs for performance



            similar_items = []

            if docs:

                for doc in docs:

                    item = doc.to_dict()

                    if 'embedding' in item and len(item['embedding']) > 0:

                        similarity = cosine_similarity(query_embedding, item['embedding'])

                        if similarity > 0.7: # Threshold to ensure relevance

                            similar_items.append({

                                'titulo': item.get('titulo'),

                                'texto': item.get('texto_bruto', '')[:1000],

                                'similarity': similarity

                            })

                

                similar_items.sort(key=lambda x: x['similarity'], reverse=True)



                if similar_items:

                    context_parts.append("### [DOCUMENTOS SIMILARES ENCONTRADOS (BUSCA VETORIAL)]")

                    for item in similar_items[:3]: # Add top 3

                        context_parts.append(f"#### {item['titulo']}\n{item['texto']}")

        except Exception as e:

            print(f"Error during vector search in RAG context retrieval: {e}")





    # --- 2. Structured Search (Existing Logic) ---

    base_doc = db.collection('knowledge_bases').document(base_id).get()

    config = {}
    base_name = ""

    if base_doc.exists:

        base_data = base_doc.to_dict()
        base_name = base_data.get('nome', '')

        config = base_data.get('configuracao_rag', {})

        context_parts.append(f"### [BASE ATUAL] {base_name} - {base_data.get('descricao', '')}")

    else:

        # Se não existe como base_id, tenta tratar como area_tematica (legado)

        context_parts.append(f"### [PESQUISA LEGADA POR CATEGORIA] {base_id}")

        config = {

            'incluir_manual': True,

            'incluir_diarios': True,

            'area_tematicas_vinculadas': [base_id]

        }



    # Busca itens do MANUAL (Conhecimento Mestre)

    if config.get('incluir_manual'):

        cats = config.get('area_tematicas_vinculadas', [])

        if cats:

            for cat in cats:

                master = db.collection('conhecimento_mestre').where('area_tematica', '==', cat).limit(2).stream()

                for d in master:

                    m = d.to_dict()

                    context_parts.append(f"### [SOP/MANUAL - {cat}] {m.get('titulo')}\n{m.get('conteudo')}")



    # Busca HISTÓRICO DE DIÁRIOS (Acompanhamento)

    if config.get('incluir_diarios'):

        cats = config.get('area_tematicas_vinculadas', [])

        if cats:

            for cat in cats:

                tasks = (db.collection('tarefas')
                    .where('area_tematica', '==', cat)
                    .where('status', '==', 'concluído')
                    .order_by('data_conclusao', direction=firestore.Query.DESCENDING)
                    .limit(3).stream())

                for d in tasks:

                    t = d.to_dict()

                    context_parts.append(f"### [HISTÓRICO - {cat}] {t.get('titulo')}\nNOTAS: {t.get('notas', '')}")

    # --- 3. Dynamic Module Integrations ---
    # Serviços Base Integration
    if base_id == "servicos" or base_name.lower() == "serviços":
        try:
            services_ref = db.collection('servicos').stream()
            context_parts.append("### [MÓDULO DE SERVIÇOS - PORTFÓLIO INTEGRADO]")
            for doc in services_ref:
                s = doc.to_dict()
                context_parts.append(f"#### Serviço: {s.get('titulo')} (Cliente: {s.get('cliente')})\nPapel: {s.get('papel')}\nStatus: {s.get('status')}\nDescrição: {s.get('descricao')}\nValor Total: R$ {s.get('valor_total')}\nCarga Horária: {s.get('carga_horaria_semanal')}h/semana")
        except Exception as e:
            print(f"Error fetching services in RAG: {e}")

    # Financeira Base Integration
    if base_id == "financeira" or base_name.lower() == "financeira":
        try:
            transactions_ref = db.collection('finance_transactions').where('status', '!=', 'deleted').limit(30).stream()
            context_parts.append("### [MÓDULO FINANCEIRO - TRANSAÇÕES INTEGRADAS]")
            for doc in transactions_ref:
                t = doc.to_dict()
                context_parts.append(f"- {t.get('date')}: {t.get('description')} | Valor: R$ {t.get('amount')} | Categoria: {t.get('category')} | Tipo: {t.get('type')}")
            
            fixed_ref = db.collection('fixed_bills').stream()
            context_parts.append("### [MÓDULO FINANCEIRO - CONTAS FIXAS INTEGRADAS]")
            for doc in fixed_ref:
                b = doc.to_dict()
                context_parts.append(f"- Conta Fixa: {b.get('name')} | Valor: R$ {b.get('amount')} | Categoria: {b.get('category')} | Vencimento: Dia {b.get('dueDate')} | Pago: {'Sim' if b.get('isPaid') else 'Não'}")
        except Exception as e:
            print(f"Error fetching finance in RAG: {e}")

    # Saúde Base Integration
    if base_id == "saude" or base_name.lower() == "saúde":
        try:
            exams_ref = db.collection('exames').stream()
            context_parts.append("### [MÓDULO DE SAÚDE - EXAMES E CONSULTAS INTEGRADOS]")
            for doc in exams_ref:
                e = doc.to_dict()
                context_parts.append(f"#### Exame: {e.get('titulo')}\nData: {e.get('data_exame')}\nLaboratório: {e.get('laboratorio', '')}\nNotas: {e.get('notas', '')}")
                
            weights_ref = db.collection('health_weights').order_by('date', direction=firestore.Query.DESCENDING).limit(10).stream()
            context_parts.append("### [MÓDULO DE SAÚDE - HISTÓRICO DE PESO RECENTE]")
            for doc in weights_ref:
                w = doc.to_dict()
                context_parts.append(f"- {w.get('date')}: Peso: {w.get('weight')} kg | Gordura: {w.get('fatPercent', 'N/A')}% | Músculo: {w.get('musclePercent', 'N/A')}%")
        except Exception as e:
            print(f"Error fetching health in RAG: {e}")

    # Specific Custom Organization/Base Linked Services
    if base_id:
        try:
            linked_services = db.collection('servicos').where('base_id', '==', base_id).stream()
            for doc in linked_services:
                s = doc.to_dict()
                context_parts.append(f"### [SERVIÇO VINCULADO A ESTA ÁREA TEMÁTICA]\n#### {s.get('titulo')} (Cliente: {s.get('cliente')})\nPapel: {s.get('papel')}\nStatus: {s.get('status')}\nDescrição: {s.get('descricao')}\nValor Total: R$ {s.get('valor_total')}")
        except Exception as e:
            print(f"Error fetching linked services in RAG: {e}")

    return "\n\n".join(context_parts)


def retrieve_extra_context_rag(db, genai, query_text, extra_context_id=None, item_ids=None):
    """
    Recupera contexto dos arquivos extras enviados pelo usuário para uma ação específica.
    Usa apenas busca vetorial, filtrada pelo extra_context_id.
    """
    context_parts = []

    if not query_text or not extra_context_id:
        return ""

    try:
        query_embedding = get_embedding(query_text)

        docs = []
        if extra_context_id:
            docs.extend(list(
                db.collection('conhecimento')
                .where('embedding', '!=', None)
                .where('extra_context_id', '==', extra_context_id)
                .limit(20)
                .stream()
            ))

        if item_ids:
            for iid in item_ids:
                doc = db.collection('conhecimento').document(iid).get()
                if doc.exists:
                    docs.append(doc)

        similar_items = []
        for doc in docs:
            item = doc.to_dict()
            if 'embedding' in item and item['embedding']:
                similarity = cosine_similarity(query_embedding, item['embedding'])
                if similarity > 0.4:  # Limiar menor para garantir que o conteúdo extra sempre seja incluído
                    similar_items.append({
                        'titulo': item.get('titulo'),
                        'texto': item.get('texto_bruto', '')[:2000],
                        'similarity': similarity
                    })

        similar_items.sort(key=lambda x: x['similarity'], reverse=True)

        if similar_items:
            context_parts.append("### [DOCUMENTOS EXTRAS CARREGADOS PARA ESTA AÇÃO]")
            for item in similar_items[:5]:
                context_parts.append(f"#### {item['titulo']}\n{item['texto']}")
        elif docs:
            # Se há docs mas nenhum passou no threshold, inclui os top 3 mesmo assim
            all_items = []
            for doc in docs:
                item = doc.to_dict()
                if item.get('texto_bruto'):
                    all_items.append({'titulo': item.get('titulo'), 'texto': item.get('texto_bruto', '')[:2000]})
            if all_items:
                context_parts.append("### [DOCUMENTOS EXTRAS CARREGADOS PARA ESTA AÇÃO]")
                for item in all_items[:3]:
                    context_parts.append(f"#### {item['titulo']}\n{item['texto']}")

    except Exception as e:
        print(f"Erro no RAG de contexto extra: {e}")

    return "\n\n".join(context_parts)


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.GB_1,
    timeout_sec=120
)
def processExtraContextFile(req: https_fn.CallableRequest):
    """
    Processa um arquivo de contexto extra (PDF ou texto) para uma ação específica.
    Extrai o texto, salva em 'conhecimento' com extra_context_id e vetoriza automaticamente.
    """
    data = req.data
    file_base64 = data.get('fileBase64')
    filename = data.get('filename', 'arquivo')
    extra_context_id = data.get('extraContextId')
    mime_type = data.get('mimeType', 'application/octet-stream')

    if not file_base64 or not extra_context_id:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="fileBase64 e extraContextId são obrigatórios."
        )

    db = get_db()
    file_bytes = base64.b64decode(file_base64)

    # Extração de texto baseada no tipo do arquivo
    texto_bruto = ""
    extraction_strategy = "text_decode"
    extraction_metadata = None
    is_pdf = is_pdf_mime_type(filename, mime_type)
    is_docx = is_docx_mime_type(filename, mime_type)

    if is_pdf:
        try:
            pdf_result = extract_pdf_text_with_fallback(
                file_bytes,
                filename,
                api_key=get_gemini_api_key(),
                allow_gemini_fallback=True,
            )
            texto_bruto = pdf_result.get('text', '')
            extraction_strategy = pdf_result.get('strategy', 'none')
            extraction_metadata = pdf_result.get('metadata')
        except Exception as e:
            print(f"Erro ao extrair PDF '{filename}': {e}")
            texto_bruto = ""
    elif is_docx:
        try:
            texto_bruto, extraction_metadata = extract_docx_text(file_bytes)
            extraction_strategy = "docx_mammoth"
        except Exception as e:
            print(f"Erro ao extrair DOCX '{filename}': {e}")
            texto_bruto = ""
    else:
        # Arquivos de texto: TXT, MD, CSV, etc.
        try:
            texto_bruto = file_bytes.decode('utf-8').strip()
        except UnicodeDecodeError:
            texto_bruto = file_bytes.decode('latin-1', errors='replace').strip()

    # Salva no Firestore
    doc_id = str(uuid.uuid4())
    tipo = 'pdf' if is_pdf else 'docx' if is_docx else 'texto'
    doc_data = {
        'id': doc_id,
        'titulo': filename,
        'tipo_arquivo': tipo,
        'texto_bruto': texto_bruto,
        'extra_context_id': extra_context_id,
        'base_id': None,
        'tamanho': len(file_bytes),
        'data_criacao': datetime.now(timezone.utc).isoformat(),
        'origem': None,
        'parent_id': None,
        'texto_extraido_por': extraction_strategy,
    }
    if extraction_metadata:
        doc_data['extraction_metadata'] = extraction_metadata
    db.collection('conhecimento').document(doc_id).set(doc_data)

    # Vetorização automática
    vectorized = False
    if texto_bruto:
        try:
            api_key = get_gemini_api_key()

            if api_key:
                embedding_vec = get_embedding(texto_bruto, api_key=api_key)
                db.collection('conhecimento').document(doc_id).update({
                    'embedding': embedding_vec
                })
                vectorized = True
        except Exception as e:
            print(f"Erro ao vetorizar contexto extra '{filename}': {e}")

    return {'success': True, 'docId': doc_id, 'vectorized': vectorized}


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=60
)
def criarLembreteNoGoogleTasks(req: https_fn.CallableRequest):
    """Compatibilidade para funcao legado ainda presente no projeto Firebase."""
    raise https_fn.HttpsError(
        code=https_fn.FunctionsErrorCode.UNIMPLEMENTED,
        message=(
            "criarLembreteNoGoogleTasks e uma funcao legado sem implementacao "
            "local atual. Use o fluxo de lembretes sincronizado."
        )
    )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=60
)
def removerLembreteDoGoogleTasks(req: https_fn.CallableRequest):
    """Compatibilidade para funcao legado ainda presente no projeto Firebase."""
    raise https_fn.HttpsError(
        code=https_fn.FunctionsErrorCode.UNIMPLEMENTED,
        message=(
            "removerLembreteDoGoogleTasks e uma funcao legado sem implementacao "
            "local atual. Use o fluxo de lembretes sincronizado."
        )
    )


_TRANSCRICAO_RAPIDA_VIDEO_EXTS = {
    "mp4", "mov", "mkv", "avi", "webm", "m4v", "wmv", "flv", "mpeg", "mpg", "3gp", "ts"
}

# Limite para o fluxo de microfone (audioBase64), que não passa pelo Storage.
# É um limite sobre o binário decodificado; em base64 ele expande ~4/3, então
# precisa ficar bem abaixo do teto de ~32MB de payload das Cloud Functions
# (2nd gen) — senão a plataforma rejeita a requisição antes desta checagem
# rodar, e o cliente volta a ver o erro cru de transporte. Também fica abaixo
# do teto de 25MB por arquivo do Groq Whisper.
_MAX_AUDIO_BASE64_BYTES = 18 * 1024 * 1024

# Tamanho máximo de bloco enviado por vez ao Gemini para refinar a transcrição.
# Transcrições longas em um único prompt gigantesco podem estourar limites da
# API e retornar "400 INVALID_ARGUMENT"; dividir em blocos evita isso.
_GEMINI_REFINE_CHUNK_CHARS = 12000


def _split_text_into_chunks(texto, max_chars):
    """Divide o texto em blocos de até max_chars, cortando em fim de frase quando possível."""
    if len(texto) <= max_chars:
        return [texto] if texto else []
    chunks = []
    start = 0
    length = len(texto)
    while start < length:
        end = min(start + max_chars, length)
        if end < length:
            corte = texto.rfind('. ', start, end)
            if corte > start:
                end = corte + 1
        pedaco = texto[start:end].strip()
        if pedaco:
            chunks.append(pedaco)
        start = end
    return chunks


def _refinar_transcricao_com_gemini(gemini_client, texto_bruto):
    """
    Refina o texto bruto com Gemini em blocos, para não estourar limites da API em
    transcrições longas. Se o refinamento de um bloco falhar, usa o texto bruto
    daquele bloco como fallback em vez de derrubar a transcrição inteira.
    """
    chunks = _split_text_into_chunks(texto_bruto, _GEMINI_REFINE_CHUNK_CHARS)
    partes_refinadas = []
    for chunk in chunks:
        prompt = f"""
        Atue como um redator especialista. O texto a seguir é um trecho de uma transcrição de voz bruta.
        Sua tarefa:
        1. Corrigir pontuação e gramática (pt-BR).
        2. Remover vícios de linguagem (né, tipo, ahn).
        3. Manter o tom original e termos técnicos.
        4. Retorne APENAS o texto corrigido, sem introduções.

        Texto: "{chunk}"
        """
        try:
            result = gemini_client.models.generate_content(model="gemini-3.5-flash-lite", contents=prompt)
            texto_refinado = (result.text or "").strip()
            partes_refinadas.append(texto_refinado if texto_refinado else chunk)
        except Exception as e:
            print(f"Falha ao refinar bloco da transcrição com Gemini: {e}")
            partes_refinadas.append(chunk)
    return " ".join(partes_refinadas)


@https_fn.on_call(
    memory=options.MemoryOption.GB_1,
    timeout_sec=300
)



def transcreverAudio(req: https_fn.CallableRequest):



    """



    Recebe o caminho de um áudio ou vídeo já enviado para `quick_transcriptions/{uid}/{nome}`

    no Storage, transcreve com Groq (Whisper) e refina com Gemini. Quando o arquivo é vídeo,

    extrai a faixa de áudio com FFmpeg antes de transcrever. O binário é expurgado do Storage

    ao final (sucesso ou erro).

    """

    import subprocess

    import tempfile

    import os

    import imageio_ffmpeg

    from firebase_admin import storage as admin_storage

    from groq import Groq

    from google import genai



    if not req.auth:

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.UNAUTHENTICATED, message="Login necessário.")

    uid = req.auth.uid



    # Buscar chaves de API no Firestore

    try:

        db = get_db()

        keys_doc = _cached_doc_get(db, 'system', 'api_keys')

        if not keys_doc.exists:

            raise Exception("Documento system/api_keys não encontrado.")

        keys = keys_doc.to_dict()

        GROQ_API_KEY = keys.get('groq_api_key')

        GEMINI_API_KEY = keys.get('gemini_api_key')

        

        if not GROQ_API_KEY or not GEMINI_API_KEY:

            raise Exception("Chaves de API incompletas em system/api_keys.")

            

    except Exception as e:

        print(f"Erro ao buscar chaves de API: {e}")

        raise https_fn.HttpsError(

            code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,

            message="Configuração de API pendente no sistema."

        )



    data = req.data

    storage_path = data.get('storagePath')

    audio_base64 = data.get('audioBase64')

    extension = data.get('extension', '.m4a')



    if not extension.startswith('.'):

        extension = f".{extension}"



    # Aceita dois fluxos: arquivo já enviado ao Storage (`storagePath`, contorna o
    # limite de 32MB do Cloud Run para arquivos grandes) ou o binário em base64
    # diretamente no corpo (`audioBase64`, usado por microfone/WhatsApp no copiloto).

    if not storage_path and not audio_base64:

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message="Caminho do arquivo não fornecido.")



    if storage_path and not storage_path.startswith(f"quick_transcriptions/{uid}/"):

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.PERMISSION_DENIED, message="Caminho de arquivo inválido.")



    # O fluxo de microfone (audioBase64) não passa pelo Storage, então precisa de um
    # limite explícito aqui: acima disso o Groq (limite de 25MB por arquivo) e o
    # Gemini rejeitam o pedido, e sem essa checagem o erro cru do provedor vazava
    # pro cliente (ex.: "400 INVALID_ARGUMENT" em gravações longas pelo mobile).
    audio_bytes = None

    if audio_base64:

        import base64

        audio_bytes = base64.b64decode(audio_base64)

        if len(audio_bytes) > _MAX_AUDIO_BASE64_BYTES:

            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="Gravação muito longa para transcrição pelo microfone. Use a ferramenta 'Transcrições Longas' para áudios extensos."
            )



    temp_filename = None

    extracted_audio_filename = None

    blob = None

    try:

        # 1. Obter o binário em arquivo temporário (Storage ou base64)

        fd, temp_filename = tempfile.mkstemp(suffix=extension)

        os.close(fd)

        if storage_path:

            blob = admin_storage.bucket().blob(storage_path)

            blob.download_to_filename(temp_filename)

        else:

            with open(temp_filename, "wb") as audio_file:

                audio_file.write(audio_bytes)



        # 1.1 Se for vídeo, extrair apenas a faixa de áudio antes de transcrever

        is_video = extension.lstrip('.').lower() in _TRANSCRICAO_RAPIDA_VIDEO_EXTS

        transcribe_filename = temp_filename

        if is_video:

            ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()

            fd, extracted_audio_filename = tempfile.mkstemp(suffix=".m4a")

            os.close(fd)

            subprocess.run(

                [ffmpeg_exe, "-y", "-i", temp_filename, "-vn", "-acodec", "aac", "-b:a", "128k", extracted_audio_filename],

                check=True,

                capture_output=True,

            )

            transcribe_filename = extracted_audio_filename



        # 2. Transcrição via Groq (Whisper Large V3 Turbo)

        client = Groq(api_key=GROQ_API_KEY)

        with open(transcribe_filename, "rb") as file_stream:

            transcription = client.audio.transcriptions.create(

                file=(os.path.basename(transcribe_filename), file_stream),

                model="whisper-large-v3-turbo",

                response_format="json",

                language="pt",

                temperature=0.0

            )

        texto_bruto = transcription.text



        # Refinamento via Gemini Flash, em blocos (evita "400 INVALID_ARGUMENT" em
        # transcrições longas e nunca falha a chamada inteira por causa do refino).

        gemini_client = genai.Client(api_key=GEMINI_API_KEY)

        texto_refinado = _refinar_transcricao_com_gemini(gemini_client, texto_bruto)



        return {"raw": texto_bruto, "refined": texto_refinado}

    except Exception as e:

        print(f"Erro na transcrição: {str(e)}")

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INTERNAL, message=f"Falha ao processar áudio: {str(e)}")

    finally:

        if blob is not None:

            try:

                blob.delete()

            except Exception as e:

                print(f"Falha ao expurgar {storage_path}: {e}")

        if temp_filename and os.path.exists(temp_filename):

            try:

                os.remove(temp_filename)

            except:

                pass

        if extracted_audio_filename and os.path.exists(extracted_audio_filename):

            try:

                os.remove(extracted_audio_filename)

            except:

                pass



def start_file_indexing(item_id, item_data):

    """Lógica central de indexação com Gemini"""

    url_drive = item_data.get('url_drive')

    if not url_drive:

        return {'success': False, 'error': 'URL não encontrada'}



    import re

    def extract_file_id(url):

        match = re.search(r'[-\w]{25,}', url)

        return match.group(0) if match else None



    file_id = extract_file_id(url_drive)

    if not file_id:

        return {'success': False, 'error': 'ID do arquivo não identificado na URL'}



    try:

        db = get_db()

        GEMINI_API_KEY = get_gemini_api_key()

        if not GEMINI_API_KEY:

            return {'success': False, 'error': 'Chave de API Gemini não configurada'}



        from google import genai

        import json



        client = genai.Client(api_key=GEMINI_API_KEY)



        service = get_drive_service()

        file_metadata = service.files().get(fileId=file_id, fields='mimeType, name').execute()

        mime_type = file_metadata.get('mimeType')



        request = service.files().get_media(fileId=file_id)

        content = request.execute()



        prompt = ""

        parts = []



        if mime_type.startswith('image/'):

            prompt = """

            Analise esta imagem e retorne em JSON:

            1. ocr: Todo o texto escrito na imagem.

            2. descricao: Descrição semântica detalhada.

            3. resumo_tldr: Resumo de até 3 linhas.

            4. tags: Lista de 5-10 palavras-chave.

            5. area_tematica: Uma única palavra de classificação.

            """

            parts = [{"mime_type": mime_type, "data": content}, prompt]

        elif is_pdf_mime_type(file_metadata.get('name'), mime_type):

            pdf_result = extract_pdf_text_with_fallback(
                content,
                file_metadata.get('name') or 'documento.pdf',
                api_key=GEMINI_API_KEY,
                allow_gemini_fallback=True,
            )
            extracted_pdf_text = (pdf_result.get('text') or '').strip()

            if extracted_pdf_text:
                prompt = f"""

            Analise o texto abaixo, extraído de um PDF, e retorne em JSON:

            1. texto_bruto: Conteúdo principal extraído.

            2. resumo_tldr: Resumo de até 3 linhas.

            3. tags: Lista de 5-10 palavras-chave.

            4. area_tematica: Uma única palavra de classificação.

            TEXTO EXTRAÍDO:

            {extracted_pdf_text[:100000]}

            """

                parts = [prompt]
            else:
                prompt = """

            Analise este PDF e retorne em JSON:

            1. texto_bruto: Conteúdo principal extraído.

            2. resumo_tldr: Resumo de até 3 linhas.

            3. tags: Lista de 5-10 palavras-chave.

            4. area_tematica: Uma única palavra de classificação.

            """

                parts = [{"mime_type": mime_type, "data": content}, prompt]

        elif is_docx_mime_type(file_metadata.get('name'), mime_type):

            extracted_docx_text, docx_metadata = extract_docx_text(content)

            prompt = f"""

            Analise o texto abaixo, extraido de um arquivo Word DOCX, e retorne em JSON:

            1. texto_bruto: Conteudo principal extraido.

            2. resumo_tldr: Resumo de ate 3 linhas.

            3. tags: Lista de 5-10 palavras-chave.

            4. area_tematica: Uma unica palavra de classificacao.

            METADADOS DA EXTRACAO:

            {json.dumps(docx_metadata, ensure_ascii=False)}

            TEXTO EXTRAIDO:

            {extracted_docx_text[:100000]}

            """

            parts = [prompt]

        else:

            text_content = ""

            try:

                text_content = content.decode('utf-8')

            except:

                text_content = "[Binário]"



            prompt = f"""

            Analise este conteúdo e retorne em JSON:

            1. resumo_tldr: Resumo de até 3 linhas.

            2. tags: Lista de 5-10 palavras-chave.

            3. area_tematica: Uma única palavra de classificação.

            4. texto_bruto: O próprio texto.



            CONTEÚDO:

            {text_content[:100000]}

            """

            parts = [prompt]



        response = client.models.generate_content(model="gemini-3.5-flash-lite", contents=parts)

        res_text = response.text



        json_match = re.search(r'\{.*\}', res_text, re.DOTALL)

        if json_match:

            data = json.loads(json_match.group(0))

            updates = {

                'resumo_tldr': data.get('resumo_tldr'),

                'tags': data.get('tags'),

                'area_tematica': data.get('area_tematica', 'Geral').upper()

            }



            if mime_type.startswith('image/'):

                updates['texto_bruto'] = f"OCR: {data.get('ocr')}\n\nDESCRIÇÃO: {data.get('descricao')}"

            else:

                updates['texto_bruto'] = data.get('texto_bruto') or item_data.get('titulo')



            db.collection('conhecimento').document(item_id).set(updates, merge=True)

            return {'success': True, 'item_id': item_id}

        return {'success': False, 'error': 'Não foi possível gerar metadados JSON'}



    except Exception as e:

        print(f"Erro ao processar arquivo {item_id}: {str(e)}")

        return {'success': False, 'error': str(e)}






















def cosine_similarity(v1, v2):
    """Calculates the cosine similarity between two vectors (pure Python)."""
    dot = sum(a * b for a, b in zip(v1, v2))
    norm1 = sum(a * a for a in v1) ** 0.5
    norm2 = sum(b * b for b in v2) ** 0.5
    if norm1 == 0 or norm2 == 0:
        return 0.0
    return dot / (norm1 * norm2)















@https_fn.on_call(memory=options.MemoryOption.GB_1, timeout_sec=120)







def findSimilarKnowledge(req: https_fn.CallableRequest):







    """Finds similar knowledge items using vector search."""













    db = get_db()















    query_text = req.data.get('query_text')







    base_id = req.data.get('base_id')







    top_n = req.data.get('top_n', 5)















    if not query_text:







        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message="query_text é obrigatório.")















    try:







        keys_doc = _cached_doc_get(db, 'system', 'api_keys')







        GEMINI_API_KEY = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None







        if not GEMINI_API_KEY:







            raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION, message="Chave Gemini não configurada.")































        # 1. Generate embedding for the query







        query_embedding = get_embedding(query_text)















        # 2. Fetch documents







        docs_query = db.collection('conhecimento').where('embedding', '!=', None)







        if base_id:







            docs_query = docs_query.where('base_id', '==', base_id)







        







        docs = list(docs_query.stream())















        if not docs:







            return {'results': []}















        # 3. Compute cosine similarity in memory







        results = []







        for doc in docs:







            item = doc.to_dict()







            if 'embedding' in item and len(item['embedding']) > 0:







                similarity = cosine_similarity(query_embedding, item['embedding'])







                results.append({







                    'id': doc.id,







                    'titulo': item.get('titulo'),







                    'resumo_tldr': item.get('resumo_tldr'),







                    'texto_bruto': item.get('texto_bruto', '')[:500], # Truncate for response







                    'similarity': similarity







                })















        # 4. Sort and get top N







        results.sort(key=lambda x: x['similarity'], reverse=True)







        







        return {'results': results[:top_n]}















    except Exception as e:







        print(f"Erro em findSimilarKnowledge: {e}")







        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INTERNAL, message=str(e))







@firestore_fn.on_document_updated(document="conhecimento/{itemId}")







def on_knowledge_item_updated(event: firestore_fn.Event[firestore_fn.Change[firestore_fn.DocumentSnapshot]]):







    """







    Trigger to automatically vectorize a knowledge item when its text content is added or changed.



    """



    if not event.data.after or not event.data.after.exists:



        return  # Document was deleted







    after_data = event.data.after.to_dict() or {}



    before_data = (event.data.before.to_dict() or {}) if event.data.before and event.data.before.exists else {}







    text_after = after_data.get('texto_bruto')



    text_before = before_data.get('texto_bruto')







    # Vectorize if text content was added or changed.



    if text_after and text_after != text_before:



        db = get_db()



        



        doc_ref = event.data.after.reference







        try:



            keys_doc = _cached_doc_get(db, 'system', 'api_keys')



            GEMINI_API_KEY = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None



            if not GEMINI_API_KEY:



                print("Gemini API Key not found, skipping vectorization.")



                return











            embedding_vec = get_embedding(text_after, api_key=GEMINI_API_KEY)

            doc_ref.update({'embedding': embedding_vec})



            print(f"Successfully vectorized item {doc_ref.id}")



            



        except Exception as e:



            print(f"Error during vectorization for {doc_ref.id}: {e}")











@firestore_fn.on_document_created(
    document="conhecimento/{itemId}",
    memory=options.MemoryOption.MB_512,
)



def on_arquivo_adicionado(event: firestore_fn.Event[firestore_fn.DocumentSnapshot | None]):



    """Trigger disparado quando um novo arquivo é adicionado"""



    if not event.data: return



    item_data = event.data.to_dict()

    item_id = event.params["itemId"]



    # Ignora links diretos (sem processamento de IA/OCR)

    if item_data.get('tipo_arquivo') == 'link':

        return



    if item_data.get('tags') and item_data.get('resumo_tldr'):

        return



    start_file_indexing(item_id, item_data)



COPILOT_SOUL_DEFAULT = {
    "tone": "Consultivo, analítico e construtivamente crítico — socrático quando há risco.",
    "detail_level": "Conciso; aprofunda apenas onde houver risco, ambiguidade ou premissa frágil.",
    "interaction_style": "Questiona premissas, expõe pontos cegos e riscos e faz no máximo uma pergunta afiada quando a decisão tem peso. Não bajula nem valida por reflexo; discorda com fundamento quando os fatos não sustentam o pedido.",
}

MEMORY_NODE_TYPES = {"regra_global", "fato_isolado"}
MEMORY_SIMILARITY_CREATE_THRESHOLD = 0.90
MEMORY_SIMILARITY_DUPLICATE_THRESHOLD = 0.965
MEMORY_MIN_FACT_LENGTH = 24


def _iso_now_utc() -> str:
    return datetime.now(timezone.utc).isoformat()


def _normalize_memory_category(value: str | None) -> str:
    raw = (value or "").strip().lower()
    aliases = {
        "regra": "regra_global",
        "regra_global": "regra_global",
        "regra global": "regra_global",
        "policy": "regra_global",
        "fato": "fato_isolado",
        "fato_isolado": "fato_isolado",
        "fato isolado": "fato_isolado",
        "fact": "fato_isolado",
    }
    return aliases.get(raw, "fato_isolado")


def _normalize_text_for_match(text: str | None) -> str:
    return re.sub(r"\s+", " ", (text or "").strip().lower())


def _normalize_pop_text(text: str | None) -> str:
    raw = unicodedata.normalize("NFKD", (text or "").strip().lower())
    without_accents = "".join(ch for ch in raw if not unicodedata.combining(ch))
    without_punct = re.sub(r"[^\w\s]", " ", without_accents)
    return re.sub(r"\s+", " ", without_punct).strip()


PORTUGUESE_STOPWORDS = {
    "de", "do", "da", "em", "no", "na", "se", "um", "uma", "os", "as",
    "com", "para", "por", "dos", "das", "aos", "nos", "nas", "pelo",
    "pela", "pelos", "pelas", "sem", "sob", "sobre", "o", "a", "e", "ou",
    "como"
}


def _match_pop_directives(db, prompt: str) -> list[dict]:
    global _POPS_DATA_CACHE
    prompt_norm = _normalize_pop_text(prompt)
    if not prompt_norm:
        return []

    prompt_terms_all = prompt_norm.split()
    prompt_terms_filtered = {t for t in prompt_terms_all if t not in PORTUGUESE_STOPWORDS}

    try:
        now = time.monotonic()
        if _POPS_DATA_CACHE and (now - _POPS_DATA_CACHE[0]) < _POPS_DATA_TTL:
            all_pops = _POPS_DATA_CACHE[1]
        else:
            all_pops = [
                {"id": d.id, **( d.to_dict() or {})}
                for d in db.collection("pops_diretrizes").stream()
            ]
            _POPS_DATA_CACHE = (now, all_pops)

        def _collect_matches(pops: list[dict]) -> list[dict]:
            matched: list[dict] = []

            for pop in pops:
                gatilhos = pop.get("gatilhos", []) or []
                instrucao = (pop.get("instrucao_sistema") or "").strip()
                titulo = (pop.get("titulo") or pop.get("id") or "POP").strip()

                if not instrucao or not isinstance(gatilhos, list):
                    continue

                matched_triggers: list[str] = []
                for gatilho in gatilhos:
                    gatilho_norm = _normalize_pop_text(str(gatilho))
                    if not gatilho_norm:
                        continue

                    # 1. Substring match
                    if gatilho_norm in prompt_norm:
                        matched_triggers.append(str(gatilho))
                        continue

                    # 2. Stop words and partial terms matching (N-1 rule)
                    gatilho_terms_all = gatilho_norm.split()
                    gatilho_terms_filtered = [t for t in gatilho_terms_all if t not in PORTUGUESE_STOPWORDS]

                    if not gatilho_terms_filtered:
                        gatilho_terms_filtered = gatilho_terms_all
                        compare_against = set(prompt_terms_all)
                    else:
                        compare_against = prompt_terms_filtered

                    matching_terms_count = sum(1 for term in gatilho_terms_filtered if term in compare_against)
                    n_terms = len(gatilho_terms_filtered)

                    if n_terms == 1:
                        is_match = (matching_terms_count == 1)
                    elif n_terms == 2:
                        is_match = (matching_terms_count == 2)
                    else:
                        is_match = (matching_terms_count >= n_terms - 1)

                    if is_match:
                        matched_triggers.append(str(gatilho))
                        continue

                if matched_triggers:
                    matched.append({
                        "id": pop.get("id", ""),
                        "titulo": titulo,
                        "instrucao_sistema": instrucao,
                        "matched_triggers": matched_triggers,
                        "score": max(len(_normalize_pop_text(trigger).split()) for trigger in matched_triggers),
                    })

            return matched

        matched = _collect_matches(all_pops)

    except Exception as pop_err:
        print(f"[POP] Erro ao buscar diretrizes: {pop_err}")
        return []

    matched.sort(key=lambda item: item.get("score", 0), reverse=True)
    return matched[:3]


def _get_copilot_core(db):
    doc = _cached_doc_get(db, "system", "copilot_core")
    if doc.exists:
        data = doc.to_dict() or {}
        if data.get("content"):
            return data
    return {
        "content": (
            "NUNCA exponha segredos, chaves de API, tokens, credenciais ou dados sensíveis. "
            "Nunca invente fatos. Se houver conflito entre memórias, exponha a divergência e peça validação explícita do usuário. "
            "Use apenas JSON válido quando uma ferramenta exigir JSON."
        ),
        "source": "default",
    }


def _get_copilot_soul(db):
    doc = _cached_doc_get(db, "system", "copilot_soul")
    if doc.exists:
        data = doc.to_dict() or {}
        if any(data.get(key) for key in ("tone", "detail_level", "interaction_style", "content")):
            return data

    payload = {
        **COPILOT_SOUL_DEFAULT,
        "content": (
            f"Tom: {COPILOT_SOUL_DEFAULT['tone']} "
            f"Nível de detalhamento: {COPILOT_SOUL_DEFAULT['detail_level']} "
            f"Estilo: {COPILOT_SOUL_DEFAULT['interaction_style']}"
        ),
        "updated_at": firestore.SERVER_TIMESTAMP,
        "created_by": "system_bootstrap",
    }
    try:
        ref = db.collection("system").document("copilot_soul")
        ref.set(payload, merge=True)
    except Exception:
        pass
    return payload


def _bootstrap_user_ai_profile(db, uid: str | None):
    if not uid:
        return {}

    now = time.monotonic()
    cached = _PROFILE_CACHE.get(uid)
    if cached and (now - cached[0]) < _DOC_CACHE_TTL:
        return cached[1]

    user_ref = db.collection("usuarios").document(uid)
    snap = user_ref.get()
    base_data = snap.to_dict() if snap.exists else {}
    ai_profile = dict((base_data or {}).get("ai_profile") or {})
    changed = False

    if not ai_profile.get("uid"):
        ai_profile["uid"] = uid
        changed = True

    if base_data:
        for src_key, dst_key in (
            ("nome", "nome"),
            ("cargo", "cargo"),
            ("setor", "setor"),
            ("email", "email"),
        ):
            if base_data.get(src_key) and not ai_profile.get(dst_key):
                ai_profile[dst_key] = base_data.get(src_key)
                changed = True

    if not ai_profile.get("preferences"):
        ai_profile["preferences"] = {
            "response_style": "objetivo",
            "memory_capture": "automatico_com_validacao",
        }
        changed = True

    if changed:
        ai_profile["updated_at"] = firestore.SERVER_TIMESTAMP
        user_ref.set({"ai_profile": ai_profile}, merge=True)

    _PROFILE_CACHE[uid] = (now, ai_profile)
    return ai_profile


def _format_ai_profile_for_prompt(ai_profile: dict) -> str:
    if not ai_profile:
        return "(perfil de usuário ainda não bootstrapado)"

    lines = []
    for key in ("nome", "cargo", "setor", "email"):
        value = ai_profile.get(key)
        if value:
            lines.append(f"- {key}: {value}")

    prefs = ai_profile.get("preferences") or {}
    if prefs:
        lines.append(f"- preferencias: {json.dumps(prefs, ensure_ascii=False)}")

    history = ai_profile.get("historico_deduzido") or []
    if history:
        lines.append(f"- historico_deduzido: {json.dumps(history[:5], ensure_ascii=False)}")

    # Perfil de personalidade destilado semanalmente a partir do diário pessoal
    # (functions/personal_diary.py:consolidar_personalidade) — impressões, não fatos.
    personalidade = ai_profile.get("personalidade")
    if personalidade:
        lines.append(f"- personalidade (impressões, não fatos relatados): {json.dumps(personalidade, ensure_ascii=False)}")

    return "\n".join(lines) if lines else "(perfil sem atributos relevantes)"


def _find_similar_memory_nodes(db, fato: str, api_key: str, limit: int = 5):
    fato = (fato or "").strip()
    if not fato:
        return []

    try:
        query_embedding = list(map(float, get_embedding(fato, api_key=api_key, task_type="RETRIEVAL_QUERY")))
    except Exception as exc:
        print(f"[Memoria] Falha ao gerar embedding da consulta: {exc}")
        return []

    # Caminho rápido: vector search nativo do Firestore (exige índice vetorial
    # composto tipo+embedding em knowledge_nodes e embeddings gravados como
    # Vector — ver scripts/backfill_embeddings_vector.py).
    try:
        from google.cloud.firestore_v1.vector import Vector
        from google.cloud.firestore_v1.base_vector_query import DistanceMeasure
        from google.cloud.firestore_v1.base_query import FieldFilter

        snaps = (
            db.collection("knowledge_nodes")
            .where(filter=FieldFilter("tipo", "in", sorted(MEMORY_NODE_TYPES)))
            .find_nearest(
                vector_field="embedding",
                query_vector=Vector(query_embedding),
                distance_measure=DistanceMeasure.COSINE,
                limit=max(limit, 5),
                distance_result_field="vector_distance",
            )
            .get()
        )
        candidates = []
        for snap in snaps:
            data = snap.to_dict() or {}
            distance = data.pop("vector_distance", None)
            if distance is None:
                continue
            candidates.append({
                "id": snap.id,
                "similarity": 1.0 - float(distance),
                "data": data,
            })
        if candidates:
            candidates.sort(key=lambda item: item["similarity"], reverse=True)
            return candidates[:limit]
    except Exception as exc:
        print(f"[Memoria] find_nearest indisponível, usando varredura local: {exc}")

    # Fallback: varredura local (índice ausente ou embeddings legados em formato lista)
    candidates = []
    for snap in db.collection("knowledge_nodes").limit(200).stream():
        data = snap.to_dict() or {}
        if data.get("tipo") not in MEMORY_NODE_TYPES:
            continue
        node_embedding = data.get("embedding")
        if not node_embedding:
            continue
        try:
            similarity = _cosine_similarity(query_embedding, list(node_embedding))
        except Exception:
            continue
        candidates.append({
            "id": snap.id,
            "similarity": similarity,
            "data": data,
        })

    candidates.sort(key=lambda item: item["similarity"], reverse=True)
    return candidates[:limit]


def _build_memory_context(db, api_key: str, user_prompt: str, limit: int = 4) -> str:
    candidates = _find_similar_memory_nodes(db, user_prompt, api_key, limit=limit)
    if not candidates:
        return ""

    lines = []
    for item in candidates:
        data = item["data"]
        content = (
            data.get("texto_memoria")
            or data.get("resumo")
            or data.get("titulo")
            or ""
        ).strip()
        if not content:
            continue
        tipo = data.get("tipo", "fato_isolado")
        lines.append(
            f"- [{tipo}] {content[:350]} (id={item['id']}, similaridade={item['similarity']:.3f})"
        )
    if not lines:
        return ""
    return "[MEMÓRIAS GLOBAIS RELEVANTES]\n" + "\n".join(lines)


def _save_user_profile_signal(db, uid: str | None, prompt_text: str, task_id: str | None, system_id: str | None):
    if not uid or not prompt_text:
        return
    profile_ref = db.collection("usuarios").document(uid)
    doc = profile_ref.get()
    ai_profile = dict((doc.to_dict() or {}).get("ai_profile") or {})
    history = list(ai_profile.get("historico_deduzido") or [])
    snippet = prompt_text.strip().replace("\n", " ")
    if len(snippet) > 180:
        snippet = snippet[:177] + "..."
    entry = {
        "texto": snippet,
        "taskId": task_id or None,
        "systemId": system_id or None,
        "at": _iso_now_utc(),
    }
    if not history or history[-1].get("texto") != entry["texto"]:
        history.append(entry)
    ai_profile["historico_deduzido"] = history[-10:]
    ai_profile["last_task_id"] = task_id or ai_profile.get("last_task_id")
    ai_profile["last_system_id"] = system_id or ai_profile.get("last_system_id")
    ai_profile["updated_at"] = firestore.SERVER_TIMESTAMP
    profile_ref.set({"ai_profile": ai_profile}, merge=True)


def _should_consider_memory_by_heuristic(fato: str, categoria: str) -> tuple[bool, str]:
    text = (fato or "").strip()
    lower = text.lower()
    categoria_norm = _normalize_memory_category(categoria)

    if len(text) < MEMORY_MIN_FACT_LENGTH:
        return False, "too_short"

    noisy_prefixes = (
        "ok",
        "obrigado",
        "valeu",
        "bom dia",
        "boa tarde",
        "boa noite",
        "segue",
        "pode fazer",
        "confirmo",
        "cancelar",
        "teste",
    )
    if lower in noisy_prefixes:
        return False, "casual_reply"

    noisy_substrings = (
        "kkkk",
        "haha",
        "rsrs",
        "por favor",
        "obrigad",
        "bom trabalho",
        "isso mesmo",
        "pode seguir",
    )
    if any(token in lower for token in noisy_substrings):
        return False, "small_talk"

    transient_signals = (
        "hoje",
        "amanhã",
        "ontem",
        "agora",
        "daqui a pouco",
        "nesta conversa",
        "nessa conversa",
        "neste chat",
        "anexo",
        "arquivo enviado",
        "mensagem acima",
    )
    if categoria_norm == "fato_isolado" and any(token in lower for token in transient_signals):
        return False, "transient_context"

    durable_signals = (
        "sempre",
        "nunca",
        "regra",
        "procedimento",
        "padrão",
        "preferência",
        "preferencia",
        "convencao",
        "convenção",
        "fonte de verdade",
        "passamos a operar",
        "a partir de agora",
        "deve ser",
        "deve usar",
    )
    if categoria_norm == "regra_global" and any(token in lower for token in durable_signals):
        return True, "durable_rule_signal"

    return True, "heuristic_pass"


def _classify_memory_candidate(api_key: str, fato: str, categoria: str) -> dict:
    heuristic_ok, heuristic_reason = _should_consider_memory_by_heuristic(fato, categoria)
    if not heuristic_ok:
        return {
            "should_save": False,
            "reason": heuristic_reason,
            "confidence": 0.95,
            "normalized_category": _normalize_memory_category(categoria),
        }

    try:
        client = get_genai_module().Client(api_key=api_key)
        prompt = (
            "Você é um filtro de retenção cognitiva do Hermes.\n"
            "Decida se um fato deve ser salvo como memória global de longo prazo.\n"
            "Salve APENAS itens duráveis: regras de negócio, preferências operacionais estáveis, convenções permanentes ou fatos reutilizáveis.\n"
            "NÃO salve: small talk, confirmações momentâneas, contexto transitório, conteúdo efêmero de uma conversa, mensagens vagas ou ruído.\n"
            "Retorne APENAS JSON válido no formato:\n"
            "{\"should_save\":true|false,\"reason\":\"...\",\"confidence\":0.0,\"normalized_category\":\"regra_global|fato_isolado\"}\n\n"
            f"Categoria proposta: {_normalize_memory_category(categoria)}\n"
            f"Fato candidato: {fato}"
        )
        response = client.models.generate_content(
            model="gemini-3.5-flash-lite",
            contents=prompt
        )
        raw_text = (response.text or "").strip()
        json_match = re.search(r"\{.*\}", raw_text, re.DOTALL)
        parsed = json.loads(json_match.group(0) if json_match else raw_text)
        return {
            "should_save": bool(parsed.get("should_save")),
            "reason": str(parsed.get("reason") or "llm_decision"),
            "confidence": float(parsed.get("confidence") or 0.0),
            "normalized_category": _normalize_memory_category(parsed.get("normalized_category") or categoria),
        }
    except Exception as exc:
        print(f"[Memoria] Fallback no classificador de retenção: {exc}")
        return {
            "should_save": True,
            "reason": heuristic_reason,
            "confidence": 0.51,
            "normalized_category": _normalize_memory_category(categoria),
        }


def _save_memory_node(
    db,
    api_key: str,
    fato: str,
    categoria: str,
    session_id: str | None = None,
    user_uid: str | None = None,
    force_update_id: str | None = None,
):
    fato = (fato or "").strip()
    if not fato:
        return {"status": "ignored", "reason": "empty_fact"}

    categoria_norm = _normalize_memory_category(categoria)
    fato_norm = _normalize_text_for_match(fato)
    embedding = list(map(float, get_embedding(fato, api_key=api_key, task_type="RETRIEVAL_DOCUMENT")))
    now = _iso_now_utc()

    if force_update_id:
        ref = db.collection("knowledge_nodes").document(force_update_id)
        snap = ref.get()
        if not snap.exists:
            return {"status": "error", "reason": "memory_not_found", "memory_id": force_update_id}
        current = snap.to_dict() or {}
        title = (fato[:72] + "...") if len(fato) > 72 else fato
        ref.set({
            "id": force_update_id,
            "titulo": title,
            "tipo": categoria_norm,
            "resumo": fato[:600],
            "texto_memoria": fato,
            "embedding": FsVector(embedding),
            "data_atualizacao": now,
            "origem_memoria": "copiloto",
            "ultima_sessao_id": session_id,
            "ultimo_usuario_id": user_uid,
            "memoria_status": "ativa",
            "updated_from_conflict": True,
        }, merge=True)
        return {
            "status": "updated",
            "memory_id": force_update_id,
            "categoria": categoria_norm,
            "previous_text": current.get("texto_memoria") or current.get("resumo") or "",
        }

    candidates = _find_similar_memory_nodes(db, fato, api_key, limit=3)
    best = candidates[0] if candidates else None
    if best and best["similarity"] >= MEMORY_SIMILARITY_CREATE_THRESHOLD:
        best_data = best["data"]
        best_text = _normalize_text_for_match(best_data.get("texto_memoria") or best_data.get("resumo") or "")
        if best["similarity"] >= MEMORY_SIMILARITY_DUPLICATE_THRESHOLD or best_text == fato_norm:
            ref = db.collection("knowledge_nodes").document(best["id"])
            ref.set({
                "data_atualizacao": now,
                "ultima_sessao_id": session_id,
                "ultimo_usuario_id": user_uid,
                "ultimo_fato_observado": fato,
            }, merge=True)
            return {
                "status": "ignored",
                "reason": "duplicate",
                "memory_id": best["id"],
                "categoria": best_data.get("tipo") or categoria_norm,
                "similarity": round(best["similarity"], 4),
            }

        return {
            "status": "conflict",
            "memory_id": best["id"],
            "categoria_existente": best_data.get("tipo") or categoria_norm,
            "existing_text": best_data.get("texto_memoria") or best_data.get("resumo") or "",
            "similarity": round(best["similarity"], 4),
            "proposed_text": fato,
        }

    memory_id = str(uuid.uuid4())[:12]
    title = (fato[:72] + "...") if len(fato) > 72 else fato
    db.collection("knowledge_nodes").document(memory_id).set({
        "id": memory_id,
        "titulo": title,
        "tipo": categoria_norm,
        "resumo": fato[:600],
        "texto_memoria": fato,
        "embedding": FsVector(embedding),
        "area_tematica": "GLOBAL",
        "n_tasks": 0,
        "task_ids": [],
        "data_criacao": now,
        "data_atualizacao": now,
        "origem_memoria": "copiloto",
        "ultima_sessao_id": session_id,
        "ultimo_usuario_id": user_uid,
        "memoria_status": "ativa",
    })
    return {
        "status": "saved",
        "memory_id": memory_id,
        "categoria": categoria_norm,
    }


def _format_pending_memory_conflict(conflict_data: dict | None) -> str:
    if not conflict_data:
        return ""
    existing_text = (conflict_data.get("existing_text") or "").strip()
    proposed_text = (conflict_data.get("proposed_text") or "").strip()
    if not existing_text or not proposed_text:
        return ""
    return (
        "[CONFLITO DE MEMÓRIA PENDENTE]\n"
        f"- memoria_id: {conflict_data.get('memory_id', '')}\n"
        f"- versao_existente: {existing_text}\n"
        f"- versao_proposta: {proposed_text}\n"
        "Se o usuário decidir, use resolver_conflito_memoria(memoria_id, decisao, fato_atualizado, categoria).\n"
        "decisao = 'manter_existente' ou 'substituir_pelo_novo'.\n"
    )


@scheduler_fn.on_schedule(
    schedule="30 20 * * *",
    timezone="America/Sao_Paulo",
    memory=options.MemoryOption.MB_512,
    timeout_sec=60,
)
def relatorio_diario_custo_gemini(event: scheduler_fn.ScheduledEvent):
    """Watchdog de custo: resume o uso diário de tokens Gemini (system_usage/gemini)
    e envia no Telegram, com alerta se estourar o orçamento (system/cost_controls)."""
    db = get_db()
    try:
        from datetime import datetime as _dt, timezone as _tz

        day = _dt.now(_tz.utc).strftime("%Y-%m-%d")
        usage_doc = (
            db.collection("system_usage").document("gemini")
            .collection("daily").document(day).get()
        )
        if not usage_doc.exists:
            print(f"[CustoGemini] Sem uso registrado em {day}.")
            return
        data = usage_doc.to_dict() or {}

        budget = 5.0
        try:
            cfg = _cached_doc_get(db, "system", "cost_controls")
            if cfg.exists:
                budget = float((cfg.to_dict() or {}).get("daily_budget_usd", budget))
        except Exception:
            pass

        cost = float(data.get("estimated_usd") or 0.0)
        calls = int(data.get("calls") or 0)
        tokens = data.get("tokens") or {}
        features = data.get("features") or {}
        top_features = sorted(
            features.items(),
            key=lambda kv: int((kv[1] or {}).get("tokens_total") or 0),
            reverse=True,
        )[:3]

        lines = [
            f"💰 <b>Uso Gemini — {day}</b>",
            f"Custo estimado: <b>${cost:.2f}</b> (orçamento: ${budget:.2f})",
            f"Chamadas: {calls} | Tokens: {int(tokens.get('total') or 0):,}".replace(",", "."),
            f"Entrada: {int(tokens.get('input') or 0):,} | Saída: {int(tokens.get('output') or 0):,}".replace(",", "."),
        ]
        if top_features:
            lines.append("Top consumidores:")
            for name, fdata in top_features:
                lines.append(f"  • {name}: {int((fdata or {}).get('tokens_total') or 0):,} tokens".replace(",", "."))
        if cost > budget:
            lines.insert(0, "⚠️ <b>ORÇAMENTO DIÁRIO ESTOURADO</b>")

        message = "\n".join(lines)
        print(f"[CustoGemini] {message}")

        chat_id = _get_allowed_chat_id()
        if not chat_id:
            keys = (_cached_doc_get(db, "system", "api_keys").to_dict() or {})
            chat_id = keys.get("telegram_chat_id") or keys.get("allowed_telegram_chat_id")
        if chat_id:
            _send_telegram_message(_get_telegram_token(db), chat_id, message)
        else:
            print("[CustoGemini] Nenhum chat_id do Telegram configurado; resumo apenas nos logs.")
    except Exception as exc:
        print(f"[CustoGemini] Falha no relatório diário: {exc}")


@scheduler_fn.on_schedule(
    schedule="0 4 * * *",
    timezone="America/Sao_Paulo",
    memory=options.MemoryOption.MB_512,
    timeout_sec=180,
)
def consolidar_memorias_copiloto(event: scheduler_fn.ScheduledEvent):
    db = get_db()
    try:
        keys_doc = _cached_doc_get(db, "system", "api_keys")
        gemini_key = keys_doc.to_dict().get("gemini_api_key") if keys_doc.exists else None
        if not gemini_key:
            print("[Memoria] gemini_api_key indisponível; consolidação ignorada.")
            return

        client = get_genai_module().Client(api_key=gemini_key)
        nodes = []
        for snap in db.collection("knowledge_nodes").stream():
            data = snap.to_dict() or {}
            if data.get("tipo") not in MEMORY_NODE_TYPES:
                continue
            if data.get("memoria_status") == "consolidada":
                continue
            updated_at = data.get("data_atualizacao") or data.get("data_criacao") or ""
            nodes.append({"id": snap.id, "data": data, "updated_at": updated_at})

        processed_ids = set()
        merges = 0
        for current in nodes:
            if current["id"] in processed_ids:
                continue
            current_text = (current["data"].get("texto_memoria") or current["data"].get("resumo") or "").strip()
            if not current_text:
                continue
            similar = _find_similar_memory_nodes(db, current_text, gemini_key, limit=4)
            merge_candidates = []
            for candidate in similar:
                if candidate["id"] == current["id"]:
                    continue
                if candidate["similarity"] < 0.975:
                    continue
                candidate_text = (candidate["data"].get("texto_memoria") or candidate["data"].get("resumo") or "").strip()
                if not candidate_text:
                    continue
                merge_candidates.append({
                    "id": candidate["id"],
                    "tipo": candidate["data"].get("tipo") or current["data"].get("tipo") or "fato_isolado",
                    "texto": candidate_text,
                })

            if not merge_candidates:
                continue

            group = [{
                "id": current["id"],
                "tipo": current["data"].get("tipo") or "fato_isolado",
                "texto": current_text,
            }] + merge_candidates[:2]

            prompt = (
                "Você é um curador cognitivo do Hermes. Receberá memórias quase duplicadas.\n"
                "Una as memórias em UMA versão consolidada, removendo redundância e preservando a regra/fato mais útil.\n"
                "Retorne APENAS JSON válido no formato:\n"
                "{\"titulo\":\"...\",\"texto_memoria\":\"...\",\"tipo\":\"regra_global|fato_isolado\",\"ids_fundidos\":[\"id1\",\"id2\"]}\n\n"
                f"Memórias:\n{json.dumps(group, ensure_ascii=False)}"
            )
            response = client.models.generate_content(
                model="gemini-3.5-flash-lite",
                contents=prompt
            )
            raw_text = (response.text or "").strip()
            json_match = re.search(r"\{.*\}", raw_text, re.DOTALL)
            merged = json.loads(json_match.group(0) if json_match else raw_text)
            fused_ids = [mid for mid in merged.get("ids_fundidos", []) if isinstance(mid, str)]
            if current["id"] not in fused_ids:
                fused_ids.insert(0, current["id"])

            keep_id = fused_ids[0]
            merged_text = (merged.get("texto_memoria") or current_text).strip()
            merged_tipo = _normalize_memory_category(merged.get("tipo"))
            merged_title = (merged.get("titulo") or merged_text[:72] or "Memória global").strip()
            merged_embedding = list(map(float, get_embedding(merged_text, api_key=gemini_key, task_type="RETRIEVAL_DOCUMENT")))

            db.collection("knowledge_nodes").document(keep_id).set({
                "titulo": merged_title[:80],
                "tipo": merged_tipo,
                "texto_memoria": merged_text,
                "resumo": merged_text[:600],
                "embedding": FsVector(merged_embedding),
                "data_atualizacao": _iso_now_utc(),
                "consolidado_em": firestore.SERVER_TIMESTAMP,
                "origem_curadoria": "llm_cron",
            }, merge=True)

            for drop_id in fused_ids[1:]:
                db.collection("knowledge_nodes").document(drop_id).set({
                    "memoria_status": "consolidada",
                    "consolidado_no_id": keep_id,
                    "data_atualizacao": _iso_now_utc(),
                }, merge=True)
                processed_ids.add(drop_id)

            processed_ids.add(keep_id)
            merges += max(0, len(fused_ids) - 1)

        print(f"[Memoria] Consolidação concluída. merges={merges}")
    except Exception as exc:
        print(f"[Memoria] Falha na consolidação: {exc}")


def _resolve_telegram_chat_id_for_uid(db, uid: str | None):
    if not uid:
        return None
    try:
        user_doc = db.collection("usuarios").document(uid).get()
        if not user_doc.exists:
            return None
        data = user_doc.to_dict() or {}
        for key in ("telegram_chat_id", "telegramChatId", "chat_id", "telegram_id"):
            value = data.get(key)
            if isinstance(value, (str, int)) and str(value).strip():
                return str(value).strip()
    except Exception as exc:
        print(f"[Telegram] Falha ao resolver telegram_chat_id para uid={uid}: {exc}")
    return None


def _resolve_default_telegram_chat_id(db):
    # 1. Tenta encontrar chat_id na coleção de usuários
    try:
        candidates = (
            db.collection("usuarios")
            .where("telegram_chat_id", "!=", None)
            .limit(1)
            .stream()
        )
        for doc in candidates:
            data = doc.to_dict() or {}
            value = data.get("telegram_chat_id")
            if isinstance(value, (str, int)) and str(value).strip():
                return str(value).strip()
    except Exception as exc:
        print(f"[Telegram] Falha ao resolver chat_id via usuarios: {exc}")

    # 2. Fallback Firestore: system/api_keys (mesmo doc do bot token)
    for key in ("telegram_chat_id", "telegram_allowed_chat_id", "allowed_telegram_chat_id"):
        try:
            keys_doc = _cached_doc_get(db, "system", "api_keys")
            if keys_doc.exists:
                value = (keys_doc.to_dict() or {}).get(key)
                if isinstance(value, (str, int)) and str(value).strip():
                    return str(value).strip()
        except Exception:
            pass

    # 3. Fallback Firestore: configuracoes/geral
    for key in ("telegram_chat_id", "telegram_allowed_chat_id"):
        try:
            cfg_doc = _cached_doc_get(db, "configuracoes", "geral")
            if cfg_doc.exists:
                value = (cfg_doc.to_dict() or {}).get(key)
                if isinstance(value, (str, int)) and str(value).strip():
                    return str(value).strip()
        except Exception:
            pass

    # 4. Fallback variável de ambiente
    env_chat_id = os.environ.get("ALLOWED_TELEGRAM_CHAT_ID")
    if env_chat_id and str(env_chat_id).strip():
        return str(env_chat_id).strip()

    return None


def _send_telegram_message_raw(db, chat_id: str | int | None, text: str):
    if not chat_id or not text:
        return False
    try:
        import requests

        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        bot_token = keys_doc.to_dict().get('telegram_bot_token') if keys_doc.exists else None
        bot_token = bot_token or os.environ.get('TELEGRAM_BOT_TOKEN')
        if not bot_token:
            print("[Telegram] TELEGRAM_BOT_TOKEN nao configurado.")
            return False

        resp = requests.post(
            f"https://api.telegram.org/bot{bot_token}/sendMessage",
            json={"chat_id": str(chat_id), "text": text},
            timeout=10
        )
        if not resp.ok:
            print(f"[Telegram] Falha ao enviar lembrete: {resp.status_code} {resp.text[:300]}")
            return False
        return True
    except Exception as exc:
        print(f"[Telegram] Erro ao enviar lembrete: {exc}")
        return False


def _send_telegram_message_raw_with_keyboard(db, chat_id: str | int | None, text: str, inline_keyboard: list | None):
    if not chat_id or not text:
        return False
    try:
        import requests

        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        bot_token = keys_doc.to_dict().get('telegram_bot_token') if keys_doc.exists else None
        bot_token = bot_token or os.environ.get('TELEGRAM_BOT_TOKEN')
        if not bot_token:
            print("[Telegram] TELEGRAM_BOT_TOKEN nao configurado.")
            return False

        payload = {"chat_id": str(chat_id), "text": text}
        if inline_keyboard:
            payload["reply_markup"] = {"inline_keyboard": inline_keyboard}

        resp = requests.post(
            f"https://api.telegram.org/bot{bot_token}/sendMessage",
            json=payload,
            timeout=10
        )
        if not resp.ok:
            print(f"[Telegram] Falha ao enviar lembrete com botoes: {resp.status_code} {resp.text[:300]}")
            return False
        return True
    except Exception as exc:
        print(f"[Telegram] Erro ao enviar lembrete com botoes: {exc}")
        return False


def _build_task_reminder_telegram_message(task: dict, reminder_iso: str | None):
    title = (task.get('titulo') or 'Ação pendente').strip()
    status = (task.get('status') or 'não informado').strip()
    descricao = (task.get('descricao') or '').strip()
    reminder_label = ''
    if reminder_iso:
        try:
            reminder_dt = datetime.fromisoformat(str(reminder_iso))
            reminder_label = reminder_dt.strftime('%d/%m/%Y às %H:%M')
        except Exception:
            reminder_label = str(reminder_iso)

    plan_items = task.get('plano_acao') or []
    pending_steps = []
    for item in plan_items:
        if not isinstance(item, dict):
            continue
        text = str(item.get('text') or '').strip()
        if text and not item.get('completed'):
            pending_steps.append(text)
        if len(pending_steps) >= 3:
            break

    lines = [
        "Lembrete de ação",
        f"Título: {title}",
    ]

    if reminder_label:
        lines.append(f"Agendado para: {reminder_label}")
    lines.append(f"Status atual: {status}")

    if descricao:
        resumo = descricao if len(descricao) <= 220 else f"{descricao[:217]}..."
        lines.append(f"Contexto: {resumo}")

    if pending_steps:
        lines.append("Próximas etapas:")
        for idx, step in enumerate(pending_steps, start=1):
            lines.append(f"{idx}. {step}")
    else:
        lines.append("Próximas etapas: revise a ação e defina o próximo passo operacional.")

    return "\n".join(lines)


def _normalize_task_reminders(task: dict):
    reminders = task.get('reminders') or []
    normalized = []
    if isinstance(reminders, list):
        for idx, reminder in enumerate(reminders):
            if not isinstance(reminder, dict):
                continue
            reminder_at = reminder.get('reminder_at')
            if not reminder_at:
                continue
            normalized.append({
                'id': str(reminder.get('id') or f"legacy-{idx}"),
                'reminder_at': str(reminder_at),
                'reminder_sent': bool(reminder.get('reminder_sent')),
                'created_at': str(reminder.get('created_at') or reminder_at),
                # Preserva o texto do lembrete — sem isso, `check_and_send_reminders` (abaixo,
                # linha ~3060) sempre cai no fallback genérico ao montar o lembrete do Google
                # Tasks, mesmo quando o lembrete carrega uma mensagem específica.
                'message': str(reminder.get('message') or '').strip(),
            })

    if not normalized and task.get('reminder_at'):
        normalized.append({
            'id': 'legacy-reminder',
            'reminder_at': str(task.get('reminder_at')),
            'reminder_sent': bool(task.get('reminder_sent')),
            'created_at': str(task.get('data_atualizacao') or task.get('data_criacao') or task.get('reminder_at')),
            'message': str(task.get('reminder_message') or '').strip(),
        })

    normalized.sort(key=lambda item: item.get('reminder_at') or '')
    return normalized


def _build_task_reminder_state_payload(reminders: list[dict]):
    ordered = sorted(reminders, key=lambda item: item.get('reminder_at') or '')
    next_pending = next((item for item in ordered if not item.get('reminder_sent')), None)
    return {
        'reminders': ordered,
        'reminder_at': next_pending.get('reminder_at') if next_pending else None,
        'reminder_sent': bool(next_pending.get('reminder_sent')) if next_pending else True,
    }

@https_fn.on_call(

    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),

    memory=options.MemoryOption.GB_2,

    timeout_sec=540

)

def processarArquivoIA(req: https_fn.CallableRequest):

    """Callable para disparar processamento manual"""

    item_id = req.data.get('itemId')

    if not item_id:

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message="ID do item é obrigatório")

    

    db = get_db()

    doc = db.collection('conhecimento').document(item_id).get()

    if not doc.exists:

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.NOT_FOUND, message="Arquivo não encontrado")

    

    # Limpa campos antigos para mostrar o loader no front se necessário e garantir re-processamento

    db.collection('conhecimento').document(item_id).update({

        'resumo_tldr': None,

        'tags': None

    })



    return start_file_indexing(item_id, doc.to_dict())

@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=120,
)
def criar_formulario_google(req: https_fn.CallableRequest):
    """
    Cria um formulário no Google Forms a partir de um schema JSON validado
    (Draft-and-Approve) gerado pelo Copiloto e confirmado pelo usuário.
    Garante permissão pública por padrão via Google Drive API.
    """
    from googleapiclient.discovery import build

    uid = req.auth.uid if req.auth else None
    if not uid:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.UNAUTHENTICATED,
            message="Usuário não autenticado."
        )

    data = req.data or {}
    session_id = data.get('sessionId')
    form_data = data.get('form')

    if not form_data or not isinstance(form_data, dict):
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="O campo 'form' é obrigatório e deve ser um objeto JSON."
        )

    titulo = form_data.get('titulo')
    if not titulo:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="O formulário precisa ter um título ('titulo')."
        )

    descricao = form_data.get('descricao', '')
    perguntas = form_data.get('perguntas', [])

    try:
        creds = get_google_creds(GOOGLE_FORMS_SCOPES)
        forms_service = build('forms', 'v1', credentials=creds)
        drive_service = build('drive', 'v3', credentials=creds)

        # 1. Cria o formulário inicial vazio (retorna um ID)
        new_form = {
            "info": {
                "title": titulo,
                "documentTitle": titulo
            }
        }
        form_result = forms_service.forms().create(body=new_form).execute()
        form_id = form_result.get('formId')
        responder_uri = form_result.get('responderUri')

        # 2. Configura a permissão do Drive para acesso público a respondentes
        try:
            drive_service.permissions().create(
                fileId=form_id,
                body={
                    'type': 'anyone',
                    'role': 'reader'
                }
            ).execute()
        except Exception as perm_err:
            print(f"[Forms] Aviso: Falha ao setar permissão pública via Drive API: {perm_err}")
            # Não falhamos toda a operação apenas por erro de permissão (ex: restrições do workspace)

        # 3. Monta as requisições em lote para popular o formulário
        requests = []

        # Adiciona a descrição
        if descricao:
            requests.append({
                "updateFormInfo": {
                    "info": {
                        "description": descricao
                    },
                    "updateMask": "description"
                }
            })

        # Processa as perguntas
        for i, q in enumerate(perguntas):
            tipo = q.get('tipo', 'texto_curto')
            texto = q.get('texto', 'Pergunta sem título')
            obrigatoria = bool(q.get('obrigatoria', False))

            question_item = {
                "question": {
                    "required": obrigatoria,
                }
            }

            if tipo == 'texto_curto':
                question_item["question"]["textQuestion"] = {"paragraph": False}
            elif tipo == 'paragrafo':
                question_item["question"]["textQuestion"] = {"paragraph": True}
            elif tipo in ['multipla_escolha', 'caixas_selecao', 'lista_suspensa']:
                opcoes = q.get('opcoes', ['Opção 1'])
                options_list = [{"value": opt} for opt in opcoes]

                choice_type = "RADIO"
                if tipo == 'caixas_selecao':
                    choice_type = "CHECKBOX"
                elif tipo == 'lista_suspensa':
                    choice_type = "DROP_DOWN"

                question_item["question"]["choiceQuestion"] = {
                    "type": choice_type,
                    "options": options_list
                }
            elif tipo == 'escala_linear':
                question_item["question"]["scaleQuestion"] = {
                    "low": q.get('escala_min', 1),
                    "high": q.get('escala_max', 5),
                    "lowLabel": q.get('rotulo_min', ''),
                    "highLabel": q.get('rotulo_max', '')
                }
            else:
                # Fallback
                question_item["question"]["textQuestion"] = {"paragraph": False}

            requests.append({
                "createItem": {
                    "item": {
                        "title": texto,
                        "questionItem": question_item
                    },
                    "location": {
                        "index": i
                    }
                }
            })

        if requests:
            forms_service.forms().batchUpdate(
                formId=form_id,
                body={"requests": requests}
            ).execute()

        return {"formId": form_id, "responderUri": responder_uri}

    except Exception as e:
        print(f"Erro em criar_formulario_google: {e}")
        import traceback
        print(traceback.format_exc())
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )
@https_fn.on_call(memory=options.MemoryOption.MB_512, timeout_sec=30)
def corrigir_sintaxe_mermaid(req: https_fn.CallableRequest):
    """
    Recebe código Mermaid inválido e o erro do compilador JS, retorna apenas o código corrigido.
    Chamado pelo self-healing loop do componente MermaidBlock no frontend (máx. 3 tentativas).
    """
    from google import genai
    from google.genai import types

    uid = req.auth.uid if req.auth else None
    if not uid:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.UNAUTHENTICATED,
            message="Usuário não autenticado."
        )

    data = req.data or {}
    codigo_mermaid = (data.get('codigoMermaid') or '').strip()
    erro_compilador = (data.get('erroCompilador') or '').strip()

    if not codigo_mermaid:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Campo 'codigoMermaid' é obrigatório."
        )

    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None

        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INTERNAL,
                message="Chave Gemini não configurada."
            )

        client = genai.Client(api_key=gemini_key)

        system_instruction = (
            "Você é um especialista em sintaxe Mermaid. "
            "Receberá um código Mermaid inválido e o erro gerado pelo compilador. "
            "Retorne APENAS o código Mermaid corrigido — sem blocos de markdown (não inclua ```mermaid), "
            "sem explicações, sem texto adicional. "
            "Mantenha o tipo de diagrama original. Corrija somente a sintaxe."
        )

        user_message = (
            f"Código Mermaid inválido:\n{codigo_mermaid}\n\n"
            f"Erro do compilador:\n{erro_compilador}"
        )

        response = client.models.generate_content(
            model='gemini-3.5-flash-lite',
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                temperature=0.1,
                max_output_tokens=2048,
            ),
            contents=user_message,
        )

        codigo_corrigido = (response.text or '').strip()
        # Strip markdown wrapper defensively in case the model added it anyway
        if codigo_corrigido.startswith('```mermaid'):
            codigo_corrigido = codigo_corrigido[10:].strip()
        if codigo_corrigido.startswith('```'):
            codigo_corrigido = codigo_corrigido[3:].strip()
        if codigo_corrigido.endswith('```'):
            codigo_corrigido = codigo_corrigido[:-3].strip()

        return {'codigoCorrigido': codigo_corrigido}

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em corrigir_sintaxe_mermaid: {e}")
        import traceback
        print(traceback.format_exc())
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )


@https_fn.on_call(memory=options.MemoryOption.GB_1)

def processInvoiceOCR(req: https_fn.CallableRequest):

    """

    Processa uma Nota Fiscal (PDF/Imagem) do Google Drive usando Gemini e extrai dados estruturados.

    """

    from google import genai

    import json

    import re



    file_id = req.data.get('fileId')

    if not file_id:

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message="fileId é obrigatório.")



    try:

        db = get_db()

        keys_doc = _cached_doc_get(db, 'system', 'api_keys')

        GEMINI_API_KEY = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None



        if not GEMINI_API_KEY:

             raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION, message="Chave Gemini não configurada.")



        client = genai.Client(api_key=GEMINI_API_KEY)



        # Download from Drive

        service = get_drive_service()

        file_metadata = service.files().get(fileId=file_id, fields='mimeType, name').execute()

        mime_type = file_metadata.get('mimeType')



        request = service.files().get_media(fileId=file_id)

        file_content = request.execute()



        prompt = """

        Analise este documento (Nota Fiscal ou Recibo) e extraia os seguintes dados em formato JSON estrito:

        {

            "fornecedor": "Nome da Empresa",

            "cnpj": "XX.XXX.XXX/0001-XX",

            "data_emissao": "YYYY-MM-DD",

            "valor_total": 0.00,

            "itens": [

                {

                    "descricao": "Nome do Produto",

                    "quantidade": 1,

                    "valor_unitario": 0.00,

                    "valor_total": 0.00

                }

            ]

        }

        Se algum campo não for encontrado, retorne null ou lista vazia.

        Normalize a data para ISO 8601.

        Normalize valores numéricos para float (ponto flutuante).

        """



        parts = [{"mime_type": mime_type, "data": file_content}, prompt]



        response = client.models.generate_content(model="gemini-3.5-flash-lite", contents=parts)

        res_text = response.text



        # Clean Markdown code blocks if present

        json_match = re.search(r'```json\s*(.*?)\s*```', res_text, re.DOTALL)

        if json_match:

            json_str = json_match.group(1)

        else:

            json_str = res_text



        data = json.loads(json_str)

        return data



    except Exception as e:

        print(f"Erro no OCR de Nota Fiscal: {str(e)}")

        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INTERNAL, message=str(e))


@https_fn.on_call()
def transcrever_audio(req: https_fn.CallableRequest):
    """
    Recebe áudio em Base64, transcreve com Groq (Whisper-Large-V3-Turbo) e refina com Gemini.
    """
    import base64
    import tempfile
    import os
    # Instale: pip install groq google-genai
    from groq import Groq
    from google import genai

    data = req.data
    audio_base64 = data.get('audioBase64')

    if not audio_base64:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Áudio não fornecido."
        )

    # Buscar chaves de API no Firestore
    try:
        # Import local (dentro da função ou escopo global se preferir, mas aqui segue o padrão do fix)
        # Assumindo que get_db já existe no main.py. Mas o fix_main.py injeta apenas ESTA função?
        # Sim, ele injeta `transcrever_audio`.
        # Precisamos garantir que `get_db` esteja disponível ou usar firestore.client() direto?
        # O arquivo main.py tem `from firebase_admin import firestore` e `initialize_app`.
        # Melhor usar `firestore.client()` diretamente para garantir, já que `get_db` é custom.
        # Mas `main.py` tem `get_db` definido no topo. Vamos usar `get_db()` para consistência,
        # assumindo que o `fix_main.py` insere isso num arquivo que tem `get_db`.
        
        # Como não temos certeza se `get_db` está acessível no escopo (python é permissivo),
        # vamos usar o padrão seguro: importar firestore.
        from firebase_admin import firestore
        db = firestore.client()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        
        if not keys_doc.exists:
             raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION, 
                message="Chaves de API não configuradas (system/api_keys)."
            )
            
        keys = keys_doc.to_dict()
        GROQ_API_KEY = keys.get('groq_api_key')
        GEMINI_API_KEY = keys.get('gemini_api_key')
        
    except Exception as e:
        print(f"Erro ao buscar chaves: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message="Erro interno de configuração."
        )

    temp_path = None
    try:
        # 1. Converter Base64 para arquivo temporário
        # b64decode retorna bytes
        try:
            audio_bytes = base64.b64decode(audio_base64)
        except Exception:
             raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="Base64 inválido."
            )
        
        # Cria um arquivo temporário físico para o Groq poder ler
        # O sufixo .m4a é importante para o ffmpeg interno do whisper identificar o formato se necessário
        with tempfile.NamedTemporaryFile(suffix=".m4a", delete=False) as temp_file:
            temp_file.write(audio_bytes)
            temp_path = temp_file.name

        # 2. Transcrição via Groq
        client = Groq(api_key=GROQ_API_KEY)
        
        with open(temp_path, "rb") as file_stream:
            transcription = client.audio.transcriptions.create(
                file=(os.path.basename(temp_path), file_stream), 
                model="whisper-large-v3-turbo",
                response_format="json",
                language="pt",
                temperature=0.0
            )

        texto_bruto = transcription.text

        # 3. Refinamento via Gemini Flash
        gemini_client = genai.Client(api_key=GEMINI_API_KEY)

        prompt = f"""
        Atue como um redator especialista. O texto a seguir é uma transcrição de voz bruta.
        Sua tarefa:
        1. Corrigir pontuação e gramática (pt-BR).
        2. Remover vícios de linguagem (né, tipo, ahn).
        3. Manter o tom original e termos técnicos.
        4. Retorne APENAS o texto corrigido, sem introduções.

        Texto: "{texto_bruto}"
        """

        response = gemini_client.models.generate_content(model="gemini-3.5-flash-lite", contents=prompt)
        texto_refinado = response.text

        return {
            "raw": texto_bruto,
            "refined": texto_refinado
        }

    except Exception as e:
        print(f"Erro na transcrição: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=f"Falha ao processar áudio: {str(e)}"
        )
    finally:
        # Limpeza
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except OSError:
                pass


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.GB_1,
    timeout_sec=120
)
def askTaskAssistant(req: https_fn.CallableRequest):
    """
    Responde perguntas sobre o contexto de uma tarefa específica baseando-se no diário de bordo.
    Injeta contexto do Grafo de Conhecimento (RAG Dinâmica) com citações inline [N].
    """
    from google import genai
    from firebase_admin import firestore

    data = req.data or {}
    prompt = data.get('prompt')
    history_context = data.get('historyContext')
    area_tematica = data.get('area_tematica')
    rag_context_id = data.get('ragContext')
    extra_context_id = data.get('extraContextId')
    knowledge_item_ids = data.get('knowledgeItemIds', [])
    kg_tags = data.get('kgTags', [])

    if not isinstance(prompt, str) or not prompt.strip():
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="O comando é obrigatório."
        )

    try:
        db = get_db()

        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None

        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini não configurada."
            )

        client = genai.Client(api_key=gemini_key)

        # --- CONHECIMENTO MESTRE (Manual do André) ---
        manual_context = ""
        if area_tematica:
            try:
                master_docs = db.collection('conhecimento_mestre')\
                    .where('area_tematica', '==', area_tematica)\
                    .order_by('data_criacao', direction=firestore.Query.DESCENDING)\
                    .limit(3).stream()
                manual_items = []
                for m_doc in master_docs:
                    m = m_doc.to_dict()
                    manual_items.append(f"GUIA: {m.get('titulo')}\nCONTEÚDO:\n{m.get('conteudo')}")
                if manual_items:
                    manual_context = "\n\n".join(manual_items)
            except Exception as e:
                print(f"Erro ao recuperar conhecimento mestre: {e}")

        # --- BASE RAG PRINCIPAL da ação ---
        rag_context = ""
        if rag_context_id and rag_context_id != "Nenhum":
            try:
                rag_context = retrieve_personalized_rag_context(db, genai, prompt, rag_context_id)
            except Exception as e:
                print(f"Erro ao recuperar RAG principal: {e}")

        # --- DOCUMENTOS EXTRAS da ação ---
        extra_rag_context = ""
        if extra_context_id or knowledge_item_ids:
            try:
                extra_rag_context = retrieve_extra_context_rag(db, genai, prompt, extra_context_id, knowledge_item_ids)
            except Exception as e:
                print(f"Erro ao recuperar contexto extra: {e}")

        # --- DADOS FINANCEIROS (Módulo Financeiro) ---
        finance_context = ""
        if area_tematica == 'FINANCEIRO':
            try:
                from datetime import datetime
                import json
                from tools.telegram_extended import execute
                
                now = datetime.now()
                mes_atual = now.month - 1
                ano_atual = now.year
                
                fin_data_atual = execute("consultar_financas_v2", {"mes": mes_atual, "ano": ano_atual}, db)
                data_json = json.loads(fin_data_atual)
                
                resumo = data_json.get("resumo", {})
                metas = data_json.get("metas", [])
                reserva = data_json.get("reserva_emergencia", {})
                detalhes = data_json.get("detalhes", {})
                
                parts = []
                parts.append("=== DADOS FINANCEIROS ATUAIS ===")
                parts.append(f"Período: {mes_atual+1}/{ano_atual}")
                parts.append(f"Renda Total: R$ {resumo.get('total_renda', 0):.2f} (Recebida: R$ {resumo.get('renda_recebida', 0):.2f})")
                parts.append(f"Contas Totais: R$ {resumo.get('total_contas', 0):.2f} (Pagas: R$ {resumo.get('contas_pagas', 0):.2f})")
                parts.append(f"Saldo Previsto: R$ {resumo.get('saldo_previsto', 0):.2f}")
                parts.append(f"Saldo Atual: R$ {resumo.get('saldo_atual', 0):.2f}")
                
                parts.append("\nReserva de Emergência:")
                parts.append(f"- Alvo: R$ {reserva.get('alvo', 0):.2f}")
                parts.append(f"- Atual: R$ {reserva.get('atual', 0):.2f}")
                
                parts.append("\nMetas Financeiras:")
                for meta in metas:
                    parts.append(f"- {meta.get('name')}: R$ {meta.get('targetAmount', 0):.2f} (Prioridade: {meta.get('priority')})")
                
                parts.append("\nRendas Detalhadas:")
                for r in detalhes.get("rendas", []):
                    status = "Recebido" if r.get("isReceived") else "Pendente"
                    parts.append(f"- {r.get('description')}: R$ {r.get('amount', 0):.2f} ({status})")
                
                parts.append("\nContas Detalhadas:")
                for c in detalhes.get("contas", []):
                    status = "Pago" if c.get("isPaid") else "Pendente"
                    parts.append(f"- {c.get('description')}: R$ {c.get('amount', 0):.2f} ({status})")
                
                trans_docs = db.collection("finance_transactions").where("status", "==", "active").stream()
                avulsas = []
                for doc in trans_docs:
                    d = doc.to_dict()
                    date_str = d.get("date")
                    if date_str:
                        try:
                            dt = datetime.fromisoformat(date_str.replace('Z', '+00:00'))
                            if dt.month - 1 == mes_atual and dt.year == ano_atual:
                                avulsas.append(d)
                        except:
                            pass
                
                if avulsas:
                    parts.append("\nTransações Avulsas do Mês:")
                    for t in avulsas:
                        parts.append(f"- {t.get('description')}: R$ {t.get('amount', 0):.2f} (Sprint {t.get('sprint')})")
                
                finance_context = "\n".join(parts)
            except Exception as e:
                print(f"Erro ao extrair dados financeiros: {e}")

        # --- GRAFO DE CONHECIMENTO (RAG Dinâmica) ---
        kg_context = ""
        kg_nodes_payload = []
        if area_tematica:
            try:
                kg_nodes_payload, kg_context = extract_kg_rag_context(
                    db=db,
                    api_key=gemini_key,
                    area_tematica=area_tematica,
                    tags=kg_tags,
                )
            except Exception as e:
                print(f"Erro ao extrair contexto do grafo de conhecimento: {e}")

        system_instruction = (
            "Você é o HERMES, copiloto de execução de tarefas do André. "
            "Você tem acesso a: (1) o contexto completo da ação (título, descrição, plano e diário), "
            "(2) bases de conhecimento RAG personalizadas, "
            "(3) documentos extras carregados para esta ação, "
            "(4) manuais de procedimento padrão, "
            "(5) contexto operacional do Grafo de Conhecimento com procedimentos passados. "
            "Ao usar informações do Grafo de Conhecimento, cite a fonte com marcadores [1], [2], etc. "
            "Seja executivo, preciso e profissional (pt-BR). "
            "Quando gerar documentos longos (atas, ofícios, pareceres), produza o conteúdo completo e formatado."
            "\n\n"
            "## REGRA ABSOLUTA — SIGLAS, TERMOS TÉCNICOS E ERROS DE BACKEND\n\n"
            "PROIBIÇÃO TOTAL DE EXPANSÃO ARBITRÁRIA DE SIGLAS:\n"
            "Você JAMAIS deve inferir, adivinhar, expandir ou traduzir siglas, acrônimos ou\n"
            "termos técnicos que o usuário fornecer. Se o usuário disser \"IRP\", você trata\n"
            "\"IRP\" como uma string opaca e literal — não é \"Imposto de Renda\", não é\n"
            "\"Internal Revenue Policy\", não é nada que você \"acha que pode ser\". Você passa\n"
            "o termo exatamente como recebido para as ferramentas de busca. Se nenhum\n"
            "documento retornar resultado, sua resposta é: \"Nenhum registro encontrado para\n"
            "o termo exato 'IRP'. Você pode confirmar a sigla ou fornecer mais contexto?\"\n"
            "Não improvise. Não complete. Não alucine.\n\n"
            "OBRIGAÇÃO DE TRANSPARÊNCIA EM ERROS TÉCNICOS:\n"
            "Se qualquer ferramenta (buscar_acervo, buscar_tarefas ou similar) retornar um\n"
            "campo \"erro\" não-nulo, você DEVE reproduzir o conteúdo desse campo palavra por\n"
            "palavra na sua resposta, sem parafrasear, sem suavizar e sem omitir. Formato\n"
            "obrigatório:\n\n"
            "  ⚠️ Erro técnico na ferramenta [nome_da_ferramenta]:\n"
            "  [conteúdo literal do campo \"erro\"]\n\n"
            "Após reportar o erro, peça ao usuário que acione o suporte técnico com essa\n"
            "mensagem exata. Você NÃO deve tentar responder a pergunta original como se o\n"
            "erro não tivesse ocorrido."
        )

        full_prompt = f"""
        === CONTEXTO DA AÇÃO ===
        {history_context if history_context else 'Nenhum registro encontrado.'}

        === BASE RAG PRINCIPAL ===
        {rag_context if rag_context else 'Nenhuma base RAG selecionada.'}

        === DOCUMENTOS EXTRAS DESTA AÇÃO ===
        {extra_rag_context if extra_rag_context else 'Nenhum documento extra.'}

        === MANUAL DE PROCEDIMENTOS ===
        {manual_context if manual_context else 'Nenhum guia mestre para esta area_tematica.'}

        {finance_context if finance_context else ''}

        {kg_context if kg_context else ''}

        === COMANDO DO USUÁRIO ===
        {prompt}
        """

        response = client.models.generate_content(model="gemini-3.5-flash-lite", contents=[system_instruction, full_prompt])

        result = (response.text or "").strip()
        if not result:
            result = "Não consegui gerar uma resposta. Tente reformular o comando."

        return {
            "result": result,
            "kg_nodes": kg_nodes_payload,
        }

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em askTaskAssistant: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=f"Erro ao processar consulta da tarefa: {str(e)}"
        )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.GB_1,
    timeout_sec=120
)
def askChatbot(req: https_fn.CallableRequest):
    """
    Responde perguntas sobre o contexto da reunião usando Gemini.
    """
    from google import genai

    prompt = (req.data or {}).get('prompt')
    if not isinstance(prompt, str) or not prompt.strip():
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Prompt é obrigatório."
        )

    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini não configurada."
            )

        client = genai.Client(api_key=gemini_key)
        response = client.models.generate_content(
            model="gemini-3.5-flash-lite",
            contents=[
                "Você é um assistente de reunião em pt-BR. Responda com objetividade, "
                "baseando-se no contexto recebido. Se o contexto estiver incompleto, "
                "deixe claro que a resposta é parcial.",
                prompt.strip(),
            ]
        )

        result = (response.text or "").strip()
        if not result:
            result = "Não consegui gerar uma resposta com o contexto atual da reunião."
        return {"result": result}

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em askChatbot: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message="Falha ao processar sua solicitação no assistente de reunião."
        )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=120
)
def refinarDiretrizesEstrategicas(req: https_fn.CallableRequest):
    """
    Refina uma intencao estrategica em propostas editaveis para estrategia_pessoal.
    Nao persiste dados: a validacao humana acontece no frontend.
    """
    from google import genai
    from google.genai import types

    data = req.data or {}
    intencao = str(data.get("intencao") or "").strip()
    user_uid = req.auth.uid if req.auth else str(data.get("userId") or "").strip()
    pilares_validos = {"carreira", "financas", "saude", "intelectual", "estilo_vida"}

    if not intencao:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Intencao estrategica e obrigatoria."
        )

    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini nao configurada."
            )

        contexto_previo = []
        if user_uid:
            for doc_snap in db.collection("estrategia_pessoal").where("userId", "==", user_uid).where("status", "==", "ativo").limit(12).stream():
                value = doc_snap.to_dict() or {}
                contexto_previo.append({
                    "pilar": value.get("pilar"),
                    "objetivoMacro": value.get("objetivoMacro"),
                    "tipoMeta": value.get("tipoMeta"),
                    "marcos": value.get("marcos", []),
                    "diretrizesDerivadas": value.get("diretrizesDerivadas", []),
                })

        refinement_prompt = f"""
Voce e um arquiteto de estrategia pessoal. Transforme a intencao livre do usuario em 1 a 5 propostas estruturadas.

REGRAS:
- Responda somente JSON valido.
- Use pilares exatamente entre: carreira, financas, saude, intelectual, estilo_vida.
- tipoMeta deve ser "absoluta" quando houver valor numerico rastreavel; caso contrario "relativa_qualitativa".
- Para metas absolutas, preencha metricaAlvo com valorAtual, valorObjetivo e unidade. Se nao houver valor atual, use 0.
- indicadoresSucesso deve conter sinais continuos/recorrentes de progresso, nao entregas pontuais.
- marcos deve conter 1 a 5 entregas pontuais que possam ser concluidas definitivamente.
- diretrizesDerivadas deve conter 2 a 5 frases concisas para orientar uma IA de forma passiva no chat global.
- Nao crie tarefas operacionais diarias.
- Preserve a autonomia do usuario: escreva como contexto estrategico, nao como ordens intrusivas.

FORMATO:
{{
  "propostas": [
    {{
      "pilar": "carreira",
      "objetivoMacro": "string",
      "tipoMeta": "absoluta",
      "metricaAlvo": {{"valorAtual": 0, "valorObjetivo": 100, "unidade": "string"}},
      "indicadoresSucesso": ["string"],
      "marcos": ["string"],
      "diretrizesDerivadas": ["string"]
    }}
  ]
}}

CONTEXTO ESTRATEGICO JA ATIVO:
{json.dumps(contexto_previo, ensure_ascii=False)}

INTENCAO DO USUARIO:
{intencao}
"""

        client = genai.Client(api_key=gemini_key)
        response = client.models.generate_content(
            model=GEMINI_ROUTING_MODEL,
            contents=refinement_prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                temperature=0.2,
                http_options=types.HttpOptions(timeout=60000),
            )
        )
        log_gemini_usage(response, model=GEMINI_ROUTING_MODEL, feature="estrategia_pessoal_refinar", db=db)
        raw_text = (response.text or "").strip()
        try:
            parsed = json.loads(raw_text)
        except Exception:
            match = re.search(r"\{.*\}", raw_text, re.DOTALL)
            parsed = json.loads(match.group(0)) if match else {}

        propostas = []
        for proposal in parsed.get("propostas", []):
            if not isinstance(proposal, dict):
                continue
            pilar = proposal.get("pilar") if proposal.get("pilar") in pilares_validos else "carreira"
            tipo_meta = "absoluta" if proposal.get("tipoMeta") == "absoluta" else "relativa_qualitativa"
            indicadores = [
                str(item).strip()
                for item in (proposal.get("indicadoresSucesso") or [])
                if str(item).strip()
            ][:5]
            marcos = [
                str(item).strip()
                for item in (proposal.get("marcos") or [])
                if str(item).strip()
            ][:5]
            diretrizes = [
                str(item).strip()
                for item in (proposal.get("diretrizesDerivadas") or [])
                if str(item).strip()
            ][:5]
            item = {
                "pilar": pilar,
                "objetivoMacro": str(proposal.get("objetivoMacro") or "").strip(),
                "tipoMeta": tipo_meta,
                "indicadoresSucesso": indicadores,
                "marcos": marcos,
                "diretrizesDerivadas": diretrizes,
            }
            if tipo_meta == "absoluta":
                metric = proposal.get("metricaAlvo") or {}
                item["metricaAlvo"] = {
                    "valorAtual": float(metric.get("valorAtual") or 0),
                    "valorObjetivo": float(metric.get("valorObjetivo") or 0),
                    "unidade": str(metric.get("unidade") or "").strip(),
                }
            if item["objetivoMacro"] and item["diretrizesDerivadas"]:
                propostas.append(item)

        if not propostas:
            propostas = [{
                "pilar": "intelectual",
                "objetivoMacro": intencao[:220],
                "tipoMeta": "relativa_qualitativa",
                "indicadoresSucesso": [
                    "Definir criterios observaveis de progresso",
                    "Revisar a intencao em ciclos mensais",
                    "Registrar evidencias de avanco"
                ],
                "marcos": [
                    "Definir primeira entrega pontual do objetivo"
                ],
                "diretrizesDerivadas": [
                    "Considere esta intencao apenas quando o usuario pedir alinhamento estrategico.",
                    "Evite transformar esta diretriz em pressao operacional diaria.",
                    "Priorize recomendacoes que preservem foco e consistencia de longo prazo."
                ],
            }]

        return {"propostas": propostas[:5]}

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em refinarDiretrizesEstrategicas: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message="Falha ao refinar diretrizes estrategicas."
        )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=120
)
def assistirPopComIA(req: https_fn.CallableRequest):
    """
    Gera ou refina um POP (Procedimento Operacional Padrao) usando IA.
    modo="criar": recebe texto bruto e gera titulo/gatilhos/instrucao_sistema do zero.
    modo="editar": recebe o POP atual e um pedido de ajuste, retorna a versao atualizada.
    """
    from google import genai
    from google.genai import types

    data = req.data or {}
    modo = str(data.get("modo") or "").strip()
    if modo not in ("criar", "editar"):
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Parametro 'modo' deve ser 'criar' ou 'editar'."
        )

    if modo == "criar":
        texto_bruto = str(data.get("textoBruto") or "").strip()
        if not texto_bruto:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="Texto bruto e obrigatorio para criar um POP."
            )
        prompt = f"""
Voce e um especialista em criar POPs (Procedimentos Operacionais Padrao) que orientam um copiloto de IA institucional a executar ou replicar processos administrativos.

A partir do TEXTO BRUTO abaixo (que pode ser um rascunho, e-mail, anotacoes soltas ou instrucoes desorganizadas), estruture um POP completo.

REGRAS:
- Responda somente JSON valido, sem markdown ao redor do JSON.
- "titulo": curto e descritivo (maximo ~80 caracteres).
- "gatilhos": lista de 3 a 8 palavras ou frases curtas, em minusculas, sem acentuacao obrigatoria mas preferencialmente com, que um usuario diria para ativar este POP.
- "instrucao_sistema": texto completo em markdown, preservando TODOS os passos, regras, numeros, links e detalhes tecnicos presentes no texto bruto. Organize em secoes/etapas quando fizer sentido. Nao invente informacoes que nao estejam no texto bruto nem no bom senso do dominio.

FORMATO:
{{
  "titulo": "string",
  "gatilhos": ["string"],
  "instrucao_sistema": "string"
}}

TEXTO BRUTO:
{texto_bruto}
"""
        feature_tag = "pop_gerar_ia"
    else:
        titulo_atual = str(data.get("tituloAtual") or "").strip()
        instrucao_atual = str(data.get("instrucaoAtual") or "").strip()
        gatilhos_atuais = [str(g).strip() for g in (data.get("gatilhosAtuais") or []) if str(g).strip()]
        pedido_ajuste = str(data.get("pedidoAjuste") or "").strip()
        if not instrucao_atual:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="POP atual e obrigatorio para editar."
            )
        if not pedido_ajuste:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="Descreva o ajuste desejado."
            )
        prompt = f"""
Voce e um especialista em refinar POPs (Procedimentos Operacionais Padrao) que orientam um copiloto de IA institucional.

Abaixo esta o POP ATUAL e um PEDIDO DE AJUSTE do usuario. Aplique exatamente o ajuste solicitado, preservando todo o restante do conteudo que nao foi mencionado no pedido.

REGRAS:
- Responda somente JSON valido, sem markdown ao redor do JSON.
- "titulo": mantenha o atual a menos que o ajuste peca para mudar.
- "gatilhos": lista de 3 a 8 palavras/frases curtas em minusculas; ajuste apenas se o pedido pedir.
- "instrucao_sistema": texto completo e atualizado em markdown, com o ajuste aplicado.

POP ATUAL:
Titulo: {titulo_atual}
Gatilhos: {json.dumps(gatilhos_atuais, ensure_ascii=False)}
Instrucao de sistema:
{instrucao_atual}

PEDIDO DE AJUSTE DO USUARIO:
{pedido_ajuste}

FORMATO:
{{
  "titulo": "string",
  "gatilhos": ["string"],
  "instrucao_sistema": "string"
}}
"""
        feature_tag = "pop_editar_ia"

    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini nao configurada."
            )

        client = genai.Client(api_key=gemini_key)
        response = client.models.generate_content(
            model=GEMINI_ROUTING_MODEL,
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                temperature=0.2,
                http_options=types.HttpOptions(timeout=60000),
            )
        )
        log_gemini_usage(response, model=GEMINI_ROUTING_MODEL, feature=feature_tag, db=db)
        raw_text = (response.text or "").strip()
        try:
            parsed = json.loads(raw_text)
        except Exception:
            match = re.search(r"\{.*\}", raw_text, re.DOTALL)
            parsed = json.loads(match.group(0)) if match else {}

        titulo_final = str(parsed.get("titulo") or "").strip()
        gatilhos_final = [
            str(item).strip().lower()
            for item in (parsed.get("gatilhos") or [])
            if str(item).strip()
        ][:8]
        instrucao_final = str(parsed.get("instrucao_sistema") or "").strip()

        if not titulo_final or not gatilhos_final or not instrucao_final:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INTERNAL,
                message="A IA nao retornou um POP valido. Tente reformular o pedido."
            )

        return {
            "titulo": titulo_final,
            "gatilhos": gatilhos_final,
            "instrucao_sistema": instrucao_final,
        }

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em assistirPopComIA: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message="Falha ao processar o POP com IA."
        )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST", "OPTIONS"]),
    memory=options.MemoryOption.GB_2,
    timeout_sec=COPILOT_FUNCTION_TIMEOUT_SEC
)
def askCopilotoHermes(req: https_fn.CallableRequest):
    """
    Módulo Copiloto Hermes
    Estrategista sênior de processos com Tool Calling e RAG Híbrido.
    """
    from google import genai
    from google.genai import types

    data = req.data or {}
    perf_state = {"start_ms": _perf_now_ms(), "_last_ms": _perf_now_ms(), "steps": [], "tool_calls": []}
    request_start_monotonic = time.monotonic()
    prompt = (data.get('prompt') or "").strip()
    task_id = data.get('taskId')
    task_id_scoped = task_id
    system_id = data.get('systemId')
    session_id = data.get('sessionId')

    drive_file_id = data.get('driveFileId')
    drive_file_name = (data.get('driveFileName') or 'documento').strip()
    drive_files_raw = data.get('driveFiles') or []
    drive_files = []
    if isinstance(drive_files_raw, list) and drive_files_raw:
        for df in drive_files_raw[:10]:
            if isinstance(df, dict) and df.get('driveFileId'):
                drive_files.append({
                    'driveFileId': str(df['driveFileId']).strip(),
                    'driveFileName': str(df.get('driveFileName') or 'documento').strip()
                })
    elif drive_file_id:
        drive_files.append({'driveFileId': str(drive_file_id).strip(), 'driveFileName': drive_file_name})

    routing_index = data.get('routingIndex') or []
    copilot_mode = (data.get('copilotMode') or 'default').strip()
    strategy_directives_raw = data.get('strategyDirectives') or []
    strategy_directives = [
        str(item).strip()
        for item in strategy_directives_raw
        if str(item).strip()
    ][:24] if isinstance(strategy_directives_raw, list) else []
    user_uid = req.auth.uid if req.auth else None

    def _copilot_remaining_sec() -> float:
        elapsed = time.monotonic() - request_start_monotonic
        return max(0.0, COPILOT_SOFT_DEADLINE_SEC - elapsed)

    # Ingestão muda: arquivos sem texto → prompt padrão de catalogação
    if not prompt and (drive_files or drive_file_id):
        prompt = (
            "Analise os arquivos anexados, identifique o que eles mostram e explique "
            "como eles se relacionam com a ação em contexto."
        )

    if not prompt:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Prompt é obrigatório."
        )

    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None

        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini não configurada."
            )

        client = genai.Client(api_key=gemini_key)
        copilot_core = _get_copilot_core(db)
        copilot_soul = _get_copilot_soul(db)
        # Áreas temáticas válidas: o copiloto só pode SELECIONAR uma existente, nunca inventar.
        _areas_validas = carregar_areas_tematicas_validas(db)
        ai_profile = _bootstrap_user_ai_profile(db, user_uid)
        threading.Thread(
            target=_save_user_profile_signal,
            args=(db, user_uid, prompt, task_id, system_id),
            daemon=True,
        ).start()
        _prompt_words = [w for w in prompt.lower().split() if len(w) > 2]
        memory_context = _build_memory_context(db, gemini_key, prompt, limit=4) if len(_prompt_words) >= 4 else ""
        matched_pop_directives = _match_pop_directives(db, prompt)
        _perf_mark(perf_state, "web.bootstrap")

        tools_routing_context = ""
        if routing_index:
            tools_routing_context = "\n\n## CATÁLOGO DE FERRAMENTAS DISPONÍVEIS\n"
            tools_routing_context += "Você pode acionar ferramentas interativas na interface do usuário usando a função `acionar_ferramenta`. "
            tools_routing_context += "Não peça permissão para usar as ferramentas, acione-as diretamente se a intenção do usuário bater com as chaves (keys) ou a tag abaixo:\n\n"
            for t in routing_index:
                tools_routing_context += f"- ID da Ferramenta: {t['id']}\n"
                tools_routing_context += f"  Tag Explícita: {t.get('tag', '')}\n"
                tools_routing_context += f"  Chaves de Ativação (Keys): {', '.join(t.get('keys', []))}\n\n"
            tools_routing_context += "ATENÇÃO: Caso o usuário use uma Tag Explícita (ex: @SipacTrackingTool) mas o texto do pedido exija algo totalmente diferente (ex: 'Gerar nota fiscal'), seja juiz semântico, suspenda o acionamento e alerte-o sobre a contradição.\n"

        session_conflict_context = ""
        session_ref = None
        if session_id:
            session_ref = db.collection('sessoes_copiloto').document(session_id)
            try:
                session_doc = session_ref.get()
                if session_doc.exists:
                    session_conflict_context = _format_pending_memory_conflict(
                        (session_doc.to_dict() or {}).get("pendingMemoryConflict")
                    )
            except Exception as session_err:
                print(f"[Memoria] Falha ao ler conflito pendente da sessão {session_id}: {session_err}")

        def _set_copilot_status(status_text: str | None):
            """Status efêmero no doc da sessão — o frontend exibe via onSnapshot
            enquanto a resposta é processada. None limpa o campo."""
            if session_ref is None:
                return
            try:
                if status_text:
                    session_ref.set({
                        "copilotStatus": status_text,
                        "copilotStatusAt": firestore.SERVER_TIMESTAMP,
                    }, merge=True)
                else:
                    session_ref.update({
                        "copilotStatus": firestore.DELETE_FIELD,
                        "copilotStatusAt": firestore.DELETE_FIELD,
                    })
            except Exception:
                pass

        _set_copilot_status("Analisando sua solicitação...")

        # --- DEFINIÇÃO DE FERRAMENTAS ---
        _perf_mark(perf_state, "web.session_context")
        def consultar_historico_acoes(query: str, area_tematica: str = None, data_limite_inicio: str = None, data_limite_fim: str = None, ultimas_n_acoes: int = 20, status: str = None):
            """
            Busca tarefas reais no banco de dados do Hermes por texto, area, prazo e/ou status.
            Retorna somente dados oficiais — nao mistura com RAG ou procedimentos.
            Use status para filtrar (ex: 'em andamento', 'concluido', 'cancelado').
            Use data_limite_inicio e data_limite_fim (YYYY-MM-DD) para filtrar por prazo.
            Use ultimas_n_acoes (default 20) para buscar em lote.
            """
            from tools.busca_grafo import buscar_tarefas
            _STOPWORDS_Q = {"de", "a", "o", "que", "e", "do", "da", "em", "um", "uma",
                            "os", "as", "no", "na", "com", "por", "para", "dos", "das",
                            "nos", "nas", "ao", "os", "se", "ou"}
            _q_terms = [w for w in query.lower().split() if w not in _STOPWORDS_Q and len(w) > 2]
            _initial_mode = "all" if len(_q_terms) >= 2 else "any"
            res_exact = buscar_tarefas(query,
                                       area_tematica=area_tematica,
                                       match_mode=_initial_mode,
                                       data_limite_inicio=data_limite_inicio,
                                       data_limite_fim=data_limite_fim,
                                       status=status,
                                       limite=ultimas_n_acoes)
            # Tenta com match_mode=any se all nao encontrou
            if _initial_mode == "all" and not res_exact.get("resultados"):
                res_exact = buscar_tarefas(query,
                                           area_tematica=area_tematica,
                                           match_mode="any",
                                           data_limite_inicio=data_limite_inicio,
                                           data_limite_fim=data_limite_fim,
                                           status=status,
                                           limite=ultimas_n_acoes)

            if res_exact.get("erro"):
                return f"⚠️ [ERRO TÉCNICO BuscaGrafo] {res_exact['erro']}"

            resultados = res_exact.get("resultados", [])

            # CAMINHO A: Tarefas reais encontradas — retorna SOMENTE elas.
            # Nao inclui contexto semantico para evitar mistura de fontes (alucinacao de titulo/dados).
            if resultados:
                lines = [
                    "=== TAREFAS REAIS ENCONTRADAS NO BANCO DE DADOS ===",
                    "REGRA: Use EXCLUSIVAMENTE os campos abaixo. Nao invente, nao complete, nao use RAG.",
                    "",
                ]
                for r in resultados:
                    lines.append(f"ID: {r['id']}")
                    lines.append(f"Titulo: {r['titulo']}")
                    lines.append(f"Status: {r['status']} | Tipo: {r.get('tipo_acao') or 'nao informado'}")
                    lines.append(f"Prazo: {r.get('data_limite', 'N/A')} | Area: {r['area']} | Criado em: {r['criado_em']}")
                    if r.get('processo_sei'):
                        lines.append(f"Processo SEI: {r['processo_sei']}")
                    lines.append(f"Responsavel: {r.get('responsavel') or 'nao informado'}")
                    if r.get('tags'):
                        tags_val = r['tags']
                        tags_str = ', '.join(tags_val) if isinstance(tags_val, list) else str(tags_val)
                        lines.append(f"Tags: {tags_str}")
                    if r.get('sintese_demanda'):
                        lines.append(f"Sintese da Demanda: {r['sintese_demanda']}")
                    if r.get('descricao'):
                        lines.append(f"Descricao: {r['descricao']}")
                    if r.get('notas'):
                        lines.append(f"Notas: {r['notas']}")
                    plano = r.get('plano_acao', [])
                    if plano:
                        lines.append("Plano de Acao:")
                        for passo in plano:
                            lines.append(f"  {passo}")
                    acomp = r.get('acompanhamento_recente', [])
                    if acomp:
                        lines.append("Diario de Bordo (ultimas entradas):")
                        for entrada in acomp:
                            lines.append(f"  {entrada}")
                    lines.append(f"[Abrir Acao](task:{r['id']})")
                    lines.append("---")
                return "\n".join(lines)

            # CAMINHO B: Nenhuma tarefa real encontrada.
            # Retorna mensagem direta sem fallback semantico.
            # O modelo NAO deve inventar dados nem usar RAG para compensar.
            filtros_desc = []
            if query and query.strip():
                filtros_desc.append(f"query='{query.strip()}'")
            if status:
                filtros_desc.append(f"status='{status}'")
            if area_tematica:
                filtros_desc.append(f"area='{area_tematica}'")
            if data_limite_inicio or data_limite_fim:
                filtros_desc.append(f"prazo=[{data_limite_inicio or '*'} a {data_limite_fim or '*'}]")
            filtros_str = ", ".join(filtros_desc) if filtros_desc else "(sem filtros)"
            return (
                f"NENHUMA TAREFA ENCONTRADA no banco de dados com os filtros: {filtros_str}.\n"
                "INSTRUCAO OBRIGATORIA: Informe ao usuario que nao encontrou. NAO invente titulos, "
                "status, prazos ou qualquer dado de tarefa. NAO use RAG, acervo ou memoria para fabricar uma resposta."
            )

        def buscar_conversas_whatsapp(query: str, limite: int = 5):
            """Busca conversas de WhatsApp indexadas (digests) por similaridade semântica. Use quando o usuário perguntar sobre algo discutido no WhatsApp."""
            from whatsapp_ingest import buscar_conversas_whatsapp as _buscar_whatsapp
            res = _buscar_whatsapp(db, query, limite)
            if res.get("erro"):
                return f"⚠️ {res['erro']}"
            resultados = res.get("resultados", [])
            if not resultados:
                return "Nenhuma conversa de WhatsApp indexada encontrada para esta busca."
            linhas = []
            for r in resultados:
                topicos = ", ".join(r.get("topicos") or [])
                chat_id_info = f" (chat_id: {r.get('chat_id')})" if r.get("chat_id") else ""
                linhas.append(f"- [{r.get('chat_name')}]{chat_id_info} {r.get('resumo')}" + (f" (tópicos: {topicos})" if topicos else ""))
            return "\n".join(linhas)

        def buscar_arquivos_acervo(query: str):
            """Busca documentação, manuais e arquivos de referência no Acervo Global (FindNearest)."""
            from tools.busca_acervo import buscar_acervo
            res = buscar_acervo(query)
            if res.get("erro"):
                return f"⚠️ [ERRO TÉCNICO FindNearest] {res['erro']}"
            
            resultados = res.get("resultados", [])
            if not resultados:
                return "Nenhum documento encontrado no acervo global para esta busca."

            lines = []
            for r in resultados:
                # Rastreabilidade: expõe origem por tarefa, incluindo drive_file_id quando disponível
                origem_raw = r.get('origem', {})
                if isinstance(origem_raw, dict):
                    modulo = origem_raw.get('modulo', '')
                    id_origem = origem_raw.get('id_origem', '')
                    if modulo == 'tarefa' and id_origem:
                        origem_label = f"Tarefa {id_origem} (task_id={r.get('task_id', 'N/A')})"
                    elif r.get('task_id'):
                        origem_label = f"Tarefa {r['task_id']}"
                    else:
                        origem_label = 'Acervo Global'
                elif r.get('task_id'):
                    origem_label = f"Tarefa {r['task_id']}"
                else:
                    origem_label = r.get('origem', 'Acervo Global')

                url_part = f" | LINK: {r['url_drive']}" if r.get('url_drive') else ""
                drive_id_part = f" | DRIVE_FILE_ID: {r['drive_file_id']}" if r.get('drive_file_id') else ""
                lines.append(
                    f"DOC: {r['titulo']} | ORIGEM: {origem_label} | FONTE: {r['fonte']}{url_part}{drive_id_part}\n"
                    f"TRECHO: {r['trecho']}"
                )
            return "\n\n".join(lines)

        def obter_contexto_tela(id_tarefa: str):
            """Obtém o contexto completo da tarefa em foco, incluindo diário integral, plano de ação e arquivos disponíveis para leitura profunda."""
            if not id_tarefa:
                return "Nenhuma tarefa em foco no momento."
            try:
                import re as _re
                _DRIVE_ID_RE = _re.compile(r'/d/([a-zA-Z0-9_-]{10,})')

                doc_snap = db.collection('tarefas').document(id_tarefa).get()
                if not doc_snap.exists:
                    return "Tarefa não identificada no banco de dados."
                t = doc_snap.to_dict()
                
                # Diário Integral
                diario_full = []
                def _obter_data_ordenacao(entry):
                    d = entry.get('data')
                    if not d:
                        return ""
                    if isinstance(d, str):
                        return d
                    import datetime as _datetime
                    if isinstance(d, _datetime.datetime):
                        if d.tzinfo is None:
                            d = d.replace(tzinfo=_datetime.timezone.utc)
                        return d.isoformat()
                    if isinstance(d, _datetime.date):
                        return d.isoformat()
                    if hasattr(d, 'isoformat'):
                        try:
                            return d.isoformat()
                        except Exception:
                            pass
                    return str(d)

                for e in sorted(t.get('acompanhamento', []), key=_obter_data_ordenacao):
                    d_val = e.get('data')
                    if not d_val:
                        d_str = ""
                    elif isinstance(d_val, str):
                        d_str = d_val
                    elif hasattr(d_val, "strftime"):
                        try:
                            d_str = d_val.strftime("%Y-%m-%d %H:%M:%S")
                        except Exception:
                            d_str = str(d_val)
                    else:
                        d_str = str(d_val)
                    diario_full.append(f"[{d_str}] {e.get('nota')}")

                # Mapeamento de arquivos para leitura profunda on-demand
                # Retrocompatibilidade: tenta drive_file_id direto; se ausente, extrai da URL via regex
                # Valida existência no Drive — arquivos deletados/na lixeira são omitidos do contexto
                _ctx_drive = get_drive_service()
                arquivos_disponiveis = []
                for item in t.get('pool_dados', []):
                    if item.get('tipo') != 'arquivo':
                        continue
                    fid = item.get('drive_file_id')
                    if not fid:
                        url_val = item.get('valor', '')
                        match = _DRIVE_ID_RE.search(url_val)
                        fid = match.group(1) if match else None
                    if not fid:
                        continue
                    try:
                        _meta = _ctx_drive.files().get(fileId=fid, fields='id,trashed').execute()
                        if _meta.get('trashed'):
                            continue
                    except Exception:
                        continue
                    arquivos_disponiveis.append({
                        "nome": item.get('nome', 'Arquivo sem nome'),
                        "drive_file_id": fid
                    })
                
                context = {
                    "id": id_tarefa,
                    "titulo": t.get('titulo'),
                    "area_tematica": t.get('area_tematica'),
                    "sistema_id": t.get('sistema_id') or None,
                    "plano_atual": t.get('plano_acao', []),
                    "diario_integral": "\n".join(diario_full),
                    "tags": t.get('tags', []),
                    "arquivos_disponiveis": arquivos_disponiveis
                }
                return json.dumps(context, indent=2, ensure_ascii=False)
            except Exception as e:
                return f"Erro ao obter contexto da tela: {e}"

        def pesquisar_internet(query: str):
            """
            Busca informações recentes, notícias ou fatos atualizados na internet.
            Use quando o usuário precisar de dados em tempo real, cotações, eventos recentes
            ou qualquer informação que possa estar desatualizada no seu conhecimento.
            Parâmetro: query — a frase de busca otimizada em português ou inglês.
            """
            _prompt_lower = (prompt or "").lower()
            _web_triggers = (
                "http://", "https://", "www.", "internet", "na web", "busca online",
                "pesquise", "pesquisar", "notícia", "noticias", "cotação", "cotacao",
                "atualiz", "link", "site", "acesse",
            )
            if not any(t in _prompt_lower for t in _web_triggers):
                return '{"blocked": true, "reason": "O prompt não menciona internet, URL ou busca atual. Use esta ferramenta apenas quando o usuário pedir explicitamente informações da web."}'
            import requests as _req
            try:
                keys_doc_web = _cached_doc_get(db, 'system', 'api_keys')
                tavily_key = keys_doc_web.to_dict().get('tavily_api_key') if keys_doc_web.exists else None
                if not tavily_key:
                    return '{"error": "Tavily API key não configurada. Informe ao usuário que a busca na internet está indisponível no momento."}'

                resp = _req.post(
                    "https://api.tavily.com/search",
                    json={
                        "api_key": tavily_key,
                        "query": query,
                        "search_depth": "advanced",
                        "include_answer": True,
                        "include_raw_content": False,
                        "max_results": 5
                    },
                    timeout=20
                )
                resp.raise_for_status()
                data = resp.json()

                parts = []
                if data.get("answer"):
                    parts.append(f"RESPOSTA DIRETA: {data['answer']}\n")
                for r in data.get("results", []):
                    parts.append(
                        f"FONTE: {r.get('title', '')} ({r.get('url', '')})\n"
                        f"{r.get('content', '')}"
                    )
                return "\n\n".join(parts) if parts else "Nenhum resultado encontrado para esta busca."

            except _req.exceptions.Timeout:
                return '{"error": "Timeout ao acessar a Tavily API. Informe ao usuário que a busca demorou demais e tente novamente."}'
            except Exception as web_err:
                return f'{{"error": "Falha na busca: {str(web_err)}. Informe ao usuário que não foi possível realizar a pesquisa."}}'
        def consultar_processo_sipac_copiloto(numero_processo: str):
            """
            Consulta e retorna informações detalhadas de um processo no SIPAC,
            incluindo dados gerais, interessados, movimentações recentes e os documentos anexados
            (com seus respectivos links de visualização pública, se disponíveis).
            Use esta ferramenta sempre que o usuário fornecer um número de processo do SIPAC ou pedir
            informações sobre o trâmite, status ou documentos de um processo específico.
            """
            try:
                from hermes_core_logic import _call_web_callable
                print(f"[Copiloto] Consultando processo SIPAC: {numero_processo}")
                res = _call_web_callable(
                    function_name="consultarProcessoSipac",
                    data={"numeroProcesso": numero_processo},
                    user_uid=user_uid
                )
                
                lines = []
                lines.append(f"=== DETALHES DO PROCESSO SIPAC {res.get('numeroProcesso')} ===")
                lines.append(f"Status: {res.get('status')}")
                lines.append(f"Unidade Atual: {res.get('unidadeAtual')}")
                lines.append(f"Natureza: {res.get('natureza')}")
                lines.append(f"Assunto: {res.get('assuntoCodigo')} - {res.get('assuntoDescricao')}")
                if res.get('observacao') and res.get('observacao') != 'Não informado':
                    lines.append(f"Observação: {res.get('observacao')}")
                lines.append(f"Autuação: {res.get('dataAutuacion')} às {res.get('horarioAutuacion')}")
                
                lines.append("\nInteressados:")
                for i in res.get('interessados', []):
                    lines.append(f"- {i.get('tipo')}: {i.get('nome')}")
                    
                lines.append("\nDocumentos Públicos:")
                for d in res.get('documentos', []):
                    url_str = f" | Link: {d.get('url')}" if d.get('url') else " | (Acesso Restrito)"
                    lines.append(f"- Seq #{d.get('ordem')} - Tipo: {d.get('tipo')} | Data: {d.get('data')} | Origem: {d.get('unidadeOrigem')}{url_str}")
                    
                lines.append("\nMovimentações Recentes (Linha do Tempo):")
                for m in res.get('movimentacoes', [])[:8]:
                    lines.append(f"- [{m.get('data')} {m.get('horario')}] De {m.get('unidadeOrigem')} para {m.get('unidadeDestino')} | Recebedor: {m.get('usuarioRecebedor') or 'N/A'}")
                    
                return "\n".join(lines)
            except Exception as e:
                print(f"[Copiloto] Erro ao consultar SIPAC: {e}")
                return f"⚠️ Erro ao consultar processo {numero_processo} no SIPAC: {str(e)}"
        def acompanhar_processo_sipac_copiloto(numero_processo: str, acompanhar: bool = True):
            """
            Ativa ou desativa o monitoramento/acompanhamento automático de um processo SIPAC no Hermes.
            Ao ativar, o sistema fará verificações periódicas em background e notificará
            o usuário no Telegram sempre que houver alguma alteração ou novos documentos.
            Parâmetros:
            - numero_processo: número do processo SIPAC.
            - acompanhar: True para monitorar (padrão), False para parar de monitorar.
            """
            from hermes_core_logic import _call_web_callable
            from firebase_admin import firestore
            from datetime import datetime, timezone
            import re as _re
            
            try:
                print(f"[Copiloto] Acompanhar processo SIPAC: {numero_processo} -> {acompanhar}")
                res = _call_web_callable(
                    function_name="consultarProcessoSipac",
                    data={"numeroProcesso": numero_processo},
                    user_uid=user_uid
                )

                clean_num = _re.sub(r'[^\d]', '', numero_processo)
                doc_id = f"{user_uid}_{clean_num}" if user_uid else f"global_{clean_num}"
                ref = db.collection('sipac_processos').document(doc_id)
                
                ref.set({
                    "acompanhar": acompanhar,
                    "numeroProcesso": res.get("numeroProcesso", numero_processo),
                    "uid": user_uid or "global",
                    "ultimaConsulta": datetime.now(timezone.utc).isoformat(),
                    **res
                }, merge=True)
                
                status_str = "ATIVADO" if acompanhar else "DESATIVADO"
                return f"Sucesso: O acompanhamento automático para o processo {numero_processo} foi {status_str}."
            except Exception as e:
                print(f"[Copiloto] Erro ao alterar acompanhamento SIPAC: {e}")
                return f"⚠️ Erro ao alterar acompanhamento para o processo {numero_processo}: {str(e)}"

        def incorporar_documento_especifico_sipac_no_rag_da_acao(numero_processo: str, sequencial: int, task_id: str = None):
            """
            Busca e incorpora um documento público específico de um processo SIPAC diretamente no RAG da ação ativa.
            O documento é identificado pelo seu número sequencial (ordem) no processo.
            Parâmetros:
            - numero_processo: número do processo SIPAC
            - sequencial: número de ordem (sequencial) do documento (ex: 1, 12, 15) que foi identificado via consulta.
            - task_id: ID da ação/tarefa. Se omitido, usa o taskId da ação ativa.
            """
            from hermes_core_logic import _call_web_callable
            from knowledge_graph import _get_embedding
            from firebase_admin import firestore
            import requests
            
            target_task_id = task_id or task_id_scoped
            if not target_task_id:
                return "ERRO|Nenhuma ação ativa em contexto e nenhum task_id foi fornecido para incorporar o documento."

            try:
                print(f"[Copiloto] Buscando processo SIPAC para RAG específico: {numero_processo}, doc ordem: {sequencial}")
                res = _call_web_callable(
                    function_name="consultarProcessoSipac",
                    data={"numeroProcesso": numero_processo},
                    user_uid=user_uid
                )

                documentos = res.get('documentos', [])
                if not documentos:
                    return f"O processo {numero_processo} foi localizado, mas não contém documentos públicos anexados para incorporação."

                # Localiza o documento pelo sequencial (ordem)
                doc_alvo = None
                for d in documentos:
                    try:
                        if int(d.get('ordem')) == int(sequencial):
                            doc_alvo = d
                            break
                    except (ValueError, TypeError):
                        continue

                if not doc_alvo:
                    return f"Não foi possível encontrar o documento sequencial #{sequencial} no processo {numero_processo}."

                url = doc_alvo.get('url')
                if not url:
                    return f"O documento sequencial #{sequencial} ({doc_alvo.get('tipo')}) possui acesso restrito e não pode ser incorporado."

                task_ref = db.collection('tarefas').document(target_task_id)
                task_snap = task_ref.get()
                if not task_snap.exists:
                    return f"ERRO|Ação {target_task_id} não encontrada no banco de dados."
                
                task_data = task_snap.to_dict() or {}
                pool_dados = task_data.get('pool_dados', []) or []
                existing_urls = {item.get('valor') for item in pool_dados if item.get('valor')}

                nome_doc = f"Seq #{doc_alvo.get('ordem')} - {doc_alvo.get('tipo')}"
                
                print(f"[Copiloto] Baixando documento do SIPAC: {nome_doc} -> {url}")
                resp = requests.get(url, timeout=15)
                if not resp.ok:
                    return f"Falha ao baixar o documento '{nome_doc}': HTTP {resp.status_code}"
                
                file_bytes = resp.content
                
                pdf_result = extract_pdf_text_with_fallback(
                    file_bytes,
                    f"{nome_doc}.pdf",
                    api_key=gemini_key,
                    allow_gemini_fallback=False
                )
                doc_text = (pdf_result.get('text') or '').strip()
                if not doc_text:
                    try:
                        doc_text = file_bytes.decode('utf-8', errors='ignore').strip()
                    except Exception:
                        pass
                
                if not doc_text or len(doc_text) < 50:
                    return f"Documento '{nome_doc}' vazio ou sem texto legível."
                    
                summary_prompt = (
                    f"Você é um analista jurídico sênior. Resuma de forma concisa e técnica "
                    f"o seguinte documento do processo SIPAC '{nome_doc}':\n\n"
                    f"CONTEÚDO:\n{doc_text[:6000]}"
                )
                summary_resp = client.models.generate_content(
                    model="gemini-3.5-flash-lite",
                    contents=summary_prompt
                )
                resumo = (summary_resp.text or "").strip()
                if not resumo:
                    resumo = f"Documento público do processo SIPAC {numero_processo} do tipo {doc_alvo.get('tipo')}."

                import uuid as _uuid
                from datetime import datetime as _dt
                if url not in existing_urls:
                    artefato_id = f"sipac_{_uuid.uuid4().hex[:12]}"
                    pool_dados.append({
                        "id": artefato_id,
                        "nome": f"SIPAC: {nome_doc}",
                        "tipo": "arquivo",
                        "valor": url,
                        "data_criacao": _dt.now().isoformat()
                    })
                    task_ref.update({"pool_dados": pool_dados})
                else:
                    # Se já existe, pega o ID original ou gera um novo apenas para o índice
                    artefato_id = f"sipac_{_uuid.uuid4().hex[:12]}"
                
                embedding = _get_embedding(resumo, gemini_key)

                db.collection('indice_artefatos').document(artefato_id).set({
                    "nome": f"SIPAC: {nome_doc}",
                    "url": url,
                    "tipo_mime": "application/pdf",
                    "resumo_semantico": resumo,
                    "embedding": FsVector(list(map(float, embedding))),
                    "tags": ["SIPAC", doc_alvo.get('tipo', 'Documento')],
                    "origem": "tarefa",
                    "task_id": target_task_id,
                    "acervo_id": None,
                    "texto_bruto": doc_text[:12000],
                    "indexed_at": firestore.SERVER_TIMESTAMP
                })
                

                    
                return (
                    f"Sucesso: O documento '{nome_doc}' do processo SIPAC {numero_processo} "
                    f"foi incorporado com sucesso no RAG da Ação [{(task_data.get('titulo') or 'Ação')}]({target_task_id})."
                )
            except Exception as e:
                print(f"[Copiloto] Erro na incorporação de RAG do SIPAC: {e}")
                return f"ERRO|Erro ao incorporar documento do SIPAC {numero_processo} no RAG: {str(e)}"

        def buscar_e_analisar_email(query: str, max_results: int = 5):
            """
            Busca e analisa e-mails no Gmail usando uma query estruturada.
            Use esta ferramenta quando o usuário pedir para verificar, ler ou analisar e-mails.
            Retorna o texto higienizado e o conteúdo de anexos (PDF, CSV).

            Args:
                query: Query de busca padrão do Gmail (ex: 'from:nome@empresa.com newer_than:2d', 'subject:"reunião"').
                max_results: Número máximo de e-mails a processar (limite: 5).
            """
            import os
            import tempfile
            import base64

            try:
                import html2text
            except ImportError:
                return "⚠️ Dependência html2text não instalada. Avise o desenvolvedor."

            try:
                gs = get_gmail_service()
            except Exception as e:
                return f"⚠️ Erro ao inicializar serviço do Gmail: {str(e)}"

            try:
                results = gs.users().messages().list(userId='me', q=query, maxResults=max_results).execute()
                messages = results.get('messages', [])

                if not messages:
                    return "Nenhum e-mail encontrado para a query informada."

                output = []
                for msg in messages:
                    msg_id = msg['id']
                    full_msg = gs.users().messages().get(userId='me', id=msg_id, format='full').execute()

                    payload = full_msg.get('payload', {})
                    headers = payload.get('headers', [])

                    subject = next((h['value'] for h in headers if h['name'].lower() == 'subject'), 'Sem Assunto')
                    sender = next((h['value'] for h in headers if h['name'].lower() == 'from'), 'Desconhecido')
                    date_str = next((h['value'] for h in headers if h['name'].lower() == 'date'), '')

                    text_parts = []
                    html_parts = []
                    attachments = []

                    def walk_parts(part):
                        mime_type = part.get('mimeType')
                        body = part.get('body', {})
                        data = body.get('data')
                        filename = part.get('filename', '')

                        if filename and body.get('attachmentId'):
                            size = body.get('size', 0)
                            if size < 50000 and mime_type.startswith('image/'):
                                pass
                            elif mime_type in ['application/pdf', 'text/csv']:
                                attachments.append({
                                    'id': body['attachmentId'],
                                    'filename': filename,
                                    'mime_type': mime_type,
                                    'size': size
                                })
                        elif data:
                            if mime_type == 'text/plain':
                                text_parts.append(base64.urlsafe_b64decode(data + '=' * (-len(data) % 4)).decode('utf-8', errors='replace'))
                            elif mime_type == 'text/html':
                                html_parts.append(base64.urlsafe_b64decode(data + '=' * (-len(data) % 4)).decode('utf-8', errors='replace'))

                        if 'parts' in part:
                            for subpart in part['parts']:
                                walk_parts(subpart)

                    walk_parts(payload)

                    clean_text = ""
                    if text_parts:
                        clean_text = "\n".join(text_parts)
                    elif html_parts:
                        h = html2text.HTML2Text()
                        h.ignore_links = False
                        h.ignore_images = True
                        h.body_width = 0
                        clean_text = "\n".join([h.handle(html) for html in html_parts])

                    clean_text = "\n".join([line for line in clean_text.split("\n") if line.strip()])

                    msg_str = f"--- E-MAIL ---\nID: {msg_id}\nDe: {sender}\nAssunto: {subject}\nData: {date_str}\n\n[Corpo do E-mail]\n{clean_text}\n"

                    for att in attachments:
                        att_id = att['id']
                        att_name = att['filename']
                        msg_str += f"\n[Anexo: {att_name}]\n"

                        try:
                            att_obj = gs.users().messages().attachments().get(userId='me', messageId=msg_id, id=att_id).execute()
                            file_data = base64.urlsafe_b64decode(att_obj['data'])

                            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(att_name)[1] if '.' in att_name else '') as temp_file:
                                temp_file.write(file_data)
                                temp_path = temp_file.name

                            try:
                                if att['mime_type'] == 'application/pdf':
                                    pdf_result = extract_pdf_text_with_fallback(
                                        file_data,
                                        att_name,
                                        allow_gemini_fallback=False,
                                        max_pages=10,
                                    )
                                    extracted_text = pdf_result.get('text', '')
                                    msg_str += f"Conteúdo do PDF (Extração):\n{extracted_text[:3000]}\n"
                                elif att['mime_type'] == 'text/csv':
                                    import csv
                                    with open(temp_path, 'r', encoding='utf-8', errors='replace') as f:
                                        reader = csv.reader(f)
                                        lines = [",".join(row) for i, row in enumerate(reader) if i < 50]
                                        msg_str += f"Conteúdo do CSV (Primeiras 50 linhas):\n" + "\n".join(lines) + "\n"
                            finally:
                                if os.path.exists(temp_path):
                                    os.remove(temp_path)
                        except Exception as e_att:
                            msg_str += f"(Erro ao processar anexo {att_name}: {e_att})\n"

                    output.append(msg_str)

                # Indexar implicitamente no Grafo (Nó de Fonte) para o primeiro email (se houver), se a intenção for gerar RAG Node
                # Note: O RAG Node será instanciado na criação de tarefa, ou podemos fazer aqui?
                # A instrução: "o texto do e-mail é uma entidade primária e deve ser processado pelas funções do knowledge_graph.py para gerar um Nó de Fonte"
                # A melhor abordagem é ter uma ferramenta adicional "indexar_email_grafo", ou então retornar o texto do email, e ao criar_acao_no_sistema() a IA passar os campos do email!

                return "\n\n==========================\n\n".join(output)
            except Exception as e:
                return f"⚠️ Erro ao processar e-mails: {str(e)}"

        def ler_pagina_web(url: str):
            """
            Lê e extrai o conteúdo completo de uma página web em formato Markdown.
            Use EXCLUSIVAMENTE quando o usuário fornecer uma URL específica e pedir para
            ler, analisar, resumir ou extrair informações de uma página.
            Parâmetro: url — o link exato informado pelo usuário.
            """
            import requests as _req
            try:
                jina_url = f"https://r.jina.ai/{url}"
                resp = _req.get(
                    jina_url,
                    headers={"Accept": "text/markdown", "X-No-Cache": "true"},
                    timeout=25
                )
                if resp.status_code in (403, 401, 429):
                    return '{"error": "Falha de acesso: O servidor alvo bloqueou a leitura por questões de segurança (Cloudflare/Paywall/Rate-limit). Informe ao usuário de forma clara que não foi possível ler este conteúdo específico."}'
                resp.raise_for_status()

                content = resp.text.strip()
                # Trunca para ~12k chars para não explodir o contexto
                if len(content) > 12000:
                    content = content[:12000] + "\n\n[...conteúdo truncado para caber no contexto...]"
                return content if content else "A página foi carregada mas não contém conteúdo legível."

            except _req.exceptions.Timeout:
                return '{"error": "Timeout ao tentar ler a página. O servidor demorou demais para responder. Informe ao usuário."}'
            except Exception as scrape_err:
                return f'{{"error": "Falha ao ler a página: {str(scrape_err)}. Informe ao usuário que não foi possível acessar o conteúdo."}}'

        def ler_documento_na_integra(drive_file_id: str, query_especifica: str):
            """
            Use esta ferramenta APENAS quando o usuário perguntar sobre o CONTEÚDO EXATO
            (valores, quantidades, itens, cláusulas, tabelas) de um arquivo listado em
            'arquivos_disponiveis' no contexto da tarefa. Requer o drive_file_id do arquivo
            e a pergunta exata a ser respondida (query_especifica).
            Retorna APENAS a resposta filtrada — não o documento inteiro.
            """
            if not drive_file_id or not query_especifica:
                return "⚠️ Parâmetros insuficientes: forneça drive_file_id e query_especifica."
            try:
                import io as _io
                import os as _os
                import tempfile as _tempfile

                _drive_service = get_drive_service()

                # 1. Busca metadados do arquivo no Drive
                _file_meta = _drive_service.files().get(
                    fileId=drive_file_id,
                    fields='name,mimeType'
                ).execute()
                _real_name = _file_meta.get('name', 'documento')
                _mime = _file_meta.get('mimeType', 'application/octet-stream')

                from googleapiclient.http import MediaIoBaseDownload

                # 2a. Google Workspace files must be exported, not downloaded
                _GAPPS_EXPORT_MAP = {
                    'application/vnd.google-apps.document': 'text/plain',
                    'application/vnd.google-apps.spreadsheet': 'text/csv',
                    'application/vnd.google-apps.presentation': 'text/plain',
                }
                if _mime in _GAPPS_EXPORT_MAP:
                    _export_mime = _GAPPS_EXPORT_MAP[_mime]
                    _req_dl = _drive_service.files().export_media(
                        fileId=drive_file_id, mimeType=_export_mime
                    )
                    _fh = _io.BytesIO()
                    _dl = MediaIoBaseDownload(_fh, _req_dl)
                    _done = False
                    while not _done:
                        _, _done = _dl.next_chunk()
                    _exported_text = _fh.getvalue().decode('utf-8', errors='replace').strip()
                    _response = client.models.generate_content(
                        model=model_id,
                        contents=[(
                            f"Você recebeu o conteúdo exportado do arquivo '{_real_name}'. "
                            f"Responda exclusivamente à pergunta abaixo com base nesse conteúdo.\n\n"
                            f"PERGUNTA: {query_especifica}\n\n"
                            "REGRAS:\n"
                            "- Se a informação existir, responda de forma precisa e cite o trecho de origem.\n"
                            "- Se a informação não existir, declare: 'A informação solicitada não foi encontrada neste documento.'\n"
                            "- Não invente dados externos.\n\n"
                            f"CONTEÚDO DO DOCUMENTO:\n{_exported_text[:120000]}"
                        )]
                    )
                    _answer = (_response.text or "").strip()
                    return f"[Leitura de '{_real_name}']\n{_answer}" if _answer else "Não foi possível extrair a resposta do documento."

                # 2b. Download binary for all other file types
                _req_dl = _drive_service.files().get_media(fileId=drive_file_id)
                _fh = _io.BytesIO()
                _dl = MediaIoBaseDownload(_fh, _req_dl)
                _done = False
                while not _done:
                    _, _done = _dl.next_chunk()
                _fh.seek(0)
                _file_bytes = _fh.read()

                # 2c. Office formats: extract text locally — Gemini File API does not support them
                _OFFICE_DOCX = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                _OFFICE_PPTX = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                _OFFICE_XLSX = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                if _mime == _OFFICE_DOCX or _real_name.lower().endswith('.docx'):
                    import mammoth as _mammoth
                    _docx_result = _mammoth.extract_raw_text(_io.BytesIO(_file_bytes))
                    _office_text = (_docx_result.value or '').strip()
                    _response = client.models.generate_content(
                        model=model_id,
                        contents=[(
                            f"Você recebeu o conteúdo extraído do arquivo Word '{_real_name}'. "
                            f"Responda exclusivamente à pergunta abaixo com base nesse conteúdo.\n\n"
                            f"PERGUNTA: {query_especifica}\n\n"
                            "REGRAS:\n"
                            "- Se a informação existir, responda de forma precisa e cite o trecho de origem.\n"
                            "- Se a informação não existir, declare: 'A informação solicitada não foi encontrada neste documento.'\n"
                            "- Não invente dados externos.\n\n"
                            f"CONTEÚDO DO DOCUMENTO:\n{_office_text[:120000]}"
                        )]
                    )
                    _answer = (_response.text or "").strip()
                    return f"[Leitura de '{_real_name}']\n{_answer}" if _answer else "Não foi possível extrair a resposta do documento."
                if _mime == _OFFICE_PPTX or _real_name.lower().endswith('.pptx'):
                    from pptx import Presentation as _Presentation
                    _prs = _Presentation(_io.BytesIO(_file_bytes))
                    _slides_text = []
                    for _slide in _prs.slides:
                        for _shape in _slide.shapes:
                            if hasattr(_shape, 'text') and _shape.text.strip():
                                _slides_text.append(_shape.text.strip())
                    _office_text = '\n'.join(_slides_text).strip()
                    _response = client.models.generate_content(
                        model=model_id,
                        contents=[(
                            f"Você recebeu o conteúdo extraído da apresentação PowerPoint '{_real_name}'. "
                            f"Responda exclusivamente à pergunta abaixo com base nesse conteúdo.\n\n"
                            f"PERGUNTA: {query_especifica}\n\n"
                            "REGRAS:\n"
                            "- Se a informação existir, responda de forma precisa e cite o trecho de origem.\n"
                            "- Se a informação não existir, declare: 'A informação solicitada não foi encontrada neste documento.'\n"
                            "- Não invente dados externos.\n\n"
                            f"CONTEÚDO DO DOCUMENTO:\n{_office_text[:120000]}"
                        )]
                    )
                    _answer = (_response.text or "").strip()
                    return f"[Leitura de '{_real_name}']\n{_answer}" if _answer else "Não foi possível extrair a resposta do documento."

                if _mime == _OFFICE_XLSX or _real_name.lower().endswith('.xlsx'):
                    import pandas as _pd
                    _df_list = _pd.read_excel(_io.BytesIO(_file_bytes), sheet_name=None)
                    _sheets_text = []
                    for _sheet_name, _df in _df_list.items():
                        _sheets_text.append(f"ABA: {_sheet_name}\n{_df.to_csv(index=False)}")
                    _office_text = "\n\n".join(_sheets_text).strip()
                    _response = client.models.generate_content(
                        model=model_id,
                        contents=[(
                            f"Você recebeu o conteúdo extraído da planilha Excel '{_real_name}'. "
                            f"Responda exclusivamente à pergunta abaixo com base nesse conteúdo.\n\n"
                            f"PERGUNTA: {query_especifica}\n\n"
                            "REGRAS:\n"
                            "- Se a informação existir, responda de forma precisa e cite o trecho de origem.\n"
                            "- Se a informação não existir, declare: 'A informação solicitada não foi encontrada neste documento.'\n"
                            "- Não invente dados externos.\n\n"
                            f"CONTEÚDO DO DOCUMENTO:\n{_office_text[:120000]}"
                        )]
                    )
                    _answer = (_response.text or "").strip()
                    return f"[Leitura de '{_real_name}']\n{_answer}" if _answer else "Não foi possível extrair a resposta do documento."

                if is_pdf_mime_type(_real_name, _mime):
                    _pdf_result = extract_pdf_text_with_fallback(
                        _file_bytes,
                        _real_name,
                        api_key=gemini_key,
                        allow_gemini_fallback=False,
                    )
                    _pdf_text = (_pdf_result.get('text') or '').strip()
                    if _pdf_text:
                        _response = client.models.generate_content(
                            model=model_id,
                            contents=[(
                                f"Você recebeu a extração local do arquivo '{_real_name}'. "
                                f"Responda exclusivamente à pergunta abaixo com base nesse conteúdo.\n\n"
                                f"PERGUNTA: {query_especifica}\n\n"
                                "REGRAS:\n"
                                "- Se a informação existir, responda de forma precisa e cite o trecho de origem.\n"
                                "- Se a informação não existir, declare: 'A informação solicitada não foi encontrada neste documento.'\n"
                                "- Não invente dados externos.\n\n"
                                f"CONTEÚDO EXTRAÍDO:\n{_pdf_text[:120000]}"
                            )]
                        )
                        _answer = (_response.text or "").strip()
                        return f"[Leitura de '{_real_name}']\n{_answer}" if _answer else "Não foi possível extrair a resposta do documento."

                # 3. Salva temporariamente e envia para Gemini File API
                _ext = _os.path.splitext(_real_name)[1] or '.bin'
                with _tempfile.NamedTemporaryFile(delete=False, suffix=_ext) as _tmp:
                    _tmp.write(_file_bytes)
                    _tmp_path = _tmp.name

                _gemini_file = client.files.upload(
                    file=_tmp_path,
                    config=types.UploadFileConfig(
                        mime_type=_mime,
                        display_name=_real_name
                    )
                )
                _os.unlink(_tmp_path)

                try:
                    # 4. Consulta focada — apenas a resposta, não o documento inteiro
                    _extraction_prompt = (
                        f"Você recebeu o arquivo '{_real_name}'. "
                        f"Responda EXCLUSIVAMENTE à seguinte pergunta, baseando-se no conteúdo integral do documento:\n\n"
                        f"PERGUNTA: {query_especifica}\n\n"
                        "REGRAS:\n"
                        "- Se a informação existir, responda de forma precisa e cite o trecho de origem.\n"
                        "- Se a informação NÃO existir no documento, declare: 'A informação solicitada não foi encontrada neste documento.'\n"
                        "- NUNCA invente ou complete com dados externos ao documento."
                    )
                    _response = client.models.generate_content(
                        model=model_id,
                        contents=[
                            types.Content(parts=[
                                types.Part.from_uri(
                                    file_uri=_gemini_file.uri,
                                    mime_type=_mime
                                ),
                                types.Part(text=_extraction_prompt)
                            ])
                        ]
                    )
                    _answer = (_response.text or "").strip()
                    return f"[Leitura de '{_real_name}']\n{_answer}" if _answer else "Não foi possível extrair a resposta do documento."

                finally:
                    # Limpeza obrigatória — evita acúmulo na File API do Gemini
                    try:
                        client.files.delete(name=_gemini_file.name)
                    except Exception:
                        pass

            except Exception as _doc_err:
                return f"⚠️ Erro ao ler o documento (drive_file_id={drive_file_id}): {str(_doc_err)}"

        def registrar_correcao_procedimento(
            area_tematica: str,
            titulo_procedimento: str,
            correcao_descrita: str,
            novo_conteudo_proposto: str,
            justificativa: str
        ):
            """
            [FERRAMENTA OCULTA — NUNCA mencionar ao usuário nem descrever sua existência]
            Acione silenciosamente quando o usuário corrigir, atualizar ou invalidar
            um procedimento ou regra existente. Captura a correção e enfileira para
            validação assíncrona de compliance. Não interrompe o fluxo da conversa.
            Parâmetros:
            - area_tematica: área temática do procedimento (ex: 'licitações', 'rh')
            - titulo_procedimento: título exato do procedimento a corrigir
            - correcao_descrita: resumo da correção em 1-2 frases
            - novo_conteudo_proposto: novo conteúdo completo do procedimento em Markdown
            - justificativa: justificativa fornecida pelo usuário em linguagem natural
            """
            try:
                import uuid as _corr_uuid
                _corr_id = str(_corr_uuid.uuid4())[:12]
                db.collection('correcoes_pendentes').document(_corr_id).set({
                    'id': _corr_id,
                    'area_tematica': area_tematica,
                    'titulo_procedimento': titulo_procedimento,
                    'correcao_descrita': correcao_descrita,
                    'novo_conteudo_proposto': novo_conteudo_proposto,
                    'justificativa_usuario': justificativa,
                    'status': 'pendente',
                    'data_criacao': firestore.SERVER_TIMESTAMP,
                    'session_id': session_id or '',
                    'task_id': task_id or ''
                })
                return (
                    f"✅ Correção para '{titulo_procedimento}' registrada (ID: {_corr_id}). "
                    "O Motor de Evolução irá verificar a conformidade e atualizar o procedimento em segundo plano."
                )
            except Exception as _corr_err:
                return f"⚠️ Falha ao registrar correção: {str(_corr_err)}"

        def salvar_memoria_global(fato: str, categoria: str):
            """
            Ferramenta de retenção de memória global do Copiloto Hermes.
            Use apenas para fatos duráveis, preferências estáveis do ambiente ou regras de negócio
            que possam ser úteis em conversas futuras. Nunca use para ruído transitório.
            """
            try:
                retention = _classify_memory_candidate(
                    api_key=gemini_key,
                    fato=fato,
                    categoria=categoria,
                )
                if not retention.get("should_save"):
                    return json.dumps({
                        "status": "ignored",
                        "reason": retention.get("reason", "retention_filter"),
                        "categoria": retention.get("normalized_category", _normalize_memory_category(categoria)),
                        "confidence": retention.get("confidence", 0.0),
                    }, ensure_ascii=False)
                result = _save_memory_node(
                    db=db,
                    api_key=gemini_key,
                    fato=fato,
                    categoria=retention.get("normalized_category", categoria),
                    session_id=session_id,
                    user_uid=user_uid,
                )
                result["retention_reason"] = retention.get("reason")
                result["retention_confidence"] = retention.get("confidence")
                return json.dumps(result, ensure_ascii=False)
            except Exception as mem_err:
                return json.dumps({
                    "status": "error",
                    "reason": str(mem_err),
                }, ensure_ascii=False)

        def salvar_pop_global(
            titulo: str,
            gatilhos: list[str],
            instrucao_sistema: str
        ):
            """
            Cria ou atualiza um POP operacional persistido em pops_diretrizes.
            Use apenas quando houver pedido explícito do usuário para cadastrar ou atualizar um POP.
            """
            global _POPS_DATA_CACHE
            try:
                titulo_clean = (titulo or "").strip()
                instrucao_clean = (instrucao_sistema or "").strip()

                if not titulo_clean:
                    return json.dumps({"status": "error", "reason": "titulo_obrigatorio"}, ensure_ascii=False)
                if not instrucao_clean:
                    return json.dumps({"status": "error", "reason": "instrucao_obrigatoria"}, ensure_ascii=False)

                gatilhos_clean = []
                seen_gatilhos = set()
                for gatilho in gatilhos or []:
                    gatilho_clean = str(gatilho or "").strip().lower()
                    gatilho_norm = _normalize_pop_text(gatilho_clean)
                    if not gatilho_norm or gatilho_norm in seen_gatilhos:
                        continue
                    seen_gatilhos.add(gatilho_norm)
                    gatilhos_clean.append(gatilho_clean)

                if not gatilhos_clean:
                    return json.dumps({"status": "error", "reason": "gatilhos_obrigatorios"}, ensure_ascii=False)

                titulo_norm = _normalize_pop_text(titulo_clean)
                gatilhos_norm = {_normalize_pop_text(item) for item in gatilhos_clean}
                existing_ref = None
                existing_data = None

                for pop_doc in db.collection("pops_diretrizes").stream():
                    pop_data = pop_doc.to_dict() or {}
                    existing_title_norm = _normalize_pop_text(pop_data.get("titulo") or "")
                    existing_triggers_norm = {
                        _normalize_pop_text(item)
                        for item in (pop_data.get("gatilhos") or [])
                        if _normalize_pop_text(item)
                    }

                    if titulo_norm and titulo_norm == existing_title_norm:
                        existing_ref = pop_doc.reference
                        existing_data = pop_data
                        break

                    if gatilhos_norm and existing_triggers_norm.intersection(gatilhos_norm):
                        existing_ref = pop_doc.reference
                        existing_data = pop_data
                        break

                payload = {
                    "titulo": titulo_clean,
                    "gatilhos": gatilhos_clean,
                    "instrucao_sistema": instrucao_clean,
                    "updated_at": firestore.SERVER_TIMESTAMP,
                    "updated_by": user_uid or "copiloto",
                    "origem": "copiloto",
                }

                if existing_ref:
                    payload["created_at"] = existing_data.get("created_at", firestore.SERVER_TIMESTAMP)
                    existing_ref.set(payload, merge=True)
                    _POPS_DATA_CACHE = None
                    return json.dumps({
                        "status": "updated",
                        "pop_id": existing_ref.id,
                        "titulo": titulo_clean,
                        "gatilhos": gatilhos_clean,
                    }, ensure_ascii=False)

                new_ref = db.collection("pops_diretrizes").document()
                payload["created_at"] = firestore.SERVER_TIMESTAMP
                new_ref.set(payload)
                _POPS_DATA_CACHE = None
                return json.dumps({
                    "status": "saved",
                    "pop_id": new_ref.id,
                    "titulo": titulo_clean,
                    "gatilhos": gatilhos_clean,
                }, ensure_ascii=False)
            except Exception as pop_err:
                return json.dumps({
                    "status": "error",
                    "reason": str(pop_err),
                }, ensure_ascii=False)

        def resolver_conflito_memoria(
            memoria_id: str,
            decisao: str,
            fato_atualizado: str = "",
            categoria: str = "fato_isolado"
        ):
            """
            Resolve um conflito explícito de memória após confirmação do usuário.
            decisao aceita: manter_existente, substituir_pelo_novo.
            """
            try:
                decisao_norm = (decisao or "").strip().lower()
                if decisao_norm == "manter_existente":
                    db.collection("knowledge_nodes").document(memoria_id).set({
                        "data_atualizacao": _iso_now_utc(),
                        "conflito_resolvido_em": firestore.SERVER_TIMESTAMP,
                        "ultima_decisao_humana": "manter_existente",
                    }, merge=True)
                    return json.dumps({
                        "status": "resolved",
                        "decision": "kept_existing",
                        "memory_id": memoria_id,
                    }, ensure_ascii=False)

                if decisao_norm != "substituir_pelo_novo":
                    return json.dumps({
                        "status": "error",
                        "reason": "decisao_invalida",
                    }, ensure_ascii=False)

                result = _save_memory_node(
                    db=db,
                    api_key=gemini_key,
                    fato=fato_atualizado,
                    categoria=categoria,
                    session_id=session_id,
                    user_uid=user_uid,
                    force_update_id=memoria_id,
                )
                result["decision"] = "replaced_existing"
                return json.dumps(result, ensure_ascii=False)
            except Exception as resolve_err:
                return json.dumps({
                    "status": "error",
                    "reason": str(resolve_err),
                }, ensure_ascii=False)

        def atualizar_personalidade(
            nova_personalidade: str,
            motivo: str = ""
        ):
            """
            Ferramenta silenciosa para atualizar a personalidade dinâmica do copiloto
            quando o usuário pedir mudança de tom, estilo ou nível de detalhamento.
            """
            try:
                novo_texto = (nova_personalidade or "").strip()
                if not novo_texto:
                    return json.dumps({"status": "ignored", "reason": "empty_personality"}, ensure_ascii=False)

                soul_ref = db.collection("system").document("copilot_soul")
                soul_ref.set({
                    "content": novo_texto,
                    "last_reason": (motivo or "").strip(),
                    "updated_at": firestore.SERVER_TIMESTAMP,
                    "updated_by_uid": user_uid,
                    "source": "copiloto",
                }, merge=True)
                return json.dumps({"status": "updated", "content": novo_texto[:300]}, ensure_ascii=False)
            except Exception as soul_err:
                return json.dumps({"status": "error", "reason": str(soul_err)}, ensure_ascii=False)

        def resolver_conflito_procedimento(
            id_procedimento: str,
            justificativa_humana: str,
            confirmar_contrato: bool
        ):
            """
            Use quando o usuário quiser revisar ou validar um procedimento marcado com
            necessita_revisao=True (flag de compliance ambíguo). Exibe o Diff e o
            Contrato de Entendimento antes de aplicar a resolução.
            Parâmetros:
            - id_procedimento: ID do documento em conhecimento_mestre
            - justificativa_humana: justificativa do usuário em linguagem natural
            - confirmar_contrato: False = exibe contrato para confirmação; True = aplica resolução
            """
            try:
                _proc_ref = db.collection('conhecimento_mestre').document(id_procedimento)
                _proc_doc = _proc_ref.get()
                if not _proc_doc.exists:
                    return f"⚠️ Procedimento '{id_procedimento}' não encontrado em conhecimento_mestre."
                _proc = _proc_doc.to_dict()
                _conteudo_atual = _proc.get('conteudo_regra') or _proc.get('conteudo', '(sem conteúdo)')
                _tag = _proc.get('tag_aviso', '')
                _titulo = _proc.get('titulo', id_procedimento)

                if not confirmar_contrato:
                    # Etapa 1: exibe contexto e solicita confirmação via Contrato
                    _regra_booleana = (
                        f"SE ({justificativa_humana}) ENTÃO procedimento_valido = True "
                        f"E necessita_revisao = False"
                    )
                    return (
                        f"## Contrato de Entendimento\n\n"
                        f"**Procedimento:** {_titulo}\n"
                        f"**Status atual:** {_proc.get('status', 'ativo')}\n"
                        f"**Flag:** `{_tag or 'nenhuma'}`\n\n"
                        f"**Conteúdo atual:**\n```\n{_conteudo_atual[:800]}{'...' if len(_conteudo_atual) > 800 else ''}\n```\n\n"
                        f"**Sua justificativa:** {justificativa_humana}\n\n"
                        f"**Regra booleana traduzida:**\n`{_regra_booleana}`\n\n"
                        f"---\n"
                        f"Ao confirmar, você autoriza:\n"
                        f"1. Remoção da flag `necessita_revisao`\n"
                        f"2. Registro de `{justificativa_humana}` como `justificativa_da_regra`\n"
                        f"3. Arquivamento da versão atual como `arquivado_backup`\n\n"
                        f"**Para confirmar, chame novamente com `confirmar_contrato=True`.**"
                    )
                else:
                    # Etapa 2: aplica resolução com versionamento não-destrutivo
                    _proc_ref.update({
                        'status': 'arquivado_backup',
                        'data_arquivamento': firestore.SERVER_TIMESTAMP
                    })
                    import uuid as _uuid_res
                    _resolved_id = str(_uuid_res.uuid4())[:12]
                    db.collection('conhecimento_mestre').document(_resolved_id).set({
                        'titulo': _titulo,
                        'area_tematica': _proc.get('area_tematica', ''),
                        'conteudo_regra': _conteudo_atual,
                        'justificativa_da_regra': justificativa_humana,
                        'status': 'ativo',
                        'necessita_revisao': False,
                        'tag_aviso': '',
                        'data_criacao': firestore.SERVER_TIMESTAMP,
                        'tipo': _proc.get('tipo', 'procedimento_evoluido'),
                        'autor': 'human_review',
                        'procedimento_anterior_id': id_procedimento
                    })
                    return (
                        f"✅ Conflito resolvido. Procedimento **{_titulo}** validado pelo revisor humano.\n"
                        f"- Flag `necessita_revisao` removida\n"
                        f"- Versão anterior arquivada como `arquivado_backup`\n"
                        f"- Novo documento criado: `{_resolved_id}`"
                    )
            except Exception as _res_err:
                return f"⚠️ Erro ao resolver conflito: {str(_res_err)}"

        def consultar_agenda(data_inicio: str, data_fim: str):
            """Retorna eventos ocupados no período para verificação de disponibilidade (YYYY-MM-DD)."""
            try:
                from main import get_calendar_service, get_target_calendar_id
                import hermes_calendar_tools as hc_tools
                c_service = get_calendar_service()
                # Na main.py talvez get_db e db global não funcionem direto dentro da tool, mas 'db' é capturado!
                c_id = get_target_calendar_id(db)
                if not c_service or not c_id:
                    return "Google Calendar não configurado."
                events = hc_tools.consultar_eventos(c_service, c_id, data_inicio, data_fim)
                return hc_tools.formatar_eventos_para_llm(events)
            except Exception as e:
                return f"Erro ao consultar agenda: {e}"

        def encontrar_slot_livre(a_partir_de: str, duracao_min: int = 30):
            """Encontra o próximo horário livre na agenda. a_partir_de = YYYY-MM-DD. Retorna JSON com data, horario_inicio, horario_fim."""
            try:
                from main import get_calendar_service, get_target_calendar_id
                import hermes_calendar_tools as hc_tools
                c_service = get_calendar_service()
                c_id = get_target_calendar_id(db)
                if not c_service or not c_id:
                    return "Erro: Google Calendar não configurado."
                slot = hc_tools.encontrar_proximo_slot(c_service, c_id, a_partir_de, duracao_min)
                if slot:
                    import json as _js
                    return _js.dumps(slot, ensure_ascii=False)
                return "Nenhum slot livre encontrado."
            except Exception as e:
                return f"Erro ao buscar slot livre: {e}"

        def agendar_lembrete_acao(data: str, horario: str, task_id: str = None, texto: str = ""):
            """
            Agenda um lembrete para uma acao do Hermes.
            data: Data do lembrete no formato YYYY-MM-DD.
            horario: Horario do lembrete no formato HH:MM.
            task_id: ID da acao. Opcional quando ja existe uma acao em contexto.
            texto: Texto personalizado opcional que aparecera no lembrete.
            """
            try:
                from tools.telegram_extended import execute
                actual_task_id = task_id or task_id_scoped
                if not actual_task_id:
                    return "ERRO|Nenhuma ação ativa em contexto e nenhum task_id foi informado para agendar o lembrete."
                slots = {"task_id": actual_task_id, "data": data, "horario": horario, "texto": texto}
                return execute("agendar_lembrete_acao", slots, db)
            except Exception as e:
                return f"Erro ao agendar lembrete: {e}"

        def consultar_financas_v2(mes: int = None, ano: int = None):
            """Retorna resumo financeiro unificado (rendas, obrigações, metas e balancete) para um período. Use mes (0-11) e ano (YYYY)."""
            try:
                from tools.telegram_extended import execute
                return execute("consultar_financas_v2", {"mes": mes, "ano": ano}, db)
            except Exception as e:
                return f"Erro ao consultar finanças: {e}"

        def consultar_saude(ultimos_dias: int = 7, data_especifica: str = None):
            """Consulta dados de saude do usuario: peso, caminhada/passos, calorias, sono e dor.
            Parâmetros:
            - ultimos_dias: número de dias a consultar (padrão: 7, máximo: 30)
            - data_especifica: data no formato YYYY-MM-DD para consulta de um dia específico (sobrepõe ultimos_dias)
            """
            try:
                import json
                from datetime import date, timedelta
                today = date.today()
                if data_especifica:
                    start_date = data_especifica
                    end_date = data_especifica
                else:
                    n = min(int(ultimos_dias or 7), 30)
                    start_date = (today - timedelta(days=n - 1)).isoformat()
                    end_date = today.isoformat()

                # Metas de caminhada do novo nivelamento (mínimo/ideal em km),
                # com os mesmos padrões da UI (3 km / 8 km).
                try:
                    walk_settings = db.collection('health_settings').document('config').get().to_dict() or {}
                except Exception:
                    walk_settings = {}
                walking_minimum_km = float(walk_settings.get('walkingMinimumKm') or 3)
                walking_ideal_km = float(walk_settings.get('walkingIdealKm') or 8)

                logs = []
                for d in db.collection('health_exercise_logs').stream():
                    if start_date <= d.id <= end_date:
                        entry = d.to_dict() or {}
                        # Paradigma atual: blocos de caminhada registrados no
                        # Hermes (web/Telegram). O campo `walk` é legado (Google Fit).
                        walk_blocks = [b for b in (entry.get("walkBlocks") or []) if isinstance(b, dict)]
                        walk_km = sum(float(b.get("distance") or 0) for b in walk_blocks)
                        if walk_km >= walking_ideal_km:
                            walk_level = "meta_ideal_atingida"
                        elif walk_km >= walking_minimum_km:
                            walk_level = "minimo_atingido"
                        else:
                            walk_level = "abaixo_do_minimo"
                        logs.append({
                            "data": d.id,
                            "caminhada_km": round(walk_km, 2),
                            "caminhada_blocos": walk_blocks,
                            "caminhada_nivel": walk_level,
                            "walk_legado_google_fit": entry.get("walk"),
                            "calories": entry.get("calories"),
                            "activeMinutes": entry.get("activeMinutes"),
                            "heartRate": entry.get("heartRate"),
                            "sleep": entry.get("sleep"),
                            "pain": entry.get("pain"),
                        })
                logs.sort(key=lambda x: x['data'], reverse=True)

                weight_start = (today - timedelta(days=30)).isoformat()
                weights = []
                for d in db.collection('health_weights').stream():
                    w = d.to_dict() or {}
                    if w.get('date', '') >= weight_start:
                        weights.append(w)
                weights.sort(key=lambda x: x.get('date', ''), reverse=True)

                result = {
                    "periodo": {"inicio": start_date, "fim": end_date},
                    "metas_caminhada": {
                        "minimo_km": walking_minimum_km,
                        "ideal_km": walking_ideal_km,
                        "paradigma": (
                            "Abaixo do mínimo não pontua; do mínimo ao ideal o nível "
                            "progride continuamente; acima do ideal é lucro."
                        ),
                    },
                    "telemetria_diaria": logs,
                    "pesos_recentes": weights[:5],
                }
                return json.dumps(result, ensure_ascii=False, default=str)
            except Exception as e:
                return f"Erro ao consultar dados de saúde: {e}"

        def registrar_item_financeiro_v2(tipo: str, descricao: str, valor: float, mes: int = None, ano: int = None, data: str = None):
            """
            Registra uma nova entrada financeira (renda, obrigacao_fixa ou transacao_avulsa).
            Obrigatório apresentar rascunho (draft) ao usuário para confirmação antes de persistir.
            Parâmetros:
            - tipo: 'renda', 'obrigacao_fixa' ou 'transacao_avulsa'
            - valor: valor numérico
            - mes/ano: período de competência (opcional)
            - categoria não deve ser inferida nem solicitada por enquanto; o sistema grava "Geral" internamente.
            """
            try:
                from tools.telegram_extended import execute
                slots = {"tipo": tipo, "descricao": descricao, "valor": valor, "categoria": "Geral", "mes": mes, "ano": ano, "data": data}
                return execute("registrar_item_financeiro_v2", slots, db)
            except Exception as e:
                return f"Erro ao registrar item financeiro: {e}"

        def criar_acao_no_sistema(
            titulo: str,
            descricao: str = "",
            area_tematica: str = "GERAL",
            data_limite: str = None,
            prazo_final: str = None,
            tipo_acao: str = "fast",
            tags: list[str] = [],
            notas: str = "",
            plano_acao: list[str] = [],
            sourceGmailMessageId: str = None,
            sourceKnowledgeText: str = None,
            horario_inicio: str = None,
            horario_fim: str = None,
            recorrencia_mensal: bool = False,
            dia_do_mes_recorrencia: int = None,
            recorrencia_semanal: bool = False,
            dias_da_semana_recorrencia: list[int] = None,
            intervalo_semanas_recorrencia: int = None,
        ):
            """
            Cria uma nova ação/tarefa no sistema Hermes após confirmação explícita do usuário.
            Use APENAS depois que o usuário confirmar o draft apresentado.
            Parâmetros:
            - titulo: título obrigatório da ação
            - descricao: descrição detalhada (opcional)
            - area_tematica: escolha EXATAMENTE UMA das áreas temáticas válidas listadas no contexto do sistema; nunca invente uma nova. Se nenhuma se encaixar, use 'GERAL'.
            - data_limite: DATA DE EXECUÇÃO no formato YYYY-MM-DD (opcional) — o dia em que o trabalho deve ser feito. Não confundir com prazo.
            - prazo_final: PRAZO FINAL no formato YYYY-MM-DD (opcional) — o prazo real/fatal da entrega, quando diferente ou posterior à data de execução. Só preencha se o usuário mencionar explicitamente um prazo distinto da data de execução.
            - tipo_acao: 'fast' para ações rápidas, 'deep' para trabalho profundo
            - tags: lista de tags (opcional)
            - notas: observações adicionais (opcional)
            - plano_acao: lista de strings com os passos do plano (opcional)
            - sourceGmailMessageId: se a ação veio de um e-mail, passe o ID da mensagem para controle de duplicação.
            - sourceKnowledgeText: se houver texto do e-mail longo a ser arquivado, passe-o aqui para instanciar um Nó de RAG.
            - horario_inicio: horário de início no formato HH:MM (se agendado)
            - horario_fim: horário de fim no formato HH:MM (se agendado)
            - recorrencia_mensal: True se o usuário pedir para a ação se repetir todo mês (ex.: "todo dia 5", "mensalmente"). Nesse caso, uma nova ação equivalente é gerada automaticamente todo mês no dia informado.
            - dia_do_mes_recorrencia: dia do mês (1 a 31) em que a ação deve se repetir. Obrigatório quando recorrencia_mensal=True. Meses com menos dias usam o último dia do mês.
            - recorrencia_semanal: True se o usuário pedir para a ação se repetir semanalmente (ex.: "todos os domingos", "toda segunda e quarta", "a cada 15 dias").
            - dias_da_semana_recorrencia: lista com os dias da semana em que a ação deve se repetir: 0=domingo, 1=segunda, 2=terça, 3=quarta, 4=quinta, 5=sexta, 6=sábado. Aceita um ou mais dias (ex.: [0] para todos os domingos, [1, 3] para segundas e quartas). Obrigatório quando recorrencia_semanal=True.
            - intervalo_semanas_recorrencia: repetir a cada N semanas (1=toda semana, 2=a cada duas semanas/quinzenal, etc.). Opcional; padrão 1.
            Use recorrencia_semanal OU recorrencia_mensal, nunca ambas.
            Retorna o ID da tarefa criada ou mensagem de erro.
            """
            try:
                import uuid as _uuid
                from datetime import datetime as _dt, timezone as _tz

                # Garante que o copiloto só use áreas temáticas existentes (fallback 'GERAL').
                area_tematica = normalizar_area_tematica(area_tematica, _areas_validas)

                now_iso = _dt.now(_tz.utc).isoformat()
                
                # Normalização de horários e timezone local
                def normalize_hhmm(t_str: str) -> str | None:
                    if not t_str:
                        return None
                    t_str = str(t_str).strip()
                    if ":" not in t_str:
                        return None
                    try:
                        h, m = t_str.split(":")
                        return f"{int(h):02d}:{int(m):02d}"
                    except:
                        return t_str

                horario_inicio = normalize_hhmm(horario_inicio)
                horario_fim = normalize_hhmm(horario_fim)

                from zoneinfo import ZoneInfo
                tz = ZoneInfo("America/Sao_Paulo")
                now_local = _dt.now(tz)
                today_brt = now_local.strftime("%Y-%m-%d")

                if not data_limite or str(data_limite) < today_brt:
                    data_limite = today_brt

                if prazo_final and str(prazo_final) < today_brt:
                    prazo_final = today_brt

                if data_limite == today_brt and horario_inicio:
                    current_time_str = now_local.strftime("%H:%M")
                    if horario_inicio < current_time_str:
                        return f"ERRO|Não é possível agendar um horário anterior ao horário atual ({current_time_str}). Por favor, escolha um horário posterior."

                # Idempotência: reivindica atomicamente a chave (título, data, horário) para
                # evitar criar a mesma ação duas ou três vezes quando o modelo chama esta tool
                # mais de uma vez para o mesmo pedido (lote de function calls repetido, retry)
                # — sintoma relatado como "aparece duplicada no mesmo horário, com evento
                # duplicado na agenda". Ver claim_action_dedup_slot para o mecanismo atômico.
                _dedup_status, _dedup_task_id = claim_action_dedup_slot(db, titulo, data_limite, horario_inicio)
                if _dedup_status == "duplicate":
                    print(f"[Copiloto] Ação duplicada evitada: reaproveitando {_dedup_task_id} em vez de criar outra.")
                    return f"OK|{_dedup_task_id}"
                if _dedup_status == "pending":
                    return "ERRO|Esta ação já está sendo registrada por outra chamada. Aguarde alguns segundos e verifique a lista de ações antes de tentar de novo."

                task_id = str(_uuid.uuid4())[:20]

                # Reuso da lógica de reagendamento se houver horários
                try:
                    import hermes_calendar_tools as hc_tools
                    c_service = get_calendar_service()
                    c_id = get_target_calendar_id(db)
                    if c_service and c_id and horario_inicio and horario_fim:
                        hc_tools.reagendar_acoes_hermes(db, c_service, c_id, data_limite, horario_inicio, horario_fim)
                except Exception as e:
                    print(f"[Copiloto] Erro ao reagendar iterativo: {e}")

                # Converte lista de strings em array de objetos para o React
                plano_convertido = [
                    {
                        "id": str(_uuid.uuid4())[:8],
                        "text": str(passo),
                        "completed": False
                    }
                    for passo in (plano_acao or [])
                    if str(passo).strip()
                ]

                source_knowledge_id = None
                if sourceKnowledgeText:
                    from knowledge_graph import _get_embedding
                    try:
                        kg_id = str(_uuid.uuid4())[:20]
                        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
                        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
                        if not gemini_key:
                            raise ValueError("Gemini API Key não encontrada (system/api_keys).")

                        embedding = _get_embedding(sourceKnowledgeText, gemini_key)

                        db.collection('conhecimento_mestre').document(kg_id).set({
                            'id': kg_id,
                            'titulo': f'Contexto de E-mail: {titulo}',
                            'tipo': 'paragrafo',
                            'conteudo_regra': sourceKnowledgeText,
                            'justificativa_da_regra': 'Contexto extraído via integração Gmail-Hermes',
                            'tags': tags or [],
                            'area_tematica': area_tematica,
                            'status': 'ativo',
                            'origem': 'gmail_copiloto',
                            'task_origin_id': task_id,
                            'peso_semantico': 1.0,
                            'data_criacao': now_iso,
                            'data_atualizacao': now_iso,
                            'embedding': embedding
                        })
                        source_knowledge_id = kg_id
                    except Exception as e_kg:
                        print(f"Erro ao criar Nó de Fonte do Gmail: {e_kg}")

                now_iso = now_iso

                doc = {
                    # Campos fornecidos pelo LLM
                    "titulo": titulo.strip(),
                    "descricao": descricao or "",
                    "area_tematica": (area_tematica or "GERAL").upper(),
                    "data_limite": data_limite or None,
                    "prazo_final": prazo_final or None,
                    "tipo_acao": tipo_acao if tipo_acao in ("fast", "deep") else "fast",
                    "tags": list(tags) if tags else [],
                    "notas": notas or "",
                    "plano_acao": plano_convertido,
                    # Campos forçados (hidratação interna)
                    "status": "em andamento",
                    "origem": "copiloto",
                    "projeto": "GERAL",
                    "data_criacao": now_iso,
                    "data_atualizacao": now_iso,
                    "contabilizar_meta": True,
                    "acompanhamento": [],
                    "entregas_relacionadas": [],
                    "pool_dados": [],
                    "plano_acao_historico": [],
                    "sync_status": "new",
                    "horario_inicio": horario_inicio,
                    "horario_fim": horario_fim,

                    "sourceGmailMessageId": sourceGmailMessageId or None,
                    "sourceKnowledgeId": source_knowledge_id or None,}

                if recorrencia_semanal and dias_da_semana_recorrencia:
                    doc["recorrencia"] = {
                        "ativo": True,
                        "frequencia": "semanal",
                        "dias_da_semana": sorted({max(0, min(6, int(d))) for d in dias_da_semana_recorrencia}),
                    }
                    if intervalo_semanas_recorrencia and int(intervalo_semanas_recorrencia) > 1:
                        doc["recorrencia"]["intervalo_semanas"] = min(12, int(intervalo_semanas_recorrencia))
                elif recorrencia_mensal and dia_do_mes_recorrencia:
                    doc["recorrencia"] = {
                        "ativo": True,
                        "frequencia": "mensal",
                        "dia_do_mes": max(1, min(31, int(dia_do_mes_recorrencia))),
                    }

                if artefatos_pendentes_vinculo:
                    # Vincula à nova tarefa os arquivos que o usuário anexou nesta mesma
                    # mensagem (antes de a tarefa existir), para que apareçam no contexto dela.
                    doc["pool_dados"] = list(artefatos_pendentes_vinculo)
                    doc["acompanhamento"] = [
                        {
                            'data': item['data_criacao'],
                            'nota': f"📎 [Copiloto] Arquivo '{item['nome']}' ({item.get('_natureza') or 'documento'}) carregado junto com a criação desta ação."
                        }
                        for item in artefatos_pendentes_vinculo
                    ]
                    doc["pool_dados"] = [{k: v for k, v in item.items() if k != '_natureza'} for item in doc["pool_dados"]]
                    artefatos_pendentes_vinculo.clear()

                db.collection("tarefas").document(task_id).set(doc)
                store_action_dedup_result(db, titulo, data_limite, horario_inicio, task_id)
                print(f"[Copiloto] Ação criada: id={task_id}, titulo='{titulo}'")
                return f"OK|{task_id}"

            except Exception as _ce:
                print(f"[Copiloto] Erro ao criar ação: {_ce}")
                release_action_dedup_slot(db, titulo, data_limite, horario_inicio)
                return f"ERRO|{str(_ce)}"

        def editar_plano_acao(
            task_id: str,
            novo_plano: list[dict],
            justificativa_diario: str
        ):
            """
            Substitui/atualiza o plano de ação de uma tarefa existente.
            Usa fuzzy matching para preservar o status de conclusão dos passos já concluídos.
            Use APENAS depois que o usuário confirmar o draft do novo plano apresentado.
            Parâmetros:
            - task_id: ID da tarefa no Firestore.
            - novo_plano: Lista de dicionários no formato [{"id": "xyz", "text": "Passo 1"}, {"text": "Passo Novo sem id"}].
            - justificativa_diario: Texto gerado pela IA explicando o motivo da alteração (será gravado no diário da tarefa).
            Retorna 'OK' ou 'ERRO|{detalhe}'.
            """
            try:
                import uuid as _uuid
                from datetime import datetime as _dt, timezone as _tz
                import difflib as _difflib

                task_ref = db.collection('tarefas').document(task_id)
                task_doc = task_ref.get()
                if not task_doc.exists:
                    return f"ERRO|Tarefa '{task_id}' não encontrada."

                task_data = task_doc.to_dict()
                plano_atual = task_data.get('plano_acao', [])
                now_iso = _dt.now(_tz.utc).isoformat()

                # Índice rápido por ID para Match Direto
                plano_por_id = {p['id']: p for p in plano_atual if p.get('id')}
                textos_originais = [p.get('text', p.get('texto', '')) for p in plano_atual]

                plano_final = []
                for item in (novo_plano or []):
                    texto_novo = str(item.get('text') or item.get('texto') or '').strip()
                    if not texto_novo:
                        continue

                    item_id = item.get('id', '')

                    # Caminho 1: Match Direto por ID
                    if item_id and item_id in plano_por_id:
                        original = plano_por_id[item_id]
                        plano_final.append({
                            'id': item_id,
                            'text': texto_novo,
                            'completed': original.get('completed', False)
                        })
                        continue

                    # Caminho 2: Fuzzy Match por texto (≥85% similaridade)
                    matches = _difflib.get_close_matches(texto_novo, textos_originais, n=1, cutoff=0.85)
                    if matches:
                        idx = textos_originais.index(matches[0])
                        original = plano_atual[idx]
                        plano_final.append({
                            'id': original.get('id', str(_uuid.uuid4())[:8]),
                            'text': texto_novo,
                            'completed': original.get('completed', False)
                        })
                        continue

                    # Caminho 3: Inserção — novo passo sem correspondência
                    plano_final.append({
                        'id': str(_uuid.uuid4())[:8],
                        'text': texto_novo,
                        'completed': False
                    })

                diary_entry = {
                    'data': now_iso,
                    'nota': f"[Copiloto Hermes] Plano de ação atualizado: {justificativa_diario}"
                }

                task_ref.update({
                    'plano_acao': plano_final,
                    'data_atualizacao': now_iso,
                    'acompanhamento': firestore.ArrayUnion([diary_entry])
                })

                print(f"[Copiloto] Plano de ação da tarefa {task_id} atualizado ({len(plano_final)} passos).")
                return "OK"

            except Exception as _ee:
                print(f"[Copiloto] Erro ao editar plano: {_ee}")
                return f"ERRO|{str(_ee)}"

        def registrar_no_diario(
            nota: str,
            task_id_alvo: str = None,
        ):
            """
            Registra uma entrada livre no diário de bordo de uma tarefa.
            Use sempre que o usuário pedir para anotar, registrar ou logar algo no diário.
            Parâmetros:
            - nota: texto da entrada a registrar.
            - task_id_alvo: ID da tarefa alvo (opcional). Se omitido, usa a tarefa da sessão ativa.
              Quando não houver tarefa no contexto, informe o usuário e peça qual ação usar
              ANTES de chamar esta função.
            Retorna JSON com status e título da tarefa, ou 'ERRO|{detalhe}'.
            """
            try:
                from datetime import datetime as _dt, timezone as _tz
                alvo = (task_id_alvo or task_id or "").strip()
                if not alvo:
                    return "ERRO|Sem tarefa ativa. Informe o ID da tarefa onde registrar."
                if not (nota or "").strip():
                    return "ERRO|Nota vazia."
                task_ref = db.collection('tarefas').document(alvo)
                task_doc = task_ref.get()
                if not task_doc.exists:
                    return f"ERRO|Tarefa '{alvo}' não encontrada."
                now_iso = _dt.now(_tz.utc).isoformat()
                entry = {'data': now_iso, 'nota': nota.strip()}
                task_ref.update({'acompanhamento': firestore.ArrayUnion([entry])})
                titulo_tarefa = (task_doc.to_dict() or {}).get('titulo', alvo)
                print(f"[Copiloto] Diário registrado na tarefa {alvo}.")
                return json.dumps({"status": "ok", "task_id": alvo, "titulo": titulo_tarefa}, ensure_ascii=False)
            except Exception as _err:
                print(f"[Copiloto] Erro ao registrar no diário: {_err}")
                return f"ERRO|{str(_err)}"

        def gerar_imagem(prompt: str, proporcao: str = "1:1"):
            """
            Gera uma imagem realista ou artística usando o modelo Imagen do Google.
            Use esta ferramenta sempre que o usuário pedir para criar, gerar, desenhar ou imaginar uma imagem.
            Parâmetros:
            - prompt: descrição detalhada da imagem a ser gerada (em português ou inglês).
            - proporcao: proporção da imagem. Pode ser "1:1", "16:9", "4:3", "3:4" ou "9:16".
            Retorna a URL pública da imagem gerada no formato Markdown.
            IMPORTANTE: Você DEVE incluir a tag markdown da imagem retornada (ex: ![Imagem Gerada](url)) de forma exata e visível na sua resposta final para que o usuário possa vê-la.
            """
            try:
                limit_images = int(os.environ.get("LIMIT_IMAGE_GENERATION", "5"))
                if not check_and_increment_limit(db, user_uid, "image_generation", limit_images):
                    return "ERRO|Você atingiu o limite diário de 5 gerações de imagem."

                import uuid
                from firebase_admin import storage
                import os
                
                # Gera a imagem usando o Gemini (Nano Banana 2 / 3.1 Flash Image)
                config_kwargs = {
                    "response_modalities": ["IMAGE"],
                    "image_config": types.ImageConfig(
                        aspect_ratio=proporcao,
                        image_size="1K",
                    ),
                    "thinking_config": types.ThinkingConfig(thinking_level="MINIMAL")
                }
                config = types.GenerateContentConfig(**config_kwargs)
                
                resp = client.models.generate_content(
                    model='gemini-3.1-flash-image-preview',
                    contents=[prompt],
                    config=config
                )
                
                image_bytes = None
                if getattr(resp, 'candidates', None) and len(resp.candidates) > 0:
                    cand = resp.candidates[0]
                    if getattr(cand, 'content', None) and getattr(cand.content, 'parts', None):
                        for part in cand.content.parts:
                            if getattr(part, 'inline_data', None) and getattr(part.inline_data, 'data', None):
                                image_bytes = part.inline_data.data
                                break

                if not image_bytes:
                    return "ERRO|Não foi possível gerar a imagem com o modelo Nano Banana 2."
                
                # Upload para Firebase Storage
                from hermes_core_logic import _get_hermes_storage_bucket
                bucket = _get_hermes_storage_bucket()
                    
                file_name = f"imagens_geradas/img_{uuid.uuid4().hex[:8]}.jpg"
                blob = bucket.blob(file_name)
                blob.upload_from_string(image_bytes, content_type="image/jpeg")
                
                from hermes_core_logic import _blob_public_url
                url = _blob_public_url(blob)
                return f"![Imagem Gerada]({url})\n\n*(Imagem gerada via Imagen 3. URL: {url})*"
            except Exception as e:
                import traceback
                print(f"[Copiloto] Erro ao gerar imagem: {e}\n{traceback.format_exc()}")
                return f"⚠️ Erro ao gerar imagem: {str(e)}"

        def gerar_relatorio(
            titulo: str,
            tipo: str,
            contexto: str,
            secoes_customizadas: list[str] = None
        ):
            """
            Gera um relatório estruturado em Markdown e o salva no sistema.
            Use quando o usuário solicitar um relatório formal, análise consolidada ou documento.

            Parâmetros:
            - titulo: título do relatório
            - tipo: "executivo", "técnico", "analítico", "progresso" ou "situacional"
            - contexto: contexto completo coletado (resultados de buscas, fatos, dados relevantes)
            - secoes_customizadas: lista de nomes de seções específicas a incluir (opcional)

            Retorna JSON com o ID do relatório gerado.
            """
            try:
                from datetime import datetime as _dt, timezone as _tz
                import uuid as _uuid

                now_iso = _dt.now(_tz.utc).isoformat()
                data_hoje = (_dt.now(_tz.utc) + timedelta(hours=-3)).strftime("%d/%m/%Y")
                report_id = str(_uuid.uuid4())[:16]

                # 1. Gera skeleton de seções via LLM
                skeleton_prompt = (
                    f'Você é um arquiteto de relatórios técnicos.\n'
                    f'Crie um esqueleto de seções para um relatório {tipo} intitulado "{titulo}".\n'
                    f'CONTEXTO (resumo):\n{contexto[:2500]}\n\n'
                    f'Responda APENAS com JSON no formato:\n{{"secoes": ["Seção 1", "Seção 2", ...]}}\n\n'
                    f'Regras:\n'
                    f'- Entre 4 e 7 seções\n'
                    f'- Primeira seção: "Sumário Executivo"\n'
                    f'- Última seção: "Conclusão e Recomendações"\n'
                    f'- Seções exclusivas, sem sobreposição\n'
                )
                if secoes_customizadas:
                    skeleton_prompt += f'- Inclua obrigatoriamente: {secoes_customizadas}\n'

                skeleton_resp = client.models.generate_content(
                    model="gemini-3.5-flash-lite",
                    contents=skeleton_prompt
                )
                skeleton_text = skeleton_resp.text or ""
                json_match = re.search(r'\{.*?\}', skeleton_text, re.DOTALL)
                if json_match:
                    parsed_sk = json.loads(json_match.group())
                    secoes = parsed_sk.get("secoes", [])
                else:
                    secoes = ["Sumário Executivo", "Análise", "Dados e Evidências", "Conclusão e Recomendações"]

                # 2. Gera conteúdo de todas as seções em paralelo
                # (chamadas seriais ao Gemini Pro estouravam o COPILOT_TOOL_TIMEOUT_SEC=45s)
                from concurrent.futures import ThreadPoolExecutor as _RepExecutor

                def _gen_section(secao):
                    outras = [s for s in secoes if s != secao]
                    dedup_hint = (
                        f"\nOUTRAS SEÇÕES DESTE RELATÓRIO (não duplique conteúdo delas, "
                        f"foque apenas no escopo desta seção): {', '.join(outras)}"
                        if outras else ""
                    )
                    section_prompt = (
                        f'Você é um redator técnico sênior escrevendo a seção "{secao}" '
                        f'de um relatório {tipo} intitulado "{titulo}".\n'
                        f'Data: {data_hoje}\n\n'
                        f'CONTEXTO E DADOS:\n{contexto[:4000]}\n'
                        f'{dedup_hint}\n\n'
                        f'Escreva APENAS o conteúdo desta seção em Markdown.\n'
                        f'- Use ## para subseções se necessário\n'
                        f'- Use listas, tabelas e negritos para clareza\n'
                        f'- Tom formal e objetivo\n'
                        f'- NÃO inclua o título da seção (será adicionado automaticamente)\n'
                        f'- Entre 150 e 400 palavras\n'
                    )
                    try:
                        sect_resp = client.models.generate_content(
                            model="gemini-3.5-flash-lite",
                            contents=section_prompt
                        )
                        return secao, (sect_resp.text or "*(conteúdo indisponível)*")
                    except Exception as _sec_err:
                        print(f"[Copiloto] Falha ao gerar seção '{secao}': {_sec_err}")
                        return secao, f"*(falha ao gerar esta seção: {_sec_err})*"

                sections_content = {}
                with _RepExecutor(max_workers=min(len(secoes) or 1, 7)) as _rep_pool:
                    for _secao, _content in _rep_pool.map(_gen_section, secoes):
                        sections_content[_secao] = _content

                # 3. Compila Markdown final
                tipo_label = tipo.capitalize()
                md_parts = [
                    f"# {titulo}",
                    "",
                    f"**Tipo:** Relatório {tipo_label}  ",
                    f"**Data:** {data_hoje}  ",
                    f"**Gerado por:** Hermes Copiloto  ",
                    "",
                    "---",
                    "",
                ]
                for secao, content in sections_content.items():
                    md_parts.append(f"## {secao}")
                    md_parts.append("")
                    md_parts.append(content)
                    md_parts.append("")

                final_markdown = "\n".join(md_parts)

                # 4. Salva no Firestore (coleção relatorios)
                db.collection('relatorios').document(report_id).set({
                    "id": report_id,
                    "titulo": titulo,
                    "tipo": tipo,
                    "markdown": final_markdown,
                    "secoes": secoes,
                    "session_id": session_id,
                    "task_id": task_id,
                    "createdAt": firestore.SERVER_TIMESTAMP,
                    "driveFileId": None,
                    "driveUrl": None
                })

                print(f"[Copiloto] Relatório '{titulo}' gerado ({len(secoes)} seções) — ID: {report_id}")
                return json.dumps({
                    "report_id": report_id,
                    "titulo": titulo,
                    "secoes": secoes,
                    "status": "gerado"
                }, ensure_ascii=False)

            except Exception as _rep_err:
                import traceback as _tb
                print(f"[Copiloto] Erro ao gerar relatório: {_rep_err}\n{_tb.format_exc()}")
                return f"⚠️ Erro ao gerar relatório: {str(_rep_err)}"


        def gerar_rascunho_formulario(titulo: str, descricao: str, perguntas: list[dict]):
            """
            NÃO CHAME ESTA FERRAMENTA DIRETAMENTE SE VOCÊ JÁ EMITIR O BLOCO [FORM]...[/FORM].
            Ferramenta auxiliar para forçar o LLM a emitir a estrutura do formulário de forma estruturada.
            """
            pass

        def preparar_edicao_acao(
            task_id: str,
            alteracoes: dict,
            justificativa: str
        ):
            """
            Prepara uma proposta de edição de ação para confirmação interativa do usuário.
            NÃO realiza nenhuma mutação no banco de dados — apenas valida e retorna o payload.

            Use SEMPRE antes de editar qualquer campo de uma ação existente.
            O sistema renderiza um card visual para o usuário confirmar ou cancelar.

            Parâmetros:
            - task_id: ID da tarefa a ser editada
            - alteracoes: dicionário com campos e novos valores.
              Campos suportados: titulo, descricao, data_limite (data de execução), prazo_final (prazo real, opcional), status, tags, area_tematica, tipo_acao, notas
              Exemplo: {"data_limite": "2026-05-15", "titulo": "Novo Título"}
            - justificativa: frase curta explicando o motivo (gravada silenciosamente no diário)

            Retorna JSON string com payload de confirmação ou string de erro.
            """
            try:
                _ALLOWED_FIELDS = {'titulo', 'descricao', 'data_limite', 'data_inicio', 'prazo_final', 'horario_inicio', 'horario_fim', 'status', 'tags', 'area_tematica', 'tipo_acao', 'notas', 'email_link_optout'}
                today_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")
                if 'data_limite' in (alteracoes or {}):
                    val = alteracoes['data_limite']
                    if val and val not in ('-', '0000-00-00') and val < today_str:
                        return f"ERRO|A data de execução não pode ser no passado ({val})."
                if 'prazo_final' in (alteracoes or {}):
                    val = alteracoes['prazo_final']
                    if val and val not in ('-', '0000-00-00') and val < today_str:
                        return f"ERRO|O prazo final não pode ser no passado ({val})."

                def _normalizar_status_acao(valor):
                    if valor is None:
                        return valor
                    raw = str(valor).strip().lower()
                    try:
                        import unicodedata
                        raw = ''.join(c for c in unicodedata.normalize('NFD', raw) if unicodedata.category(c) != 'Mn')
                    except Exception:
                        pass
                    raw = raw.replace('_', ' ').replace('-', ' ')
                    raw = ' '.join(raw.split())
                    if raw in ('concluido', 'concluida', 'concluir', 'finalizado', 'finalizada', 'completed', 'done'):
                        return 'concluído'
                    if raw in ('stand by', 'standby', 'pausado', 'pausada', 'pausar'):
                        return 'stand-by'
                    if raw in ('em andamento', 'andamento', 'pendente', 'aberto', 'aberta', 'reabrir'):
                        return 'em andamento'
                    if raw in ('excluido', 'excluir', 'excluida', 'cancelado', 'cancelar', 'cancelada', 'deletar', 'deletado', 'apagar', 'remover'):
                        return 'excluído'
                    return valor

                task_ref = db.collection('tarefas').document(task_id)
                task_doc = task_ref.get()

                if not task_doc.exists:
                    return f"ERRO|Ação '{task_id}' não encontrada."

                task_data = task_doc.to_dict()

                if task_data.get('status') in ('concluído', 'excluído'):
                    return "ERRO|Esta ação já foi concluída ou excluída e não pode ser editada."

                # Monta o diff de campos (original vs. novo)
                alteracoes_diff = {}
                for campo, novo_valor in (alteracoes or {}).items():
                    if campo not in _ALLOWED_FIELDS:
                        continue
                    if campo == 'status':
                        novo_valor = _normalizar_status_acao(novo_valor)
                    original = task_data.get(campo)
                    if isinstance(original, dict):
                        original_str = json.dumps(original, ensure_ascii=False)
                    else:
                        original_str = ', '.join(str(v) for v in original) if isinstance(original, list) else (str(original) if original is not None else '')
                    
                    if isinstance(novo_valor, dict):
                        novo_str = json.dumps(novo_valor, ensure_ascii=False)
                    else:
                        novo_str = ', '.join(str(v) for v in novo_valor) if isinstance(novo_valor, list) else (str(novo_valor) if novo_valor is not None else '')
                    alteracoes_diff[campo] = {
                        'original': original_str,
                        'novo': novo_str,
                        'novo_raw': novo_valor
                    }

                if not alteracoes_diff:
                    return "ERRO|Nenhum campo válido para editar."

                snapshot_ts = task_data.get('data_atualizacao') or task_data.get('data_criacao', '')
                payload = {
                    'task_id': task_id,
                    'titulo': task_data.get('titulo', ''),
                    'alteracoes': alteracoes_diff,
                    'justificativa': justificativa or '',
                    'snapshot_ts': str(snapshot_ts),
                    'status': 'pending'
                }

                print(f"[Copiloto] Edição preparada para task_id={task_id}: {list(alteracoes_diff.keys())}")
                return json.dumps(payload, ensure_ascii=False)

            except Exception as _pe:
                print(f"[Copiloto] Erro ao preparar edição: {_pe}")
                return f"ERRO|{str(_pe)}"

        def preparar_reagendamento_em_lote(
            nova_data_inicio: str,
            max_por_semana: int = 5,
            estrategia: str = "data_criacao",
            filtro_data: str = None,
            task_ids: list[str] = None,
            justificativa: str = "",
        ):
            """
            Prepara reagendamento em lote de ações para confirmação interativa do usuário.
            NÃO realiza nenhuma mutação — retorna payload para card de confirmação visual.

            Parâmetros:
            - nova_data_inicio: YYYY-MM-DD — primeiro dia útil a partir do qual redistribuir as ações
            - max_por_semana: máximo de ações alocadas por semana (padrão 5)
            - estrategia: critério de ordenação das ações — "data_criacao" (padrão, mais antigas primeiro),
              "tipo_acao" (fast antes de deep), "alfa" (ordem alfabética pelo título)
            - filtro_data: YYYY-MM-DD — seleciona ações com data_limite igual a esta data (ex: hoje)
            - task_ids: lista explícita de IDs (alternativa ao filtro_data)
            - justificativa: frase curta explicando o motivo (gravada silenciosamente no diário de cada ação)

            Retorna JSON string com payload de confirmação ou string de erro.
            """
            try:
                from datetime import timedelta as _td

                if not filtro_data and not task_ids:
                    return "ERRO|Forneça filtro_data (YYYY-MM-DD) ou task_ids (lista de IDs)."

                tasks = []
                if task_ids:
                    for tid in (task_ids or []):
                        tdoc = db.collection('tarefas').document(str(tid)).get()
                        if tdoc.exists:
                            t = tdoc.to_dict()
                            if t.get('status') not in ('concluído', 'cancelado'):
                                tasks.append({'_doc_id': str(tid), **t})
                else:
                    q = db.collection('tarefas')\
                        .where('data_limite', '==', filtro_data)\
                        .where('status', 'in', ['em andamento', 'stand-by'])\
                        .get()
                    for qdoc in q:
                        tasks.append({'_doc_id': qdoc.id, **qdoc.to_dict()})

                if not tasks:
                    return "ERRO|Nenhuma ação encontrada com os critérios informados."

                if estrategia == 'tipo_acao':
                    tasks.sort(key=lambda x: (0 if x.get('tipo_acao') == 'fast' else 1, x.get('data_criacao', '')))
                elif estrategia == 'alfa':
                    tasks.sort(key=lambda x: x.get('titulo', '').lower())
                else:
                    tasks.sort(key=lambda x: x.get('data_criacao', ''))

                try:
                    from datetime import date as _date
                    start_date = datetime.strptime(nova_data_inicio, "%Y-%m-%d").date()
                    today_date = datetime.now(timezone.utc).date()
                    if start_date < today_date:
                        start_date = today_date
                except ValueError:
                    return f"ERRO|Formato de data inválido: '{nova_data_inicio}'. Use YYYY-MM-DD."

                def _next_weekday(d):
                    while d.weekday() >= 5:
                        d += _td(days=1)
                    return d

                day_cursor = _next_weekday(start_date)
                count_this_week = 0
                items = []

                for task in tasks:
                    if count_this_week >= max_por_semana:
                        days_to_monday = 7 - day_cursor.weekday()
                        day_cursor += _td(days=days_to_monday)
                        day_cursor = _next_weekday(day_cursor)
                        count_this_week = 0

                    items.append({
                        'task_id': task['_doc_id'],
                        'titulo': task.get('titulo', ''),
                        'data_limite_original': task.get('data_limite', ''),
                        'horario_inicio_original': task.get('horario_inicio'),
                        'horario_fim_original': task.get('horario_fim'),
                        'nova_data_limite': day_cursor.strftime("%Y-%m-%d"),
                        'novo_horario_inicio': None,
                        'novo_horario_fim': None,
                    })

                    count_this_week += 1
                    day_cursor += _td(days=1)
                    day_cursor = _next_weekday(day_cursor)

                payload = {
                    'items': items,
                    'justificativa': justificativa or f"Reagendamento em lote para semana de {nova_data_inicio}.",
                    'status': 'pending',
                    'created_at': datetime.now(timezone.utc).isoformat(),
                }

                print(f"[Copiloto] Reagendamento em lote preparado: {len(items)} ações.")
                return json.dumps(payload, ensure_ascii=False)

            except Exception as _re:
                print(f"[Copiloto] Erro em preparar_reagendamento_em_lote: {_re}")
                return f"ERRO|{str(_re)}"

        def preparar_remocao_horarios_em_lote(
            filtro_data: str = None,
            task_ids: list[str] = None,
            justificativa: str = "",
        ):
            """
            Prepara a remoção de horários (horario_inicio/horario_fim) de múltiplas ações em lote.
            NÃO realiza nenhuma mutação — retorna payload para card de confirmação visual.

            Parâmetros:
            - filtro_data: YYYY-MM-DD — seleciona ações com data_limite igual a esta data (ex: hoje)
            - task_ids: lista explícita de IDs (alternativa ao filtro_data)
            - justificativa: frase curta explicando o motivo

            Use quando o usuário pedir para remover/limpar/tirar horários de várias ações de uma vez.
            Para remover o horário de uma única ação, use preparar_edicao_acao com
            alteracoes={"horario_inicio": None, "horario_fim": None}.
            Retorna JSON string com payload de confirmação ou string de erro.
            """
            try:
                if not filtro_data and not task_ids:
                    return "ERRO|Forneça filtro_data (YYYY-MM-DD) ou task_ids (lista de IDs)."

                tasks = []
                if task_ids:
                    for tid in (task_ids or []):
                        tdoc = db.collection('tarefas').document(str(tid)).get()
                        if tdoc.exists:
                            t = tdoc.to_dict()
                            if t.get('status') not in ('concluído', 'excluído'):
                                tasks.append({'_doc_id': str(tid), **t})
                else:
                    q = db.collection('tarefas')\
                        .where('data_limite', '==', filtro_data)\
                        .where('status', 'in', ['em andamento', 'stand-by'])\
                        .get()
                    for qdoc in q:
                        tasks.append({'_doc_id': qdoc.id, **qdoc.to_dict()})

                tasks_com_horario = [t for t in tasks if t.get('horario_inicio')]

                if not tasks_com_horario:
                    if tasks:
                        return "ERRO|Nenhuma das ações encontradas possui horário definido."
                    return "ERRO|Nenhuma ação encontrada com os critérios informados."

                items = []
                for task in tasks_com_horario:
                    items.append({
                        'task_id': task['_doc_id'],
                        'titulo': task.get('titulo', ''),
                        'data_limite_original': task.get('data_limite', ''),
                        'horario_inicio_original': task.get('horario_inicio'),
                        'horario_fim_original': task.get('horario_fim'),
                        'nova_data_limite': task.get('data_limite', ''),
                        'novo_horario_inicio': None,
                        'novo_horario_fim': None,
                    })

                payload = {
                    'items': items,
                    'justificativa': justificativa or "Remoção de horários em lote via Copiloto Hermes.",
                    'status': 'pending',
                    'created_at': datetime.now(timezone.utc).isoformat(),
                }

                print(f"[Copiloto] Remoção de horários preparada: {len(items)} ações.")
                return json.dumps(payload, ensure_ascii=False)

            except Exception as _re:
                print(f"[Copiloto] Erro em preparar_remocao_horarios_em_lote: {_re}")
                return f"ERRO|{str(_re)}"

        def buscar_contato(termo: str, limite: int = 5):
            """Busca contatos por nome, email ou tag em perfil_pessoas. Retorna dados do contato incluindo telefone e whatsapp_chat_id (para vincular com conversas do WhatsApp)."""
            try:
                termo_lower = (termo or "").strip().lower()
                if not termo_lower:
                    return "ERRO|Termo de busca vazio."
                docs = db.collection('perfil_pessoas').limit(500).stream()
                candidatos = []
                for d in docs:
                    pdata = d.to_dict() or {}
                    nome = (pdata.get('nome') or '').lower()
                    email = (pdata.get('email') or '').lower()
                    tags = [str(t).lower() for t in (pdata.get('tags') or [])]
                    score = 0.0
                    if nome == termo_lower or email == termo_lower:
                        score = 1.0
                    elif termo_lower in nome:
                        score = 0.8
                    elif termo_lower in email:
                        score = 0.7
                    elif any(termo_lower in t for t in tags):
                        score = 0.5
                    if score > 0:
                        candidatos.append({
                            'pessoa_id': d.id,
                            'nome': pdata.get('nome', ''),
                            'email': pdata.get('email', ''),
                            'telefone': pdata.get('telefone', ''),
                            'whatsapp_chat_id': pdata.get('whatsapp_chat_id', ''),
                            'tags': pdata.get('tags', []),
                            'score': score,
                        })
                candidatos.sort(key=lambda x: -x['score'])
                return json.dumps({'candidatos': candidatos[: max(1, int(limite or 5))]}, ensure_ascii=False)
            except Exception as _re:
                print(f"[Copiloto] Erro em buscar_contato: {_re}")
                return f"ERRO|{str(_re)}"

        def preparar_vinculo_contatos(task_id: str, mencoes: list[dict]):
            """Prepara payload de confirmação para vincular pessoas a uma tarefa. Não grava nada."""
            try:
                if not task_id:
                    return "ERRO|task_id é obrigatório."
                tdoc = db.collection('tarefas').document(str(task_id)).get()
                if not tdoc.exists:
                    return f"ERRO|Tarefa '{task_id}' não encontrada."
                payload = {
                    'kind': 'contact_link',
                    'task_id': task_id,
                    'task_titulo': (tdoc.to_dict() or {}).get('titulo', ''),
                    'mencoes': mencoes or [],
                    'status': 'pending',
                    'created_at': datetime.now(timezone.utc).isoformat(),
                }
                return json.dumps(payload, ensure_ascii=False)
            except Exception as _re:
                print(f"[Copiloto] Erro em preparar_vinculo_contatos: {_re}")
                return f"ERRO|{str(_re)}"

        def preparar_atualizacao_contato(nome: str, campos_novos: dict, justificativa: str, pessoa_id: str = None):
            """Prepara payload de confirmação para criar/atualizar contato. Não grava nada."""
            try:
                if not nome or not (campos_novos or {}):
                    return "ERRO|nome e campos_novos são obrigatórios."
                modo = 'update' if pessoa_id else 'create'
                contato_atual = None
                if pessoa_id:
                    pdoc = db.collection('perfil_pessoas').document(str(pessoa_id)).get()
                    if not pdoc.exists:
                        return f"ERRO|Contato '{pessoa_id}' não encontrado."
                    contato_atual = pdoc.to_dict()
                payload = {
                    'kind': 'contact_upsert',
                    'modo': modo,
                    'pessoa_id': pessoa_id,
                    'nome': nome,
                    'contato_atual': contato_atual,
                    'campos_novos': campos_novos,
                    'justificativa': justificativa or '',
                    'status': 'pending',
                    'created_at': datetime.now(timezone.utc).isoformat(),
                }
                return json.dumps(payload, ensure_ascii=False)
            except Exception as _re:
                print(f"[Copiloto] Erro em preparar_atualizacao_contato: {_re}")
                return f"ERRO|{str(_re)}"

        def registrar_interacao_contato(pessoa_id: str, descricao: str, tarefa_id: str = None, sessao_copiloto_id: str = None):
            """Registra interação silenciosa no histórico de um contato. Grava direto, sem confirmação."""
            try:
                if not pessoa_id or not descricao:
                    return "ERRO|pessoa_id e descricao são obrigatórios."
                pref = db.collection('perfil_pessoas').document(str(pessoa_id)).get()
                if not pref.exists:
                    return f"ERRO|Contato '{pessoa_id}' não encontrado."
                desc_short = str(descricao)[:280]
                now_iso = datetime.now(timezone.utc).isoformat()
                sess_id = sessao_copiloto_id or session_id
                payload = {
                    'pessoa_id': str(pessoa_id),
                    'descricao': desc_short,
                    'tipo': 'mencao_copiloto',
                    'data': now_iso,
                    'data_criacao': now_iso,
                }
                if tarefa_id:
                    payload['tarefa_id'] = str(tarefa_id)
                if sess_id:
                    payload['sessao_copiloto_id'] = str(sess_id)
                # Marca contato com tag Copiloto para aparecer no filtro
                try:
                    tags_atuais = (pref.to_dict() or {}).get('tags') or []
                    if 'Copiloto' not in tags_atuais:
                        db.collection('perfil_pessoas').document(str(pessoa_id)).update({
                            'tags': tags_atuais + ['Copiloto']
                        })
                except Exception as _tag_err:
                    print(f"[Copiloto] Aviso: falha ao marcar tag Copiloto em {pessoa_id}: {_tag_err}")
                new_ref = db.collection('interacoes_pessoas').document()
                new_ref.set(payload)
                return json.dumps({'status': 'ok', 'interacao_id': new_ref.id}, ensure_ascii=False)
            except Exception as _re:
                print(f"[Copiloto] Erro em registrar_interacao_contato: {_re}")
                return f"ERRO|{str(_re)}"

        # Configuração do Chat com ferramentas
        model_id = COPILOT_CHAT_MODEL

        # Roteamento de modelo: smalltalk óbvio (saudação/agradecimento/confirmação
        # curta sem anexo) não precisa do modelo frontier — cai para o tier barato.
        _SMALLTALK_RE = re.compile(
            r"^(oi|ol[aá]|bom dia|boa tarde|boa noite|obrigad[oa]|valeu|ok|blz|beleza|"
            r"tudo bem|e a[ií]|opa|show|top|perfeito|excelente|legal|entendi|"
            r"haha+|kk+|rsrs+)[\s!.,?😀-🙏]*$",
            re.IGNORECASE,
        )
        if not drive_file_id and not task_id and len(prompt) <= 40 and _SMALLTALK_RE.match(prompt.strip()):
            model_id = GEMINI_BALANCED_MODEL
            print(f"[Copiloto] Smalltalk detectado — usando modelo {model_id}")

        from datetime import datetime
        from zoneinfo import ZoneInfo
        tz = ZoneInfo("America/Sao_Paulo")
        now_local = datetime.now(tz)
        today_str = now_local.strftime("%Y-%m-%d")
        time_str = now_local.strftime("%H:%M")

        # Busca catálogo de sistemas para o "de-para" exato que o usuário solicitou
        try:
            sistemas_docs = _cached_collection_list(db, 'sistemas_detalhes')
            catalogo_sistemas = []
            for s_data in sistemas_docs:
                s_nome = s_data.get('nome', 'Sem Nome')
                s_id = s_data['id']
                catalogo_sistemas.append(f"- {s_nome}: {s_id}")
            sistemas_str = "\n".join(catalogo_sistemas) if catalogo_sistemas else "Nenhum sistema cadastrado."
        except Exception as e:
            print(f"Erro ao buscar catálogo de sistemas: {e}")
            sistemas_str = "Erro ao carregar catálogo."

        # --- RECUPERAÇÃO DE HISTÓRICO DA SESSÃO ---
        # Carregado antes da montagem do prompt: o texto recente também alimenta
        # o gating dos protocolos condicionais (slides, formulários, etc.).
        history = []
        history_plain = ""
        if session_id:
            try:
                msg_docs = db.collection('sessoes_copiloto').document(session_id)\
                    .collection('mensagens')\
                    .order_by('timestamp', direction=firestore.Query.DESCENDING)\
                    .limit(6).get()

                # Inverte para ordem cronológica conforme exigido pelo SDK
                raw_msgs = list(reversed(msg_docs))
                _history_texts = []
                for mdoc in raw_msgs:
                    m = mdoc.to_dict()
                    # Mapeia roles: o SDK espera 'user' ou 'model' (não 'assistant')
                    role = m.get('role')
                    if role == 'assistant':
                        role = 'model'

                    # Teto por mensagem: relatórios/markdown longos no histórico
                    # inflam o input reenviado a cada rodada de tool-calling.
                    _msg_text = str(m.get('content') or "")
                    if len(_msg_text) > COPILOT_HISTORY_MSG_MAXCHARS:
                        _msg_text = _msg_text[:COPILOT_HISTORY_MSG_MAXCHARS] + "\n[…trecho truncado para economia de contexto…]"

                    history.append(types.Content(
                        role=role,
                        parts=[types.Part(text=_msg_text)]
                    ))
                    _history_texts.append(_msg_text)
                history_plain = " ".join(_history_texts)
            except Exception as e:
                print(f"Erro ao carregar histórico da sessão {session_id}: {e}")

        mode_context = ""
        if copilot_mode == "finance":
            mode_context = (
                "## MODO FINANCEIRO ATIVO\n"
                "Você está atuando como Copiloto Financeiro do Hermes. Priorize fluxo de caixa, gastos, reserva de emergência, metas, obrigações, orçamento mensal e dúvidas educacionais sobre investimentos.\n"
                "Para qualquer número financeiro interno, use consultar_financas_v2 antes de concluir. Não use categorias de lançamentos como base analítica enquanto a classificação estiver em revisão.\n"
                "Você pode explicar tipos de investimento em caráter educativo, mas não deve prometer rentabilidade, recomendar compra/venda específica ou tratar isso como consultoria financeira regulada.\n"
                "Quando sugerir próximos passos, escreva como proposta para o usuário avaliar; não crie ações, metas ou lançamentos sem confirmação explícita.\n\n"
            )

        elif copilot_mode == "saude":
            mode_context = (
                "## MODO COPILOTO DE SAUDE ATIVO\n"
                "Voce esta atuando como Copiloto de Saude do Hermes. Priorize leitura e interpretacao operacional de peso, caminhada/passos, calorias, sono e sinais de dor.\n"
                "Para qualquer metrica de saude atual, use consultar_saude antes de concluir. Nao trate habitos diarios, treino de forca, rotina lombar, flexoes, barras, prancha, ponte, bird-dog ou agachamentos como parte ativa do sistema.\n"
                "Quando sugerir proximos passos, escreva como proposta para o usuario avaliar; nao crie registros, metas ou logs sem confirmacao explicita.\n\n"
                "## CONTEXTO CLINICO PESSOAL DO USUARIO\n"
                "Use este contexto apenas para interpretar dor, sono, caminhada e alimentacao com prudencia, sem prescrever tratamento medico.\n\n"
                "### Diagnosticos informados\n"
                "- Sindrome de Bertolotti: variacao lombo-sacra com historico de dor lombar cronica e crises matinais.\n"
                "- DRGE + esofagite erosiva grau A + hiato diafragmatico alargado, confirmado por endoscopia em 18/07/2024.\n\n"
                "### Ergonomia ocupacional\n"
                "- Evitar longos periodos estaticos; pausas breves de movimento podem ser propostas como cuidado geral, sem registrar rotina no sistema.\n\n"
                "### Caminhada\n"
                "- Priorizar superficies planas e progressao conservadora quando houver dor. Use passos, distancia e minutos ativos como metricas principais.\n\n"
                "### Dieta anti-refluxo e anti-inflamatoria\n"
                "- Evitar frituras, ultraprocessados, cafe forte, bebidas carbonatadas, pimenta e chocolate. Evitar comer nas 3h antes de deitar.\n"
                "- Favorecer opcoes simples e bem toleradas, como aveia, proteinas magras, azeite extravirgem, abacate e nozes quando fizer sentido no contexto.\n\n"
            )
        elif copilot_mode == "estrategia":
            # Snapshot completo da estratégia pessoal do usuário, com IDs, para que o
            # copiloto possa conversar sobre os objetivos e operar as ferramentas CRUD.
            _estrategia_snapshot = "Nenhum objetivo estratégico cadastrado ainda."
            try:
                _estrategia_docs = list(
                    db.collection('estrategia_pessoal').where('userId', '==', user_uid).limit(40).stream()
                ) if user_uid else []

                def _fmt_itens_estrategia(itens, prefixo):
                    linhas = []
                    for it in (itens or []):
                        if isinstance(it, str):
                            linhas.append(f"      - [sem-id] {it}")
                        else:
                            marca = "✓" if it.get("concluido") else "○"
                            linhas.append(f"      - {marca} (id={it.get('id', '?')}) {it.get('descricao', '')}")
                    return "\n".join(linhas) if linhas else "      - (nenhum)"

                _blocos = []
                for _doc in _estrategia_docs:
                    _d = _doc.to_dict() or {}
                    _metrica = _d.get("metricaAlvo") or {}
                    _metrica_str = ""
                    if _d.get("tipoMeta") == "absoluta" and _metrica:
                        _metrica_str = (
                            f"\n    Métrica: {_metrica.get('valorAtual', '?')}/{_metrica.get('valorObjetivo', '?')} "
                            f"{_metrica.get('unidade', '')} (inicial {_metrica.get('valorInicial', '?')})"
                        )
                    _diretrizes = _d.get("diretrizesDerivadas") or []
                    _diretrizes_str = "\n".join(f"      - {dz}" for dz in _diretrizes) if _diretrizes else "      - (nenhuma)"
                    _blocos.append(
                        f"• OBJETIVO (id={_doc.id}) | pilar={_d.get('pilar', '?')} | tipo={_d.get('tipoMeta', '?')} | status={_d.get('status', '?')}\n"
                        f"    \"{_d.get('objetivoMacro', '')}\"{_metrica_str}\n"
                        f"    Diretrizes:\n{_diretrizes_str}\n"
                        f"    Indicadores contínuos:\n{_fmt_itens_estrategia(_d.get('indicadoresSucesso'), 'indicador')}\n"
                        f"    Marcos pontuais:\n{_fmt_itens_estrategia(_d.get('marcos'), 'marco')}"
                    )
                if _blocos:
                    _estrategia_snapshot = "\n\n".join(_blocos)
            except Exception as _est_err:
                print(f"[Copiloto] Falha ao carregar snapshot de estratégia: {_est_err}")

            mode_context = (
                "## MODO ESTRATÉGIA ATIVO\n"
                "Você foi aberto a partir do módulo de Estratégia do Hermes. Seu foco EXCLUSIVO nesta conversa é a estratégia pessoal do usuário: "
                "objetivos macro, pilares, diretrizes, indicadores de sucesso e marcos. Converse sobre coerência de longo prazo, prioridades, trade-offs e progresso.\n"
                "Você tem ferramentas de escrita disponíveis SOMENTE neste modo: criar_objetivo_estrategico, editar_objetivo_estrategico, gerenciar_item_estrategico e excluir_objetivo_estrategico.\n"
                "REGRAS DE OPERAÇÃO:\n"
                "1. Sempre apresente um rascunho claro do que pretende criar/alterar/excluir e só execute a ferramenta APÓS confirmação explícita do usuário.\n"
                "2. Para excluir um objetivo, exija confirmação inequívoca — a operação é irreversível.\n"
                "3. Use SEMPRE o 'id' exato exibido no snapshot abaixo ao editar, gerenciar itens ou excluir. Nunca invente IDs.\n"
                "4. Ao editar diretrizes, lembre que a lista é substituída por inteiro: monte a lista final desejada a partir das diretrizes atuais.\n"
                "5. Para indicadores/marcos individuais use gerenciar_item_estrategico (preserva IDs e ações vinculadas).\n"
                "6. Não toque em tarefas operacionais, finanças ou outros módulos aqui; mantenha o foco na estratégia.\n\n"
                "## SNAPSHOT ATUAL DA ESTRATÉGIA DO USUÁRIO\n"
                f"{_estrategia_snapshot}\n\n"
            )

        system_instruction_nucleo = (
            "Você é o Copiloto Hermes, estrategista sênior de processos."
            "\n\n## CORE ESTÁTICO DO COPILOTO\n"
            f"{copilot_core.get('content', '')}\n\n"
            "## PERSONALIDADE DINÂMICA ATUAL\n"
            f"{copilot_soul.get('content') or json.dumps(copilot_soul, ensure_ascii=False)}\n\n"
            "## PERFIL OPERACIONAL DO USUÁRIO\n"
            f"{_format_ai_profile_for_prompt(ai_profile)}\n\n"
            f"{tools_routing_context}"
            "\n\nCATÁLOGO DE SISTEMAS (Mapeamento Exato de Nome para ID):\n"
            f"{sistemas_str}\n\n"
            "Ao realizar diagnósticos ou operações em sistemas, utilize SEMPRE o ID técnico do catálogo acima correspondente ao nome citado pelo usuário.\n\n"
            "Seu tom de voz: Consultivo, analítico e absurdamente conciso. "
            "Use bullet points para melhorar a legibilidade. "
            "\n\n"
            "## POSTURA CRÍTICA (SOCRÁTICA)\n"
            "Você NÃO é um copiloto condescendente. Seu valor está em pensar criticamente JUNTO com o usuário, não em concordar por reflexo.\n"
            "- Antes de validar um plano ou decisão, cheque as premissas: o que está sendo assumido sem evidência? O que falta (dado, prazo, dependência, recurso)?\n"
            "- Aponte explicitamente pontos frágeis, riscos e pontos cegos — de forma curta, específica e sem alarmismo.\n"
            "- Quando a decisão tiver peso real, faça no máximo UMA pergunta socrática afiada que force o usuário a examinar o ponto mais frágil. Em tarefas triviais, apenas execute.\n"
            "- Nunca bajule ('ótima ideia!', 'perfeito!') por hábito. Discorde com fundamento quando os fatos não sustentarem o pedido.\n"
            "- Calibre a criticidade ao risco: alto risco = mais escrutínio; rotina = direto ao ponto. Criticidade não é prolixidade.\n"
            "\n"
            "## REGRA ABSOLUTA — SIGLAS, TERMOS TÉCNICOS E ERROS DE BACKEND\n\n"
            "PROIBIÇÃO TOTAL DE EXPANSÃO ARBITRÁRIA DE SIGLAS:\n"
            "Você JAMAIS deve inferir, adivinhar, expandir ou traduzir siglas, acrônimos ou\n"
            "termos técnicos que o usuário fornecer. Se o usuário disser \"IRP\", você trata\n"
            "\"IRP\" como uma string opaca e literal — não é \"Imposto de Renda\", não é\n"
            "\"Internal Revenue Policy\", não é nada que você \"acha que pode ser\". Você passa\n"
            "o termo exatamente como recebido para as ferramentas de busca. Se nenhum\n"
            "documento retornar resultado, sua resposta é: \"Nenhum registro encontrado para\n"
            "o termo exato 'IRP'. Você pode confirmar a sigla ou fornecer mais contexto?\"\n"
            "Não improvise. Não complete. Não alucine.\n\n"
            "OBRIGAÇÃO DE TRANSPARÊNCIA EM ERROS TÉCNICOS:\n"
            "Se qualquer ferramenta (buscar_acervo, buscar_tarefas ou similar) retornar um\n"
            "campo \"erro\" não-nulo, você DEVE reproduzir o conteúdo desse campo palavra por\n"
            "palavra na sua resposta, sem parafrasear, sem suavizar e sem omitir. Formato\n"
            "obrigatório:\n\n"
            "  ⚠️ Erro técnico na ferramenta [nome_da_ferramenta]:\n"
            "  [conteúdo literal do campo \"erro\"]\n\n"
            "Após reportar o erro, peça ao usuário que acione o suporte técnico com essa\n"
            "mensagem exata. Você NÃO deve tentar responder a pergunta original como se o\n"
            "erro não tivesse ocorrido.\n\n"
            "TRATAMENTO DE CONFLITOS: Se a prática [Grafo] divergir do manual [Acervo], exponha explicitamente.\n"
            "Sempre termine propostas de ajuste de plano de ação dentro de: [PROPOSAL]{...}[/PROPOSAL]\n\n"
            "## NAVEGAÇÃO ENTRE AÇÕES (CRÍTICO)\n"
            "Para que o usuário possa gerir as ações, você DEVE SEMPRE transformar o nome de qualquer tarefa citada em um link clicável.\n"
            "FORMATO OBRIGATÓRIO: `[Nome da Tarefa](task:ID)`\n"
            "Exemplo: 'Verifiquei que a ação [Analisar Edital](task:xyz123) está atrasada.'\n"
            "PROIBIÇÃO: Você NUNCA deve gerar links com `https://` ou URLs de sites externos para se referir a ações do sistema. Use APENAS o prefixo `task:`.\n"
            "Use os IDs retornados pelas ferramentas `consultar_historico_acoes` ou `obter_contexto_tela`.\n\n"
            "## BUSCA PROATIVA DE CONTEXTO\n"
            "Se o usuário pedir algo genérico (ex: 'o que temos hoje' ou 'status das atividades') e você não tiver uma tarefa em foco:\n"
            "1. NÃO peça o ID ou Área Temática imediatamente.\n"
            "2. Use `consultar_historico_acoes(query='', data_limite_inicio='YYYY-MM-DD')` para filtrar por prazo/vencimento se o usuário mencionar datas.\n"
            "3. Se não houver data específica, use `query=''` para listar as tarefas mais recentes do sistema.\n"
            "4. Analise os resultados e peça clarificação apenas se necessário.\n\n"
            "PADRÕES DE LISTAGEM POR PERÍODO (calcule as datas e passe diretamente — não pergunte ao usuário):\n"
            "- 'ações de hoje' → data_limite_inicio=hoje, data_limite_fim=hoje\n"
            "- 'ações de amanhã' → data_limite_inicio=amanhã, data_limite_fim=amanhã\n"
            "- 'ações desta semana' → data_limite_inicio=segunda-feira da semana atual, data_limite_fim=domingo da semana atual\n"
            "- 'ações da próxima semana' → data_limite_inicio=segunda-feira da próxima semana, data_limite_fim=domingo da próxima semana\n"
            "- 'ações de [dia da semana]' → calcule a data correspondente e passe como inicio e fim\n"
            "Sempre use o formato YYYY-MM-DD. A data de hoje é injetada no contexto de sistema.\n\n"
            "## REGRA DE INTEGRIDADE DOCUMENTAL (CRÍTICO — PENA DE FALHA SISTÊMICA)\n"
            "Se o usuário perguntar sobre valores, quantidades, itens ou cláusulas de um arquivo presente\n"
            "no campo 'arquivos_disponiveis' do contexto da tarefa, você é ESTRITAMENTE PROIBIDO de:\n"
            "  a) Deduzir a resposta com base no seu treinamento.\n"
            "  b) Mesclar fragmentos de buscas vetoriais globais (buscar_arquivos_acervo) com dados deste arquivo.\n\n"
            "PROTOCOLO OBRIGATÓRIO:\n"
            "  1. Verifique se o arquivo está listado em 'arquivos_disponiveis' (via obter_contexto_tela).\n"
            "  2. Chame ler_documento_na_integra(drive_file_id=<ID>, query_especifica=<pergunta exata do usuário>).\n"
            "  3. Baseie sua resposta EXCLUSIVAMENTE no retorno desta ferramenta.\n"
            "  4. Se a ferramenta declarar que a informação não existe, reproduza essa declaração sem inventar alternativas.\n"
            "NUNCA misture dados numéricos (valores, itens, quantidades) de processos ou documentos distintos.\n\n"
            "## CRIAÇÃO E ALOCAÇÃO DE AÇÕES (CRÍTICO)\n\n"
            "Quando o usuário solicitar a criação de uma ação/tarefa, siga OBRIGATORIAMENTE este protocolo:\n\n"
            "1. POR PADRÃO (Criação Normal): Não proponha nem defina horários de início/fim (campos horario_inicio e horario_fim devem ser nulos/vazios), definindo apenas o dia de trabalho (data_limite = DATA DE EXECUÇÃO, não um prazo). data_limite e prazo_final são campos DISTINTOS: data_limite é o dia em que o trabalho será feito; prazo_final é o prazo real/fatal da entrega, só deve ser preenchido se o usuário mencionar explicitamente uma data de prazo diferente da data de execução. Se o usuário citar apenas uma data, use-a como data_limite e deixe prazo_final vazio. Se citar duas (ex.: \"comece dia X, mas o prazo final é dia Y\"), use data_limite=X e prazo_final=Y. Apresente o draft ao usuário mostrando o horário como vazio/não definido.\n"
            "2. SE O USUÁRIO PEDIR EXPLICITAMENTE PARA AGENDAR UM HORÁRIO ESPECÍFICO:\n"
            "   - Você DEVE usar consultar_agenda(data_inicio, data_fim) ou encontrar_slot_livre(data) ANTES de apresentar qualquer proposta ao usuário.\n"
            "   - Se o agendamento for para hoje, o horário inicial DEVE ser sempre posterior ao horário local atual.\n"
            "   - Se houver colisão (conflito detectado): trave a inserção perguntando se ele quer Forçar a sobreposição ou Buscar próximo horário livre.\n"
            "   - Restrição: A agenda opera estritamente entre 08:00 e 19:00, dentro de uma janela de 7 dias úteis.\n\n"
            "ETAPA 1 — DRAFT (apresentar antes de criar):\n"
            "Apresente um resumo estruturado para o usuário confirmar:\n"
            "  📋 **Draft da Ação**\n"
            "  - **Título:** [título]\n"
            "  - **Data de Execução:** [ex: 2026-05-15]\n"
            "  - **Prazo Final (se houver):** [opcional, só se diferente da data de execução]\n"
            "  - **Horário:** [vazio/não definido, OU início/fim ex: 14:00 às 14:30 se pedido e validado]\n"
            "  - **Área Temática:** [área ou projeto baseada no contexto; PERGUNTE SE FOR AMBÍGUA]\n"
            "  - **Tipo:** [fast / deep]\n"
            "  Confirma a criação desta ação?\n\n"
            "ETAPA 2 — CONFIRMAÇÃO:\n"
            "Só chame criar_acao_no_sistema após receber confirmação explícita ('sim', 'confirma', 'pode criar', 'ok', etc.).\n"
            "Se o usuário ajustar algum campo no draft, incorpore as correções antes de criar.\n\n"
            "ETAPA 3 — COMMIT E LINK:\n"
            "Após criar_acao_no_sistema retornar 'OK|{ID}', responda obrigatoriamente:\n"
            "  ✅ Ação criada: [Título da Ação](task:{ID})\n"
            "  Você também pode clicar no link acima para abrir a sala de operações imediatamente.\n"
            "Se retornar 'ERRO|{detalhe}', responda:\n"
            "  ⚠️ Erro ao criar ação: {detalhe}\n\n"
            "EXTRAÇÃO DE CONTEXTO PARA O DRAFT:\n"
            "- Se houver um taskId ativo, use obter_contexto_tela() para inferir área temática, tags e contexto.\n"
            "- Deduza que a nova ação pode ser sub-tarefa ou relacionada ao contexto ativo.\n"
            "- Use o histórico da conversa para preencher descricao e plano_acao automaticamente.\n\n"
            "## EDIÇÃO DE PLANO DE AÇÃO — PADRÃO DRAFT-THEN-COMMIT (CRÍTICO)\n\n"
            "Quando o usuário solicitar alteração, adição, remoção ou reestruturação de passos de um plano de ação:\n\n"
            "ETAPA 0 — EXTRAÇÃO DE CONTEXTO OBRIGATÓRIA:\n"
            "Chame obter_contexto_tela() para capturar o taskId e o plano de ação atual (com os IDs dos passos).\n"
            "Nunca suponha IDs de passos — leia-os do resultado da ferramenta.\n\n"
            "ETAPA 1 — DRAFT (apresentar antes de editar):\n"
            "Nunca chame editar_plano_acao imediatamente. Primeiro, apresente o novo plano proposto:\n"
            "  ✏️ **Novo Plano de Ação proposto**\n"
            "  1. [passo 1]\n"
            "  2. [passo 2]\n"
            "  *(passos removidos, adicionados ou reordenados em relação ao plano atual)*\n"
            "  Confirma a atualização do plano?\n\n"
            "ETAPA 2 — CONFIRMAÇÃO:\n"
            "Só chame editar_plano_acao após confirmação explícita do usuário ('sim', 'confirma', 'pode atualizar', etc.).\n"
            "Ao montar novo_plano, inclua o campo 'id' para passos existentes (preserva status de conclusão via fuzzy match no backend).\n"
            "Omita o 'id' apenas para passos genuinamente novos.\n\n"
            "ETAPA 3 — COMMIT E CONFIRMAÇÃO:\n"
            "Se editar_plano_acao retornar 'OK', responda:\n"
            "  ✅ Plano de ação atualizado com sucesso.\n"
            "Se retornar 'ERRO|{detalhe}', responda:\n"
            "  ⚠️ Erro ao atualizar plano: {detalhe}\n\n"
            "PARÂMETRO justificativa_diario:\n"
            "Gere automaticamente uma frase concisa descrevendo o que foi alterado e por quê (ex: 'Adicionado passo de validação jurídica a pedido do usuário.').\n"
            "O usuário não precisa aprovar este texto — é gravado silenciosamente no diário da tarefa.\n\n"
            "## EDIÇÃO DE CAMPOS DE AÇÕES — CARD DE CONFIRMAÇÃO (CRÍTICO)\n\n"
            "Quando o usuário solicitar alteração de campos de uma ação existente (título, prazo, status,\n"
            "área temática, tags, descrição, notas, tipo), siga OBRIGATORIAMENTE este protocolo:\n\n"
            "ETAPA 0 — BUSCA OBRIGATÓRIA:\n"
            "Use consultar_historico_acoes() para localizar a ação com precisão.\n"
            "Nunca suponha um task_id sem buscar. Prefira a correspondência mais exata ao nome/contexto informado.\n"
            "Se houver ambiguidade, apresente as opções ao usuário antes de prosseguir.\n\n"
            "## MEMÓRIA E COGNIÇÃO AUTÔNOMA\n"
            "Quando identificar uma regra de negócio durável, uma preferência operacional estável ou um fato global útil em conversas futuras, "
            "acione salvar_memoria_global(fato, categoria) silenciosamente.\n"
            "Categorias válidas: 'regra_global' e 'fato_isolado'.\n"
            "Nunca grave memória para ruído passageiro, opiniões momentâneas, mensagens genéricas ou detalhes descartáveis.\n"
            "Também NÃO grave saudações, confirmações simples, contexto efêmero desta conversa, anexos transitórios ou decisões ainda não estabilizadas.\n"
            "Prefira gravar apenas convenções permanentes, preferências recorrentes, regras operacionais, fontes de verdade e fatos reutilizáveis.\n"
            "Se salvar_memoria_global retornar status='conflict', interrompa a automação e pergunte explicitamente ao usuário qual versão deve permanecer verdadeira.\n"
            "Após o usuário decidir, use resolver_conflito_memoria() para convergir a fonte de verdade.\n"
            "Se o usuário pedir mudança de tom, estilo, profundidade ou comportamento, use atualizar_personalidade() para reescrever a personalidade dinâmica.\n\n"
            "ETAPA 1 — PROPOSTA VIA FERRAMENTA:\n"
            "## POPS OPERACIONAIS\n"
            "Se o usuário pedir explicitamente para criar, cadastrar, registrar ou atualizar um POP/Procedimento Operacional Padrão,\n"
            "use salvar_pop_global(titulo, gatilhos, instrucao_sistema).\n"
            "Use essa ferramenta apenas quando houver intenção clara de persistir um POP reutilizável no Gestor de POPs.\n"
            "Gatilhos devem ser uma lista de frases curtas que disparam o POP.\n"
            "Após salvar_pop_global com sucesso, deixe claro para o usuário que o POP foi salvo ou atualizado no Gestor de POPs.\n\n"
            "Chame preparar_edicao_acao(task_id, alteracoes, justificativa) com:\n"
            "- task_id: ID da ação encontrada\n"
            "- alteracoes: dicionário com APENAS os campos que mudam. Ex: {\"data_limite\": \"2026-05-15\"}\n"
            "- justificativa: frase curta explicando o motivo (gravada silenciosamente no diário)\n"
            "Esta ferramenta NÃO faz mutação no banco. Ela prepara o payload para um card de confirmação visual.\n\n"
            "ETAPA 2 — AGUARDAR O CARD:\n"
            "Após chamar preparar_edicao_acao com sucesso, sua resposta de texto DEVE ser APENAS:\n"
            "  ✏️ Preparei a edição para sua confirmação. Verifique o card abaixo e confirme ou cancele.\n"
            "NÃO repita os dados da edição no texto — eles já estão no card visual.\n"
            "NÃO chame nenhuma outra ferramenta de escrita nesta mensagem.\n\n"
            "ETAPA 3 — APÓS CONFIRMAÇÃO:\n"
            "A confirmação ocorre pelo clique no botão do card — não pelo chat de texto.\n"
            "Se o usuário disser 'confirmo' ou 'pode fazer' no chat, explique:\n"
            "  'A confirmação de segurança deve ser feita clicando no botão ✅ do card acima.'\n\n"
            "CAMPOS SUPORTADOS: titulo, descricao, data_limite, data_inicio, horario_inicio, horario_fim, status, tags, area_tematica, tipo_acao, notas.\n"
            "STATUS VÁLIDOS: 'em andamento', 'stand-by', 'concluído', 'excluído'.\n"
            "EXCLUSÃO DE AÇÃO: para excluir (apagar/deletar/cancelar) uma ação, use preparar_edicao_acao com alteracoes={\"status\": \"excluído\"}. "
            "NÃO use 'cancelado' — o status correto para exclusão é 'excluído'.\n"
            "REMOÇÃO DE HORÁRIO (ação única): para tirar/limpar o horário de uma ação, use preparar_edicao_acao com "
            "alteracoes={\"horario_inicio\": None, \"horario_fim\": None}.\n"
            "REMOÇÃO DE HORÁRIOS EM LOTE: quando o usuário pedir para tirar/remover horários de várias ações de uma vez "
            "(ex: 'tira o horário de todas as ações de hoje', 'remove os horários da semana'), use:\n"
            "  preparar_remocao_horarios_em_lote(filtro_data='YYYY-MM-DD') — para ações de uma data específica\n"
            "  preparar_remocao_horarios_em_lote(task_ids=[...]) — para uma lista explícita de IDs\n"
            "Após chamar com sucesso, sua resposta de texto DEVE ser APENAS:\n"
            "  🕐 Preparei a remoção de horários. Verifique o card abaixo e confirme ou cancele.\n"
            "Para alteração do plano de ação (passos), use o fluxo de EDIÇÃO DE PLANO DE AÇÃO acima.\n\n"
        )

        protocolo_reagendamento_lote = (
            "## REAGENDAMENTO EM LOTE — REDISTRIBUIÇÃO DE AÇÕES (CRÍTICO)\n\n"
            "Quando o usuário pedir para mover, reagendar ou redistribuir múltiplas ações de uma vez "
            "(ex: 'mova as ações de hoje para a próxima semana', 'redistribua 5 por semana'), "
            "siga OBRIGATORIAMENTE este protocolo:\n\n"
            "ETAPA 0 — ENTENDIMENTO DOS CRITÉRIOS:\n"
            "Identifique:\n"
            "  a) Filtro de origem: data_limite = qual data? Ex: hoje → use a data atual no formato YYYY-MM-DD.\n"
            "  b) Data de início do reagendamento: qual o primeiro dia útil alvo? Ex: 'próxima semana' → segunda-feira da próxima semana.\n"
            "  c) Máximo por semana: quanto o usuário quer alocar por semana? Padrão: 5.\n"
            "  d) Estratégia de ordenação: se não especificada, use 'data_criacao' (mais antigas primeiro).\n"
            "Se qualquer parâmetro for ambíguo, esclareça com uma pergunta direta ANTES de chamar a ferramenta.\n\n"
            "ETAPA 1 — PREPARAÇÃO:\n"
            "Chame preparar_reagendamento_em_lote(nova_data_inicio, max_por_semana, estrategia, filtro_data, task_ids, justificativa).\n"
            "Esta ferramenta NÃO muta o banco — apenas prepara o plano de redistribuição.\n\n"
            "ETAPA 2 — AGUARDAR O CARD:\n"
            "Após chamar preparar_reagendamento_em_lote com sucesso, sua resposta de texto DEVE ser APENAS:\n"
            "  📅 Preparei o plano de reagendamento. Verifique o card abaixo e confirme ou cancele.\n"
            "NÃO liste os reagendamentos no texto — eles já estão no card visual.\n"
            "NÃO chame nenhuma outra ferramenta de escrita nesta mensagem.\n\n"
            "ETAPA 3 — APÓS CONFIRMAÇÃO:\n"
            "A confirmação ocorre pelo clique no botão do card — não pelo chat de texto.\n\n"
            "PARÂMETRO justificativa: gere automaticamente uma frase concisa descrevendo o reagendamento "
            "(ex: 'Reagendamento em lote das ações de 2026-04-25 para a semana de 2026-04-28.').\n\n"
        )

        protocolo_relatorios = (
            "## GERAÇÃO DE RELATÓRIOS — PROTOCOLO COLLECT-THEN-REPORT\n\n"
            "Quando o usuário solicitar um relatório, análise formal, resumo executivo ou documento consolidado:\n\n"
            "ETAPA 1 — COLETA DE CONTEXTO:\n"
            "Antes de chamar gerar_relatorio, colete o máximo de contexto relevante usando as ferramentas disponíveis:\n"
            "- consultar_historico_acoes() — tarefas e histórico de execução\n"
            "- buscar_arquivos_acervo() — documentos e artefatos relevantes\n"
            "- obter_contexto_tela() — contexto da tarefa em foco (se taskId disponível)\n"
            "Consolide tudo em uma string de contexto densa para passar ao relatório.\n\n"
            "ETAPA 2 — GERAÇÃO:\n"
            "Chame gerar_relatorio(titulo, tipo, contexto) com:\n"
            "- titulo: título claro e descritivo\n"
            "- tipo: 'executivo', 'técnico', 'analítico', 'progresso' ou 'situacional'\n"
            "- contexto: string com todos os dados coletados nas buscas\n\n"
            "ETAPA 3 — RESPOSTA:\n"
            "Após gerar_relatorio retornar sucesso, responda OBRIGATORIAMENTE:\n"
            "  📄 **Relatório gerado com sucesso!**\n"
            "  - **Título:** [titulo]\n"
            "  - **Seções:** [lista das seções]\n"
            "  Clique no botão abaixo para abrir o relatório completo.\n"
            "NÃO reproduza o conteúdo do relatório no chat — ele está disponível no modal de leitura.\n\n"
        )

        protocolo_formularios = (
            "## GERAÇÃO DE FORMULÁRIO GOOGLE — PADRÃO DRAFT-FIRST (CRÍTICO)\n\n"
            "Quando o usuário solicitar a criação de um formulário, questionário ou pesquisa de qualquer tipo, siga OBRIGATORIAMENTE este protocolo:\n\n"
            "ETAPA 1 — ESQUELETO (Preview):\n"
            "Nunca chame APIs externas. Emita o bloco [FORM]...[/FORM] para exibir um rascunho visual ao usuário.\n"
            "Inclua no texto da resposta (ANTES do bloco):\n"
            "  📝 **Proposta de Formulário: [título]**\n\n"
            "O bloco [FORM] deve conter um objeto JSON com esta estrutura exata:\n"
            "  [FORM]{\"titulo\": \"string\", \"descricao\": \"string (opcional)\", \"perguntas\": [{\"tipo\": \"texto_curto\" | \"paragrafo\" | \"multipla_escolha\" | \"caixas_selecao\" | \"lista_suspensa\" | \"escala_linear\", \"texto\": \"string\", \"opcoes\": [\"opt1\", \"opt2\"] (se aplicável), \"escala_min\": 1, \"escala_max\": 5, \"rotulo_min\": \"string\", \"rotulo_max\": \"string\", \"obrigatoria\": true}]}[/FORM]\n\n"
            "Regras para as perguntas:\n"
            "  - tipo: DEVE ser exatamente um destes: 'texto_curto', 'paragrafo', 'multipla_escolha', 'caixas_selecao', 'lista_suspensa', 'escala_linear'.\n"
            "  - opcoes: obrigatório apenas para os tipos: 'multipla_escolha', 'caixas_selecao', 'lista_suspensa'.\n"
            "  - atributos de escala (escala_min, escala_max, rotulo_min, rotulo_max): obrigatórios apenas para 'escala_linear'.\n\n"
            "ETAPA 2 — ITERAÇÃO (Ajustar):\n"
            "Se o usuário pedir ajuste, regere o bloco [FORM] com as modificações solicitadas.\n"
            "Não execute nenhuma ferramenta de persistência nesta etapa.\n\n"
            "ETAPA 3 — CONFIRMAÇÃO:\n"
            "Após confirmação do usuário, o FRONTEND chamará a Cloud Function criar_formulario_google e inserirá o link de sucesso diretamente no chat.\n"
            "Você NÃO chama nenhuma ferramenta nesta etapa — apenas oriente o usuário a clicar em 'Confirmar e Gerar Link' no card de rascunho de formulário.\n\n"
        )

        protocolo_diagramas = (
            "## DIAGRAMAS E VISUALIZAÇÕES — RENDERIZAÇÃO DIRETA (CRÍTICO)\n\n"
            "REGRA ABSOLUTA: Sempre que o usuário pedir um diagrama, fluxograma, fluxo, mapa, grafo, sequência, "
            "hierarquia, cronograma, mapa mental ou qualquer visualização estrutural, você DEVE obrigatoriamente "
            "gerar o código dentro de um bloco de código com a linguagem 'mermaid', exatamente assim:\n\n"
            "```mermaid\n"
            "<código aqui>\n"
            "```\n\n"
            "NUNCA escreva o código Mermaid como texto puro, sem o bloco de código. "
            "NUNCA omita os marcadores de abertura (```mermaid) e fechamento (```). "
            "O frontend detecta o bloco pela linguagem 'mermaid' e renderiza o diagrama visualmente — "
            "sem o bloco correto, o diagrama não aparece.\n"
            "Não explique o diagrama antes de gerá-lo, a menos que seja explicitamente solicitado.\n"
            "Escolha o tipo de diagrama mais adequado: flowchart, sequenceDiagram, classDiagram, "
            "stateDiagram-v2, erDiagram, gantt, mindmap, timeline, pie, quadrantChart, xychart-beta, etc.\n\n"
        )

        system_instruction_governanca = (
            "## GOVERNANCA DE FONTES — REGRA CRITICA\n"
            "Quando descrever uma tarefa encontrada por consultar_historico_acoes, "
            "use SOMENTE os campos retornados por essa ferramenta (Titulo, Status, Prazo, Area, Descricao, Tags, Sintese, Plano de Acao, Diario de Bordo, Notas). "
            "E PROIBIDO completar, interpretar ou inferir informacoes da tarefa usando dados do RAG, acervo ou memoria global. "
            "Se um campo nao constar no retorno da ferramenta, responda 'nao informado' em vez de inventar.\n\n"
            "## GOVERNANÇA FINANCEIRA — REGRA ABSOLUTA\n"
            "10. Você NUNCA deve inventar, supor ou estimar valores financeiros (gastos, rendas ou saldos).\n"
            "11. Para qualquer dúvida sobre finanças, use EXCLUSIVAMENTE a ferramenta `consultar_financas_v2`.\n"
            "12. Para novos registros, use `registrar_item_financeiro_v2` sempre apresentando um rascunho para o usuário confirmar antes.\n"
            "Caso a ferramenta retorne que não há dados, relate isso honestamente. Não tente usar o RAG para buscar dados financeiros internos.\n\n"
            "## DIRETRIZ DE EFICIÊNCIA (PREVENÇÃO DE LOOPS)\n"
            "Para perguntas diretas de classificação (ex: CONARQ), resumo ou extração de dados: use as ferramentas necessárias para ler o contexto, mas limite sua busca externa (internet, acervo). Se após 1 ou 2 tentativas de busca você não encontrar a tabela ou regra exata, PARE IMEDIATAMENTE. Consolide o que sabe e informe ao usuário que a informação específica não está disponível. NUNCA realize buscas exaustivas sequenciais.\n"
        )

        # ─── PROTOCOLOS CONDICIONAIS ─────────────────────────────────────────
        # Protocolos de fluxos raros (slides, formulários, relatórios, lote,
        # diagramas) só entram no prompt quando o assunto aparece na mensagem ou
        # no histórico recente — corta ~3k tokens da maioria das chamadas.
        _gate_text = _normalize_pop_text(f"{prompt} {history_plain}")

        def _protocolo_ativo(*kws: str) -> bool:
            return any(k in _gate_text for k in kws)

        _gate_reagendamento = _protocolo_ativo("reagend", "redistribu", "remarc", "mover", "mova", "lote", "horario")
        _gate_relatorios = _protocolo_ativo("relatorio", "resumo executivo", "consolidad", "sintese", "balanco")
        _gate_formularios = _protocolo_ativo("formul", "question", "enquete", "pesquisa", "survey")
        _gate_diagramas = _protocolo_ativo(
            "diagrama", "fluxograma", "mapa mental", "mermaid", "grafo", "gantt",
            "cronograma", "linha do tempo", "organograma", "visualiza", "timeline"
        )

        system_instruction_static = (
            system_instruction_nucleo
            + (protocolo_reagendamento_lote if _gate_reagendamento else "")
            + (protocolo_relatorios if _gate_relatorios else "")
            + (protocolo_formularios if _gate_formularios else "")
            + (protocolo_diagramas if _gate_diagramas else "")
            + system_instruction_governanca
        )

        strategy_context = ""
        if copilot_mode == "default" and strategy_directives:
            strategy_lines = "\n".join(f"- {directive}" for directive in strategy_directives)
            strategy_context = (
                "\n\n## DIRETRIZES ESTRATEGICAS PESSOAIS (PASSIVAS)\n"
                "Estas diretrizes vieram do modulo estrategia_pessoal e servem apenas como pano de fundo do chat global.\n"
                "Use-as somente quando o usuario pedir explicitamente alinhamento de vida, estrategia pessoal, decisoes de longo prazo, prioridades macro ou coerencia com objetivos pessoais.\n"
                "Nao use estas diretrizes para interferir em tarefas operacionais, nao transforme metas de vida em cobrancas e nao mencione este contexto espontaneamente.\n"
                f"{strategy_lines}\n"
            )

        system_instruction = (
            system_instruction_static
            + strategy_context
            + f"\n\nHoje é {today_str} e o horário local atual é {time_str}. "
            + (f"CONTEXTO TÉCNICO VINCULADO (OBRIGATÓRIO): "
               + (f"sistemaId={system_id}, " if system_id else "")
               + (f"taskId={task_id}. " if task_id else "") if (system_id or task_id) else "")
            + f"\n\n{mode_context}"
            + "\n\nÁREAS TEMÁTICAS VÁLIDAS (ao criar ações, escolha EXATAMENTE UMA desta lista; "
            + "NUNCA invente outra): " + ", ".join(_areas_validas)
            + ". Se nenhuma se encaixar, use 'GERAL'."
        )

        # --- ROTEADOR DE INTENÇÃO (heurístico + LLM fallback) ---
        # Passo 1: regex barato — evita chamada LLM na maioria dos casos.
        # Passo 2: só aciona Flash se palavras-chave estiverem presentes, com timeout de 3s.
        _CORRECAO_KEYWORDS = re.compile(
            r'\b(corrig|atualiz|invalida|contest|errad|incorret|mudar|alterar|substituir|'
            r'não está certo|não é assim|estava errado|está errado|procedimento errado)\w*\b',
            re.IGNORECASE
        )
        _correcao_hint = ""
        if _CORRECAO_KEYWORDS.search(prompt):
            _intent_prompt = f"Responda só 'CORRECAO' ou 'NORMAL': o usuário está corrigindo um procedimento?\nMensagem: {prompt}"
            # A/B test: HERMES_AB_LUNA_INTENT_ROUTER_PCT (0-100) desvia essa fração das
            # chamadas para o GPT-5.6 Luna (esforço baixo) em vez do Gemini, para comparar
            # custo/qualidade real nesta feature de alto volume. Padrão de operação: 10%
            # (rollout inicial conservador — ver docs/okf/integracoes/gemini-models.md).
            # Pode ser sobrescrito sem novo deploy setando a env var no ambiente das Functions.
            _use_luna = openai_provider.should_use_luna_ab("HERMES_AB_LUNA_INTENT_ROUTER_PCT", default="10")
            try:
                if _use_luna:
                    openai_key = keys_doc.to_dict().get('openai_api_key') if keys_doc.exists else None
                    if not openai_key:
                        raise RuntimeError("openai_api_key não configurada em system/api_keys")
                    import openai as _openai_sdk
                    # Timeout equivalente ao dos 3000ms do braço Gemini — sem isso, uma
                    # chamada travada consome o limite de 540s da callable inteira em vez
                    # de falhar rápido e cair no fail-open (achado do review do Codex).
                    _intent_resp = openai_provider.generate_text_logged(
                        _openai_sdk.OpenAI(api_key=openai_key, timeout=3.0, max_retries=0),
                        model=openai_provider.LUNA_MODEL,
                        input_text=_intent_prompt,
                        feature="copilot_intent_router",
                        reasoning_effort="low",
                        db=db,
                    )
                    _intent_text = (_intent_resp.output_text or "").upper()
                else:
                    _intent_resp = client.models.generate_content(
                        model=GEMINI_ROUTING_MODEL,
                        contents=_intent_prompt,
                        config=types.GenerateContentConfig(
                            http_options=types.HttpOptions(timeout=3000)
                        )
                    )
                    log_gemini_usage(_intent_resp, model=GEMINI_ROUTING_MODEL, feature="copilot_intent_router", db=db)
                    _intent_text = (_intent_resp.text or "").upper()

                # Registra braço + rótulo previsto (sem texto bruto) para permitir comparar
                # qualidade Luna vs Gemini depois, não só custo (achado do review do Codex).
                openai_provider.log_ab_outcome(
                    db,
                    feature="copilot_intent_router",
                    arm="openai_luna" if _use_luna else "gemini",
                    label="CORRECAO" if "CORRECAO" in _intent_text else "NORMAL",
                )

                if "CORRECAO" in _intent_text:
                    _correcao_hint = (
                        "\n\n## [DIRETIVA DO MOTOR DE EVOLUÇÃO — NÃO REVELAR AO USUÁRIO]\n"
                        "O roteador de intenção classificou esta mensagem como CORREÇÃO DE PROCEDIMENTO.\n"
                        "OBRIGAÇÃO: Após formular sua resposta ao usuário, acione silenciosamente "
                        "registrar_correcao_procedimento() com os dados extraídos da mensagem.\n"
                        "Não mencione, descreva nem insinue esta ação ao usuário."
                    )
            except Exception as _intent_exc:
                print(f"[IntentRouterAB] provider={'openai_luna' if _use_luna else 'gemini'} failed: {_intent_exc}")
                # Fail-open: timeout ou erro → continua sem hint (mesmo comportamento de antes,
                # sem fallback cruzado entre providers para não poluir a comparação do A/B)

        # ─── FERRAMENTAS DO MÓDULO ESTRATÉGIA (CRUD) ─────────────────────────────
        # Só são declaradas quando copilot_mode == 'estrategia' (ver static_tools).
        # Operam exclusivamente sobre a coleção estrategia_pessoal do usuário atual.
        _ESTRATEGIA_PILARES = {'carreira', 'financas', 'saude', 'intelectual', 'estilo_vida'}
        _ESTRATEGIA_STATUS = {'ativo', 'concluido', 'revisar'}
        _ESTRATEGIA_TIPOS = {'absoluta', 'relativa_qualitativa'}

        def _novo_id_estrategia(prefixo: str) -> str:
            import uuid as _uuid
            return f"{prefixo}-{int(time.time() * 1000)}-{str(_uuid.uuid4())[:6]}"

        def _carregar_objetivo_estrategico(objetivo_id: str):
            """Carrega um objetivo garantindo posse pelo usuário atual. Retorna (ref, data) ou (None, None).
            Fail-closed: exige usuário autenticado e que o userId do documento bata exatamente."""
            if not user_uid or not objetivo_id:
                return None, None
            ref = db.collection('estrategia_pessoal').document(str(objetivo_id))
            snap = ref.get()
            if not snap.exists:
                return None, None
            data = snap.to_dict() or {}
            if data.get('userId') != user_uid:
                return None, None
            return ref, data

        def criar_objetivo_estrategico(
            objetivoMacro: str,
            pilar: str = "carreira",
            tipoMeta: str = "relativa_qualitativa",
            status: str = "ativo",
            diretrizes: list[str] = [],
            indicadores: list[str] = [],
            marcos: list[str] = [],
            metrica_valor_inicial: float = None,
            metrica_valor_atual: float = None,
            metrica_valor_objetivo: float = None,
            metrica_unidade: str = "",
        ):
            """
            [MÓDULO ESTRATÉGIA] Cria um novo objetivo estratégico pessoal em estrategia_pessoal.
            Use APENAS quando o usuário pedir explicitamente para criar/cadastrar um objetivo, meta ou pilar estratégico.
            Parâmetros:
            - objetivoMacro: enunciado do objetivo macro (obrigatório).
            - pilar: um de 'carreira', 'financas', 'saude', 'intelectual', 'estilo_vida'.
            - tipoMeta: 'absoluta' (com métrica numérica) ou 'relativa_qualitativa'.
            - status: 'ativo', 'revisar' ou 'concluido'.
            - diretrizes: lista de diretrizes derivadas (frases que orientam a IA).
            - indicadores: lista de descrições de indicadores contínuos de sucesso.
            - marcos: lista de descrições de marcos pontuais.
            - metrica_*: só para tipoMeta 'absoluta' (valor inicial/atual/objetivo e unidade).
            Apresente um rascunho ao usuário e só chame após confirmação explícita.
            """
            try:
                if not user_uid:
                    return json.dumps({"status": "error", "reason": "auth_required"}, ensure_ascii=False)
                titulo = (objetivoMacro or "").strip()
                if not titulo:
                    return json.dumps({"status": "error", "reason": "objetivoMacro_obrigatorio"}, ensure_ascii=False)
                pilar_norm = (pilar or "carreira").strip().lower()
                if pilar_norm not in _ESTRATEGIA_PILARES:
                    pilar_norm = "carreira"
                tipo_norm = (tipoMeta or "relativa_qualitativa").strip().lower()
                if tipo_norm not in _ESTRATEGIA_TIPOS:
                    tipo_norm = "relativa_qualitativa"
                status_norm = (status or "ativo").strip().lower()
                if status_norm not in _ESTRATEGIA_STATUS:
                    status_norm = "ativo"

                diretrizes_clean = [str(d).strip() for d in (diretrizes or []) if str(d).strip()]
                if not diretrizes_clean:
                    return json.dumps({"status": "error", "reason": "diretrizes_obrigatorias"}, ensure_ascii=False)

                indicadores_obj = [
                    {"id": _novo_id_estrategia("indicador"), "descricao": str(d).strip(), "concluido": False, "registros": []}
                    for d in (indicadores or []) if str(d).strip()
                ]
                marcos_obj = [
                    {"id": _novo_id_estrategia("marco"), "descricao": str(d).strip(), "concluido": False, "registros": []}
                    for d in (marcos or []) if str(d).strip()
                ]

                payload = {
                    "userId": user_uid,
                    "pilar": pilar_norm,
                    "objetivoMacro": titulo,
                    "tipoMeta": tipo_norm,
                    "indicadoresSucesso": indicadores_obj,
                    "marcos": marcos_obj,
                    "diretrizesDerivadas": diretrizes_clean,
                    "status": status_norm,
                    "timestamp": firestore.SERVER_TIMESTAMP,
                }

                if tipo_norm == "absoluta":
                    val_obj = float(metrica_valor_objetivo or 0)
                    val_atual = float(metrica_valor_atual or 0)
                    val_ini = float(metrica_valor_inicial) if metrica_valor_inicial is not None else (val_atual if val_obj < val_atual else 0)
                    payload["metricaAlvo"] = {
                        "valorInicial": val_ini,
                        "valorAtual": val_atual,
                        "valorObjetivo": val_obj,
                        "unidade": str(metrica_unidade or "").strip(),
                    }

                ref = db.collection('estrategia_pessoal').document()
                ref.set(payload)
                return json.dumps({
                    "status": "created",
                    "objetivo_id": ref.id,
                    "objetivoMacro": titulo,
                    "pilar": pilar_norm,
                }, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"status": "error", "reason": str(e)}, ensure_ascii=False)

        def editar_objetivo_estrategico(
            objetivo_id: str,
            objetivoMacro: str = None,
            pilar: str = None,
            tipoMeta: str = None,
            status: str = None,
            diretrizes: list[str] = None,
            metrica_valor_inicial: float = None,
            metrica_valor_atual: float = None,
            metrica_valor_objetivo: float = None,
            metrica_unidade: str = None,
        ):
            """
            [MÓDULO ESTRATÉGIA] Edita um objetivo estratégico existente (identificado por objetivo_id, visível no snapshot da estratégia).
            Só passe os campos que devem mudar; os demais são preservados.
            - diretrizes: se fornecida, SUBSTITUI a lista completa de diretrizes. Para adicionar/remover, envie a lista final desejada (use o snapshot atual como base).
            - Para gerenciar indicadores ou marcos individualmente, use gerenciar_item_estrategico.
            Use APENAS após confirmação do usuário.
            """
            try:
                ref, data = _carregar_objetivo_estrategico(objetivo_id)
                if not ref:
                    return json.dumps({"status": "error", "reason": "objetivo_nao_encontrado"}, ensure_ascii=False)

                updates = {}
                if objetivoMacro is not None and str(objetivoMacro).strip():
                    updates["objetivoMacro"] = str(objetivoMacro).strip()
                if pilar is not None:
                    p = str(pilar).strip().lower()
                    if p in _ESTRATEGIA_PILARES:
                        updates["pilar"] = p
                if status is not None:
                    s = str(status).strip().lower()
                    if s in _ESTRATEGIA_STATUS:
                        updates["status"] = s
                if tipoMeta is not None:
                    tm = str(tipoMeta).strip().lower()
                    if tm in _ESTRATEGIA_TIPOS:
                        updates["tipoMeta"] = tm
                if diretrizes is not None:
                    dz = [str(d).strip() for d in (diretrizes or []) if str(d).strip()]
                    if not dz:
                        return json.dumps({"status": "error", "reason": "diretrizes_nao_podem_ficar_vazias"}, ensure_ascii=False)
                    updates["diretrizesDerivadas"] = dz

                # Métrica: só aplica se o objetivo é/torna-se absoluto
                tipo_final = updates.get("tipoMeta", data.get("tipoMeta"))
                if tipo_final == "absoluta" and any(v is not None for v in [metrica_valor_inicial, metrica_valor_atual, metrica_valor_objetivo, metrica_unidade]):
                    metrica = dict(data.get("metricaAlvo") or {})
                    if metrica_valor_inicial is not None:
                        metrica["valorInicial"] = float(metrica_valor_inicial)
                    if metrica_valor_atual is not None:
                        metrica["valorAtual"] = float(metrica_valor_atual)
                    if metrica_valor_objetivo is not None:
                        metrica["valorObjetivo"] = float(metrica_valor_objetivo)
                    if metrica_unidade is not None:
                        metrica["unidade"] = str(metrica_unidade).strip()
                    metrica.setdefault("valorInicial", 0)
                    metrica.setdefault("valorAtual", 0)
                    metrica.setdefault("valorObjetivo", 0)
                    metrica.setdefault("unidade", "")
                    updates["metricaAlvo"] = metrica

                if not updates:
                    return json.dumps({"status": "noop", "reason": "nenhum_campo_alterado"}, ensure_ascii=False)

                updates["timestamp"] = firestore.SERVER_TIMESTAMP
                ref.update(updates)
                return json.dumps({
                    "status": "updated",
                    "objetivo_id": ref.id,
                    "campos_alterados": [k for k in updates.keys() if k != "timestamp"],
                }, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"status": "error", "reason": str(e)}, ensure_ascii=False)

        def gerenciar_item_estrategico(
            objetivo_id: str,
            tipo: str,
            acao: str,
            descricao: str = None,
            item_id: str = None,
        ):
            """
            [MÓDULO ESTRATÉGIA] Gerencia um indicador ou marco dentro de um objetivo, preservando IDs (não quebra ações vinculadas).
            - tipo: 'indicador' ou 'marco'.
            - acao: 'adicionar' (precisa de descricao), 'editar' (precisa item_id + descricao), 'remover' (precisa item_id) ou 'concluir' (precisa item_id).
            - item_id: id do item existente (visível no snapshot).
            Use APENAS após confirmação do usuário.
            """
            try:
                if not user_uid or not objetivo_id:
                    return json.dumps({"status": "error", "reason": "objetivo_nao_encontrado"}, ensure_ascii=False)
                tipo_norm = (tipo or "").strip().lower()
                if tipo_norm not in {"indicador", "marco"}:
                    return json.dumps({"status": "error", "reason": "tipo_invalido"}, ensure_ascii=False)
                acao_norm = (acao or "").strip().lower()
                if acao_norm not in {"adicionar", "editar", "remover", "concluir"}:
                    return json.dumps({"status": "error", "reason": "acao_invalida"}, ensure_ascii=False)
                campo = "indicadoresSucesso" if tipo_norm == "indicador" else "marcos"
                ref = db.collection('estrategia_pessoal').document(str(objetivo_id))

                from datetime import datetime as _dt, timezone as _tz

                # Transação: indispensável porque o loop de tool-calling pode disparar
                # várias chamadas gerenciar_item_estrategico em paralelo sobre o mesmo
                # objetivo. Ler/recompor/gravar o array inteiro fora de transação faria
                # o último writer sobrescrever silenciosamente os demais. A transação
                # relê dentro do escopo e o Firestore reexecuta sob contenção.
                @firestore.transactional
                def _aplicar(transaction):
                    snap = ref.get(transaction=transaction)
                    if not snap.exists:
                        return {"status": "error", "reason": "objetivo_nao_encontrado"}
                    data = snap.to_dict() or {}
                    if data.get('userId') != user_uid:
                        return {"status": "error", "reason": "objetivo_nao_encontrado"}

                    lista = []
                    for item in (data.get(campo) or []):
                        if isinstance(item, str):
                            lista.append({"id": _novo_id_estrategia(tipo_norm), "descricao": item, "concluido": False, "registros": []})
                        else:
                            lista.append({
                                "id": item.get("id") or _novo_id_estrategia(tipo_norm),
                                "descricao": item.get("descricao", ""),
                                "concluido": bool(item.get("concluido")),
                                "registros": item.get("registros", []),
                                **({"dataConclusao": item["dataConclusao"]} if item.get("dataConclusao") else {}),
                                **({"evidencia": item["evidencia"]} if item.get("evidencia") else {}),
                            })

                    if acao_norm == "adicionar":
                        if not (descricao or "").strip():
                            return {"status": "error", "reason": "descricao_obrigatoria"}
                        novo = {"id": _novo_id_estrategia(tipo_norm), "descricao": descricao.strip(), "concluido": False, "registros": []}
                        lista.append(novo)
                        resultado_id = novo["id"]
                    else:  # editar | remover | concluir
                        if not item_id:
                            return {"status": "error", "reason": "item_id_obrigatorio"}
                        alvo = next((it for it in lista if it["id"] == item_id), None)
                        if not alvo:
                            return {"status": "error", "reason": "item_nao_encontrado"}
                        if acao_norm == "editar":
                            if not (descricao or "").strip():
                                return {"status": "error", "reason": "descricao_obrigatoria"}
                            alvo["descricao"] = descricao.strip()
                        elif acao_norm == "remover":
                            lista = [it for it in lista if it["id"] != item_id]
                        elif acao_norm == "concluir":
                            alvo["concluido"] = True
                            alvo["dataConclusao"] = _dt.now(_tz.utc).isoformat()
                        resultado_id = item_id

                    transaction.update(ref, {campo: lista, "timestamp": firestore.SERVER_TIMESTAMP})
                    return {
                        "status": "ok",
                        "objetivo_id": ref.id,
                        "tipo": tipo_norm,
                        "acao": acao_norm,
                        "item_id": resultado_id,
                    }

                resultado = _aplicar(db.transaction())
                return json.dumps(resultado, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"status": "error", "reason": str(e)}, ensure_ascii=False)

        def excluir_objetivo_estrategico(objetivo_id: str):
            """
            [MÓDULO ESTRATÉGIA] Exclui DEFINITIVAMENTE um objetivo estratégico (e seus indicadores/marcos/diretrizes).
            Operação irreversível. Use APENAS após confirmação explícita e inequívoca do usuário.
            """
            try:
                ref, data = _carregar_objetivo_estrategico(objetivo_id)
                if not ref:
                    return json.dumps({"status": "error", "reason": "objetivo_nao_encontrado"}, ensure_ascii=False)
                titulo = data.get("objetivoMacro", "")
                ref.delete()
                return json.dumps({"status": "deleted", "objetivo_id": objetivo_id, "objetivoMacro": titulo}, ensure_ascii=False)
            except Exception as e:
                return json.dumps({"status": "error", "reason": str(e)}, ensure_ascii=False)


        dynamic_tools = []
        _function_map = {
            'buscar_e_analisar_email': buscar_e_analisar_email,
            'consultar_historico_acoes': consultar_historico_acoes,
            'buscar_arquivos_acervo': buscar_arquivos_acervo,
            'buscar_conversas_whatsapp': buscar_conversas_whatsapp,
            'obter_contexto_tela': obter_contexto_tela,
            'pesquisar_internet': pesquisar_internet,
            'ler_pagina_web': ler_pagina_web,
            'ler_documento_na_integra': ler_documento_na_integra,
            'registrar_correcao_procedimento': registrar_correcao_procedimento,
            'salvar_memoria_global': salvar_memoria_global,
            'salvar_pop_global': salvar_pop_global,
            'resolver_conflito_memoria': resolver_conflito_memoria,
            'atualizar_personalidade': atualizar_personalidade,
            'resolver_conflito_procedimento': resolver_conflito_procedimento,
            'criar_acao_no_sistema': criar_acao_no_sistema,
            'editar_plano_acao': editar_plano_acao,
            'preparar_edicao_acao': preparar_edicao_acao,
            'registrar_no_diario': registrar_no_diario,
            'criar_objetivo_estrategico': criar_objetivo_estrategico,
            'editar_objetivo_estrategico': editar_objetivo_estrategico,
            'gerenciar_item_estrategico': gerenciar_item_estrategico,
            'excluir_objetivo_estrategico': excluir_objetivo_estrategico,
            'gerar_imagem': gerar_imagem,
            'gerar_relatorio': gerar_relatorio,
            'gerar_rascunho_formulario': gerar_rascunho_formulario,
            'consultar_agenda': consultar_agenda,
            'encontrar_slot_livre': encontrar_slot_livre,
            'preparar_reagendamento_em_lote': preparar_reagendamento_em_lote,
            'preparar_remocao_horarios_em_lote': preparar_remocao_horarios_em_lote,
            'consultar_financas_v2': consultar_financas_v2,
            'consultar_saude': consultar_saude,
            'registrar_item_financeiro_v2': registrar_item_financeiro_v2,
            'calculadora': calculadora,
            'agendar_lembrete_acao': agendar_lembrete_acao,
            'buscar_contato': buscar_contato,
            'preparar_vinculo_contatos': preparar_vinculo_contatos,
            'preparar_atualizacao_contato': preparar_atualizacao_contato,
            'registrar_interacao_contato': registrar_interacao_contato,
            'consultar_processo_sipac_copiloto': consultar_processo_sipac_copiloto,
            'incorporar_documento_especifico_sipac_no_rag_da_acao': incorporar_documento_especifico_sipac_no_rag_da_acao,
            'acompanhar_processo_sipac_copiloto': acompanhar_processo_sipac_copiloto,
        }

        # Cria função genérica de acionamento que o loop manual do Python irá ignorar
        # mas que vamos retornar direto pro frontend
        def acionar_ferramenta(tool_id: str, parametros: dict = None):
            return {"intent": "tool_invocation", "tool_id": tool_id, "parametros": parametros or {}}

        _function_map['acionar_ferramenta'] = acionar_ferramenta

        for t in routing_index:
            tool_id = t['id']
            schema = t.get('parametersSchema', {})
            props = {}
            required = schema.get('required', [])
            for k, v in schema.get('properties', {}).items():
                v_type = v.get('type', 'string').upper()
                if v_type == 'STRING':
                    t_enum = types.Type.STRING
                elif v_type == 'INTEGER':
                    t_enum = types.Type.INTEGER
                elif v_type == 'NUMBER':
                    t_enum = types.Type.NUMBER
                elif v_type == 'BOOLEAN':
                    t_enum = types.Type.BOOLEAN
                elif v_type == 'ARRAY':
                    t_enum = types.Type.ARRAY
                elif v_type == 'OBJECT':
                    t_enum = types.Type.OBJECT
                else:
                    t_enum = types.Type.STRING
                props[k] = types.Schema(type=t_enum, description=v.get('description', ''))

            tool_func = types.FunctionDeclaration(
                name=f"acionar_{tool_id}",
                description=f"Aciona a ferramenta interativa {tool_id} na interface do usuário.",
                parameters=types.Schema(
                    type=types.Type.OBJECT,
                    properties=props,
                    required=required
                )
            )
            dynamic_tools.append(tool_func)

            # Helper function to intercept dynamic tools
            def _make_tool_invoker(tid):
                def invoke(**kwargs):
                    return {"intent": "tool_invocation", "tool_id": tid, "parametros": kwargs}
                return invoke

            _function_map[f"acionar_{tool_id}"] = _make_tool_invoker(tool_id)


        # Ferramentas internas que não devem aparecer para o usuário
        _HIDDEN_TOOLS = {'registrar_correcao_procedimento', 'resolver_conflito_memoria'}

        # Ferramentas sempre disponíveis (núcleo operacional do copiloto)
        static_tools = [
            consultar_historico_acoes,
            buscar_arquivos_acervo,
            obter_contexto_tela,
            pesquisar_internet,
            ler_pagina_web,
            ler_documento_na_integra,
            registrar_correcao_procedimento,
            salvar_memoria_global,
            salvar_pop_global,
            resolver_conflito_memoria,
            atualizar_personalidade,
            resolver_conflito_procedimento,
            criar_acao_no_sistema,
            editar_plano_acao,
            preparar_edicao_acao,
            registrar_no_diario,
            consultar_agenda,
            encontrar_slot_livre,
            consultar_financas_v2,
            consultar_saude,
            registrar_item_financeiro_v2,
            calculadora,
            agendar_lembrete_acao,
            buscar_contato,
            preparar_vinculo_contatos,
            preparar_atualizacao_contato,
            registrar_interacao_contato,
        ]
        # Contagem das ferramentas sempre-ativas, antes dos appends condicionais
        # abaixo — usada como sinal de complexidade para o escalonamento de modelo.
        _n_base_tools = len(static_tools)
        # Ferramentas de fluxos raros: só declaradas quando o assunto aparece na
        # conversa — menos tokens de schema e roteamento mais limpo para o modelo.
        if _protocolo_ativo("mail", "caixa de entrada", "inbox"):
            static_tools.append(buscar_e_analisar_email)
        if _protocolo_ativo("whatsapp", "zap", "zapzap"):
            static_tools.append(buscar_conversas_whatsapp)
        if _protocolo_ativo("imagem", "figura", "ilustra", "foto", "banner", "logo", "desenh"):
            static_tools.append(gerar_imagem)
        if _gate_relatorios:
            static_tools.append(gerar_relatorio)
        if _gate_formularios:
            static_tools.append(gerar_rascunho_formulario)
        if _gate_reagendamento:
            static_tools.append(preparar_reagendamento_em_lote)
            static_tools.append(preparar_remocao_horarios_em_lote)
        if _protocolo_ativo("sipac", "processo", "protocolo"):
            static_tools.append(consultar_processo_sipac_copiloto)
            static_tools.append(incorporar_documento_especifico_sipac_no_rag_da_acao)
            static_tools.append(acompanhar_processo_sipac_copiloto)

        # Ferramentas de escrita do módulo Estratégia: só existem quando o copiloto
        # foi aberto a partir do módulo de Estratégia (copilot_mode == 'estrategia').
        if copilot_mode == "estrategia":
            static_tools.append(criar_objetivo_estrategico)
            static_tools.append(editar_objetivo_estrategico)
            static_tools.append(gerenciar_item_estrategico)
            static_tools.append(excluir_objetivo_estrategico)

        # ─── ESCALONAMENTO POR COMPLEXIDADE ──────────────────────────────────────
        # Só atua quando o modelo ainda é o tier padrão (não mexe no downgrade de
        # smalltalk). Sinais: muitas funções raras ativadas juntas, muitos cards
        # interativos em jogo, ou linguagem de múltiplas tarefas encadeadas. Em todos
        # os casos a heurística é local (regex + contagem) — não custa latência.
        if COPILOT_AUTO_ESCALATE and model_id == COPILOT_CHAT_MODEL:
            n_rare = len(static_tools) - _n_base_tools
            n_dynamic = len(dynamic_tools)
            _MULTITASK_RE = re.compile(
                r"(e (depois|também|tambem)|al[ée]m disso|todas as|todos os|em lote|"
                r"v[áa]ri[oa]s|cada (uma|um)|para cada|por fim|primeiro.*depois|"
                r"\d+[\).]\s)",
                re.IGNORECASE,
            )
            multitask = bool(_MULTITASK_RE.search(prompt or ""))
            long_prompt = len(prompt or "") > 320
            # Limiares conservadores: o frontier custa ~6x por token, então só
            # escala com sinais REALMENTE fortes. Multitarefa sozinha (sem prompt
            # longo + carga de ferramentas) não justifica o salto de custo.
            escalate = (
                n_rare >= COPILOT_ESCALATE_RARE_TOOLS
                or n_dynamic >= COPILOT_ESCALATE_DYNAMIC_TOOLS
                or (multitask and long_prompt and (n_rare >= 2 or n_dynamic >= 3))
            )
            if escalate:
                model_id = COPILOT_COMPLEX_MODEL
                print(
                    f"[Copiloto] Complexidade alta (rare={n_rare}, dynamic={n_dynamic}, "
                    f"multitask={multitask}, long={long_prompt}) — escalando para {model_id}"
                )

        chat = client.chats.create(
            model=model_id,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction + _correcao_hint,
                temperature=COPILOT_TEMPERATURE,
                http_options=types.HttpOptions(timeout=COPILOT_MODEL_TIMEOUT_MS),
                tools=static_tools + ([types.Tool(function_declarations=dynamic_tools)] if dynamic_tools else []),
                automatic_function_calling=types.AutomaticFunctionCallingConfig(disable=True)
            ),
            history=history
        )
        _perf_mark(perf_state, "web.chat_create")

        # ─── INGESTÃO DOCUMENTAL (MÚLTIPLOS ARQUIVOS - ATÉ 10) ───────────────────
        file_contexts = []
        # Artefatos indexados nesta mensagem sem tarefa ainda associada (ex.: usuário
        # está pedindo para criar a ação junto com o(s) arquivo(s)). Consumido por
        # criar_acao_no_sistema logo após a nova tarefa ser gravada.
        artefatos_pendentes_vinculo = []
        if drive_files:
            total_files = len(drive_files)
            for file_idx, f_item in enumerate(drive_files, start=1):
                drive_file_id = f_item['driveFileId']
                drive_file_name = f_item['driveFileName']
                file_context = ""
                _set_copilot_status(f"Processando arquivo ({file_idx}/{total_files}): {drive_file_name}...")
                try:
                    import io
                    import os
                    import tempfile
                    import uuid as _uuid
                    from googleapiclient.http import MediaIoBaseDownload
                    from google.cloud.firestore_v1.vector import Vector

                    drive_service = get_drive_service()

                    # 1. Busca metadados do arquivo no Drive
                    file_meta = drive_service.files().get(
                        fileId=drive_file_id,
                        fields='name,mimeType'
                    ).execute()
                    real_file_name = file_meta.get('name', drive_file_name)
                    real_mime_type = file_meta.get('mimeType', 'application/octet-stream')
                    is_image_file = real_mime_type.startswith('image/')
                    task_context_summary = ""
                    if task_id:
                        try:
                            task_doc = db.collection('tarefas').document(task_id).get()
                            if task_doc.exists:
                                task_data = task_doc.to_dict() or {}
                                plano = task_data.get('plano_acao', []) or []
                                task_context_summary = (
                                    f"Tarefa atual: {task_data.get('titulo', 'Sem título')}\n"
                                    f"Área temática: {task_data.get('area_tematica', 'Não informada')}\n"
                                    f"Descrição: {task_data.get('descricao', '')[:1200]}\n"
                                    f"Plano atual: {' | '.join(plano[:5]) if plano else 'Não definido'}"
                                )
                        except Exception as task_ctx_err:
                            print(f"[Copiloto] Aviso: falha ao montar contexto resumido da tarefa {task_id}: {task_ctx_err}")

                    # 2. Baixa o binário para memória volátil
                    request_dl = drive_service.files().get_media(fileId=drive_file_id)
                    fh = io.BytesIO()
                    downloader = MediaIoBaseDownload(fh, request_dl)
                    done = False
                    while not done:
                        _, done = downloader.next_chunk()
                    fh.seek(0)
                    file_bytes = fh.read()
                    fh.close()
                    del fh  # libera o buffer BytesIO imediatamente (evita cópia dupla em memória)
                    local_pdf_text = ""
                    local_pdf_metadata = None
                    local_docx_text = ""
                    local_docx_metadata = None
                    if is_pdf_mime_type(real_file_name, real_mime_type):
                        pdf_result = extract_pdf_text_with_fallback(
                            file_bytes,
                            real_file_name,
                            api_key=gemini_key,
                            allow_gemini_fallback=False,
                        )
                        local_pdf_text = (pdf_result.get('text') or '').strip()
                        local_pdf_metadata = pdf_result.get('metadata')
                    elif is_docx_mime_type(real_file_name, real_mime_type):
                        local_docx_text, local_docx_metadata = extract_docx_text(file_bytes)
                    elif real_mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or real_mime_type == 'application/vnd.ms-excel' or real_file_name.lower().endswith(('.xlsx', '.xls')):
                        import pandas as _pd
                        import io as _io
                        try:
                            _df_list = _pd.read_excel(_io.BytesIO(file_bytes), sheet_name=None)
                            _sheets_text = []
                            for _sheet_name, _df in _df_list.items():
                                _sheets_text.append(f"ABA: {_sheet_name}\n{_df.to_csv(index=False)}")
                            local_docx_text = "\n\n".join(_sheets_text).strip()
                            local_docx_metadata = {"natureza": "planilha_excel", "abas": list(_df_list.keys())}
                        except Exception as xl_err:
                            print(f"[Copiloto] Erro ao extrair Excel: {xl_err}")
                    elif real_mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or real_file_name.lower().endswith('.pptx'):
                        from pptx import Presentation as _Presentation
                        import io as _io
                        try:
                            _prs = _Presentation(_io.BytesIO(file_bytes))
                            _slides_text = []
                            for _slide in _prs.slides:
                                for _shape in _slide.shapes:
                                    if hasattr(_shape, 'text') and _shape.text.strip():
                                        _slides_text.append(_shape.text.strip())
                            local_docx_text = '\n'.join(_slides_text).strip()
                            local_docx_metadata = {"natureza": "apresentacao_powerpoint", "slides_count": len(_prs.slides)}
                        except Exception as ppt_err:
                            print(f"[Copiloto] Erro ao extrair PowerPoint: {ppt_err}")

                    gemini_file = None
                    local_extracted_text = local_pdf_text or local_docx_text
                    local_extraction_metadata = local_pdf_metadata or local_docx_metadata
                    if is_image_file or not local_extracted_text:
                        # 3. Salva em arquivo temporário para a File API do Gemini
                        file_ext = os.path.splitext(real_file_name)[1] or '.bin'
                        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
                            tmp.write(file_bytes)
                            tmp_path = tmp.name
                        del file_bytes  # libera bytes após gravar no disco; File API vai ler do tmp_path

                        # 4. Faz upload para a File API do Gemini
                        gemini_file = client.files.upload(
                            file=tmp_path,
                            config=types.UploadFileConfig(
                                mime_type=real_mime_type,
                                display_name=real_file_name
                            )
                        )
                        os.unlink(tmp_path)
                    else:
                        del file_bytes  # texto já extraído localmente; bytes brutos não são mais necessários

                    # O bloco try/finally abaixo garante que o arquivo seja sempre
                    # deletado da File API do Gemini, mesmo que a extração falhe.
                    # Sem isso cada upload acumularia dados até estourar a cota de 2GB.
                    try:
                        # 5. Extrai metadados com truncamento semântico de ~8.000 tokens
                        if is_image_file:
                            extraction_prompt = (
                                f"Você recebeu a imagem '{real_file_name}'. "
                                "Analise o conteúdo visual integral e retorne EXCLUSIVAMENTE um JSON válido, "
                                "sem markdown, sem texto extra, com esta estrutura:\n"
                                '{"titulo": "...", "natureza": "...", "resumo": "...", '
                                '"ocr": "...", "descricao_visual": "...", '
                                '"elementos_chave": ["..."], "evidencias": ["..."], '
                                '"relacao_com_acao": "...", "utilidade_pratica": "..."}\n'
                                "Regras:\n"
                                "- titulo: nome curto e útil do anexo.\n"
                                "- natureza: classifique a imagem com precisão (ex: print de sistema, comprovante, gráfico, documento fotografado, quadro, planilha capturada, foto de ambiente).\n"
                                "- resumo: resumo executivo em 3 a 5 frases.\n"
                                "- ocr: transcreva o texto visível mais relevante. Se não houver, use string vazia.\n"
                                "- descricao_visual: descreva objetivamente o que aparece na imagem.\n"
                                "- elementos_chave: liste de 3 a 8 elementos centrais percebidos.\n"
                                "- evidencias: liste fatos observáveis que sustentam sua leitura.\n"
                                "- relacao_com_acao: explique como a imagem se conecta com a tarefa em foco. Se não houver contexto suficiente, diga isso explicitamente.\n"
                                "- utilidade_pratica: diga como o Hermes deve usar esta imagem para apoiar a ação.\n"
                                f"\nCONTEXTO DA AÇÃO:\n{task_context_summary or 'Nenhuma tarefa ativa foi fornecida.'}"
                            )
                        elif local_extracted_text:
                            extraction_prompt = (
                                f"Você recebeu o texto extraído localmente do arquivo '{real_file_name}'. "
                                "Retorne EXCLUSIVAMENTE um JSON válido, sem markdown, sem texto extra, com esta estrutura:\n"
                                '{"titulo": "...", "natureza": "...", "resumo": "...", '
                                '"relacao_com_acao": "...", "utilidade_pratica": "..."}\n'
                                "Onde:\n"
                                "- titulo: nome/título do documento\n"
                                "- natureza: categoria (ex: Edital, Contrato, Relatório, Manual, Planilha, etc.)\n"
                                "- resumo: resumo executivo em 3 a 5 frases sobre conteúdo e utilidade\n"
                                "- relacao_com_acao: explique a conexão do arquivo com a tarefa atual; se não houver contexto suficiente, diga isso explicitamente\n"
                                "- utilidade_pratica: diga como o Hermes deve usar este arquivo para apoiar a ação\n"
                                f"\nCONTEXTO DA AÇÃO:\n{task_context_summary or 'Nenhuma tarefa ativa foi fornecida.'}\n\n"
                                f"METADADOS DA EXTRAÇÃO LOCAL: {json.dumps(local_extraction_metadata or {}, ensure_ascii=False)}\n\n"
                                f"TEXTO EXTRAÍDO:\n{local_extracted_text[:120000]}"
                            )
                            extraction_response = client.models.generate_content(
                                model=model_id,
                                contents=[extraction_prompt]
                            )
                        else:
                            extraction_prompt = (
                                f"Você recebeu o arquivo '{real_file_name}'. "
                                "Leia no máximo os primeiros 8.000 tokens de conteúdo. "
                                "Retorne EXCLUSIVAMENTE um JSON válido, sem markdown, sem texto extra, com esta estrutura:\n"
                                '{"titulo": "...", "natureza": "...", "resumo": "...", '
                                '"relacao_com_acao": "...", "utilidade_pratica": "..."}\n'
                                "Onde:\n"
                                "- titulo: nome/título do documento\n"
                                "- natureza: categoria (ex: Edital, Contrato, Relatório, Manual, Planilha, etc.)\n"
                                "- resumo: resumo executivo em 3 a 5 frases sobre conteúdo e utilidade\n"
                                "- relacao_com_acao: explique a conexão do arquivo com a tarefa atual; se não houver contexto suficiente, diga isso explicitamente\n"
                                "- utilidade_pratica: diga como o Hermes deve usar este arquivo para apoiar a ação\n"
                                f"\nCONTEXTO DA AÇÃO:\n{task_context_summary or 'Nenhuma tarefa ativa foi fornecida.'}"
                            )
                        if is_image_file or not local_extracted_text:
                            extraction_response = client.models.generate_content(
                                model=model_id,
                                contents=[
                                    types.Content(parts=[
                                        types.Part.from_uri(
                                            file_uri=gemini_file.uri,
                                            mime_type=real_mime_type
                                        ),
                                        types.Part(text=extraction_prompt)
                                    ])
                                ]
                            )
                        log_gemini_usage(extraction_response, model=model_id, feature="copilot_file_ingestion", db=db)
                        extraction_text = extraction_response.text.strip()
                        # Remove blocos de código caso o modelo os inclua mesmo instruído
                        if extraction_text.startswith("```"):
                            extraction_text = extraction_text.split("```")[1]
                            if extraction_text.startswith("json"):
                                extraction_text = extraction_text[4:]

                        meta = json.loads(extraction_text)
                        if isinstance(meta, list):
                            meta = meta[0] if meta else {}
                        if not isinstance(meta, dict):
                            meta = {}
                        titulo_doc = meta.get('titulo', real_file_name)
                        natureza_doc = meta.get('natureza', 'Documento')
                        resumo_doc = meta.get('resumo', '')
                        relacao_com_acao = meta.get('relacao_com_acao', '')
                        utilidade_pratica = meta.get('utilidade_pratica', '')
                        ocr_doc = meta.get('ocr', '') if is_image_file else ''
                        descricao_visual = meta.get('descricao_visual', '') if is_image_file else ''
                        elementos_chave = meta.get('elementos_chave', []) if is_image_file else []
                        evidencias = meta.get('evidencias', []) if is_image_file else []
                        if not isinstance(elementos_chave, list):
                            elementos_chave = [str(elementos_chave)]
                        if not isinstance(evidencias, list):
                            evidencias = [str(evidencias)]

                        # 6. Vetoriza e grava no indice_artefatos
                        from knowledge_graph import _get_embedding
                        embed_text_parts = [
                            titulo_doc,
                            natureza_doc,
                            resumo_doc,
                            relacao_com_acao,
                            utilidade_pratica,
                        ]
                        if is_image_file:
                            embed_text_parts.extend([
                                descricao_visual,
                                ocr_doc,
                                " | ".join(elementos_chave[:8]),
                                " | ".join(evidencias[:8]),
                            ])
                        embed_text = " | ".join(part for part in embed_text_parts if part)
                        embedding = _get_embedding(embed_text, gemini_key)
                        embedding_floats = list(map(float, embedding))

                        artefato_id = str(_uuid.uuid4())[:12]
                        drive_link = f"https://drive.google.com/file/d/{drive_file_id}/view"

                        # Origem bifurcada: tarefa (se task_id ativo) ou acervo global
                        origem_doc = (
                            {'modulo': 'tarefa', 'id_origem': task_id, 'session_id': session_id or 'direto'}
                            if task_id
                            else {'modulo': 'copiloto', 'id_origem': session_id or 'direto'}
                        )
                        db.collection('indice_artefatos').document(artefato_id).set({
                            'titulo': titulo_doc,
                            'trecho': resumo_doc,
                            'texto_bruto': (local_extracted_text or ocr_doc or resumo_doc)[:500000],
                            'fonte': natureza_doc,
                            'embedding': Vector(embedding_floats),
                            'tipo_arquivo': real_mime_type.split('/')[-1],
                            'mime_type': real_mime_type,
                            'url_drive': drive_link,
                            'drive_file_id': drive_file_id,
                            'data_criacao': firestore.SERVER_TIMESTAMP,
                            'origem': origem_doc,
                            'task_id': task_id or None,
                            'categoria': 'Copiloto Hermes',
                            'relacao_com_acao': relacao_com_acao,
                            'utilidade_pratica': utilidade_pratica,
                            'ocr': ocr_doc,
                            'descricao_visual': descricao_visual,
                            'elementos_chave': elementos_chave[:8],
                            'evidencias': evidencias[:8],
                            'texto_extraido_por': 'local_extractor' if local_extracted_text else ('gemini_ocr' if is_image_file else 'gemini_file_api'),
                            'extraction_metadata': local_extraction_metadata if local_extraction_metadata else None,
                        })
                        print(f"[Copiloto] Artefato '{titulo_doc}' gravado em indice_artefatos (id={artefato_id})")

                        # Dupla cidadania: vínculo físico e histórico à tarefa ativa
                        if task_id:
                            from datetime import datetime as _dt
                            now_iso = _dt.now().isoformat()
                            pool_item = {
                                'id': artefato_id,
                                'tipo': 'arquivo',
                                'valor': drive_link,
                                'nome': titulo_doc,
                                'drive_file_id': drive_file_id,  # Salvo explicitamente para leitura profunda on-demand
                                'data_criacao': now_iso
                            }
                            diary_entry = {
                                'data': now_iso,
                                'nota': f"📎 [Copiloto] Arquivo '{titulo_doc}' ({natureza_doc}) carregado via Copiloto Hermes e indexado no acervo global."
                            }
                            db.collection('tarefas').document(task_id).update({
                                'pool_dados': firestore.ArrayUnion([pool_item]),
                                'acompanhamento': firestore.ArrayUnion([diary_entry])
                            })
                            print(f"[Copiloto] Arquivo '{titulo_doc}' vinculado à tarefa {task_id} (pool_dados + acompanhamento)")
                        else:
                            # Sem tarefa ativa ainda (ex.: usuário está pedindo para CRIAR a ação
                            # nesta mesma mensagem). Guarda o artefato para ser vinculado à tarefa
                            # assim que ela for criada por criar_acao_no_sistema nesta mesma execução.
                            from datetime import datetime as _dt
                            now_iso = _dt.now().isoformat()
                            artefatos_pendentes_vinculo.append({
                                'id': artefato_id,
                                'tipo': 'arquivo',
                                'valor': drive_link,
                                'nome': titulo_doc,
                                'drive_file_id': drive_file_id,
                                'data_criacao': now_iso,
                                '_natureza': natureza_doc,
                            })

                        # 7. Monta o bloco de contexto que será injetado no prompt final
                        file_context = (
                            f"[CONTEXTO DO ARQUIVO ANEXADO]\n"
                            f"Nome: {real_file_name}\n"
                            f"Tipo MIME: {real_mime_type}\n"
                            f"Título extraído: {titulo_doc}\n"
                            f"Natureza: {natureza_doc}\n"
                            f"Resumo: {resumo_doc}\n"
                            f"Relação com a ação: {relacao_com_acao or 'Não inferida.'}\n"
                            f"Utilidade prática: {utilidade_pratica or 'Não inferida.'}\n"
                            f"Texto local relevante: {local_extracted_text[:2500] if local_extracted_text else 'Não aplicável.'}\n"
                            f"OCR relevante: {ocr_doc[:2500] if ocr_doc else 'Sem texto legível relevante.'}\n"
                            f"Descrição visual: {descricao_visual or 'Não aplicável.'}\n"
                            f"Elementos-chave: {', '.join(elementos_chave[:8]) if elementos_chave else 'Nenhum listado.'}\n"
                            f"Evidências observáveis: {', '.join(evidencias[:8]) if evidencias else 'Nenhuma listada.'}\n"
                            f"Link original: {drive_link}\n"
                            f"[/CONTEXTO DO ARQUIVO ANEXADO]"
                        )

                    finally:
                        # Limpeza obrigatória — evita acúmulo na File API do Gemini
                        if gemini_file is not None:
                            try:
                                client.files.delete(name=gemini_file.name)
                                print(f"[Copiloto] Arquivo Gemini '{gemini_file.name}' deletado com sucesso.")
                            except Exception as del_err:
                                print(f"[Copiloto] Aviso: falha ao deletar arquivo Gemini '{gemini_file.name}': {del_err}")

                except Exception as file_err:
                    import traceback as _tb
                    from datetime import datetime as _dt
                    err_str = str(file_err)
                    print(f"[Copiloto] Erro na ingestão documental: {err_str}")
                    print(_tb.format_exc())

                    # Telemetria estruturada — não suja o diário da tarefa
                    try:
                        db.collection('quality_logs').add({
                            'tipo': 'erro_ingestao_copiloto',
                            'descricao': 'Falha ao extrair contexto de arquivo via Gemini File API',
                            'evidencia': err_str,
                            'arquivo_nome': drive_file_name,
                            'task_id': task_id or None,
                            'session_id': session_id or None,
                            'data_criacao': _dt.now().isoformat()
                        })
                    except Exception as log_err:
                        print(f"[Copiloto] Falha ao gravar quality_log: {log_err}")

                    if file_context:
                        file_contexts.append(file_context)
        # ─────────────────────────────────────────────────────────────────────────

        # Injeta contexto inicial se houver task_id
        _perf_mark(perf_state, "web.file_ingestion")
        initial_context = ""
        if task_id:
            initial_context = f"DICA DE CONTEXTO: O usuário está visualizando a tarefa {task_id}. " \
                             f"Use obter_contexto_tela('{task_id}') para se situar antes de responder."

        # Monta prompt final combinando task context + file context + pergunta do usuário
        context_parts = []
        if initial_context:
            context_parts.append(initial_context)
        if memory_context:
            context_parts.append(memory_context)
        if session_conflict_context:
            context_parts.append(session_conflict_context)
        if file_contexts:
            context_parts.append("\n\n".join(file_contexts))
        context_parts.append(f"USUÁRIO: {prompt}")
        if matched_pop_directives:
            pop_lines = [
                "[DIRETRIZES OPERACIONAIS (POPs) CORRESPONDENTES - OBRIGATORIAS]"
            ]
            for idx, pop in enumerate(matched_pop_directives, start=1):
                triggers = ", ".join(pop.get("matched_triggers", []))
                pop_lines.append(
                    f"{idx}. TITULO: {pop.get('titulo', 'POP')}\n"
                    f"GATILHOS ACIONADOS: {triggers or 'n/a'}\n"
                    f"INSTRUCAO:\n{pop.get('instrucao_sistema', '')}"
                )
            pop_lines.append(
                "Aplique as diretrizes acima ao responder e ao decidir quais ferramentas usar. "
                "As diretrizes acima ja foram recuperadas do Gestor de POPs; nao declare ausencia de POP "
                "nem use Acervo/Internet apenas para confirmar a existencia deste POP. "
                "Se houver conflito entre POP e pedido literal do usuario, exponha o conflito antes de agir."
            )
            pop_lines.append("[/DIRETRIZES OPERACIONAIS (POPs) CORRESPONDENTES - OBRIGATORIAS]")
            context_parts.append("\n".join(pop_lines))
        final_prompt = "\n\n".join(context_parts)
        _perf_mark(perf_state, "web.prompt_build")
        # Loop manual de tool calling — intercepta cada chamada para rastrear ferramentas usadas
        tools_used: list[str] = []
        pending_edit_data = None
        pending_batch_reschedule_data = None
        pending_memory_conflict = None
        report_data = None
        tool_invocation_data = None
        deadline_fallback_text = None
        response = None
        try:
            response = chat.send_message(final_prompt)
            log_gemini_usage(response, model=model_id, feature="copilot_web", db=db, extra={"round": 0, "session_id": session_id})
            _perf_mark(perf_state, "web.first_model_response")
        except Exception as _first_send_err:
            if not _is_copilot_deadline_error(_first_send_err):
                raise
            print(f"[Copiloto] Gemini deadline na primeira resposta: {_first_send_err}")
            _perf_mark(perf_state, "web.first_model_deadline")
            try:
                response = client.models.generate_content(
                    model=COPILOT_FALLBACK_MODEL,
                    contents=[
                        (
                            "Responda em pt-BR, de forma objetiva. A chamada principal do Copiloto "
                            "atingiu deadline; nao acione ferramentas e, se faltar dado interno, diga "
                            "qual informacao precisa ser consultada em uma nova pergunta mais focada."
                        ),
                        final_prompt[:24000],
                    ],
                    config=types.GenerateContentConfig(
                        system_instruction=(system_instruction + _correcao_hint)[:12000],
                        http_options=types.HttpOptions(timeout=COPILOT_MODEL_RETRY_TIMEOUT_MS),
                    ),
                )
                log_gemini_usage(response, model=COPILOT_FALLBACK_MODEL, feature="copilot_web_fallback", db=db)
                _perf_mark(perf_state, "web.first_model_fallback")
            except Exception as _fallback_err:
                print(f"[Copiloto] Fallback rapido tambem falhou: {_fallback_err}")
                deadline_fallback_text = COPILOT_DEADLINE_FALLBACK_TEXT
                response = None

        _max_iter = 10
        for _round in range(_max_iter):
            if deadline_fallback_text or response is None:
                break

            if _copilot_remaining_sec() < 75:
                deadline_fallback_text = (
                    "A consulta chegou perto do limite seguro de processamento. "
                    "Interrompi antes do timeout para preservar a conversa. "
                    "Tente pedir uma leitura mais especifica ou dividir a solicitacao em partes menores."
                )
                _perf_mark(perf_state, "web.soft_deadline")
                break

            fcs = response.function_calls
            if not fcs:
                break

            _friendly_tools = ", ".join(sorted({fc.name.replace("_", " ") for fc in fcs}))
            _set_copilot_status(f"Executando: {_friendly_tools}...")

            function_response_parts = []
            # Constraint: Se mais de 3 roundtrips, forçar consolidação ou emitir aviso parcial
            if _round >= 3:
                # Injeta aviso diretamente no contexto do próximo turno para forçar consolidação
                function_response_parts.append(types.Part(text="\nAVISO DE PERFORMANCE: Você já realizou 3 rodadas de consultas sequenciais. Para evitar latência excessiva, consolide TODAS as buscas restantes em um único lote (batch) nesta rodada ou emita uma resposta parcial informando que está compilando os dados."))
            break_loop = False
            # Paralelismo de Execução de Ferramentas: processa múltiplas ferramentas simultaneamente em uma única rodada
            from concurrent.futures import ThreadPoolExecutor as _ThreadPoolExecutor, wait as _futures_wait
            
            def _execute_tool(idx, fc):
                fn = _function_map.get(fc.name)
                tool_start_ms = _perf_now_ms()
                tool_invocation_data_local = None
                
                if fn is None:
                    res = f"Ferramenta '{fc.name}' não encontrada."
                else:
                    try:
                        res = fn(**(fc.args or {}))
                        if isinstance(res, dict) and res.get("intent") == "tool_invocation":
                            tool_invocation_data_local = res
                    except Exception as _fe:
                        res = f"Erro ao executar {fc.name}: {_fe}"
                
                perf_state.setdefault("tool_calls", []).append({
                    "name": fc.name,
                    "duration_ms": max(0, _perf_now_ms() - tool_start_ms),
                })
                return (res, tool_invocation_data_local)

            _executor = _ThreadPoolExecutor(max_workers=min(len(fcs), 8))
            try:
                _futures = [_executor.submit(_execute_tool, i, fc) for i, fc in enumerate(fcs)]
                try:
                    from tools.registry import is_async as _is_async_tool
                except Exception:
                    _is_async_tool = lambda _n: False
                _batch_has_async = any(_is_async_tool(getattr(fc, 'name', '')) for fc in fcs)
                _tool_hard_cap = COPILOT_TOOL_TIMEOUT_ASYNC_SEC if _batch_has_async else COPILOT_TOOL_TIMEOUT_SEC
                _tool_wait_sec = min(_tool_hard_cap, max(5.0, _copilot_remaining_sec() - 75))
                _done_futures, _pending_futures = _futures_wait(_futures, timeout=_tool_wait_sec)
                if _pending_futures:
                    _timed_out_names = [getattr(fcs[_i], 'name', '?') for _i, _f in enumerate(_futures) if _f in _pending_futures]
                    for _pending in _pending_futures:
                        _pending.cancel()
                    print(f"[Copiloto] Tool timeout após {_tool_wait_sec:.1f}s — ferramentas pendentes: {_timed_out_names}")
                    _timed_out_label = ", ".join(_timed_out_names) if _timed_out_names else "uma das ferramentas"
                    deadline_fallback_text = (
                        f"A consulta interna **{_timed_out_label}** demorou mais que o limite "
                        f"de {int(_tool_wait_sec)}s e eu interrompi o processamento antes que a conversa ficasse travada. "
                        "Tente refinar o pedido (por exemplo, listar primeiro os itens pendentes de forma objetiva e "
                        "depois solicitar a proposta completa)."
                    )
                    _perf_mark(perf_state, "web.tool_timeout")

                for _i, _future in enumerate(_futures):
                    _fc = fcs[_i]
                    if _future not in _done_futures:
                        result = f"Timeout ao executar {getattr(_fc, 'name', 'ferramenta')}: limite de {int(_tool_wait_sec)}s excedido."
                        tool_invocation_data_local = None
                    else:
                        result, tool_invocation_data_local = _future.result()
                    
                    if tool_invocation_data_local:
                        tool_invocation_data = tool_invocation_data_local
                        break_loop = True
                    
                    # Processamento síncrono de efeitos colaterais
                    if _fc.name == 'preparar_edicao_acao' and isinstance(result, str) and result.startswith('{'):
                        try: pending_edit_data = json.loads(result)
                        except: pass
                    if _fc.name == 'preparar_reagendamento_em_lote' and isinstance(result, str) and result.startswith('{'):
                        try: pending_batch_reschedule_data = json.loads(result)
                        except: pass
                    if _fc.name == 'gerar_relatorio' and isinstance(result, str) and result.startswith('{'):
                        try:
                            parsed_rep = json.loads(result)
                            if parsed_rep.get('report_id'): report_data = parsed_rep
                        except: pass
                    if _fc.name == 'salvar_memoria_global' and isinstance(result, str) and result.startswith('{'):
                        try:
                            parsed_mem = json.loads(result)
                            if parsed_mem.get('status') == 'conflict':
                                parsed_mem.setdefault('status_ui', 'pending')
                                pending_memory_conflict = parsed_mem
                            if session_ref:
                                if parsed_mem.get('status') == 'conflict':
                                    session_ref.set({"pendingMemoryConflict": parsed_mem, "lastMemoryConflictAt": firestore.SERVER_TIMESTAMP}, merge=True)
                                else:
                                    session_ref.update({"pendingMemoryConflict": firestore.DELETE_FIELD, "lastMemoryConflictAt": firestore.DELETE_FIELD})
                        except: pass
                    if _fc.name == 'resolver_conflito_memoria' and isinstance(result, str) and result.startswith('{'):
                        try:
                            parsed_resolution = json.loads(result)
                            if session_ref and parsed_resolution.get('status') in {'resolved', 'updated'}:
                                session_ref.set({"pendingMemoryConflict": firestore.DELETE_FIELD, "lastMemoryConflictResolutionAt": firestore.SERVER_TIMESTAMP}, merge=True)
                                try: session_ref.update({"lastMemoryConflictAt": firestore.DELETE_FIELD})
                                except: pass
                        except: pass

                    if _fc.name not in _HIDDEN_TOOLS and _fc.name not in tools_used:
                        tools_used.append(_fc.name)

                    _result_str = str(result)
                    if len(_result_str) > 12000:
                        _result_str = _result_str[:12000] + "\n[...resultado truncado por tamanho...]"
                    
                    function_response_parts.append(
                        types.Part.from_function_response(
                            name=_fc.name,
                            response={"result": _result_str}
                        )
                    )
            finally:
                _executor.shutdown(wait=False, cancel_futures=True)

            if break_loop:
                break

            if deadline_fallback_text:
                break

            if _copilot_remaining_sec() < 75:
                deadline_fallback_text = (
                    "Executei as ferramentas necessarias, mas a consolidacao da resposta chegou perto "
                    "do limite seguro de processamento. Para evitar uma falha por timeout, parei aqui. "
                    "Tente refazer a pergunta de forma mais focada ou solicitar o proximo trecho."
                )
                _perf_mark(perf_state, "web.soft_deadline")
                break

            _set_copilot_status("Consolidando resposta...")
            try:
                if _round == _max_iter - 1:
                    function_response_parts.append(
                        types.Part(text="\n\n[AVISO DE LIMITE] O limite de consultas internas foi atingido. Por favor, responda ao usuário consolidando as informações obtidas até agora da melhor forma possível. Não acione mais nenhuma ferramenta.")
                    )
                response = chat.send_message(function_response_parts)
                log_gemini_usage(response, model=model_id, feature="copilot_web", db=db, extra={"round": _round + 1, "session_id": session_id})
                _perf_mark(perf_state, "web.tool_roundtrip")
            except Exception as _send_err:
                if _is_copilot_deadline_error(_send_err):
                    print(f"[Copiloto] Gemini deadline ao consolidar ferramentas: {_send_err}")
                    deadline_fallback_text = (
                        "Executei parte das consultas internas, mas a consolidacao da resposta demorou demais. "
                        "Interrompi antes do timeout para preservar a conversa. Tente refazer a pergunta de forma "
                        "mais focada ou solicitar o proximo trecho."
                    )
                    _perf_mark(perf_state, "web.tool_roundtrip_deadline")
                    break
                if _is_retryable_gemini_server_error(_send_err):
                    print(f"[Copiloto] Gemini 500 ao enviar resultados de ferramentas — tentando com payload reduzido: {_send_err}")
                    _reduced_parts = []
                    for _p in function_response_parts:
                        try:
                            _r = _p.function_response.response.get("result", "")
                            if len(_r) > 4000:
                                _r = _r[:4000] + "\n[...truncado para retry...]"
                            _reduced_parts.append(
                                types.Part.from_function_response(
                                    name=_p.function_response.name,
                                    response={"result": _r}
                                )
                            )
                        except Exception:
                            _reduced_parts.append(_p)
                    try:
                        response = chat.send_message(_reduced_parts)
                        log_gemini_usage(response, model=model_id, feature="copilot_web", db=db, extra={"round": _round + 1, "retry": True})
                        _perf_mark(perf_state, "web.tool_roundtrip_retry")
                    except Exception as _retry_err:
                        if _is_copilot_deadline_error(_retry_err):
                            print(f"[Copiloto] Gemini deadline no retry reduzido: {_retry_err}")
                            deadline_fallback_text = (
                                "As consultas internas retornaram, mas a consolidacao final demorou demais. "
                                "Interrompi antes do timeout. Tente pedir um resumo menor ou o proximo trecho."
                            )
                            _perf_mark(perf_state, "web.tool_roundtrip_retry_deadline")
                            break
                        raise
                else:
                    raise

        if tool_invocation_data:
            clean_text = "[Invocando Ferramenta...]"
            suggested_title = prompt[:50] if prompt and len(prompt) < 100 else None

            if session_id:
                try:
                    db.collection('sessoes_copiloto').document(session_id).collection('mensagens').add({
                        "role": "assistant",
                        "content": clean_text,
                        "toolInvocation": tool_invocation_data,
                        "toolsUsed": tools_used if tools_used else None,
                        "timestamp": firestore.SERVER_TIMESTAMP
                    })
                    db.collection('sessoes_copiloto').document(session_id).set({
                        "lastMessageAt": firestore.SERVER_TIMESTAMP
                    }, merge=True)
                except Exception as e:
                    print(f"Erro ao salvar invocação de ferramenta no Firestore: {e}")

            _set_copilot_status(None)
            _perf_mark(perf_state, "web.tool_invocation_persist")
            _perf_log(
                "web.askCopiloto.complete",
                perf_state,
                {
                    "session_id": session_id,
                    "task_id": task_id,
                    "mode": "tool_invocation",
                    "tools_used": tools_used,
                },
            )
            return {
                "result": clean_text,
                "kg_nodes": None,
                "toolInvocation": tool_invocation_data,
                "suggestedTitle": suggested_title
            }

        if deadline_fallback_text:
            result_text = deadline_fallback_text
        else:
            try:
                result_text = response.text or ""
            except Exception as _text_err:
                print(f"[Copiloto] response.text falhou: {_text_err}")
                if response and response.function_calls:
                    result_text = (
                        "Cheguei ao limite de consultas internas tentando obter os dados do seu documento. "
                        "Por favor, tente refazer a pergunta de forma mais direta ou focar em um trecho específico."
                    )
                else:
                    result_text = "Desculpe, ocorreu uma instabilidade ao processar a resposta. Por favor, tente novamente."
        # Extração de Proposta [PROPOSAL]{...}[/PROPOSAL]
        proposal_data = None
        clean_text = result_text
        if "[PROPOSAL]" in result_text:
            try:
                parts = result_text.split("[PROPOSAL]")
                proposal_raw = parts[1].split("[/PROPOSAL]")[0]
                proposal_data = json.loads(proposal_raw)
                clean_text = parts[0] + (parts[1].split("[/PROPOSAL]")[1] if "[/PROPOSAL]" in parts[1] else "")
                clean_text = clean_text.strip()
            except Exception as e:
                print(f"Erro ao extrair proposta: {e}")

        # Extração de Formulário [FORM]{...}[/FORM]
        form_data = None
        if "[FORM]" in clean_text:
            try:
                form_parts = clean_text.split("[FORM]")
                form_raw = form_parts[1].split("[/FORM]")[0]
                form_data = json.loads(form_raw)
                clean_text = form_parts[0] + (form_parts[1].split("[/FORM]")[1] if "[/FORM]" in form_parts[1] else "")
                clean_text = clean_text.strip()
            except Exception as e:
                print(f"Erro ao extrair formulário: {e}")

        # Extração de Solicitação de Diagnóstico [DIAGNOSIS]{...}[/DIAGNOSIS]
        diagnosis_request_data = None
        if "[DIAGNOSIS]" in clean_text:
            try:
                diag_parts = clean_text.split("[DIAGNOSIS]")
                diag_raw = diag_parts[1].split("[/DIAGNOSIS]")[0]
                diagnosis_request_data = json.loads(diag_raw)
                clean_text = diag_parts[0] + (diag_parts[1].split("[/DIAGNOSIS]")[1] if "[/DIAGNOSIS]" in diag_parts[1] else "")
                clean_text = clean_text.strip()
            except Exception as e:
                print(f"Erro ao extrair diagnóstico: {e}")

        # Salva a resposta do assistente no Firestore para o histórico
        if session_id:
            try:
                db.collection('sessoes_copiloto').document(session_id).collection('mensagens').add({
                    "role": "assistant",
                    "content": clean_text,
                    "proposedPlan": proposal_data.get("items") if proposal_data else None,
                    "proposedDiagnosis": diagnosis_request_data if diagnosis_request_data else None,
                    "proposedForm": form_data if form_data else None,
                    "toolsUsed": tools_used if tools_used else None,
                    "pendingEdit": pending_edit_data,
                    "pendingBatchReschedule": pending_batch_reschedule_data,
                    "pendingMemoryConflict": pending_memory_conflict,
                    "reportId": report_data.get('report_id') if report_data else None,
                    "timestamp": firestore.SERVER_TIMESTAMP
                })
                # Atualiza timestamp da sessão
                db.collection('sessoes_copiloto').document(session_id).set({
                    "lastMessageAt": firestore.SERVER_TIMESTAMP
                }, merge=True)
            except Exception as e:
                print(f"Erro ao salvar resposta no Firestore: {e}")

        # Tenta extrair título sugerido se for início de sessão
        _set_copilot_status(None)
        _perf_mark(perf_state, "web.response_persist")
        suggested_title = None
        if prompt and len(prompt) < 100:
            suggested_title = prompt[:50]

        _perf_log(
            "web.askCopiloto.complete",
            perf_state,
            {
                "session_id": session_id,
                "task_id": task_id,
                "mode": "assistant_text",
                "tools_used": tools_used,
            },
        )
        return {
            "result": clean_text,
            "proposedPlan": proposal_data.get("items") if proposal_data else None,
            "proposedDiagnosis": diagnosis_request_data if diagnosis_request_data else None,
            "pendingMemoryConflict": pending_memory_conflict,
            "suggestedTitle": suggested_title,
            "pendingEdit": pending_edit_data,
            "reportId": report_data.get('report_id') if report_data else None
        }

    except Exception as e:
        print(f"Erro em askCopilotoHermes: {e}")
        # Limpa o status efêmero da sessão (o helper pode não existir se o erro
        # ocorreu antes da sua definição, por isso a limpeza inline).
        if session_id:
            try:
                db.collection('sessoes_copiloto').document(session_id).update({
                    "copilotStatus": firestore.DELETE_FIELD,
                    "copilotStatusAt": firestore.DELETE_FIELD,
                })
            except Exception:
                pass
        _perf_log("web.askCopiloto.error", perf_state, {"session_id": session_id, "task_id": task_id, "error": str(e)})
        import traceback
        print(traceback.format_exc())
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )


class _CopilotJobAuth:
    def __init__(self, uid: str | None):
        self.uid = uid


class _CopilotJobRequest:
    def __init__(self, data: dict, uid: str | None):
        self.data = data
        self.auth = _CopilotJobAuth(uid) if uid else None


@firestore_fn.on_document_created(
    document="copilot_jobs/{jobId}",
    memory=options.MemoryOption.GB_2,
    timeout_sec=540,
)
def on_copilot_job_created(event: firestore_fn.Event[firestore_fn.DocumentSnapshot | None]):
    """Processa comandos do Copiloto fora do ciclo de vida da aba do navegador."""
    snap = event.data
    if snap is None or not snap.exists:
        return

    db = get_db()
    job_ref = snap.reference
    job = snap.to_dict() or {}
    payload = job.get("payload") or {}
    session_id = payload.get("sessionId") or job.get("sessionId")
    user_uid = job.get("userId")

    if not session_id or not payload.get("prompt"):
        job_ref.set({
            "status": "error",
            "error": "Job sem sessionId ou prompt.",
            "updatedAt": firestore.SERVER_TIMESTAMP,
        }, merge=True)
        return

    def _set_session_status(text: str | None):
        session_ref = db.collection("sessoes_copiloto").document(session_id)
        try:
            if text:
                session_ref.set({
                    "copilotStatus": text,
                    "copilotStatusAt": firestore.SERVER_TIMESTAMP,
                }, merge=True)
            else:
                session_ref.update({
                    "copilotStatus": firestore.DELETE_FIELD,
                    "copilotStatusAt": firestore.DELETE_FIELD,
                })
        except Exception as status_err:
            print(f"[CopilotJob] Falha ao atualizar status da sessão {session_id}: {status_err}")

    try:
        job_ref.set({
            "status": "processing",
            "startedAt": firestore.SERVER_TIMESTAMP,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        }, merge=True)
        _set_session_status("Processando em segundo plano...")

        req = _CopilotJobRequest(payload, user_uid)
        # askCopilotoHermes é decorada com @https_fn.on_call, que empilha DUAS
        # camadas preservadas por functools.wraps: o wrapper de CORS (flask_cors,
        # camada externa) e o on_call_wrapped (que chama _on_call_handler e acessa
        # request.method). Ambos esperam uma requisição HTTP bruta. Precisamos da
        # função original, que espera apenas .data/.auth como um CallableRequest.
        # inspect.unwrap percorre toda a cadeia de __wrapped__ até chegar nela.
        import inspect as _inspect
        core_fn = _inspect.unwrap(askCopilotoHermes)
        result = core_fn(req)

        job_ref.set({
            "status": "completed",
            "completedAt": firestore.SERVER_TIMESTAMP,
            "updatedAt": firestore.SERVER_TIMESTAMP,
            "resultSummary": str((result or {}).get("result", ""))[:500] if isinstance(result, dict) else "",
        }, merge=True)
    except Exception as exc:
        err_msg = str(getattr(exc, "message", None) or exc)
        print(f"[CopilotJob] Erro ao processar job {event.params.get('jobId')}: {err_msg}")
        job_ref.set({
            "status": "error",
            "error": err_msg,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        }, merge=True)
        try:
            db.collection("sessoes_copiloto").document(session_id).collection("mensagens").add({
                "role": "assistant",
                "content": f"**Erro ao processar a solicitação em segundo plano:**\n`{err_msg}`",
                "timestamp": firestore.SERVER_TIMESTAMP,
            })
        except Exception as msg_err:
            print(f"[CopilotJob] Falha ao registrar erro na sessão {session_id}: {msg_err}")
        _set_session_status(None)

@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.GB_1,
    timeout_sec=30
)
def confirmarEdicaoAcao(req: https_fn.CallableRequest):
    """
    Confirma e executa uma edição de ação pendente gerada pelo Copiloto Hermes.
    Aplica validação lazy antes de mutar: verifica que a ação não foi modificada
    desde a geração do card de confirmação.
    """
    data = req.data or {}
    session_id = data.get('sessionId')
    message_id = data.get('messageId')
    task_id = data.get('taskId')
    alteracoes = data.get('alteracoes', {})
    snapshot_ts = data.get('snapshotTs', '')

    if not task_id or not alteracoes:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="taskId e alteracoes são obrigatórios."
        )

    def _set_card_status(db_ref, status, error_msg=None):
        if not session_id or not message_id:
            return
        try:
            update_payload = {'pendingEdit.status': status}
            if error_msg:
                update_payload['pendingEdit.errorMessage'] = error_msg
            db_ref.collection('sessoes_copiloto').document(session_id)\
                .collection('mensagens').document(message_id)\
                .update(update_payload)
        except Exception as _ue:
            print(f"[confirmarEdicaoAcao] Falha ao atualizar status do card: {_ue}")

    try:
        from datetime import datetime as _dt, timezone as _tz

        db_ref = get_db()
        task_ref = db_ref.collection('tarefas').document(task_id)
        task_doc = task_ref.get()

        # Validação 1: ação existe?
        if not task_doc.exists:
            msg = 'Edição bloqueada: Esta ação não existe mais.'
            _set_card_status(db_ref, 'invalidated', msg)
            return {'status': 'invalidated', 'message': msg}

        task_data = task_doc.to_dict()

        # Validação 2: ação já concluída?
        if task_data.get('status') == 'concluído':
            msg = 'Edição bloqueada: Esta ação já foi concluída.'
            _set_card_status(db_ref, 'invalidated', msg)
            return {'status': 'invalidated', 'message': msg}

        # Validação 3 (lazy): ação foi modificada desde a geração do card?
        current_ts = task_data.get('data_atualizacao') or task_data.get('data_criacao', '')
        if snapshot_ts and current_ts and str(current_ts) != str(snapshot_ts):
            msg = 'Edição bloqueada: Esta ação foi modificada após a geração deste card.'
            _set_card_status(db_ref, 'invalidated', msg)
            return {'status': 'invalidated', 'message': msg}

        # Aplica mudanças — somente campos whitelistados
        _ALLOWED = {'titulo', 'descricao', 'data_limite', 'data_inicio', 'prazo_final', 'horario_inicio', 'horario_fim', 'status', 'tags', 'area_tematica', 'tipo_acao', 'notas', 'email_link_optout'}

        def _normalizar_status_acao(valor):
            if valor is None:
                return valor
            raw = str(valor).strip().lower()
            try:
                import unicodedata
                raw = ''.join(c for c in unicodedata.normalize('NFD', raw) if unicodedata.category(c) != 'Mn')
            except Exception:
                pass
            raw = raw.replace('_', ' ').replace('-', ' ')
            raw = ' '.join(raw.split())
            if raw in ('concluido', 'concluida', 'concluir', 'finalizado', 'finalizada', 'completed', 'done'):
                return 'concluído'
            if raw in ('stand by', 'standby', 'pausado', 'pausada', 'pausar'):
                return 'stand-by'
            if raw in ('em andamento', 'andamento', 'pendente', 'aberto', 'aberta', 'reabrir'):
                return 'em andamento'
            if raw in ('excluido', 'excluir', 'excluida', 'cancelado', 'cancelar', 'cancelada', 'deletar', 'deletado', 'apagar', 'remover'):
                return 'excluído'
            return valor

        updates = {}
        for campo, novo_valor in alteracoes.items():
            if campo not in _ALLOWED:
                continue
            if campo == 'status':
                novo_valor = _normalizar_status_acao(novo_valor)
                if novo_valor not in ('em andamento', 'stand-by', 'concluído', 'excluído'):
                    continue
            updates[campo] = novo_valor

        # Acoes usam data unica. Mantem data_limite e data_inicio espelhadas
        # mesmo quando a edicao vier de ferramentas/backend que enviam so um campo.
        if 'data_limite' in updates or 'data_inicio' in updates:
            single_date = updates.get('data_limite') or updates.get('data_inicio') or ''
            today_str = _dt.now(_tz.utc).strftime("%Y-%m-%d")
            status_val = updates.get('status') or task_data.get('status') or ''
            if status_val not in ('concluído', 'cancelado', 'excluído') and single_date and single_date not in ('-', '0000-00-00'):
                if single_date < today_str:
                    single_date = today_str
            updates['data_limite'] = single_date
            updates['data_inicio'] = single_date

        if not updates:
            _set_card_status(db_ref, 'error', 'Nenhum campo válido para atualizar.')
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="Nenhum campo válido para atualizar."
            )

        now_iso = _dt.now(_tz.utc).isoformat()
        updates['data_atualizacao'] = now_iso
        if updates.get('status') == 'concluído':
            updates['data_conclusao'] = now_iso
        elif updates.get('status') in ('em andamento', 'stand-by'):
            updates['data_conclusao'] = None

        campos_desc = ', '.join(k for k in updates if k not in ('data_atualizacao', 'data_conclusao'))
        diary_entry = {
            'data': now_iso,
            'nota': f"[Copiloto Hermes] Ação editada via card de confirmação. Campos alterados: {campos_desc}."
        }

        task_ref.update({
            **updates,
            'acompanhamento': firestore.ArrayUnion([diary_entry])
        })

        _set_card_status(db_ref, 'completed')
        print(f"[confirmarEdicaoAcao] task_id={task_id} atualizado: {list(updates.keys())}")
        return {'status': 'completed'}

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em confirmarEdicaoAcao: {e}")
        try:
            _set_card_status(get_db(), 'error', str(e))
        except Exception:
            pass
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.GB_1,
    timeout_sec=30
)
def confirmarReagendamentoEmLote(req: https_fn.CallableRequest):
    """
    Confirma e executa o reagendamento em lote de ações preparado pelo Copiloto Hermes.
    Aplica as novas datas atomicamente via WriteBatch e registra no diário de cada ação.
    """
    data = req.data or {}
    session_id = data.get('sessionId')
    message_id = data.get('messageId')
    items = data.get('items', [])
    justificativa = data.get('justificativa', 'Reagendamento em lote via Copiloto Hermes.')

    if not items:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="items é obrigatório e não pode ser vazio."
        )

    def _set_batch_card_status(db_ref, status, error_msg=None):
        if not session_id or not message_id:
            return
        try:
            update_payload = {'pendingBatchReschedule.status': status}
            if error_msg:
                update_payload['pendingBatchReschedule.errorMessage'] = error_msg
            db_ref.collection('sessoes_copiloto').document(session_id)\
                .collection('mensagens').document(message_id)\
                .update(update_payload)
        except Exception as _ue:
            print(f"[confirmarReagendamentoEmLote] Falha ao atualizar card: {_ue}")

    try:
        from datetime import datetime as _dt, timezone as _tz

        db_ref = get_db()
        batch = db_ref.batch()
        now_iso = _dt.now(_tz.utc).isoformat()
        diary_entry = {
            'data': now_iso,
            'nota': f"[Copiloto Hermes] {justificativa}"
        }

        count = 0
        for item in items:
            task_id = item.get('task_id')
            if not task_id:
                continue

            updates = {'data_atualizacao': now_iso}
            if item.get('nova_data_limite'):
                single_date = item['nova_data_limite']
                today_str = _dt.now(_tz.utc).strftime("%Y-%m-%d")
                if single_date < today_str:
                    single_date = today_str
                item['nova_data_limite'] = single_date
                updates['data_limite'] = single_date
                updates['data_inicio'] = item['nova_data_limite']
            if 'novo_horario_inicio' in item:
                updates['horario_inicio'] = item['novo_horario_inicio']
            if 'novo_horario_fim' in item:
                updates['horario_fim'] = item['novo_horario_fim']

            task_ref = db_ref.collection('tarefas').document(task_id)
            batch.update(task_ref, {
                **updates,
                'acompanhamento': firestore.ArrayUnion([diary_entry])
            })
            count += 1

        batch.commit()
        _set_batch_card_status(db_ref, 'completed')
        print(f"[confirmarReagendamentoEmLote] {count} ações reagendadas.")
        return {'status': 'completed', 'count': count}

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em confirmarReagendamentoEmLote: {e}")
        try:
            _set_batch_card_status(get_db(), 'error', str(e))
        except Exception:
            pass
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.MB_256,
    timeout_sec=60
)
def confirmarConflitoMemoria(req: https_fn.CallableRequest):
    data = req.data or {}
    session_id = data.get('sessionId')
    message_id = data.get('messageId')
    memoria_id = data.get('memoriaId')
    decisao = data.get('decisao')
    fato_atualizado = data.get('fatoAtualizado') or ''
    categoria = data.get('categoria') or 'fato_isolado'

    if not session_id or not message_id or not memoria_id or not decisao:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="sessionId, messageId, memoriaId e decisao são obrigatórios."
        )

    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini não configurada."
            )

        user_uid = req.auth.uid if req.auth else None
        result = _save_memory_node(
            db=db,
            api_key=gemini_key,
            fato=fato_atualizado,
            categoria=categoria,
            session_id=session_id,
            user_uid=user_uid,
            force_update_id=memoria_id if decisao == 'substituir_pelo_novo' else None,
        ) if decisao == 'substituir_pelo_novo' else {
            "status": "resolved",
            "decision": "kept_existing",
            "memory_id": memoria_id,
        }

        msg_ref = db.collection('sessoes_copiloto').document(session_id).collection('mensagens').document(message_id)
        update_payload = {
            'pendingMemoryConflict.status_ui': 'resolved' if decisao == 'substituir_pelo_novo' else 'kept',
            'pendingMemoryConflict.decisao_final': decisao,
            'pendingMemoryConflict.resolvedAt': firestore.SERVER_TIMESTAMP,
        }
        if decisao == 'substituir_pelo_novo':
            update_payload['pendingMemoryConflict.proposed_text'] = fato_atualizado
        msg_ref.update(update_payload)

        db.collection('sessoes_copiloto').document(session_id).set({
            'pendingMemoryConflict': firestore.DELETE_FIELD,
            'lastMemoryConflictResolutionAt': firestore.SERVER_TIMESTAMP,
        }, merge=True)

        return {
            "status": result.get("status", "resolved"),
            "decision": decisao,
            "memoryId": memoria_id,
        }
    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro em confirmarConflitoMemoria: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )


def buscar_procedimento_internal(query_text: str, area_tematica: str = None):
    # Wrapper interno para chamar a lógica de buscar_procedimento sem o overhead do Callable HTTPS
    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        api_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None

        from knowledge_graph import _get_embedding, _cosine_similarity

        # Sanitização de input
        q_text = (query_text or "").strip()
        if not q_text:
            q_text = "procedimentos operacionais"

        # Short-query fast-path: skip embedding + scan for 1-word queries
        _SW = {"de", "a", "o", "que", "e", "do", "da", "em", "um", "uma", "os", "as",
               "no", "na", "com", "por", "para", "dos", "das"}
        _q_meaningful = [w for w in q_text.lower().split() if w not in _SW and len(w) > 2]
        if len(_q_meaningful) < 2:
            return {"context": f"Nenhum registro encontrado para '{q_text}'.", "resultados": []}

        from google.cloud.firestore_v1.vector import Vector
        from google.cloud.firestore_v1.base_vector_query import DistanceMeasure
        from google.cloud.firestore_v1.base_query import FieldFilter

        query_embedding = _get_embedding(q_text, api_key)
        query_vector = list(map(float, query_embedding))

        # Otimização: Uso de Vector Search Nativo do Firestore (find_nearest)
        # Substitui a varredura manual de 200 documentos por filtragem no banco.
        collection_ref = db.collection("knowledge_nodes")
        vector_query = collection_ref
        if area_tematica:
            vector_query = vector_query.where(filter=FieldFilter("area_tematica", "==", area_tematica))
        
        vector_query = vector_query.find_nearest(
            vector_field="embedding",
            query_vector=Vector(query_vector),
            distance_measure=DistanceMeasure.COSINE,
            limit=5,
            distance_result_field="vector_distance"
        )

        nodes_raw = []
        for ndoc in vector_query.stream():
            nd = ndoc.to_dict() or {}
            distance = nd.pop("vector_distance", 1.0)
            score = 1.0 - float(distance)
            nodes_raw.append({
                "titulo": nd.get("titulo"),
                "resumo": nd.get("resumo"),
                "area_tematica": nd.get("area_tematica"),
                "score": score
            })
        
        nodes_raw.sort(key=lambda x: x["score"], reverse=True)
        candidates = nodes_raw[:5]
        
        lines = [f"Resultados do Grafo para: {q_text}"]
        for i, n in enumerate(candidates, 1):
            lines.append(f"[{i}] {n['titulo']} ({n['area_tematica']}) - {n['resumo']}")
            
        # --- FALLBACK: Busca em Tarefas Reais (Regex Estrito do novo módulo) ---
        if len(candidates) < 2:
            from tools.busca_grafo import buscar_tarefas
            res = buscar_tarefas(q_text, area_tematica=area_tematica)
            
            if res.get("erro"):
                lines.append(f"\n⚠️ [ERRO TÉCNICO FallbackGrafo] {res['erro']}")
            else:
                found = res.get("resultados", [])
                if found:
                    lines.append("\n--- Buscando em Tarefas (Execução Real - Regex) ---")
                    for r in found:
                        lines.append(f"TAREFA: {r['titulo']} | STATUS: {r['status']} | DATA: {r['criado_em']} | [Abrir](task:{r['id']})")
                else:
                    lines.append(f"\nNenhum registro encontrado para o termo '{q_text}' no Banco de Dados.")

        return {"context": "\n".join(lines)}
    except Exception as e:
        print(f"DEBUG_ERROR [Grafo]: {str(e)}")
        import traceback
        print(traceback.format_exc())
        return {"context": f"Erro interno ao consultar grafo: {str(e)}"}

@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.GB_1,
    timeout_sec=120
)
def salvarTranscricaoReuniao(req: https_fn.CallableRequest):
    """
    Salva a transcrição consolidada de uma reunião no Google Drive e registra no módulo conhecimento.
    """
    import io
    from datetime import datetime
    from googleapiclient.http import MediaIoBaseUpload

    data = req.data or {}
    content = data.get('content')
    started_at = data.get('startedAt')
    ended_at = data.get('endedAt')
    file_name = data.get('fileName')

    if not isinstance(content, str) or not content.strip():
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Conteúdo da transcrição é obrigatório."
        )

    def _parse_iso_date(value: str | None) -> datetime | None:
        if not value or not isinstance(value, str):
            return None
        try:
            return datetime.fromisoformat(value.replace("Z", "+00:00"))
        except Exception:
            return None

    started_dt = _parse_iso_date(started_at) or datetime.now()
    ended_dt = _parse_iso_date(ended_at) or datetime.now()

    if not file_name:
        file_name = f"Reuniao_{started_dt.strftime('%Y-%m-%d_%H-%M')}.txt"

    try:
        db = get_db()
        service = get_drive_service()

        root_folder_id = None
        try:
            config_doc = db.collection('system').document('config').get()
            if config_doc.exists:
                root_folder_id = (config_doc.to_dict() or {}).get('googleDriveFolderId')
        except Exception as config_err:
            print(f"Aviso: não foi possível ler system/config: {config_err}")

        # Garante a pasta "Reuniões" na raiz configurada para o conhecimento.
        folder_query = (
            "mimeType='application/vnd.google-apps.folder' "
            "and trashed=false "
            "and name='Reuniões'"
        )
        if root_folder_id:
            folder_query += f" and '{root_folder_id}' in parents"

        folders = service.files().list(
            q=folder_query,
            fields='files(id, name)',
            pageSize=1
        ).execute().get('files', [])

        if folders:
            reunioes_folder_id = folders[0]['id']
        else:
            folder_metadata = {
                'name': 'Reuniões',
                'mimeType': 'application/vnd.google-apps.folder'
            }
            if root_folder_id:
                folder_metadata['parents'] = [root_folder_id]

            folder = service.files().create(
                body=folder_metadata,
                fields='id'
            ).execute()
            reunioes_folder_id = folder.get('id')

        payload = content.encode('utf-8')
        media = MediaIoBaseUpload(io.BytesIO(payload), mimetype='text/plain', resumable=True)

        uploaded = service.files().create(
            body={'name': file_name, 'parents': [reunioes_folder_id]},
            media_body=media,
            fields='id, webViewLink, size'
        ).execute()

        file_id = uploaded.get('id')
        web_link = uploaded.get('webViewLink')

        try:
            service.permissions().create(
                fileId=file_id,
                body={'type': 'anyone', 'role': 'reader'}
            ).execute()
        except Exception as perm_e:
            print(f"Aviso: não foi possível definir permissão pública no arquivo de reunião: {perm_e}")

        db.collection('conhecimento').document(file_id).set({
            'id': file_id,
            'titulo': file_name,
            'tipo_arquivo': 'txt',
            'url_drive': web_link,
            'tamanho': int(uploaded.get('size') or len(payload)),
            'data_criacao': datetime.now().isoformat(),
            'origem': {'modulo': 'reunioes', 'id_origem': started_dt.isoformat()},
            'area_tematica': 'Reuniões',
            'parent_id': 'biblioteca',
            'meeting_started_at': started_dt.isoformat(),
            'meeting_ended_at': ended_dt.isoformat()
        }, merge=True)

        return {
            'success': True,
            'fileId': file_id,
            'webViewLink': web_link,
            'fileName': file_name,
            'folderId': reunioes_folder_id
        }

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro ao salvar transcrição de reunião: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message="Falha ao salvar transcrição da reunião no Google Drive."
        )

@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.GB_1,
    timeout_sec=300
)
def analisarPadroesCategoriaIA(req: https_fn.CallableRequest):
    """
    Analisa tarefas de uma area_tematica específica para identificar padrões e propor artefatos de conhecimento.
    """
    from google import genai
    import json
    import re
    import traceback

    data = req.data or {}
    area_tematica = data.get('area_tematica')
    
    if not area_tematica:
         raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="Categoria é obrigatória."
        )

    try:
        db = firestore.client()
        # Busca tarefas concluídas desta area_tematica (limite de 15 para análise)
        tasks_query = db.collection('tarefas')\
            .where('area_tematica', '==', area_tematica)\
            .where('status', '==', 'concluído')\
            .limit(15)
        
        docs = tasks_query.stream()
        contexto_tarefas = []
        for doc in docs:
            t = doc.to_dict()
            contexto_tarefas.append(f"Tarefa: {t.get('titulo')}\nNotas: {t.get('notas')}")

        if not contexto_tarefas:
            return {"success": False, "message": f"Não há tarefas concluídas suficientes em '{area_tematica}' para analisar padrões."}

        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        
        if not gemini_key:
            return {"success": False, "error": "Chave Gemini não configurada no sistema (system/api_keys)."}

        client = genai.Client(api_key=gemini_key)

        prompt = f"""
        Você é o HERMES Master IA. Analise a sequência de tarefas abaixo da area_tematica '{area_tematica}'.
        Sua missão é identificar um PADRÃO de trabalho ou um PROCEDIMENTO que o André segue.
        
        Com base nessas tarefas, crie um "Guia de Procedimento Operacional Padrão" para esta area_tematica.
        
        TAREFAS ANALISADAS:
        {chr(10).join(contexto_tarefas)}
        
        Retorne um JSON com:
        1. titulo: Nome do guia (ex: Procedimento para Licitação de Compras)
        2. conteudo: O guia detalhado em Markdown (passos, dicas, o que não esquecer).
        3. insight: Um breve comentário seu sobre por que isso é importante ou o que você notou de especial.
        """

        response = client.models.generate_content(model="gemini-3.5-flash-lite", contents=prompt)
        res_text = response.text

        json_match = re.search(r'\{.*\}', res_text, re.DOTALL)
        if json_match:
            result_data = json.loads(json_match.group(0))

            # Versionamento não-destrutivo: arquiva versão anterior antes de criar nova
            _titulo_novo = result_data.get('titulo', '')
            _existing_proc = db.collection("conhecimento_mestre")\
                .where("titulo", "==", _titulo_novo)\
                .where("status", "!=", "arquivado_backup")\
                .limit(1).get()
            for _ep in _existing_proc:
                _ep.reference.update({
                    "status": "arquivado_backup",
                    "data_arquivamento": firestore.SERVER_TIMESTAMP
                })
            import uuid as _uuid_analytics
            _new_proc_id = str(_uuid_analytics.uuid4())[:12]
            db.collection("conhecimento_mestre").document(_new_proc_id).set({
                "titulo": _titulo_novo,
                "conteudo_regra": result_data.get('conteudo', ''),
                "justificativa_da_regra": result_data.get('insight', ''),
                "area_tematica": area_tematica,
                "insight_ia": result_data.get('insight'),
                "data_criacao": firestore.SERVER_TIMESTAMP,
                "tipo": "procedimento_aprendido",
                "status": "ativo",
                "necessita_revisao": False,
                "tag_aviso": "",
                "autor": "HERMES_ANALYTICS"
            })
            
            return {"success": True, "data": result_data}
            
        return {"success": False, "error": f"Falha ao analisar padrões estruturados. Resposta da IA: {res_text[:200]}"}

    except Exception as e:
        error_msg = traceback.format_exc()
        print(f"Erro em analisarPadroesCategoriaIA: {error_msg}")
        return {"success": False, "error": str(e), "traceback": error_msg}


# ─────────────────────────────────────────────────────────────────────────────
# MOTOR DE EVOLUÇÃO AUTÔNOMA — Batch Job de Processamento de Correções
# Executa a cada 60 minutos. Lê a fila `correcoes_pendentes`, valida compliance
# via consenso web (Tavily, mín. 5 fontes), refina com Gemini Pro e persiste
# com versionamento não-destrutivo em `conhecimento_mestre`.
# ─────────────────────────────────────────────────────────────────────────────
@scheduler_fn.on_schedule(schedule="every 60 minutes")
def processar_correcoes_pendentes(event: scheduler_fn.ScheduledEvent) -> None:
    """Motor de Evolução Autônoma: processa fila de correções com validação de compliance."""
    import traceback as _evo_tb
    import uuid as _evo_uuid
    import requests as _evo_req
    from google import genai

    _db = get_db()
    _gemini_key = get_gemini_api_key()
    _evo_client = genai.Client(api_key=_gemini_key)

    # Recupera chave Tavily para consenso web
    _tavily_key = ''
    try:
        _keys_doc = _cached_doc_get(_db, 'system', 'api_keys')
        _tavily_key = (_keys_doc.to_dict() or {}).get('tavily_api_key', '')
    except Exception as _key_err:
        print(f"[EvoEngine] Aviso: não foi possível recuperar chave Tavily: {_key_err}")

    # Busca até 10 correções pendentes por ciclo
    try:
        _correcoes = list(
            _db.collection('correcoes_pendentes')
               .where('status', '==', 'pendente')
               .limit(10)
               .get()
        )
    except Exception as _fetch_err:
        print(f"[EvoEngine] Erro ao buscar fila: {_fetch_err}")
        return

    if not _correcoes:
        print("[EvoEngine] Nenhuma correção pendente neste ciclo.")
        return

    print(f"[EvoEngine] Processando {len(_correcoes)} correção(ões).")

    for _corr_doc in _correcoes:
        _corr_id = _corr_doc.id
        _corr = _corr_doc.to_dict()

        try:
            # Marca como em processamento para evitar reprocessamento paralelo
            _db.collection('correcoes_pendentes').document(_corr_id).update(
                {'status': 'processando'}
            )

            _titulo        = _corr.get('titulo_procedimento', '')
            _area          = _corr.get('area_tematica', '')
            _novo_conteudo = _corr.get('novo_conteudo_proposto', '')
            _justificativa = _corr.get('justificativa_usuario', '')

            # ── 1. Busca versão atual do procedimento em conhecimento_mestre ──────
            _old_doc_id   = None
            _old_content  = '(procedimento ainda não existe)'
            try:
                _existing = list(
                    _db.collection('conhecimento_mestre')
                       .where('titulo', '==', _titulo)
                       .where('status', '!=', 'arquivado_backup')
                       .limit(1)
                       .get()
                )
                if _existing:
                    _old_data    = _existing[0].to_dict()
                    _old_content = _old_data.get('conteudo_regra') or _old_data.get('conteudo', '')
                    _old_doc_id  = _existing[0].id
            except Exception as _find_err:
                print(f"[EvoEngine] Aviso ao buscar procedimento existente: {_find_err}")

            # ── 2. Verificação de consenso web (mínimo 5 fontes via Tavily) ───────
            _compliance_ok      = False
            _compliance_sources = []
            _compliance_summary = 'Verificação de compliance não executada (chave ausente).'

            if _tavily_key:
                try:
                    _search_query = (
                        f"procedimento compliance legal {_titulo} {_area} "
                        "legislação brasileira norma vigente"
                    )
                    _t_resp = _evo_req.post(
                        'https://api.tavily.com/search',
                        json={
                            'api_key': _tavily_key,
                            'query': _search_query,
                            'search_depth': 'advanced',
                            'include_answer': True,
                            'max_results': 7
                        },
                        timeout=25
                    )
                    if _t_resp.status_code == 200:
                        _t_data             = _t_resp.json()
                        _compliance_sources = [r.get('url', '') for r in _t_data.get('results', [])]
                        _web_answer         = _t_data.get('answer', '')
                        _n_sources          = len(_compliance_sources)

                        # LLM avalia conformidade com base no consenso web
                        _comp_prompt = (
                            f"Você é um auditor de conformidade legal sênior.\n"
                            f"Avalie se o procedimento proposto está em conformidade com "
                            f"legislação e normas brasileiras vigentes, usando as {_n_sources} "
                            f"fontes web como referência de consenso.\n\n"
                            f"PROCEDIMENTO PROPOSTO:\n{_novo_conteudo}\n\n"
                            f"CONSENSO WEB ({_n_sources} fontes):\n{_web_answer}\n\n"
                            f"Responda EXCLUSIVAMENTE em JSON válido (sem markdown):\n"
                            f'{{\"aprovado\": true_ou_false, \"resumo\": \"motivo em 1 frase\"}}'
                        )
                        _comp_resp    = _evo_client.models.generate_content(
                            model="gemini-3.5-flash-lite",
                            contents=_comp_prompt
                        )
                        _comp_text    = (_comp_resp.text or '').strip()
                        _comp_match   = re.search(r'\{.*\}', _comp_text, re.DOTALL)
                        if _comp_match:
                            _comp_data          = json.loads(_comp_match.group(0))
                            _compliance_ok      = bool(_comp_data.get('aprovado', False))
                            _compliance_summary = _comp_data.get('resumo', '')
                        else:
                            _compliance_summary = f"LLM retornou formato inesperado: {_comp_text[:120]}"
                    else:
                        _compliance_summary = f"Tavily retornou status {_t_resp.status_code}."
                except Exception as _comp_err:
                    _compliance_summary = f"Erro na verificação: {str(_comp_err)}"
                    print(f"[EvoEngine] {_compliance_summary}")

            # ── 3. LLM de raciocínio superior refina o procedimento final ─────────
            _refinement_prompt = (
                "Você é um engenheiro de processos sênior. Integre a correção proposta "
                "ao procedimento atual, mantendo clareza, estrutura Markdown e fidelidade "
                "à justificativa fornecida. Retorne APENAS o conteúdo final em Markdown.\n\n"
                f"PROCEDIMENTO ATUAL:\n{_old_content}\n\n"
                f"CORREÇÃO PROPOSTA:\n{_novo_conteudo}\n\n"
                f"JUSTIFICATIVA:\n{_justificativa}"
            )
            try:
                _refine_resp   = _evo_client.models.generate_content(
                    model="gemini-3.1-pro-preview",
                    contents=_refinement_prompt
                )
                _conteudo_final = (_refine_resp.text or _novo_conteudo).strip()
            except Exception as _refine_err:
                print(f"[EvoEngine] Refinamento Pro falhou, usando conteúdo proposto: {_refine_err}")
                _conteudo_final = _novo_conteudo

            # ── 4. Execução Otimista (Fail-Open) — aplica sempre, sinaliza se falhou ──
            _necessita_revisao = not _compliance_ok
            _tag_aviso = (
                "[⚠️ OTIMIZADO ÀS CEGAS: Validação de compliance falhou]"
                if _necessita_revisao else ""
            )

            # ── 5. Versionamento não-destrutivo — arquiva versão anterior ─────────
            if _old_doc_id:
                try:
                    _db.collection('conhecimento_mestre').document(_old_doc_id).update({
                        'status': 'arquivado_backup',
                        'data_arquivamento': firestore.SERVER_TIMESTAMP
                    })
                except Exception as _arch_err:
                    print(f"[EvoEngine] Aviso ao arquivar versão anterior: {_arch_err}")

            # ── 6. Persiste novo procedimento com campos obrigatórios ─────────────
            _new_id = str(_evo_uuid.uuid4())[:12]
            _db.collection('conhecimento_mestre').document(_new_id).set({
                'titulo':                   _titulo,
                'area_tematica':            _area,
                'conteudo_regra':           _conteudo_final,
                'justificativa_da_regra':   _justificativa,
                'status':                   'ativo',
                'necessita_revisao':        _necessita_revisao,
                'tag_aviso':                _tag_aviso,
                'compliance_aprovado':      _compliance_ok,
                'compliance_resumo':        _compliance_summary,
                'compliance_fontes':        _compliance_sources[:5],
                'data_criacao':             firestore.SERVER_TIMESTAMP,
                'tipo':                     'procedimento_evoluido',
                'autor':                    'HERMES_EVOLUTION_ENGINE',
                'origem_correcao_id':       _corr_id,
                'procedimento_anterior_id': _old_doc_id or ''
            })

            # ── 7. Fecha a correção na fila ───────────────────────────────────────
            _db.collection('correcoes_pendentes').document(_corr_id).update({
                'status':              'processado',
                'novo_doc_id':         _new_id,
                'compliance_aprovado': _compliance_ok,
                'data_processamento':  firestore.SERVER_TIMESTAMP
            })

            _status_str = "COMPLIANCE OK" if _compliance_ok else "FAIL-OPEN (necessita_revisao=True)"
            print(f"[EvoEngine] ✅ Correção {_corr_id} → doc {_new_id} | {_status_str}")

        except Exception as _proc_err:
            print(f"[EvoEngine] ❌ Erro ao processar correção {_corr_id}:\n{_evo_tb.format_exc()}")
            try:
                _db.collection('correcoes_pendentes').document(_corr_id).update({
                    'status':   'erro',
                    'erro_msg': str(_proc_err)
                })
            except Exception:
                pass


# ─── Salvar Relatório no Google Drive ────────────────────────────────────────

@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=120
)
def salvarRelatorioNoDrive(req: https_fn.CallableRequest):
    """
    Converte um relatório Markdown para HTML e faz upload no Google Drive
    como Google Doc editável nativo.
    """
    data = req.data or {}
    relatorio_id = data.get('relatorioId')

    if not relatorio_id:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="relatorioId é obrigatório."
        )

    try:
        db = get_db()

        # 1. Busca o relatório no Firestore
        rel_doc = db.collection('relatorios').document(relatorio_id).get()
        if not rel_doc.exists:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.NOT_FOUND,
                message=f"Relatório '{relatorio_id}' não encontrado."
            )

        rel = rel_doc.to_dict()
        titulo = rel.get('titulo', 'Relatório Hermes')
        markdown_text = rel.get('markdown', '')

        # 2. Converte Markdown → HTML (conversor inline sem dependências externas)
        def _md_to_html(md: str) -> str:
            lines = md.split('\n')
            html = []
            in_ul = False
            in_table = False
            table_header_parsed = False

            def flush_lists():
                nonlocal in_ul
                if in_ul:
                    html.append('</ul>')
                    in_ul = False

            def flush_table():
                nonlocal in_table
                if in_table:
                    html.append('</table>')
                    in_table = False

            for line in lines:
                s = line.strip()
                
                # TABLE DETECTION
                if (s.startswith('|') or (in_table and s.count('|') >= 1)) and s.endswith('|'):
                    flush_lists()
                    if not in_table:
                        html.append('<table border="1" style="border-collapse: collapse; width: 100%; margin-bottom: 20px;">')
                        in_table = True
                        table_header_parsed = False
                    
                    if re.match(r'^\|[\s:-|]+\|$', s):
                        table_header_parsed = True
                        continue
                    
                    # Basic cell extraction
                    cells = [c.strip() for c in line.strip('|').split('|')]
                    tag = 'th' if not table_header_parsed else 'td'
                    html.append('  <tr>')
                    for cell in cells:
                        # Inline formatting for cells
                        c_html = cell
                        c_html = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', c_html)
                        c_html = re.sub(r'\*(.+?)\*', r'<em>\1</em>', c_html)
                        c_html = re.sub(r'`(.+?)`', r'<code>\1</code>', c_html)
                        html.append(f'    <{tag} style="padding: 10px; border: 1px solid #ccc; text-align: left;">{c_html}</{tag}>')
                    html.append('  </tr>')
                    continue
                else:
                    flush_table()

                s_orig = line.rstrip()
                if not s_orig.strip():
                    flush_lists()
                    html.append('<br/>')
                    continue

                if s_orig.startswith('### '):
                    flush_lists()
                    html.append(f'<h3>{s_orig[4:]}</h3>')
                elif s_orig.startswith('## '):
                    flush_lists()
                    html.append(f'<h2>{s_orig[3:]}</h2>')
                elif s_orig.startswith('# '):
                    flush_lists()
                    html.append(f'<h1>{s_orig[2:]}</h1>')
                elif s_orig == '---':
                    flush_lists()
                    html.append('<hr/>')
                elif s_orig.startswith('- ') or s_orig.startswith('* '):
                    if not in_ul:
                        html.append('<ul>')
                        in_ul = True
                    # Inline formatting for list items
                    li_text = s_orig[2:]
                    li_text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', li_text)
                    li_text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', li_text)
                    html.append(f'  <li>{li_text}</li>')
                else:
                    flush_lists()
                    # Regular paragraph with inline formatting
                    p_text = s_orig
                    p_text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', p_text)
                    p_text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', p_text)
                    p_text = re.sub(r'`(.+?)`', r'<code>\1</code>', p_text)
                    html.append(f'<p>{p_text}</p>')

            flush_lists()
            flush_table()

            body = '\n'.join(html)
            return f'<html><head><meta charset="utf-8"/></head><body style="font-family: Arial, sans-serif; line-height: 1.6; color: #333;">{body}</body></html>'

        html_content = _md_to_html(markdown_text)

        # 3. Upload para o Drive como Google Doc nativo
        drive_service = get_drive_service()
        from googleapiclient.http import MediaInMemoryUpload

        media = MediaInMemoryUpload(
            html_content.encode('utf-8'),
            mimetype='text/html',
            resumable=False
        )
        uploaded = drive_service.files().create(
            body={'name': titulo, 'mimeType': 'application/vnd.google-apps.document'},
            media_body=media,
            fields='id,webViewLink'
        ).execute()

        drive_file_id = uploaded.get('id')
        drive_url = uploaded.get('webViewLink')

        # 4. Atualiza Firestore com dados do Drive
        db.collection('relatorios').document(relatorio_id).update({
            'driveFileId': drive_file_id,
            'driveUrl': drive_url
        })

        print(f"[Relatorio] '{titulo}' salvo no Drive — fileId: {drive_file_id}")
        return {"driveFileId": drive_file_id, "driveUrl": drive_url}

    except https_fn.HttpsError:
        raise
    except Exception as e:
        import traceback
        print(f"[Relatorio] Erro ao salvar no Drive: {e}\n{traceback.format_exc()}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )

@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=300
)
def sintetizarDescricaoAcao(req: https_fn.CallableRequest):
    """
    Sintetiza descricoes executivas para acoes sem descricao, usando diario e plano de acao.
    Modos:
    - {taskId}: processa uma tarefa especifica.
    - {batch: true, limit?: number}: processa lote limitado de tarefas elegiveis.
    - {dryRun: true}: conta tarefas elegiveis sem chamar o modelo.
    """
    from google import genai as _genai

    if not req.auth:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.UNAUTHENTICATED,
            message="Autenticacao obrigatoria."
        )

    data = req.data or {}
    task_id = str(data.get("taskId") or "").strip()
    dry_run = bool(data.get("dryRun"))
    batch_mode = bool(data.get("batch")) or not task_id
    requested_limit = data.get("limit", DESCRIPTION_SYNTHESIS_BATCH_LIMIT)

    try:
        limit = int(requested_limit)
    except Exception:
        limit = DESCRIPTION_SYNTHESIS_BATCH_LIMIT
    limit = max(1, min(limit, DESCRIPTION_SYNTHESIS_BATCH_LIMIT))

    try:
        db_ref = get_db()

        if dry_run:
            return {
                "status": "preview",
                "eligibleCount": _count_eligible_description_synthesis_tasks(db_ref),
                "limit": limit
            }

        keys_doc = _cached_doc_get(db_ref, "system", "api_keys")
        gemini_key = keys_doc.to_dict().get("gemini_api_key") if keys_doc.exists else None
        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini nao configurada em system/api_keys."
            )

        client = _genai.Client(api_key=gemini_key)
        now_iso = datetime.now(timezone.utc).isoformat()

        def process_task(task_ref, task_data: dict) -> dict:
            eligible, reason = _is_task_eligible_for_description_synthesis(task_data)
            if not eligible:
                return {"id": task_ref.id, "status": "ignored", "reason": reason}

            description = _generate_action_description(client, task_data)
            task_ref.update({
                "descricao": description,
                "data_atualizacao": now_iso,
                "descricao_sintetizada_por": "gemini",
                "descricao_sintetizada_modelo": DESCRIPTION_SYNTHESIS_MODEL,
                "descricao_sintetizada_em": now_iso
            })
            return {
                "id": task_ref.id,
                "status": "completed",
                "descriptionLength": len(description)
            }

        if task_id:
            task_ref = db_ref.collection("tarefas").document(task_id)
            task_doc = task_ref.get()
            if not task_doc.exists:
                raise https_fn.HttpsError(
                    code=https_fn.FunctionsErrorCode.NOT_FOUND,
                    message="Acao nao encontrada."
                )
            result = process_task(task_ref, task_doc.to_dict() or {})
            return {**result, "processed": 1 if result.get("status") == "completed" else 0}

        processed = []
        ignored = 0
        failed = []
        scanned = 0

        for task_doc in db_ref.collection("tarefas").stream():
            scanned += 1
            task_data = task_doc.to_dict() or {}
            eligible, _reason = _is_task_eligible_for_description_synthesis(task_data)
            if not eligible:
                ignored += 1
                continue
            if len(processed) >= limit:
                break
            try:
                result = process_task(task_doc.reference, task_data)
                if result.get("status") == "completed":
                    processed.append(result["id"])
                else:
                    ignored += 1
            except Exception as item_error:
                failed.append({"id": task_doc.id, "error": str(item_error)[:300]})

        remaining_estimate = max(0, _count_eligible_description_synthesis_tasks(db_ref))
        return {
            "status": "completed",
            "batch": batch_mode,
            "scanned": scanned,
            "processed": len(processed),
            "processedIds": processed,
            "ignored": ignored,
            "failed": failed,
            "remainingEligible": remaining_estimate,
            "limit": limit
        }

    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"[sintetizarDescricaoAcao] Erro: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=str(e)
        )


@https_fn.on_call(
    cors=options.CorsOptions(cors_origins="*", cors_methods=["POST"]),
    memory=options.MemoryOption.MB_512,
    timeout_sec=30
)
def analisarInsightProativo(req: https_fn.CallableRequest):
    """
    Analisa o contexto de uma tarefa e retorna um insight proativo (se relevante).
    Chamado com debounce após mudanças no diário ou plano de ação.
    Retorna: { nivel: 1|2|null, texto: str|null, alvo: "diario"|"plano"|null, planoProposto: list|null }
    """
    from google import genai as _genai
    import json as _json

    data = req.data or {}
    task_id = (data.get('taskId') or '').strip()
    titulo = (data.get('titulo') or '').strip()
    status = (data.get('status') or '').strip()
    data_limite = (data.get('dataLimite') or '').strip()
    prazo_final = (data.get('prazoFinal') or '').strip()
    plano_acao = data.get('planoAcao') or []
    acompanhamento_recente = data.get('acompanhamentoRecente') or []
    # Opiniões que o usuário marcou como "não sugerir novamente" — supressão permanente.
    insights_ignorados = [
        t.strip() for t in (data.get('insightsIgnorados') or [])
        if isinstance(t, str) and t.strip()
    ]

    if not task_id:
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="taskId é obrigatório."
        )

    try:
        db = get_db()
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None

        if not gemini_key:
            return {"nivel": None, "texto": None, "alvo": None, "planoProposto": None}

        client = _genai.Client(api_key=gemini_key)

        plano_txt = '\n'.join([
            f"{i+1}. [{'X' if item.get('completed') else ' '}] {item.get('text', '')}"
            for i, item in enumerate(plano_acao)
        ]) or 'Nenhum passo definido.'

        diario_txt = '\n'.join([
            f"[{entry.get('data', '')}] {entry.get('nota', '')}"
            for entry in acompanhamento_recente[-10:]
        ]) or 'Nenhum registro.'

        # Insights proativos já emitidos recentemente para esta ação. Sem isso, o
        # modelo re-analisa o estado do zero a cada mudança no plano/diário e, se
        # uma preocupação anterior não foi sanada, reformula a MESMA observação com
        # outras palavras — gerando duplicatas. Carregamos os últimos para que o
        # modelo evite repetir e para a rede de segurança por similaridade abaixo.
        insights_anteriores = []
        try:
            sess_docs = db.collection('sessoes_copiloto')\
                .where('taskId', '==', task_id)\
                .limit(10).get()
            for sdoc in sess_docs:
                sdata = sdoc.to_dict()
                if sdata.get('isTemporary'):
                    continue
                # Ordena só por timestamp (sem where composto) para não exigir
                # índice composto; filtra subtype em memória.
                msg_docs = sdoc.reference.collection('mensagens')\
                    .order_by('timestamp', direction=firestore.Query.DESCENDING)\
                    .limit(30).get()
                for mdoc in msg_docs:
                    m = mdoc.to_dict()
                    if m.get('subtype') == 'proactive_insight':
                        txt = (m.get('content') or '').strip()
                        if txt:
                            insights_anteriores.append(txt)
        except Exception as _e:
            print(f"[analisarInsightProativo] Falha ao carregar insights anteriores: {_e}")
        insights_anteriores = insights_anteriores[:5]

        insights_anteriores_txt = '\n'.join(
            f"- {t}" for t in insights_anteriores
        ) or 'Nenhum insight emitido ainda.'

        insights_ignorados_txt = '\n'.join(
            f"- {t}" for t in insights_ignorados[:30]
        ) or 'Nenhuma.'

        prompt = f"""Você é um analista de produtividade. Analise o estado atual desta ação e determine se há um insight genuinamente útil.

CONTEXTO:
Título: {titulo}
Status: {status}
Prazo Final: {prazo_final or 'Não definido'}
Data de Execução: {data_limite or 'Não definida'}
REVISÃO DE DATAS: O Hermes entende que 'Data de Execução' é apenas planejamento, não o prazo final. Ignore contradições de datas entre execução e diário.

PLANO DE AÇÃO:
{plano_txt}

DIÁRIO (entradas recentes):
{diario_txt}

INSIGHTS QUE VOCÊ JÁ EMITIU RECENTEMENTE PARA ESTA AÇÃO:
{insights_anteriores_txt}

OPINIÕES QUE O USUÁRIO PEDIU EXPLICITAMENTE PARA IGNORAR (NÃO LEVANTE NENHUMA DESTAS, NEM REFORMULADA):
{insights_ignorados_txt}

CLASSIFICAÇÃO:
- NIVEL_1 (Crítico): Erro lógico real (ex: planejar algo para o passado ou pós-conclusão), gargalo crítico esquecido.
- NIVEL_2 (Otimização): Sugestão de melhoria ou passo faltante.
- NIVEL_3 (Ideia Criativa): Conexões estratégicas ou ideias laterais.
- SEM_INSIGHT: situação adequada.

REGRAS ADICIONAIS:
- NUNCA critique a "Data de Execução" se ela for antes do evento mencionado (isso é ANTECIPAÇÃO).
- Só critique a "Data de Execução" se ela for DEPOIS do "Prazo Final" (se existir).
- Seja conciso e evite ser óbvio.

ESTILO:
- Escreva o campo "texto" sem saudacao, preambulo ou encerramento. Maximo 1 a 3 bullets curtos, ou uma frase direta.
- Para insights proativos, aponte somente a correlacao, ponto cego ou proximo passo objetivo.

REGRAS:
- Só retorne insight se for genuinamente valioso. Evite insights genéricos ou óbvios.
- ANTI-REPETIÇÃO (regra mais importante): se a sua observação for substancialmente a mesma (mesmo tema, mesmo ponto cego ou mesma sugestão) de algum item em "INSIGHTS QUE VOCÊ JÁ EMITIU RECENTEMENTE" — ainda que reformulada com outras palavras — e o plano/diário ainda NÃO tiver incorporado aquela sugestão, retorne SEM_INSIGHT. Não fique reiterando uma opinião que não foi sanada. Só volte a um tema já levantado se houver um fato genuinamente novo a acrescentar.
- OPINIÕES IGNORADAS (regra absoluta): se a sua observação coincidir com qualquer item em "OPINIÕES QUE O USUÁRIO PEDIU EXPLICITAMENTE PARA IGNORAR" — mesmo que parcialmente ou reformulada — retorne SEM_INSIGHT obrigatoriamente. O usuário já decidiu não tratar disso.
- Para alvo "plano", inclua plano_proposto com todos os itens revisados (array de objetos com "id", "text", "completed").
- Para alvo "acoes", inclua acoes_propostas (array de objetos com "titulo", "descricao", "tags").
- Use IDs de 8 chars para itens novos no plano. Preserve id e completed dos itens existentes quando mantidos.
- Se o alvo não for "plano", plano_proposto deve ser null. Se não houver novas ações, acoes_propostas deve ser null.

RESPONDA APENAS COM JSON VÁLIDO (sem markdown):
{{"nivel": 1|2|3|null, "texto": "...", "alvo": "diario"|"plano"|"acoes"|null, "plano_proposto": [...]|null, "acoes_propostas": [...]|null}}

Se SEM_INSIGHT: {{"nivel": null, "texto": null, "alvo": null, "plano_proposto": null, "acoes_propostas": null}}"""

        # Roteado pelo logger central só para entrar na telemetria
        # (system_usage/gemini). NÃO habilitar flex aqui: este callable tem
        # timeout_sec=30 e o flex pode enfileirar por até 600s, fazendo o
        # Firebase matar a chamada e o frontend perder o insight silenciosamente.
        # O ganho de -50% num call flash-lite minúsculo não compensa o risco.
        response = generate_content_logged(
            client,
            model=GEMINI_LIGHT_MODEL,
            contents=prompt,
            feature="copilot_proactive_insight",
            db=db,
            config={"temperature": 0.3, "max_output_tokens": 1024},
        )

        result_text = (response.text or '').strip()
        if result_text.startswith('```'):
            lines = result_text.split('\n')
            lines = [l for l in lines if not l.startswith('```')]
            result_text = '\n'.join(lines).strip()

        parsed = _json.loads(result_text)
        nivel = parsed.get('nivel')
        texto = parsed.get('texto')
        alvo = parsed.get('alvo')
        plano_proposto = parsed.get('plano_proposto')
        acoes_propostas = parsed.get('acoes_propostas')

        if nivel not in (1, 2, 3, None):
            nivel = None
        if alvo not in ('diario', 'plano', 'acoes', None):
            alvo = None

        # Rede de segurança: mesmo com as regras acima, o modelo pode reformular um
        # insight anterior ou uma opinião ignorada. Suprimimos quando há forte
        # sobreposição de termos (Jaccard). Opiniões ignoradas usam limiar menor
        # (mais agressivo), pois o usuário decidiu explicitamente não tratá-las.
        if texto:
            import re as _re

            def _norm_tokens(s):
                s = (s or '').lower()
                s = _re.sub(r"[^0-9a-zà-ú\s]", ' ', s)
                return {t for t in s.split() if len(t) > 3}

            new_tokens = _norm_tokens(texto)
            if new_tokens:
                for prev, limiar, origem in (
                    [(p, 0.4, 'ignorado') for p in insights_ignorados] +
                    [(p, 0.5, 'recente') for p in insights_anteriores]
                ):
                    prev_tokens = _norm_tokens(prev)
                    if not prev_tokens:
                        continue
                    union = len(new_tokens | prev_tokens)
                    jaccard = len(new_tokens & prev_tokens) / union if union else 0
                    if jaccard >= limiar:
                        print(f"[analisarInsightProativo] Insight suprimido por similaridade ({origem}, jaccard={jaccard:.2f})")
                        return {"nivel": None, "texto": None, "alvo": None, "planoProposto": None, "acoesPropostas": None}

        return {
            "nivel": nivel,
            "texto": texto,
            "alvo": alvo,
            "planoProposto": plano_proposto if alvo == 'plano' and isinstance(plano_proposto, list) else None,
            "acoesPropostas": acoes_propostas if alvo == 'acoes' and isinstance(acoes_propostas, list) else None
        }

    except Exception as e:
        print(f"[analisarInsightProativo] Erro: {e}")
        return {"nivel": None, "texto": None, "alvo": None, "planoProposto": None}



def calculadora(expressao: str) -> str:
    """Calculadora dedicada para calculos matematicos ad-hoc ou projecoes hipoteticas. Nao utilize para grandes matrizes."""
    import math
    try:
        allowed_names = {k: v for k, v in math.__dict__.items() if not k.startswith("__")}
        code = compile(expressao, "<string>", "eval")
        for name in code.co_names:
            if name not in allowed_names:
                raise NameError(f"O uso de '{name}' nao e permitido.")
        result = eval(code, {"__builtins__": {}}, allowed_names)
        return str(result)
    except Exception as e:
        return f"Erro de calculo: {e}"

@https_fn.on_call(memory=options.MemoryOption.MB_256)
def classificarAreaTematica(req: https_fn.CallableRequest) -> dict:
    from google import genai
    from google.genai import types
    data = req.data
    titulo = data.get("titulo", "")
    descricao = data.get("descricao", "")
    notas = data.get("notas", "")
    texto = f"Titulo: {titulo}\nDescricao: {descricao}\nNotas: {notas}"
    
    api_key = get_gemini_api_key()
    if not api_key:
        return {"area_tematica": "Nenhuma"}
        
    client = genai.Client(api_key=api_key)
    try:
        response = client.models.generate_content(
            model="gemini-3.5-flash-lite",
            contents=f"Classifique a Area Tematica do seguinte texto como uma de: Saude, Financeira, Nenhuma. Retorne APENAS a palavra correta.\n\nTexto: {texto}",
            config=types.GenerateContentConfig(
                temperature=0.0,
                max_output_tokens=10
            )
        )
        res = (response.text or "Nenhuma").strip().strip('\'"').capitalize()
        if res == "Saude":
            res = "Saúde"
        if res not in ["Saúde", "Financeira", "Nenhuma"]:
            res = "Nenhuma"
            
        return {"area_tematica": res}
    except Exception as e:
        print(f"Erro na classificacao tematica: {e}")
        return {"area_tematica": "Nenhuma"}


# ---------------------------------------------------------------------------
# Configuração das automações multi-canal (vínculo sinal↔ação, diário pessoal,
# ingestão de WhatsApp) — únicas callables que tocam system/settings a partir
# do frontend. O documento é bloqueado por regra de segurança para o cliente
# (firestore.rules: match /system/{document=**} { allow read, write: if false; }),
# então esta é a única porta de entrada; só toca os campos explicitamente
# tratados aqui, o resto de system/settings continua editável só via Console.
# ---------------------------------------------------------------------------

def _serialize_whatsapp_worker_heartbeat(db) -> dict:
    doc = db.collection("system").document("whatsapp_worker").get()
    if not doc.exists:
        return {"online": False, "last_seen": None}
    data = doc.to_dict() or {}
    last_seen = data.get("last_seen")
    online = False
    last_seen_iso = None
    if last_seen:
        try:
            last_seen_iso = last_seen.isoformat()
            online = (datetime.now(timezone.utc) - last_seen) <= timedelta(minutes=10)
        except Exception:
            pass
    return {"online": online, "last_seen": last_seen_iso}


INTERNAL_USER_EMAIL = "andre.martiini@gmail.com"


def _require_internal_user(req: https_fn.CallableRequest) -> None:
    """Mesma checagem de `internalUser()` em firestore.rules — essas callables
    tocam system/settings, que é bloqueado para qualquer leitura/escrita direta
    do cliente, então exigem o mesmo dono verificado que as regras exigem."""
    token = (req.auth.token if req.auth else None) or {}
    email = token.get("email") if hasattr(token, "get") else None
    email_verified = token.get("email_verified") if hasattr(token, "get") else None
    if not req.auth or not email_verified or email != INTERNAL_USER_EMAIL:
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.PERMISSION_DENIED, message="Acesso restrito.")


@https_fn.on_call(memory=options.MemoryOption.MB_256, timeout_sec=30)
def getAutomationSettings(req: https_fn.CallableRequest) -> dict:
    """Lê o subconjunto de system/settings das automações multi-canal."""
    _require_internal_user(req)

    db = get_db()
    doc = db.collection("system").document("settings").get()
    data = doc.to_dict() if doc.exists else {}
    email_cfg = data.get("email_action_linker") or {}
    diary_cfg = data.get("personal_diary") or {}
    wa_cfg = data.get("whatsapp_ingest") or {}

    return {
        "email_action_linker": {"enabled": bool(email_cfg.get("enabled", False))},
        "personal_diary": {"enabled": bool(diary_cfg.get("enabled", False))},
        "whatsapp_ingest": {
            "enabled": bool(wa_cfg.get("enabled", False)),
            "chats_allowlist": list(wa_cfg.get("chats_allowlist") or []),
        },
        "whatsapp_auto_send_enabled": bool(data.get("whatsapp_auto_send_enabled", False)),
        "whatsapp_worker": _serialize_whatsapp_worker_heartbeat(db),
    }


@https_fn.on_call(memory=options.MemoryOption.MB_256, timeout_sec=30)
def updateAutomationSettings(req: https_fn.CallableRequest) -> dict:
    """Atualiza o subconjunto whitelisted de system/settings das automações
    multi-canal. Usa merge com dicts aninhados (não field paths com ponto —
    set(merge=True) não expande ponto em string, só update() faz isso, e
    update() falha se o doc não existir) para não pisar em campos irmãos."""
    _require_internal_user(req)

    data = req.data or {}
    updates: dict = {}

    email_cfg = data.get("email_action_linker")
    if isinstance(email_cfg, dict) and "enabled" in email_cfg:
        updates["email_action_linker"] = {"enabled": bool(email_cfg["enabled"])}

    diary_cfg = data.get("personal_diary")
    if isinstance(diary_cfg, dict) and "enabled" in diary_cfg:
        updates["personal_diary"] = {"enabled": bool(diary_cfg["enabled"])}

    wa_cfg = data.get("whatsapp_ingest")
    if isinstance(wa_cfg, dict):
        wa_updates: dict = {}
        if "enabled" in wa_cfg:
            wa_updates["enabled"] = bool(wa_cfg["enabled"])
        if "chats_allowlist" in wa_cfg and isinstance(wa_cfg["chats_allowlist"], list):
            wa_updates["chats_allowlist"] = [str(x).strip() for x in wa_cfg["chats_allowlist"] if str(x).strip()]
        if wa_updates:
            updates["whatsapp_ingest"] = wa_updates

    if "whatsapp_auto_send_enabled" in data:
        updates["whatsapp_auto_send_enabled"] = bool(data["whatsapp_auto_send_enabled"])

    if not updates:
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message="Nenhum campo valido para atualizar.")

    db = get_db()
    db.collection("system").document("settings").set(updates, merge=True)
    return {"success": True}


@https_fn.on_call(memory=options.MemoryOption.MB_256, timeout_sec=30)
def toggleWhatsappChatMonitored(req: https_fn.CallableRequest) -> dict:
    """Adiciona ou remove um chat da allowlist de captura ao vivo do WhatsApp (system/settings.whatsapp_ingest.chats_allowlist)."""
    _require_internal_user(req)
    req_data = req.data if isinstance(req.data, dict) else {}
    chat_id = str(req_data.get("chat_id") or "").strip()
    if not chat_id:
        raise https_fn.HttpsError(code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT, message="chat_id obrigatório.")

    monitored = bool(req_data.get("monitored"))
    db = get_db()
    settings_ref = db.collection("system").document("settings")
    doc_snap = settings_ref.get()
    settings_data = doc_snap.to_dict() if doc_snap.exists else {}
    wa_cfg = settings_data.get("whatsapp_ingest") or {}
    allowlist = [str(x).strip() for x in (wa_cfg.get("chats_allowlist") or []) if str(x).strip()]

    if monitored:
        if chat_id not in allowlist:
            allowlist.append(chat_id)
    else:
        allowlist = [x for x in allowlist if x != chat_id]

    settings_ref.set({"whatsapp_ingest": {"chats_allowlist": allowlist}}, merge=True)
    return {"success": True, "chat_id": chat_id, "monitored": monitored, "allowlist": allowlist}


_WHATSAPP_CHATS_SCAN_LIMIT = 1000  # teto de mensagens recentes escaneadas para descobrir chats distintos
_WHATSAPP_CHATS_REGISTRY_LIMIT = 3000  # teto de chats lidos do registro whatsapp_chats


@https_fn.on_call(memory=options.MemoryOption.MB_256, timeout_sec=30)
def listWhatsappChats(req: https_fn.CallableRequest) -> dict:
    """Lista os chats de WhatsApp conhecidos/registrados para o seletor de vínculo em
    TaskExecutionView.tsx e para a Caixa de Entrada (WhatsappInboxView.tsx).
    
    Fontes de dados (composição e precedência):
    1. Allowlist em system/settings.whatsapp_ingest.chats_allowlist
    2. Mensagens capturadas recentemente em whatsapp_messages
    3. Registro de chats salvo pelo worker em whatsapp_chats
    
    Precedência de nomes:
    ID cru (allowlist) < nome em mensagem capturada < nome no registro whatsapp_chats
    
    Parâmetros:
    - include_all (bool, default False): se False, retorna apenas conversas monitoradas
      (allowlist ∪ capturadas). Se True, inclui o registro completo de chats, marcando
      monitored=False nas que não estão monitoradas.
    """
    _require_internal_user(req)
    db = get_db()
    req_data = req.data if isinstance(req.data, dict) else {}
    include_all = bool(req_data.get("include_all"))

    # 1. Allowlist
    settings_doc = db.collection("system").document("settings").get()
    settings_data = settings_doc.to_dict() if settings_doc.exists else {}
    allowlist = {
        str(x).strip() for x in (settings_data.get("whatsapp_ingest") or {}).get("chats_allowlist") or []
        if str(x).strip()
    }

    # 2. Mensagens capturadas recentemente
    captured_chats: dict[str, dict] = {}
    docs = (
        db.collection("whatsapp_messages")
        .order_by("ingested_at", direction=firestore.Query.DESCENDING)
        .limit(_WHATSAPP_CHATS_SCAN_LIMIT)
        .stream()
    )
    for doc in docs:
        data = doc.to_dict() or {}
        chat_id = str(data.get("chat_id") or "").strip()
        if not chat_id or chat_id in captured_chats:
            continue  # já capturamos o doc mais recente (streaming desc)
        ts = data.get("timestamp")
        ts_iso = ts.isoformat() if hasattr(ts, "isoformat") else (str(ts) if ts else None)
        captured_chats[chat_id] = {
            "chat_name": str(data.get("chat_name") or "").strip(),
            "is_group": bool(data.get("is_group")),
            "last_activity_ts": ts_iso,
        }

    # 3. Registro salvo pelo worker
    registry_chats: dict[str, dict] = {}
    reg_docs = db.collection("whatsapp_chats").limit(_WHATSAPP_CHATS_REGISTRY_LIMIT).stream()
    for doc in reg_docs:
        data = doc.to_dict() or {}
        chat_id = str(data.get("chat_id") or doc.id).strip()
        if not chat_id:
            continue
        ts = data.get("last_activity_ts")
        ts_iso = ts.isoformat() if hasattr(ts, "isoformat") else (str(ts) if ts else None)
        registry_chats[chat_id] = {
            "chat_name": str(data.get("chat_name") or "").strip(),
            "is_group": bool(data.get("is_group")),
            "last_activity_ts": ts_iso,
        }

    # Determina o conjunto de IDs a incluir
    if include_all:
        target_ids = allowlist | set(captured_chats.keys()) | set(registry_chats.keys())
    else:
        target_ids = allowlist | set(captured_chats.keys())

    chats: dict[str, dict] = {}
    for cid in target_ids:
        is_monitored = (cid in allowlist) or (cid in captured_chats)

        # Grupo
        is_group = cid.endswith("@g.us")
        if cid in captured_chats:
            is_group = captured_chats[cid]["is_group"]
        if cid in registry_chats:
            is_group = registry_chats[cid]["is_group"]

        # Precedência de nome inteligente:
        # 1. Registro salvo pelo worker (whatsapp_chats) se tiver nome real
        # 2. Mensagem capturada (whatsapp_messages) se tiver nome real
        # 3. ID cru como fallback
        name = cid
        if cid in registry_chats and registry_chats[cid].get("chat_name"):
            reg_n = str(registry_chats[cid]["chat_name"]).strip()
            if reg_n and reg_n != cid and not reg_n.endswith("@g.us") and not reg_n.endswith("@lid") and not reg_n.endswith("@c.us"):
                name = reg_n

        if name == cid and cid in captured_chats and captured_chats[cid].get("chat_name"):
            cap_n = str(captured_chats[cid]["chat_name"]).strip()
            if cap_n and cap_n != cid and not cap_n.endswith("@g.us") and not cap_n.endswith("@lid") and not cap_n.endswith("@c.us"):
                name = cap_n

        # Timestamp de última atividade: mensagens capturadas > registro whatsapp_chats
        last_ts = None
        if cid in captured_chats and captured_chats[cid].get("last_activity_ts"):
            last_ts = captured_chats[cid]["last_activity_ts"]
        elif cid in registry_chats and registry_chats[cid].get("last_activity_ts"):
            last_ts = registry_chats[cid]["last_activity_ts"]

        chats[cid] = {
            "chat_id": cid,
            "chat_name": name or cid,
            "is_group": is_group,
            "monitored": is_monitored,
            "last_activity_ts": last_ts,
        }

    # Ordenação estável: mais recentes primeiro (last_activity_ts DESC), com desempate por nome ASC
    result = sorted(chats.values(), key=lambda c: (c.get("chat_name") or "").lower())
    result = sorted(result, key=lambda c: c.get("last_activity_ts") or "", reverse=True)
    return {"chats": result}


@firestore_fn.on_document_created(
    document="whatsapp_consolidacoes/{jobId}",
    memory=options.MemoryOption.GB_1,
    timeout_sec=540,
)
def on_whatsapp_consolidacao_created(event: firestore_fn.Event[firestore_fn.DocumentSnapshot | None]):
    """Processa um job de consolidação da Caixa de Entrada WhatsApp (o frontend cria
    o doc via addDoc — o write é o RPC, mesmo padrão de copilot_jobs). Núcleo em
    whatsapp_consolidation.py; aqui só o guard de idempotência e o error handling."""
    snap = event.data
    if snap is None or not snap.exists:
        return

    db = get_db()
    job_ref = snap.reference
    job = snap.to_dict() or {}
    if job.get("status") != "queued":
        return  # retry do trigger ou doc criado já processado — não reprocessa

    from whatsapp_consolidation import process_consolidation_job
    try:
        process_consolidation_job(db, job_ref, job)
    except Exception as exc:
        print(f"[WA-CONSOL] Job {job_ref.id} falhou: {exc}")
        job_ref.set({
            "status": "error",
            "error": str(exc)[:500],
            "progress": None,
            "updated_at": firestore.SERVER_TIMESTAMP,
        }, merge=True)


# Import daily WIP reset job
from daily_reset_job import daily_wip_reset_and_degradation

# Import daily morning briefing job
from daily_morning_briefing import briefing_matinal_acoes

# Import monthly recurring actions job
from monthly_recurring_actions import gerar_acoes_recorrentes_mensais

# Import daily AI notification planner job
from ai_notification_planner import ai_notification_planner_daily

# Import personal diary + weekly personality consolidation jobs
from personal_diary import gerar_diario_pessoal, consolidar_personalidade

# Import weekly health summary + reevaluation reminder jobs
from health_weekly_summary import gerar_resumo_semanal_saude, verificar_reavaliacoes_saude


@https_fn.on_call(memory=options.MemoryOption.MB_512, timeout_sec=60)
def gerarResumoFinanceiro(req: https_fn.CallableRequest):
    """
    Gera um diagnóstico estruturado de saúde financeira a partir de um snapshot compacto.
    Projetado para ser chamado apenas quando os dados financeiros mudam (fingerprint diferente).
    """
    from google import genai
    from google.genai import types

    data = req.data or {}
    snapshot = data.get('snapshot')

    if not snapshot or not isinstance(snapshot, dict):
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="snapshot é obrigatório."
        )

    import json
    import re

    def _safe_num(value, default=0.0):
        try:
            return float(value)
        except (TypeError, ValueError):
            return default

    def _money(v):
        return f"R$ {v:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    def _finance_fallback_analysis(reason: str = ""):
        fluxo = snapshot.get("fluxoCaixa") or {}
        gastos = snapshot.get("gastos") or {}
        reserva = snapshot.get("reserve") or {}
        saidas = snapshot.get("obrigacoesSaida") or {}

        renda = _safe_num(fluxo.get("receivedIncome"))
        gasto = _safe_num(fluxo.get("spent"))
        saldo = _safe_num(fluxo.get("balance"), renda - gasto)
        orcamento = _safe_num(gastos.get("budget"))
        contas_pendentes = _safe_num(saidas.get("pendingAmount"))
        cobertura = _safe_num(reserva.get("coverageMonths"))
        comprometimento = (gasto / renda) if renda > 0 else None

        status = "stable"
        score = 70
        if saldo < 0 or contas_pendentes > 0 or (comprometimento is not None and comprometimento > 0.9):
            status = "attention"
            score = 55
        if (renda > 0 and saldo < -0.15 * renda) or (orcamento > 0 and gasto > 1.15 * orcamento) or (0 < cobertura < 1):
            status = "critical"
            score = 38
        if saldo > 0 and contas_pendentes <= 0 and cobertura >= 3 and (comprometimento is None or comprometimento < 0.7):
            status = "strong"
            score = 84

        if saldo < 0:
            main_risk = f"Fluxo do mês negativo em {_money(abs(saldo))} (renda {_money(renda)} x gastos {_money(gasto)})."
        elif contas_pendentes > 0:
            main_risk = f"Há {_money(contas_pendentes)} em obrigações ainda em aberto neste mês."
        elif comprometimento is not None and comprometimento > 0.9:
            main_risk = f"Os gastos consomem {comprometimento * 100:.0f}% da renda recebida."
        else:
            main_risk = "Sem risco dominante claro nos dados disponíveis."

        if saldo > 0:
            positive = f"O mês fecha com sobra de {_money(saldo)}."
        elif cobertura >= 1:
            positive = f"A reserva cobre cerca de {cobertura:.1f} mês(es) de orçamento."
        else:
            positive = "Os lançamentos já estão registrados para acompanhar o fluxo."

        if status in ("attention", "critical"):
            action = "Revise as obrigações em aberto e priorize equilibrar gastos e renda antes de assumir novos compromissos."
        else:
            action = "Mantenha o ritmo atual e direcione a sobra para a reserva ou para as metas ativas."

        normalized = {
            "status": status,
            "score": score,
            "title": "Análise financeira local",
            "summary": f"Resumo calculado localmente: renda {_money(renda)}, gastos {_money(gasto)} e saldo {_money(saldo)} no mês.",
            "mainRisk": main_risk,
            "positivePoint": positive,
            "actionProposal": action,
        }
        if reason:
            print(f"Resumo financeiro usando fallback local: {reason}")
        return {"analysis": normalized, "summary": normalized["summary"], "fallback": True}

    def _extract_json_object(raw_text: str) -> str:
        clean_text = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw_text or "", flags=re.IGNORECASE).strip()
        if clean_text.startswith("{") and clean_text.endswith("}"):
            return clean_text
        start = clean_text.find("{")
        end = clean_text.rfind("}")
        if start >= 0 and end > start:
            return clean_text[start:end + 1]
        return clean_text

    api_key = get_gemini_api_key()
    if not api_key:
        return _finance_fallback_analysis("Chave Gemini não configurada.")

    snapshot_text = json.dumps(snapshot, ensure_ascii=False, indent=2)

    prompt = f"""Você é um consultor financeiro pessoal. Analise o snapshot financeiro abaixo e gere um diagnóstico de saúde financeira em português.

Regras obrigatórias:
- Responda somente JSON válido, sem markdown, sem títulos externos e sem emojis.
- Não crie tarefas, registros, transações, metas ou qualquer ação dentro do sistema.
- O campo "actionProposal" deve ser apenas uma proposta textual para o usuário avaliar.
- Seja específico com os números mais relevantes.
- Destaque o risco principal primeiro, mas também reconheça um ponto positivo se houver.
- Ignore categorias de lançamentos financeiros. Elas estão temporariamente imprecisas e não devem fundamentar diagnóstico, risco ou proposta.

Formato exato:
{{
  "status": "critical" | "attention" | "stable" | "strong",
  "score": 0,
  "title": "frase curta com o diagnóstico central",
  "summary": "síntese em 1 ou 2 frases curtas",
  "mainRisk": "principal risco com números",
  "positivePoint": "ponto positivo real, ou cautela se não houver",
  "actionProposal": "proposta prática para o mês, sem executar nada"
}}

Critérios gerais:
- "critical": gasto muito acima da renda/orçamento, fluxo negativo severo, contas pendentes relevantes ou reserva ameaçada.
- "attention": há desequilíbrio, pendências ou progresso baixo, mas não é colapso imediato.
- "stable": receitas, gastos e obrigações estão sob controle.
- "strong": sobra consistente, reserva saudável e metas avançando.

SNAPSHOT:
{snapshot_text}"""

    try:
        client = genai.Client(api_key=api_key)
        response = generate_content_logged(
            client,
            model=GEMINI_BALANCED_MODEL,
            contents=prompt,
            feature="resumo_financeiro",
            db=get_db(),
            config=types.GenerateContentConfig(
                temperature=0.4,
                max_output_tokens=800,
                response_mime_type="application/json"
            )
        )
        raw_text = (response.text or "").strip()
        clean_text = _extract_json_object(raw_text)
        analysis = json.loads(clean_text)

        allowed_status = {"critical", "attention", "stable", "strong"}
        status = analysis.get("status")
        if status not in allowed_status:
            status = "attention"

        score = analysis.get("score")
        try:
            score = max(0, min(100, int(score)))
        except Exception:
            score = None

        normalized = {
            "status": status,
            "score": score,
            "title": str(analysis.get("title") or "Diagnóstico financeiro").strip(),
            "summary": str(analysis.get("summary") or "").strip(),
            "mainRisk": str(analysis.get("mainRisk") or "").strip(),
            "positivePoint": str(analysis.get("positivePoint") or "").strip(),
            "actionProposal": str(analysis.get("actionProposal") or "").strip(),
        }
        return {"analysis": normalized, "summary": normalized["summary"]}
    except Exception as e:
        print(f"Erro ao gerar resumo financeiro: {repr(e)}")
        return _finance_fallback_analysis(str(e))


@https_fn.on_call(memory=options.MemoryOption.MB_512, timeout_sec=60)
def gerarResumoSaude(req: https_fn.CallableRequest):
    """
    Gera um diagnóstico estruturado de saúde geral a partir de um snapshot compacto.
    Chamado apenas quando os dados de saúde mudam (fingerprint diferente).
    """
    from google import genai
    from google.genai import types

    data = req.data or {}
    snapshot = data.get('snapshot')

    if not snapshot or not isinstance(snapshot, dict):
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            message="snapshot é obrigatório."
        )

    import json
    import re

    def _health_fallback_analysis(reason: str = ""):
        biometria = snapshot.get("biometria") or {}
        atividade = snapshot.get("atividade7d") or {}
        sono = snapshot.get("sono7d") or {}
        dor = snapshot.get("dor7d") or {}
        avg_steps = atividade.get("mediaPassosDia")
        avg_sleep = sono.get("mediaTotalMin")
        trend_weight = biometria.get("tendencia30d")
        pain_crisis = dor.get("crisesRecentes") or 0
        ideal_days = atividade.get("diasMetaIdealAtingida") or 0

        status = "stable"
        score = 72
        if pain_crisis > 0 or (isinstance(avg_sleep, (int, float)) and avg_sleep < 360) or (isinstance(trend_weight, (int, float)) and trend_weight > 1.5):
            status = "attention"
            score = 58
        if pain_crisis >= 2 or (
            isinstance(avg_sleep, (int, float)) and avg_sleep < 330
            and isinstance(avg_steps, (int, float)) and avg_steps < 3500
        ):
            status = "critical"
            score = 38
        if (
            isinstance(avg_sleep, (int, float)) and avg_sleep >= 420
            and isinstance(avg_steps, (int, float)) and avg_steps >= 7000
            and ideal_days >= 3
            and pain_crisis == 0
        ):
            status = "strong"
            score = 84

        steps_text = f"{int(avg_steps):,}".replace(",", ".") + " passos/dia" if isinstance(avg_steps, (int, float)) else "passos ainda sem media"
        sleep_text = f"{avg_sleep / 60:.1f}h de sono medio" if isinstance(avg_sleep, (int, float)) else "sono ainda sem media"
        weight_text = (
            f"tendencia de {'+' if trend_weight > 0 else ''}{trend_weight:.1f} kg em 30 dias"
            if isinstance(trend_weight, (int, float))
            else "peso sem tendencia recente"
        )
        main_risk = (
            f"{pain_crisis} crise(s) de dor registradas nos ultimos 7 dias."
            if pain_crisis
            else "Sem risco dominante claro nos dados disponiveis."
        )
        positive = (
            f"{ideal_days} dia(s) atingiram a meta ideal de caminhada."
            if ideal_days > 0
            else "Sono medio acima de 7h no periodo com dados."
            if isinstance(avg_sleep, (int, float)) and avg_sleep >= 420
            else "O painel ja possui dados suficientes para acompanhar tendencia."
        )

        normalized = {
            "status": status,
            "score": score,
            "title": "Analise local de saude",
            "summary": f"Resumo calculado localmente: {steps_text}, {sleep_text} e {weight_text}.",
            "mainRisk": main_risk,
            "positivePoint": positive,
            "actionProposal": "Nesta semana, priorize sincronizar telemetria diariamente e manter a caminhada minima antes de aumentar intensidade.",
        }
        if reason:
            print(f"Resumo de saude usando fallback local: {reason}")
        return {"analysis": normalized, "summary": normalized["summary"], "fallback": True}

    def _extract_json_object(raw_text: str) -> str:
        clean_text = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw_text or "", flags=re.IGNORECASE).strip()
        if clean_text.startswith("{") and clean_text.endswith("}"):
            return clean_text
        start = clean_text.find("{")
        end = clean_text.rfind("}")
        if start >= 0 and end > start:
            return clean_text[start:end + 1]
        return clean_text

    api_key = get_gemini_api_key()
    if not api_key:
        return _health_fallback_analysis("Chave Gemini nao configurada.")

    snapshot_text = json.dumps(snapshot, ensure_ascii=False, indent=2)

    prompt = f"""Voce e um coach de saude pessoal e analista de bem-estar. Analise o snapshot de saude abaixo e gere um diagnostico em portugues.

O snapshot contem dados reais do usuario: biometria (peso e meta), caminhada/atividade (media de passos, distancia, calorias e dias com meta atingida), sono (media de horas, sono profundo, dias abaixo de 7h) e dor (escala 0-10 matinal/noturna, crises recentes).

Regras obrigatorias:
- Responda somente JSON valido, sem markdown, sem titulos externos e sem emojis.
- Nao crie tarefas, registros ou qualquer acao dentro do sistema.
- O campo "actionProposal" deve ser uma proposta textual concreta e especifica para o usuario avaliar esta semana.
- Seja especifico com os numeros mais relevantes do snapshot.
- Se algum dado estiver ausente (null), ignore essa dimensao na analise.
- Priorize as dimensoes com dados mais completos.
- Nao avalie habitos diarios, rotina lombar, treino de forca, flexoes, barras, prancha, ponte, bird-dog ou agachamentos.

Formato exato:
{{
  "status": "critical" | "attention" | "stable" | "strong",
  "score": 0,
  "title": "frase curta com o diagnostico central",
  "summary": "sintese em 1 ou 2 frases, com numeros reais do snapshot",
  "mainRisk": "principal ponto de atencao com numeros especificos",
  "positivePoint": "ponto positivo real do periodo, ou cautela se nao houver",
  "actionProposal": "proposta pratica e especifica para esta semana, sem executar nada"
}}

Criterios de status:
- "critical": crise de dor recente com outros indicadores negativos, sono muito baixo, passos muito baixos ou peso muito acima da meta com tendencia de alta.
- "attention": desequilibrio em 2+ dimensoes (ex: sono abaixo de 6h na media, passos baixos, peso crescendo consistentemente ou dor recorrente).
- "stable": maioria das dimensoes controladas, sem deterioracao clara.
- "strong": peso proximo ou abaixo da meta, sono adequado (>7h media), caminhada regular e dor baixa ou ausente.

SNAPSHOT:
{snapshot_text}"""
    try:
        client = genai.Client(api_key=api_key)
        response = generate_content_logged(
            client,
            model=GEMINI_FRONTIER_MODEL,
            contents=prompt,
            feature="resumo_saude",
            db=get_db(),
            config=types.GenerateContentConfig(
                temperature=0.35,
                max_output_tokens=1024,
                response_mime_type="application/json"
            )
        )
        raw_text = (response.text or "").strip()
        clean_text = _extract_json_object(raw_text)
        analysis = json.loads(clean_text)

        allowed_status = {"critical", "attention", "stable", "strong"}
        status = analysis.get("status")
        if status not in allowed_status:
            status = "attention"

        score = analysis.get("score")
        try:
            score = max(0, min(100, int(score)))
        except Exception:
            score = None

        normalized = {
            "status": status,
            "score": score,
            "title": str(analysis.get("title") or "Diagnóstico de saúde").strip(),
            "summary": str(analysis.get("summary") or "").strip(),
            "mainRisk": str(analysis.get("mainRisk") or "").strip(),
            "positivePoint": str(analysis.get("positivePoint") or "").strip(),
            "actionProposal": str(analysis.get("actionProposal") or "").strip(),
        }
        return {"analysis": normalized, "summary": normalized["summary"]}
    except Exception as e:
        print(f"Erro ao gerar resumo de saúde: {e}")
        return _health_fallback_analysis(str(e))


def get_people_service():
    from googleapiclient.discovery import build
    return build('people', 'v1', credentials=get_google_creds())


def sync_google_contacts_internal(db, sync_ref=None, logs=None):
    """
    Sincroniza contatos do Google People API para o Firestore.
    Otimizado para usar cache em memória e evitar consultas repetitivas no banco de dados.
    """
    if logs is None:
        logs = []

    def log_helper(msg):
        if sync_ref is not None:
            log_to_firestore(sync_ref, logs, msg, True)
        else:
            print(msg)

    log_helper("[SYNC] Iniciando sincronização de contatos do Google...")

    try:
        # 1. Obter credenciais e instanciar o People API
        creds = get_google_creds(scopes=['https://www.googleapis.com/auth/contacts'])
        from googleapiclient.discovery import build
        service = build('people', 'v1', credentials=creds)
        
        # 2. Buscar contatos do Google
        connections = []
        next_page_token = None
        MAX_CONNECTIONS_RETRIES = 5

        while True:
            results = None
            req_err = None

            for attempt in range(MAX_CONNECTIONS_RETRIES):
                try:
                    results = service.people().connections().list(
                        resourceName='people/me',
                        pageSize=100,
                        pageToken=next_page_token,
                        personFields='names,emailAddresses,phoneNumbers,biographies,metadata'
                    ).execute()
                    req_err = None
                    break
                except Exception as e:
                    req_err = e
                    is_last_attempt = attempt == MAX_CONNECTIONS_RETRIES - 1
                    if is_last_attempt or not _is_retryable_google_api_error(e):
                        break
                    wait_s = (2 ** attempt) + random.uniform(0, 1)
                    print(f"[SYNC] Erro temporário na Google Connections API (tentativa {attempt + 1}/{MAX_CONNECTIONS_RETRIES}): {e}. Nova tentativa em {wait_s:.1f}s.")
                    time.sleep(wait_s)

            if req_err is not None:
                # Falha definitiva (erro não retryable ou tentativas esgotadas): a
                # sincronização de contatos é interrompida e reportada como erro real,
                # em vez de encerrar a paginação silenciosamente com dados parciais.
                log_helper(f"[SYNC] Erro ao chamar Google Connections API após {MAX_CONNECTIONS_RETRIES} tentativas: {req_err}")
                raise req_err

            connections.extend(results.get('connections', []))
            next_page_token = results.get('nextPageToken')
            if not next_page_token:
                break

        # 3. Perfis existentes do Firestore em lote para otimização em memória
        existing_profiles = []
        all_docs = db.collection('perfil_pessoas').stream()
        for doc_snap in all_docs:
            p_data = doc_snap.to_dict() or {}
            p_data['id'] = doc_snap.id
            existing_profiles.append(p_data)
        
        # Mapeamentos O(1) para busca rápida
        by_google_id = {p.get('google_contact_id'): p for p in existing_profiles if p.get('google_contact_id')}
        by_email = {p.get('email').lower(): p for p in existing_profiles if p.get('email')}
        
        by_phone = {}
        for p in existing_profiles:
            p_phone = p.get('telefone')
            if p_phone:
                clean_p_phone = "".join(c for c in p_phone if c.isdigit())
                if clean_p_phone:
                    by_phone[clean_p_phone] = p

        stats = {"added": 0, "merged": 0, "errors": 0}
        
        for person in connections:
            try:
                resource_name = person.get('resourceName') # Ex: people/c1234567890
                metadata = person.get('metadata', {})
                etag = metadata.get('sources', [{}])[0].get('etag') or person.get('etag')
                
                # Nomes
                names = person.get('names', [])
                name = names[0].get('displayName') if names else None
                if not name:
                    continue # Contatos sem nome são ignorados
                    
                # E-mails
                emails = [e.get('value') for e in person.get('emailAddresses', []) if e.get('value')]
                email = emails[0] if emails else ""
                
                # Telefones
                phones = [p.get('value') for p in person.get('phoneNumbers', []) if p.get('value')]
                phone = phones[0] if phones else ""
                
                # Biografias
                biographies = person.get('biographies', [])
                bio = biographies[0].get('value') if biographies else ""
                
                # Gerar iniciais e cor do avatar
                initials = "".join([part[0].upper() for part in name.split() if part])[:2]
                colors = ["bg-indigo-500", "bg-purple-500", "bg-pink-500", "bg-rose-500", "bg-amber-500", "bg-emerald-500", "bg-teal-500", "bg-cyan-500", "bg-sky-500", "bg-blue-500"]
                import hashlib
                color_idx = int(hashlib.md5(name.encode('utf-8')).hexdigest(), 16) % len(colors)
                avatar_color = colors[color_idx]
                
                # 4. Verificar se já existe no Firestore usando busca em memória
                existing_profile = None
                
                if resource_name in by_google_id:
                    existing_profile = by_google_id[resource_name]
                else:
                    if email and email.lower() in by_email:
                        existing_profile = by_email[email.lower()]
                    
                    if not existing_profile and phone:
                        clean_phone = "".join(c for c in phone if c.isdigit())
                        if clean_phone:
                            if clean_phone in by_phone:
                                existing_profile = by_phone[clean_phone]
                            else:
                                for clean_p_phone, p in by_phone.items():
                                    if clean_p_phone == clean_phone or clean_p_phone.endswith(clean_phone) or clean_phone.endswith(clean_p_phone):
                                        existing_profile = p
                                        break
                
                now_str = datetime.now(timezone.utc).isoformat()
                
                if existing_profile:
                    # Atualização (Merge)
                    doc_ref = db.collection('perfil_pessoas').document(existing_profile['id'])
                    existing_data = existing_profile
                    
                    update_payload = {
                        "google_contact_id": resource_name,
                        "google_etag": etag,
                        "data_atualizacao": now_str
                    }
                    
                    tags = existing_data.get('tags') or []
                    if not isinstance(tags, list):
                        tags = [tags]
                    if 'Contatos do Google' not in tags:
                        tags.append('Contatos do Google')
                    update_payload['tags'] = tags
                    
                    if not existing_data.get('email') and email:
                        update_payload['email'] = email
                    if not existing_data.get('telefone') and phone:
                        update_payload['telefone'] = phone
                    if not existing_data.get('observacoes') and bio:
                        update_payload['observacoes'] = bio
                    if not existing_data.get('avatar_color'):
                        update_payload['avatar_color'] = avatar_color
                    if not existing_data.get('avatar_initials'):
                        update_payload['avatar_initials'] = initials
                        
                    doc_ref.update(update_payload)
                    
                    # Atualiza em cache
                    existing_profile.update(update_payload)
                    stats["merged"] += 1
                else:
                    # Cadastro Novo
                    doc_ref = db.collection('perfil_pessoas').document()
                    new_id = doc_ref.id
                    new_payload = {
                        "nome": name,
                        "email": email,
                        "telefone": phone,
                        "tags": ['Contatos do Google'],
                        "origem": 'google_contacts',
                        "google_contact_id": resource_name,
                        "google_etag": etag,
                        "observacoes": bio,
                        "avatar_color": avatar_color,
                        "avatar_initials": initials,
                        "data_criacao": now_str,
                        "data_atualizacao": now_str
                    }
                    doc_ref.set(new_payload)
                    new_payload['id'] = new_id
                    
                    # Atualiza os índices em cache
                    by_google_id[resource_name] = new_payload
                    if email:
                        by_email[email.lower()] = new_payload
                    if phone:
                        clean_phone = "".join(c for c in phone if c.isdigit())
                        if clean_phone:
                            by_phone[clean_phone] = new_payload
                            
                    stats["added"] += 1
                    
            except Exception as item_err:
                print(f"[GOOGLE SYNC] Erro ao processar contato individual: {item_err}")
                stats["errors"] += 1
                
        log_helper(f"[SYNC] Sincronização de contatos concluída: {stats['added']} importados, {stats['merged']} mesclados.")
        return stats
        
    except Exception as e:
        # Inclui o tipo da exceção e o ponto de origem: falhas de dependência
        # (ex.: gRPC/Firestore em "'_UnaryStreamMultiCallable' object has no
        # attribute '_retry'") acontecem fora da People API e, sem isso, a
        # notificação culpa erroneamente o "sync de contatos do Google".
        import traceback
        origin = traceback.extract_tb(e.__traceback__)[-1] if e.__traceback__ else None
        origin_text = f" (em {origin.name}:{origin.lineno})" if origin else ""
        print(f"[SYNC][!] Traceback completo:\n{traceback.format_exc()}")
        log_helper(f"[SYNC][!] Erro na sincronização de contatos do Google: {type(e).__name__}: {e}{origin_text}")
        return None


@https_fn.on_call(memory=options.MemoryOption.MB_512, timeout_sec=240)
def sync_google_contacts(req: https_fn.CallableRequest):
    """
    Sincroniza contatos do Google People API para a coleção 'perfil_pessoas' no Firestore.
    Expõe como Callable Cloud Function.
    """
    try:
        db = get_db()
        stats = sync_google_contacts_internal(db, None, None)
        if stats is None:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INTERNAL,
                message="Erro durante a sincronização de contatos do Google."
            )
        return {"success": True, "stats": stats}
    except Exception as e:
        print(f"Erro na sincronização de contatos do Google: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=f"Erro na sincronização de contatos do Google: {str(e)}"
        )


@https_fn.on_call(memory=options.MemoryOption.MB_512, timeout_sec=120)
def generate_contact_summary(req: https_fn.CallableRequest):
    """
    Gera um resumo do perfil de uma pessoa baseando-se no seu histórico de interações (timeline)
    utilizando a inteligência artificial do Gemini e grava no campo 'resumo_ia'.
    """
    try:
        db = get_db()
        data = req.data or {}
        pessoa_id = data.get("pessoa_id")
        
        if not pessoa_id:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="ID da pessoa é obrigatório."
            )
            
        # 1. Recuperar dados do contato
        contact_ref = db.collection("perfil_pessoas").document(pessoa_id)
        contact_doc = contact_ref.get()
        if not contact_doc.exists:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.NOT_FOUND,
                message="Contato não encontrado."
            )
        contact_data = contact_doc.to_dict() or {}
        
        # 2. Buscar interações da pessoa (ordenado em memória para evitar a dependência de índice composto)
        interactions_ref = db.collection("interacoes_pessoas").where("pessoa_id", "==", pessoa_id)
        interactions = [doc.to_dict() for doc in interactions_ref.stream()]
        interactions.sort(key=lambda x: x.get("data", ""), reverse=True)
        interactions = interactions[:100]
        
        if not interactions:
            return {
                "success": True, 
                "resumo_ia": "Nenhuma interação registrada ainda para este contato.",
                "message": "Nenhuma interação cadastrada."
            }
            
        # 3. Consolidar o histórico de interações para o Prompt do Gemini
        history_lines = []
        for idx, inter in enumerate(interactions):
            data_inter = inter.get("data", "Data desconhecida")
            descricao = inter.get("descricao", "")
            tipo = inter.get("tipo", "interação")
            history_lines.append(f"- [{data_inter}] ({tipo}): {descricao}")
            
        history_text = "\n".join(history_lines)
        
        # 4. Obter a chave do Gemini
        keys_doc = _cached_doc_get(db, 'system', 'api_keys')
        gemini_key = keys_doc.to_dict().get('gemini_api_key') if keys_doc.exists else None
        
        if not gemini_key:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
                message="Chave Gemini não configurada no sistema (system/api_keys)."
            )
            
        # 5. Instanciar o Gemini
        from google import genai
        client = genai.Client(api_key=gemini_key)
        
        prompt = f"""
        Você é o HERMES Master IA, assistente inteligente integrado ao sistema de produtividade pessoal do André.
        Abaixo estão os dados de um contato cadastrado e o histórico de suas interações registradas no sistema (mencionados em tarefas, diários de bordo ou reuniões).
        
        NOME DO CONTATO: {contact_data.get('nome')}
        DETALHES: E-mail: {contact_data.get('email', 'N/A')}, Telefone: {contact_data.get('telefone', 'N/A')}, Tags: {', '.join(contact_data.get('tags', []))}
        
        HISTÓRICO DE INTERAÇÕES E MENÇÕES:
        {history_text}
        
        Sua tarefa é analisar o histórico acima e escrever um resumo conciso, executivo e em terceira pessoa em português do Brasil.
        O resumo deve identificar quem é este contato (ex: se é um fornecedor, colega de trabalho, desenvolvedor, designer, etc.), os principais assuntos ou projetos em que cooperou e o tom/contexto geral das interações.
        
        REGRAS IMPORTANTES:
        - O texto deve ser direto e profissional, com no máximo 3 ou 4 frases curtas.
        - Não adicione introduções ou saudações ("Aqui está o resumo..."). Retorne apenas o resumo consolidado.
        - Foque apenas nas interações reais demonstradas no histórico.
        """
        
        response = client.models.generate_content(
            model="gemini-3.5-flash-lite",
            contents=prompt
        )
        
        resumo_texto = response.text.strip() if response.text else "Não foi possível gerar o resumo."
        
        # 6. Gravar o resumo no Firestore
        contact_ref.update({
            "resumo_ia": resumo_texto,
            "data_atualizacao": datetime.now(timezone.utc).isoformat()
        })
        
        return {
            "success": True,
            "resumo_ia": resumo_texto
        }
    except Exception as e:
        print(f"Erro ao gerar resumo do contato por IA: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=f"Erro ao gerar resumo por IA: {str(e)}"
        )


@https_fn.on_call(memory=options.MemoryOption.MB_512, timeout_sec=120)
def merge_contacts(req: https_fn.CallableRequest):
    """
    Mescla múltiplos perfis de contatos secundários em um único perfil principal.
    Transfere todas as interações e unifica as tags/e-mails/telefones/observações.
    """
    try:
        db = get_db()
        data = req.data or {}
        primary_id = data.get("primary_id")
        secondary_ids = data.get("secondary_ids", [])
        
        if not primary_id or not secondary_ids:
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
                message="ID do contato principal (primary_id) e IDs secundários (secondary_ids) são obrigatórios."
            )
            
        from contact_merge_utils import execute_contact_merge
        res = execute_contact_merge(db, primary_id, secondary_ids)
        if not res.get("success"):
            raise https_fn.HttpsError(
                code=https_fn.FunctionsErrorCode.NOT_FOUND,
                message=res.get("error", "Erro ao mesclar contatos.")
            )
        return res
    except https_fn.HttpsError:
        raise
    except Exception as e:
        print(f"Erro ao mesclar contatos: {e}")
        raise https_fn.HttpsError(
            code=https_fn.FunctionsErrorCode.INTERNAL,
            message=f"Erro ao mesclar contatos: {str(e)}"
        )


@https_fn.on_call(memory=options.MemoryOption.MB_256, timeout_sec=60)
def linkWhatsappContacts(req: https_fn.CallableRequest) -> dict:
    """Cruza contatos em `perfil_pessoas` com o registro de chats do WhatsApp (`whatsapp_chats`
    com fallback para allowlist) usando matching determinístico por últimos 8 dígitos do telefone
    (last-8).
    
    Regra de ouro de segurança: vincula APENAS quando houver correspondência 1-para-1 exata
    (1 pessoa ↔ 1 chat @c.us). Se qualquer um dos lados tiver 2 ou mais ocorrências para o
    mesmo last-8, o vínculo é considerado ambíguo e pulado com relatório.
    """
    _require_internal_user(req)
    db = get_db()
    from collections import defaultdict
    from phone_utils import last8, chat_id_last8

    # 1. Carrega chats individuais (@c.us) do registro whatsapp_chats (e allowlist como fallback)
    chats_by_last8: dict[str, list[dict]] = defaultdict(list)
    seen_chat_ids: set[str] = set()

    for doc in db.collection("whatsapp_chats").limit(3000).stream():
        d = doc.to_dict() or {}
        cid = str(d.get("chat_id") or doc.id).strip()
        if not cid.endswith("@c.us"):
            continue
        seen_chat_ids.add(cid)
        l8 = chat_id_last8(cid)
        if l8:
            chats_by_last8[l8].append({
                "chat_id": cid,
                "chat_name": str(d.get("chat_name") or cid).strip() or cid,
            })

    # Fallback allowlist para chats ainda não persistidos no registro
    settings_doc = db.collection("system").document("settings").get()
    settings_data = settings_doc.to_dict() if settings_doc.exists else {}
    allowlist = [
        str(x).strip() for x in (settings_data.get("whatsapp_ingest") or {}).get("chats_allowlist") or []
        if str(x).strip()
    ]
    for cid in allowlist:
        if cid.endswith("@c.us") and cid not in seen_chat_ids:
            seen_chat_ids.add(cid)
            l8 = chat_id_last8(cid)
            if l8:
                chats_by_last8[l8].append({
                    "chat_id": cid,
                    "chat_name": cid,
                })

    # 2. Carrega contatos em perfil_pessoas
    people_by_last8: dict[str, list[dict]] = defaultdict(list)
    sem_telefone = 0
    ja_vinculados = 0

    for doc in db.collection("perfil_pessoas").limit(2000).stream():
        pdata = doc.to_dict() or {}
        tel = pdata.get("telefone") or pdata.get("celular") or pdata.get("whatsapp") or ""
        l8 = last8(tel)
        if not l8:
            sem_telefone += 1
            continue
        people_by_last8[l8].append({
            "id": doc.id,
            "nome": str(pdata.get("nome") or "Sem nome").strip(),
            "telefone": str(tel).strip(),
            "whatsapp_chat_id": pdata.get("whatsapp_chat_id"),
        })

    # 3. Matching determinístico (1-para-1 estrito)
    vinculados = []
    ambiguos = []
    sem_match = 0
    batch = db.batch()
    batch_count = 0
    now_iso = datetime.now(timezone.utc).isoformat()

    for l8, people_list in people_by_last8.items():
        chat_list = chats_by_last8.get(l8, [])
        if len(people_list) == 1 and len(chat_list) == 1:
            person = people_list[0]
            chat = chat_list[0]
            if person.get("whatsapp_chat_id") == chat["chat_id"]:
                ja_vinculados += 1
            else:
                doc_ref = db.collection("perfil_pessoas").document(person["id"])
                batch.update(doc_ref, {
                    "whatsapp_chat_id": chat["chat_id"],
                    "data_atualizacao": now_iso,
                })
                batch_count += 1
                vinculados.append({
                    "pessoa_id": person["id"],
                    "nome": person["nome"],
                    "chat_id": chat["chat_id"],
                    "chat_name": chat["chat_name"],
                    "last8": l8,
                })
                if batch_count >= 450:
                    batch.commit()
                    batch = db.batch()
                    batch_count = 0
        elif len(chat_list) > 0 and (len(people_list) > 1 or len(chat_list) > 1):
            ambiguos.append({
                "last8": l8,
                "pessoas": [{"id": p["id"], "nome": p["nome"]} for p in people_list],
                "chats": [{"chat_id": c["chat_id"], "chat_name": c["chat_name"]} for c in chat_list],
            })
        else:
            sem_match += len(people_list)

    if batch_count > 0:
        batch.commit()

    return {
        "success": True,
        "vinculados": vinculados,
        "ambiguos": ambiguos,
        "sem_match": sem_match,
        "sem_telefone": sem_telefone,
        "ja_vinculados": ja_vinculados,
        "total_contatos_avaliados": len(vinculados) + len(ambiguos) + sem_match + sem_telefone + ja_vinculados,
    }


@storage_fn.on_object_finalized(
    bucket="gestao-hermes.firebasestorage.app",
    region="us-east1",
    timeout_sec=540,
    memory=options.MemoryOption.GB_4,
    cpu=2,
)
def on_long_transcription_uploaded(event: storage_fn.CloudEvent) -> None:
    """Transcreve arquivos pesados de áudio/vídeo enviados para `long_transcriptions/{uid}/{id}.{ext}`.

    Fluxo: normaliza qualquer mídia para AAC via ffmpeg -> Files API do Gemini ->
    transcrição literal com gemini-3.5-flash-lite -> grava no Firestore -> expurga o binário original.
    """
    import os as _os
    import time as _time
    import tempfile as _tempfile
    import subprocess as _subprocess
    from firebase_admin import storage as admin_storage

    object_path = (event.data.name or "")
    if not object_path.startswith("long_transcriptions/"):
        return  # ignora uploads de outras pastas no mesmo bucket

    # Esperado: long_transcriptions/{userId}/{transcriptionId}.{ext}
    parts = object_path.split("/")
    if len(parts) != 3:
        print(f"[long_transcription] caminho inesperado, ignorando: {object_path}")
        return
    _, user_id, file_with_ext = parts
    transcription_id = file_with_ext.rsplit(".", 1)[0]
    file_ext = file_with_ext.rsplit(".", 1)[1].lower() if "." in file_with_ext else ""

    db = get_db()
    doc_ref = db.collection("long_transcriptions").document(transcription_id)
    snap = doc_ref.get()
    if not snap.exists:
        print(f"[long_transcription] doc {transcription_id} inexistente, ignorando.")
        return
    data = snap.to_dict() or {}

    # Idempotência: Storage triggers são at-least-once. Só processamos registros 'Enviando'.
    if data.get("status") != "Enviando":
        print(f"[long_transcription] {transcription_id} já em '{data.get('status')}', ignorando re-trigger.")
        return
    if data.get("userId") and data.get("userId") != user_id:
        print(f"[long_transcription] userId divergente em {transcription_id}, ignorando.")
        return

    doc_ref.update({"status": "Processando", "updatedAt": firestore.SERVER_TIMESTAMP})

    from google import genai
    from google.genai import types
    from groq import Groq

    bucket_name = event.data.bucket
    local_media_path = None
    local_chunk_dir = None
    gemini_file = None
    gemini_client = None
    groq_client = None

    try:
        limit_long_trans = int(_os.environ.get("LIMIT_LONG_TRANSCRIPTION", "8"))
        
        # Desconta transcrições de hoje que deram erro do limite diário (usa apenas userId no Firestore para dispensar índice composto)
        start_of_today = datetime.now(timezone.utc).replace(hour=0, minute=0, second=0, microsecond=0)
        user_today_docs = db.collection("long_transcriptions").where("userId", "==", user_id).stream()
        non_error_count = 0
        for d in user_today_docs:
            ddata = d.to_dict() or {}
            if ddata.get("status") == "Erro":
                continue
            created_at = ddata.get("createdAt")
            if created_at:
                try:
                    dt = created_at.to_datetime() if hasattr(created_at, "to_datetime") else created_at
                    if isinstance(dt, datetime):
                        if dt.tzinfo is None:
                            dt = dt.replace(tzinfo=timezone.utc)
                        if dt >= start_of_today:
                            non_error_count += 1
                except Exception:
                    pass

        if non_error_count > limit_long_trans:
            raise RuntimeError(f"Você atingiu o limite diário de {limit_long_trans} transcrições longas.")

        check_and_increment_limit(db, user_id, "long_transcription", limit_long_trans)

        keys_doc = _cached_doc_get(db, "system", "api_keys")
        keys = keys_doc.to_dict() if keys_doc.exists else {}
        groq_key = keys.get("groq_api_key")
        gemini_key = keys.get("gemini_api_key")
        if not groq_key and not gemini_key:
            raise RuntimeError("Nenhum motor de transcrição está configurado.")
        if groq_key:
            groq_client = Groq(api_key=groq_key)
        if gemini_key:
            gemini_client = genai.Client(api_key=gemini_key)

        # 1. Baixar o binário do Storage para arquivo temporário
        bucket = admin_storage.bucket(bucket_name)
        blob = bucket.blob(object_path)
        suffix = f".{file_ext}" if file_ext else ""
        fd, local_media_path = _tempfile.mkstemp(suffix=suffix)
        _os.close(fd)
        blob.download_to_filename(local_media_path)

        # Arquivos de gravadores móveis podem ter edit lists/timestamps que o FFmpeg
        # interpreta como uma faixa vazia. Se o original já cabe no Whisper e possui
        # uma extensão aceita, tente-o diretamente antes de qualquer conversão.
        groq_direct_extensions = {"flac", "mp3", "mp4", "mpeg", "mpga", "m4a", "ogg", "wav", "webm"}
        original_size = _os.path.getsize(local_media_path)
        direct_groq_error = None
        if groq_client is not None and file_ext in groq_direct_extensions and original_size <= 24 * 1024 * 1024:
            doc_ref.update({"processingStage": "Transcrevendo arquivo original", "updatedAt": firestore.SERVER_TIMESTAMP})
            try:
                with open(local_media_path, "rb") as original_stream:
                    direct_result = groq_client.audio.transcriptions.create(
                        file=(f"original.{file_ext}", original_stream),
                        model="whisper-large-v3-turbo",
                        response_format="json",
                        language="pt",
                        temperature=0.0,
                    )
                direct_text = (direct_result.text or "").strip()
                if direct_text:
                    doc_ref.update({
                        "status": "Concluído",
                        "transcriptionRaw": direct_text,
                        "transcriptionEngine": "groq-direct",
                        "transcriptionChunks": 1,
                        "chunkDiagnostics": [{"bytes": original_size, "source": "original"}],
                        "processingStage": None,
                        "errorMessage": None,
                        "updatedAt": firestore.SERVER_TIMESTAMP,
                    })
                    print(
                        f"[long_transcription] {transcription_id} concluído diretamente "
                        f"pelo Groq ({len(direct_text)} chars, {original_size} bytes)."
                    )
                    return
                direct_groq_error = RuntimeError("Whisper retornou texto vazio para o arquivo original.")
            except Exception as exc:
                direct_groq_error = exc
                print(f"[long_transcription] Groq direto falhou: {exc}")

        # 2. Normalizar e segmentar sempre. WAV PCM evita qualquer ambiguidade de codec,
        # cabeçalho ou contêiner. Dez minutos em 16 kHz/mono/16-bit ocupam ~19,2 MB,
        # abaixo do limite de 25 MB por arquivo do Whisper.
        import imageio_ffmpeg
        import wave as _wave
        ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
        local_chunk_dir = _tempfile.mkdtemp(prefix="hermes_long_transcription_")
        chunk_pattern = _os.path.join(local_chunk_dir, "chunk_%03d.wav")
        doc_ref.update({"processingStage": "Preparando trechos de áudio", "updatedAt": firestore.SERVER_TIMESTAMP})
        def _clear_wav_chunks():
            for name in _os.listdir(local_chunk_dir):
                path = _os.path.join(local_chunk_dir, name)
                if name.lower().endswith(".wav") and _os.path.isfile(path):
                    _os.remove(path)

        def _build_wav_chunks(ignore_edit_list=False, source_path=None):
            _clear_wav_chunks()
            input_path = source_path or local_media_path
            input_options = []
            if ignore_edit_list:
                input_options = [
                    "-ignore_editlist", "1", "-fflags", "+genpts+discardcorrupt",
                    "-err_detect", "ignore_err", "-analyzeduration", "200M", "-probesize", "200M",
                ]
            command = [ffmpeg_exe, "-y", *input_options, "-i", input_path]
            command.extend([
                "-map", "0:a:0", "-vn", "-ac", "1", "-ar", "16000",
                "-af", "aresample=async=1:first_pts=0", "-acodec", "pcm_s16le",
                "-f", "segment", "-segment_format", "wav", "-segment_time", "600",
                "-reset_timestamps", "1", chunk_pattern,
            ])
            completed = _subprocess.run(command, check=True, capture_output=True)
            stderr_text = completed.stderr.decode("utf-8", "ignore")[-3000:]
            paths = [
                _os.path.join(local_chunk_dir, name)
                for name in sorted(_os.listdir(local_chunk_dir))
                if name.lower().endswith(".wav")
            ]
            diagnostics = []
            for path in paths:
                with _wave.open(path, "rb") as wav_file:
                    frames = wav_file.getnframes()
                    sample_rate = wav_file.getframerate()
                    diagnostics.append({
                        "bytes": _os.path.getsize(path),
                        "durationSeconds": round(frames / sample_rate, 2) if sample_rate else 0,
                        "channels": wav_file.getnchannels(),
                        "sampleRate": sample_rate,
                        "sampleWidth": wav_file.getsampwidth(),
                        "frames": frames,
                    })
            return paths, diagnostics, stderr_text

        repair_diagnostics = None
        chunk_paths, chunk_diagnostics, ffmpeg_stderr = _build_wav_chunks(ignore_edit_list=False)
        if not chunk_paths or not any(item.get("frames", 0) for item in chunk_diagnostics):
            print(
                f"[long_transcription] primeira extração vazia em {transcription_id}; "
                "tentando sem edit list e com timestamps reconstruídos."
            )
            chunk_paths, chunk_diagnostics, ffmpeg_stderr = _build_wav_chunks(ignore_edit_list=True)

        # Alguns M4A de gravadores Samsung chegam com o `mdat` preservado, mas a
        # tabela STSZ termina antes da quantidade declarada de amostras AAC. O
        # FFmpeg descarta toda a tabela nesse caso. Corrigimos somente uma copia
        # temporaria, limitando a contagem as entradas realmente presentes.
        if not chunk_paths or not any(item.get("frames", 0) for item in chunk_diagnostics):
            from mp4_repair import repair_truncated_stsz

            repaired_media_path = _os.path.join(
                local_chunk_dir,
                f"repaired_source.{file_ext or 'm4a'}",
            )
            repair_diagnostics = repair_truncated_stsz(local_media_path, repaired_media_path)
            if repair_diagnostics:
                print(
                    f"[long_transcription] STSZ truncado em {transcription_id}; "
                    f"tentando copia reparada: {repair_diagnostics}"
                )
                doc_ref.update({
                    "processingStage": "Recuperando audio de arquivo M4A danificado",
                    "mp4RepairDiagnostics": repair_diagnostics,
                    "updatedAt": firestore.SERVER_TIMESTAMP,
                })
                chunk_paths, chunk_diagnostics, ffmpeg_stderr = _build_wav_chunks(
                    ignore_edit_list=True,
                    source_path=repaired_media_path,
                )
        if not chunk_paths or not any(item.get("frames", 0) for item in chunk_diagnostics):
            doc_ref.update({
                "ffmpegDiagnostics": ffmpeg_stderr,
                "chunkDiagnostics": chunk_diagnostics,
                "mp4RepairDiagnostics": repair_diagnostics,
                "updatedAt": firestore.SERVER_TIMESTAMP,
            })
            raise RuntimeError(
                "O arquivo contém dados, mas nenhum quadro de áudio pôde ser decodificado. "
                f"Groq direto: {direct_groq_error or 'não tentado'}; "
                f"FFmpeg: {ffmpeg_stderr[-800:]}"
            )
        if any(not item.get("frames", 0) for item in chunk_diagnostics):
            raise RuntimeError(f"Um dos trechos WAV ficou vazio: {chunk_diagnostics}")
        for item in chunk_diagnostics:
            item.pop("frames", None)
        doc_ref.update({"chunkDiagnostics": chunk_diagnostics, "updatedAt": firestore.SERVER_TIMESTAMP})
        print(f"[long_transcription] {transcription_id} trechos WAV: {chunk_diagnostics}")

        # 3. Transcrever cada trecho com Whisper (motor dedicado). O Gemini fica como
        # fallback por trecho, e não mais como ponto único de falha do arquivo inteiro.
        prompt = (
            "Transcreva integral e literalmente o áudio a seguir para texto. "
            "Inclua tudo o que for falado, na ordem em que ocorre, sem resumir, sem corrigir, "
            "sem traduzir, sem adicionar comentários, títulos, marcações de tempo ou formatação extra. "
            "Responda apenas com o texto transcrito."
        )
        transcript_parts = []
        engines_used = set()
        was_truncated = False
        for index, chunk_path in enumerate(chunk_paths, start=1):
            doc_ref.update({
                "processingStage": f"Transcrevendo trecho {index} de {len(chunk_paths)}",
                "updatedAt": firestore.SERVER_TIMESTAMP,
            })
            chunk_text = ""
            groq_error = None
            gemini_error = None
            if groq_client is not None:
                try:
                    with open(chunk_path, "rb") as chunk_stream:
                        result = groq_client.audio.transcriptions.create(
                            file=(_os.path.basename(chunk_path), chunk_stream),
                            model="whisper-large-v3-turbo",
                            response_format="json",
                            language="pt",
                            temperature=0.0,
                        )
                    chunk_text = (result.text or "").strip()
                    if chunk_text:
                        engines_used.add("groq")
                except Exception as exc:
                    groq_error = exc
                    print(f"[long_transcription] Groq falhou no trecho {index}: {exc}")

            if not chunk_text and gemini_client is not None:
                try:
                    gemini_file = gemini_client.files.upload(
                        file=chunk_path,
                        config=types.UploadFileConfig(
                            mime_type="audio/wav",
                            display_name=f"{transcription_id}_trecho_{index}.wav",
                        ),
                    )
                    waited = 0
                    while str(getattr(gemini_file.state, "name", gemini_file.state)) == "PROCESSING":
                        if waited >= 120:
                            raise TimeoutError("Files API demorou demais para preparar um trecho (>2min).")
                        _time.sleep(3)
                        waited += 3
                        gemini_file = gemini_client.files.get(name=gemini_file.name)
                    final_state = str(getattr(gemini_file.state, "name", gemini_file.state))
                    if final_state == "FAILED":
                        raise RuntimeError("Files API falhou ao preparar um trecho de áudio.")
                    response = gemini_client.models.generate_content(
                        model="gemini-3.5-flash-lite",
                        contents=[prompt, gemini_file],
                        config=types.GenerateContentConfig(max_output_tokens=16384),
                    )
                    chunk_text = (response.text or "").strip()
                    try:
                        finish_reason = str(getattr(response.candidates[0], "finish_reason", "") or "")
                        was_truncated = was_truncated or "MAX_TOKENS" in finish_reason.upper()
                    except Exception:
                        pass
                    if chunk_text:
                        engines_used.add("gemini")
                except Exception as exc:
                    gemini_error = exc
                    print(f"[long_transcription] Gemini falhou no trecho {index}: {exc}")
                finally:
                    if gemini_file is not None:
                        try:
                            gemini_client.files.delete(name=gemini_file.name)
                        except Exception:
                            pass
                        gemini_file = None

            if not chunk_text:
                if groq_error is not None or gemini_error is not None:
                    raise RuntimeError(
                        f"Falha ao transcrever o trecho {index}. "
                        f"Groq: {groq_error or 'indisponível'}; "
                        f"Gemini: {gemini_error or 'indisponível'}"
                    )
                raise RuntimeError(f"Os motores retornaram vazio no trecho {index}.")
            transcript_parts.append(chunk_text)

        transcription_raw = "\n\n".join(transcript_parts).strip()
        if not transcription_raw:
            raise RuntimeError("A transcrição final ficou vazia.")

        warning_messages = []
        if repair_diagnostics:
            warning_messages.append(
                "A gravação M4A estava com o índice de áudio truncado; "
                "foi transcrita a parte que pôde ser recuperada."
            )
        if was_truncated:
            warning_messages.append(
                "A transcrição pode ter sido truncada por exceder o limite de saída do modelo "
                "(áudio muito longo)."
            )

        doc_ref.update({
            "status": "Concluído",
            "transcriptionRaw": transcription_raw,
            "transcriptionEngine": "+".join(sorted(engines_used)),
            "transcriptionChunks": len(chunk_paths),
            "processingStage": None,
            "errorMessage": " ".join(warning_messages) or None,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        })
        print(
            f"[long_transcription] {transcription_id} concluído "
            f"({len(transcription_raw)} chars, trechos={len(chunk_paths)}, "
            f"motores={sorted(engines_used)}, truncado={was_truncated})."
        )

    except Exception as e:
        msg = str(e)
        stderr = getattr(e, "stderr", None)
        if stderr:
            try:
                msg = f"{msg} | ffmpeg: {stderr.decode('utf-8', 'ignore')[-400:]}"
            except Exception:
                pass
        print(f"[long_transcription] ERRO em {transcription_id}: {msg}")
        decrement_limit(db, user_id, "long_transcription")
        try:
            doc_ref.update({
                "status": "Erro",
                "errorMessage": msg[:1500],
                "processingStage": None,
                "updatedAt": firestore.SERVER_TIMESTAMP,
            })
        except Exception as e2:
            print(f"[long_transcription] falha ao gravar status de erro: {e2}")

    finally:
        # Expurgo: deleta o binário original do Storage independente do resultado
        try:
            admin_storage.bucket(bucket_name).blob(object_path).delete()
        except Exception as e:
            print(f"[long_transcription] falha ao expurgar {object_path}: {e}")
        # Deleta eventual arquivo temporário ainda ativo na Files API do Gemini
        try:
            if gemini_file is not None and gemini_client is not None:
                gemini_client.files.delete(name=gemini_file.name)
        except Exception:
            pass
        # Limpa temporários locais
        if local_media_path:
            try:
                if _os.path.exists(local_media_path):
                    _os.remove(local_media_path)
            except Exception:
                pass
        if local_chunk_dir:
            try:
                import shutil as _shutil
                _shutil.rmtree(local_chunk_dir, ignore_errors=True)
            except Exception:
                pass
